#!/usr/bin/env python3
"""
search_engine.py — Strukturierte Suche fuer AssistantDev Memory.

Drei Hauptklassen:
  SearchIndex  — Baut und verwaltet den Index pro Agent
  QueryParser  — Versteht natuerliche Sprache (DE/EN/PT)
  HybridSearch — Mehrstufige Suche mit BM25 + Volltext

Nutzung:
  from search_engine import SearchIndex, QueryParser, HybridSearch

  index = SearchIndex(speicher_path)
  index.build_or_update()

  query = QueryParser.parse("suche die Email von Simonas von gestern")
  results = HybridSearch.search(query, speicher_path, max_results=3)
"""

import os
import re
import json
import time
import datetime
import unicodedata
import threading

# ─── UNICODE NORMALISIERUNG ──────────────────────────────────────────────────

def normalize_unicode(s):
    """Normalize unicode for fuzzy matching: Simonas -> simonas, Vysniukas -> vysniukas."""
    return unicodedata.normalize('NFKD', s.lower()).encode('ascii', 'ignore').decode()


def fuzzy_match(needle, haystack, threshold=0.70):
    """Check if needle fuzzy-matches any word in haystack.
    Returns (match_found, score_bonus).
    Uses: prefix match, edit distance, bigram overlap, compound word matching.
    All inputs should already be normalize_unicode'd (lowercase, ascii).
    """
    if not needle or len(needle) < 3:
        return False, 0

    # Normalize separators: treat _-. as spaces for word splitting
    hay_clean = haystack.replace('_', ' ').replace('-', ' ').replace('.', ' ').replace('@', ' ')
    words = hay_clean.split()
    best_score = 0

    # Also check multi-word compound: "arnevidar" matches "arne vidar"
    hay_flat = hay_clean.replace(' ', '')
    if len(needle) >= 4 and needle in hay_flat and needle not in haystack:
        best_score = max(best_score, 4)

    for word in words:
        if len(word) < 2:
            continue

        # Exact substring — already handled elsewhere, skip
        if needle in word or word in needle:
            continue

        # 1. Prefix match (min 3 chars shared prefix)
        prefix_len = min(len(needle), len(word))
        common = 0
        for i in range(prefix_len):
            if needle[i] == word[i]:
                common += 1
            else:
                break
        if common >= 3 and common >= len(needle) * 0.55:
            best_score = max(best_score, 3)
            continue

        # 2. Edit distance (Levenshtein) for similar-length words
        if abs(len(needle) - len(word)) <= 3 and len(needle) >= 4:
            dist = _edit_distance(needle, word)
            max_len = max(len(needle), len(word))
            similarity = 1.0 - (dist / max_len)
            if similarity >= threshold:
                best_score = max(best_score, int(similarity * 5))
                continue

        # 3. Bigram overlap for longer words
        if len(needle) >= 4 and len(word) >= 4:
            n_bigrams = set(needle[i:i+2] for i in range(len(needle)-1))
            w_bigrams = set(word[i:i+2] for i in range(len(word)-1))
            if n_bigrams and w_bigrams:
                overlap = len(n_bigrams & w_bigrams) / len(n_bigrams | w_bigrams)
                if overlap >= 0.35:
                    best_score = max(best_score, int(overlap * 4))

    # 4. Multi-word needle: check each part separately
    if ' ' in needle and best_score == 0:
        parts = needle.split()
        part_scores = []
        for part in parts:
            if len(part) >= 3:
                pm, ps = fuzzy_match(part, haystack, threshold)
                # Also check exact substring for each part
                if part in haystack:
                    part_scores.append(3)
                elif pm:
                    part_scores.append(ps)
        if len(part_scores) >= 1:
            best_score = max(best_score, sum(part_scores))

    return best_score > 0, best_score


def _edit_distance(s1, s2):
    """Levenshtein edit distance — optimized for short strings."""
    if len(s1) < len(s2):
        return _edit_distance(s2, s1)
    if len(s2) == 0:
        return len(s1)
    prev = list(range(len(s2) + 1))
    for i, c1 in enumerate(s1):
        curr = [i + 1]
        for j, c2 in enumerate(s2):
            cost = 0 if c1 == c2 else 1
            curr.append(min(curr[j] + 1, prev[j + 1] + 1, prev[j] + cost))
        prev = curr
    return prev[-1]


# ─── CONSTANTS ────────────────────────────────────────────────────────────────

BINARY_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.mp4', '.mov', '.mp3', '.wav',
               '.zip', '.gz', '.tar', '.dmg', '.iso'}

STOPWORDS = {
    # Deutsch
    'die', 'der', 'das', 'ein', 'eine', 'und', 'oder', 'von', 'fuer', 'für',
    'mit', 'den', 'dem', 'des', 'auf', 'aus', 'bei', 'bis',
    'ist', 'mir', 'ich', 'er', 'sie', 'wir', 'ihr',
    'ueber', 'über', 'wie', 'was', 'wer', 'zum', 'zur',
    'meine', 'meinen', 'meinem', 'dein', 'sein', 'ihre',
    'bitte', 'nochmal', 'noch', 'mal', 'auch', 'nur', 'sehr',
    'dann', 'jetzt', 'hier', 'dort', 'aber', 'wenn', 'weil',
    'kann', 'dass', 'nach', 'hat', 'haben', 'hatte', 'werden',
    'nicht', 'schon', 'gibt', 'gibt', 'alle', 'alles',
    # Englisch
    'the', 'and', 'for', 'from', 'with', 'that', 'this',
    'also', 'please', 'just', 'have', 'has', 'had', 'will',
    'would', 'could', 'should', 'are', 'was', 'were', 'been',
    'not', 'all', 'some', 'any', 'each', 'every', 'into',
    'about', 'than', 'then', 'when', 'where', 'which', 'who',
    'how', 'can', 'may', 'more', 'most', 'very', 'your', 'our',
    # Portugiesisch
    'uma', 'uns', 'umas', 'para', 'com', 'por', 'que',
    'mas', 'nao', 'não', 'ao', 'na', 'no', 'nas', 'nos',
    'pelo', 'pela', 'este', 'essa', 'esse', 'isso',
    'dos', 'tem', 'algum', 'alguma',
}

# ─── SEARCH TRIGGER KEYWORDS ─────────────────────────────────────────────────

SEARCH_ACTIONS = {
    'finde', 'suche', 'such', 'zeig', 'zeige', 'nachschlagen', 'lookup',
    'hol', 'hole', 'lad', 'lade', 'schau', 'guck', 'lies',
    'find', 'search', 'retrieve', 'fetch', 'check', 'get', 'show', 'load', 'read',
    'encontra', 'encontre', 'procura', 'procure', 'acha', 'ache',
    'mostra', 'mostre', 'busca', 'busque', 'localiza', 'localize',
}

SEARCH_PHRASES = [
    'wo ist', 'hast du', 'gibt es', 'schau nach', 'zeig mir', 'hol mir',
    'look for', 'show me', 'where is', 'do you have', 'look up', 'get me',
    'onde está', 'onde esta', 'tem algum', 'tem alguma',
]

# Question-intent markers: user asking about past/stored information.
# Together with a subject (SEARCH_OBJECTS or a proper noun) these count as search intent
# even without an explicit SEARCH_ACTIONS verb.
QUESTION_INTENT_PHRASES = [
    # Deutsch
    'was stand', 'was steht', 'was schrieb', 'was schreibt', 'was hatte', 'was habe', 'was haben',
    'was sagte', 'was sagt', 'was meinte', 'was meint', 'was war', 'was ist mit',
    'welches datum', 'welche datei', 'welcher tag', 'welche mail', 'welche email',
    'welcher termin', 'welche nachricht', 'welche info',
    'wo haben wir', 'wo hat', 'wo hatten wir', 'wo steht', 'wo stand',
    'wann haben wir', 'wann war', 'wann hatten wir', 'wann ist', 'wann wurde',
    'wer hat', 'wer war', 'wer schrieb', 'wer schickt',
    'erinnerst du dich', 'weisst du noch', 'weißt du noch', 'erinnere mich',
    'erinnerung an', 'erinner dich an',
    'nochmal die', 'noch einmal die', 'nochmal der', 'nochmal das',
    'uns nochmal', 'mir nochmal',
    # Englisch
    'what did', 'what was', 'what does', 'what is the', 'which date', 'which file',
    'which email', 'which message', 'when did', 'when was', 'who was', 'who did',
    'who sent', 'where did', 'where is the',
    'remember when', 'do you remember', 'remind me', 'remember the',
    # Portugiesisch
    'o que disse', 'o que escreveu', 'o que tinha', 'qual data', 'qual arquivo',
    'quando foi', 'quando disse', 'quem disse', 'quem enviou',
    'lembra', 'se lembra', 'lembra-se',
    # Konversationelle Personen-Phrasen ("hat sich gemeldet", "news von", ...).
    # Diese Phrasen sind selbst kein Frage-Verb, signalisieren aber klar
    # Such-Intent wenn ein Personenname mitgeliefert wird.
    'hat sich gemeldet', 'meldet sich', 'meldete sich', 'meldung von',
    'gibt es was von', 'gibts was von', 'gibts neues von', 'gibt es neues von',
    'gibt es was neues von', 'gibts was neues von',
    'was gibt es von', 'was gibts von', 'was gibts neues von',
    'was l\u00e4uft mit', 'was laeuft mit', 'was geht mit', 'was macht',
    'wie steht es um', 'wie sieht es aus mit',
    'h\u00f6re ich von', 'hoere ich von', 'h\u00f6ren wir von', 'hoeren wir von',
    'news von', 'update zu', 'update von', 'neues von', 'antwort von',
    'antwort auf', 'feedback von', 'r\u00fcckmeldung von', 'rueckmeldung von',
    'zuletzt von', 'zuletzt mit',
    # Englisch
    'heard from', 'any news from', 'any update from', 'any update on',
    'update on', 'news from', 'reached out', 'replied', 'response from',
    'latest from', 'last from',
    # Portugiesisch
    'novidades de', 'resposta de', 'falou comigo', 'me escreveu', 'me mandou',
]

# German often splits a verb phrase around the subject — "hat sich [X] gemeldet"
# instead of "hat sich gemeldet [X]". A flat substring match misses those, so
# we add regex patterns that allow up to a few words between the parts.
# Each pattern is `re.search`ed against the lowercased message.
QUESTION_INTENT_REGEXES = [
    # "hat/hatte sich [X] gemeldet/gemailt"
    r'\bhat(?:te)?\s+sich\s+\w+(?:\s+\w+){0,4}\s+(?:gemeldet|gemailt|gemailed|gemeldet)\b',
    # "meldet/meldete sich [X]"
    r'\bmelde(?:t|te)\s+sich\s+\w+',
    # "[X] hat sich gemeldet/gemailt"
    r'\b\w+(?:\s+\w+){0,3}\s+hat\s+sich\s+(?:gemeldet|gemailt|gemailed|gemeldet)\b',
    # "gibt es (was/neues/news) (von/zu) [X]" / "gibts was neues von [X]"
    r'\bgib(?:t|ts|t es)\s+(?:was\s+)?(?:neues|news)?\s*(?:von|zu)\s+\w+',
    # "was (macht|läuft|geht) [X]"
    r'\bwas\s+(?:macht|l\u00e4uft|laeuft|geht|gibt es bei)\s+\w+',
    # "wie steht es um [X]" / "wie sieht es aus mit [X]"
    r'\bwie\s+steht\s+es\s+um\s+\w+',
    r'\bwie\s+sieht\s+es\s+aus\s+(?:mit|bei)\s+\w+',
    # English forms
    r'\bany\s+(?:news|update[s]?)\s+from\s+\w+',
    r'\bheard\s+from\s+\w+',
    r'\bdid\s+\w+(?:\s+\w+){0,3}\s+(?:reply|respond|reach\s+out|message|email|write)',
    r'\bhas\s+\w+(?:\s+\w+){0,3}\s+(?:replied|responded|reached\s+out|messaged|emailed|written|gotten\s+back)',
    # Portuguese
    r'\b(?:tem|teve)\s+(?:noticia|not\u00edcia)s?\s+(?:de|do|da)\s+\w+',
    r'\b\w+(?:\s+\w+){0,3}\s+me\s+(?:escreveu|mandou|chamou|respondeu)',
]

# Phrases the user uses to signal they want the *most recent* item(s) first.
# When detected, the search ranks results by date (newest first) and uses
# scores only as a tie-break.
RECENCY_TRIGGER_PHRASES = [
    # Deutsch
    'letzte', 'letzten', 'letzter', 'letztes',
    'neueste', 'neuesten', 'neueste', 'neuste', 'neusten',
    'aktuell', 'aktuelle', 'aktuellste', 'aktuellsten',
    'j\u00fcngste', 'juengste', 'j\u00fcngsten', 'juengsten',
    'zuletzt', 'k\u00fcrzlich', 'kuerzlich', 'gerade',
    # Englisch
    'last', 'latest', 'newest', 'most recent', 'recent', 'recently', 'just now',
    # Portugiesisch
    '\u00faltima', 'ultima', '\u00faltimo', 'ultimo', '\u00faltimas', 'ultimas',
    'recente', 'recentes', 'mais recente', 'mais recentes',
]

# Conversational phrases that should NOT trigger a search even if a proper noun is present.
# Examples: greetings, acknowledgements, meta-conversation about the assistant itself.
NO_SEARCH_OVERRIDES = [
    'hallo', 'hi ', 'hey', 'guten morgen', 'guten tag', 'guten abend',
    'hello', 'good morning', 'good evening', 'bom dia', 'boa tarde', 'boa noite',
    'danke', 'vielen dank', 'thanks', 'thank you', 'obrigado', 'obrigada',
    'verstanden', 'alles klar', 'ok', 'ok.', 'okay', 'got it', 'entendido',
    # Meta about the agent
    'wer bist du', 'was kannst du', 'hilfe', 'help', 'who are you',
]

SEARCH_OBJECTS = {
    # E-Mail
    'email', 'mail', 'e-mail', 'nachricht', 'message', 'mensagem', 'correo',
    'emails', 'mails', 'nachrichten', 'messages', 'mensagens', 'correos',
    # Dokumente allgemein
    'datei', 'dokument', 'file', 'document', 'arquivo',
    'dateien', 'dokumente', 'files', 'documents', 'arquivos',
    'doc', 'bericht', 'report', 'relatorio',
    'docs', 'berichte', 'reports', 'relatorios',
    # Word
    'word', 'brief', 'schreiben', 'letter', 'carta',
    'word-dokument', 'worddokument', 'textdokument',
    # Excel / Tabellen
    'excel', 'tabelle', 'spreadsheet', 'planilha',
    'kalkulation', 'calculation', 'calculo',
    'auswertung', 'analysis', 'analise',
    'sheet', 'sheets', 'zahlen', 'numbers',
    # PowerPoint / Praesentationen
    'praesentation', 'presentation', 'apresentacao',
    'powerpoint', 'pptx', 'slides', 'deck', 'folie',
    'folien', 'pitch', 'vortrag',
    # PDF
    'pdf', 'anhang', 'attachment', 'anexo',
    # Vertraege / Rechnungen
    'vertrag', 'contract', 'contrato', 'agreement',
    'vertraege', 'contracts', 'contratos', 'agreements',
    'rechnung', 'invoice', 'fatura', 'billing',
    'rechnungen', 'invoices', 'faturas',
    'angebot', 'offer', 'proposta', 'quote',
    'angebote', 'offers', 'propostas', 'quotes',
    'bestellung', 'order', 'pedido',
    'bestellungen', 'orders', 'pedidos',
    # Notizen / Sonstiges
    'notiz', 'note', 'nota', 'memo',
    'protokoll', 'minutes', 'ata',
    'zusammenfassung', 'summary', 'resumo',
    # Temporal / date context
    'gestern', 'heute', 'letzte', 'letzten', 'letzter',
    'yesterday', 'today', 'last', 'recent', 'latest',
    'ontem', 'hoje', 'ultimo', 'ultima', 'recente',
}

# ─── SOURCE TAXONOMY ─────────────────────────────────────────────────────────

SOURCE_TAXONOMY = {
    'email': {
        'label': 'E-Mail',
        'icon': '\u2709',
        'subcategories': ['notification'],
        'patterns': ['email_*.txt', '*.eml', '*_IN_*.txt', '*_OUT_*.txt'],
        'keywords': {'email', 'mail', 'e-mail', 'nachricht', 'message', 'mensagem', 'correo', 'brief', 'letter', 'carta', 'correio'},
    },
    'notification': {
        'label': 'Notifikation',
        'icon': '\U0001f4e2',
        'parent': 'email',
        'is_noise': True,
        'auto_detect': [
            'noreply', 'no-reply', 'donotreply', 'notifications', 'mailer', 'daemon',
            'teams.mail.microsoft', 'notify.microsoft', 'posted a message', 'mentioned you',
            'automatic reply', 'automatische antwort', 'out of office', 'delivery failed',
            'invited you', 'shared with you',
        ],
        'keywords': set(),
    },
    'webclip': {
        'label': 'Web Clip',
        'icon': '\U0001f310',
        'subcategories': ['webclip_salesforce', 'webclip_slack', 'webclip_linkedin', 'webclip_general'],
        'patterns': ['web_*.txt', 'web_*.json', 'slack_*.txt', 'slack_*.json', 'salesforce_*.txt', 'salesforce_*.json', 'linkedin_*.txt', 'linkedin_*.json'],
        'keywords': {'web clip', 'webclip', 'webseite', 'lesezeichen', 'bookmark', 'clip', 'artikel', 'article', 'webpage', 'saved page', 'pagina web', 'marcador', 'clipe'},
    },
    'webclip_salesforce': {
        'label': 'Salesforce',
        'icon': '\U0001f310',
        'parent': 'webclip',
        'patterns': ['salesforce_*.txt', 'salesforce_*.json'],
        'keywords': {'salesforce', 'lead', 'account', 'opportunity', 'crm', 'contact', 'case', 'conta', 'oportunidade'},
    },
    'webclip_slack': {
        'label': 'Slack',
        'icon': '\U0001f310',
        'parent': 'webclip',
        'patterns': ['slack_*.txt', 'slack_*.json'],
        'keywords': {'slack', 'kanal', 'channel', 'slack nachricht', 'slack message', 'canal'},
    },
    'webclip_linkedin': {
        'label': 'LinkedIn',
        'icon': '\U0001f310',
        'parent': 'webclip',
        'patterns': ['linkedin_*.txt', 'linkedin_*.json'],
        'keywords': {'linkedin', 'profil', 'profile', 'connection', 'post', 'netzwerk', 'network'},
    },
    'webclip_general': {
        'label': 'Web',
        'icon': '\U0001f310',
        'parent': 'webclip',
        'patterns': ['web_*.txt', 'web_*.json'],
        'keywords': {'webseite', 'artikel', 'seite', 'internet', 'website', 'article', 'page', 'web page', 'site', 'pagina'},
    },
    'document': {
        'label': 'Dokument',
        'icon': '\U0001f4c4',
        'subcategories': ['document_word', 'document_excel', 'document_pdf', 'document_pptx'],
        'patterns': ['*.pdf', '*.docx', '*.xlsx', '*.pptx'],
        'keywords': {'dokument', 'document', 'arquivo', 'datei', 'file', 'bericht', 'report', 'relatorio', 'rechnung', 'invoice', 'fatura', 'vertrag', 'contract', 'contrato'},
    },
    'document_word': {
        'label': 'Word',
        'icon': '\U0001f4c4',
        'parent': 'document',
        'patterns': ['*.docx', '*.doc'],
        'keywords': {'word', 'word-dokument', 'textdokument', 'word document', 'documento word'},
    },
    'document_excel': {
        'label': 'Excel',
        'icon': '\U0001f4c4',
        'parent': 'document',
        'patterns': ['*.xlsx', '*.xls', '*.csv'],
        'keywords': {'excel', 'tabelle', 'kalkulation', 'spreadsheet', 'csv', 'planilha', 'tabela', 'calculation', 'calculo', 'sheet', 'sheets', 'zahlen', 'numbers'},
    },
    'document_pdf': {
        'label': 'PDF',
        'icon': '\U0001f4c4',
        'parent': 'document',
        'patterns': ['*.pdf'],
        'keywords': {'pdf', 'acrobat', 'anhang', 'attachment', 'anexo'},
    },
    'document_pptx': {
        'label': 'PowerPoint',
        'icon': '\U0001f4c4',
        'parent': 'document',
        'patterns': ['*.pptx', '*.ppt'],
        'keywords': {'praesentation', 'presentation', 'apresentacao', 'powerpoint', 'pptx', 'slides', 'deck', 'folie', 'folien', 'pitch', 'vortrag'},
    },
    'conversation': {
        'label': 'Konversation',
        'icon': '\U0001f4ac',
        'subcategories': [],
        'patterns': ['konversation_*.txt'],
        'keywords': {'konversation', 'gespraech', 'chat', 'unterhaltung', 'verlauf', 'session', 'conversation', 'discussion', 'history', 'conversa', 'sessao', 'discussao', 'historico'},
    },
    'screenshot': {
        'label': 'Screenshot',
        'icon': '\U0001f4f8',
        'subcategories': [],
        'patterns': ['*.png', '*.jpg', '*.jpeg', 'Screenshot_*.png'],
        'keywords': {'screenshot', 'bild', 'image', 'foto', 'photo', 'picture', 'imagem', 'bildschirmfoto', 'aufnahme', 'screen', 'captura', 'print'},
    },
    'whatsapp': {
        'label': 'WhatsApp',
        'icon': '\U0001f4ac',
        'subcategories': ['whatsapp_direct', 'whatsapp_group'],
        'patterns': ['whatsapp_chat_*.txt'],
        'keywords': {'whatsapp', 'wa nachricht', 'whatsapp chat', 'wa chat', 'wa message', 'whatsapp nachricht', 'whatsapp gruppe', 'whatsapp group'},
    },
    'whatsapp_direct': {
        'label': 'WhatsApp DM',
        'icon': '\U0001f4ac',
        'parent': 'whatsapp',
        'patterns': ['whatsapp_chat_*.txt'],
        'keywords': {'whatsapp dm', 'whatsapp direkt', 'wa direkt', 'whatsapp direct'},
    },
    'whatsapp_group': {
        'label': 'WhatsApp Gruppe',
        'icon': '\U0001f4ac',
        'parent': 'whatsapp',
        'patterns': ['whatsapp_chat_group_*.txt'],
        'keywords': {'whatsapp gruppe', 'wa gruppe', 'whatsapp group', 'gruppenchat'},
    },
}

# Build TYPE_KEYWORDS from taxonomy (for backward compat with QueryParser)
TYPE_KEYWORDS = {}
for _st, _info in SOURCE_TAXONOMY.items():
    if 'parent' not in _info:  # Only top-level types
        all_kws = set(_info.get('keywords', set()))
        for _sub in _info.get('subcategories', []):
            if _sub in SOURCE_TAXONOMY:
                all_kws |= SOURCE_TAXONOMY[_sub].get('keywords', set())
        TYPE_KEYWORDS[_st] = all_kws

# Build reverse lookup: keyword -> most specific source type
SOURCE_TYPE_KEYWORDS = {}
for _st, _info in SOURCE_TAXONOMY.items():
    for _kw in _info.get('keywords', set()):
        # More specific (child) type wins over parent
        if _kw not in SOURCE_TYPE_KEYWORDS or 'parent' in _info:
            SOURCE_TYPE_KEYWORDS[_kw] = _st


def detect_source_type(filename, preview=''):
    """Detect source type from filename and preview text."""
    fname_lower = filename.lower()

    # Web Clips — subcategories first
    if fname_lower.startswith('salesforce_'):
        return 'webclip_salesforce'
    if fname_lower.startswith('slack_'):
        return 'webclip_slack'
    if fname_lower.startswith('web_'):
        return 'webclip_general'
    if fname_lower.startswith('linkedin_'):
        return 'webclip_linkedin'

    # E-Mails: email_*.txt OR DATUM_IN/OUT_absender_at_domain_*.txt (email_watcher format)
    is_email = False
    if fname_lower.startswith('email_') or fname_lower.endswith('.eml'):
        is_email = True
    elif re.match(r'\d{4}-\d{2}-\d{2}_\d{2}-\d{2}-\d{2}_(in|out)_', fname_lower):
        is_email = True
    if is_email:
        # Check if notification
        if preview:
            check = preview[:500].lower()
            notif_patterns = SOURCE_TAXONOMY.get('notification', {}).get('auto_detect', [])
            if any(p in check for p in notif_patterns):
                return 'notification'
        return 'email'

    # WhatsApp Chats
    if fname_lower.startswith('whatsapp_chat_'):
        if '_group_' in fname_lower:
            return 'whatsapp_group'
        return 'whatsapp_direct'

    # Konversationen
    if fname_lower.startswith('konversation_'):
        return 'conversation'

    # Screenshots
    if fname_lower.endswith(('.png', '.jpg', '.jpeg')):
        return 'screenshot'

    # Dokumente — subcategories
    if fname_lower.endswith(('.docx', '.doc')):
        return 'document_word'
    if fname_lower.endswith(('.xlsx', '.xls', '.csv')):
        return 'document_excel'
    if fname_lower.endswith('.pdf'):
        return 'document_pdf'
    if fname_lower.endswith(('.pptx', '.ppt')):
        return 'document_pptx'

    return 'file'


def detect_source_filter(query):
    """Detect source type filter from query. Returns (source_type, effective_types).
    effective_types includes subcategories for parent types.
    """
    query_lower = query.lower()
    words = query_lower.split()

    best_type = None
    best_specificity = 0  # higher = more specific

    for w in words:
        w_clean = w.strip('.,;:!?()')
        if w_clean in SOURCE_TYPE_KEYWORDS:
            st = SOURCE_TYPE_KEYWORDS[w_clean]
            # Subcategories are more specific (specificity=2) than parents (specificity=1)
            spec = 2 if 'parent' in SOURCE_TAXONOMY.get(st, {}) else 1
            if spec > best_specificity:
                best_type = st
                best_specificity = spec

    # Check multi-word keywords
    for st, info in SOURCE_TAXONOMY.items():
        for kw in info.get('keywords', set()):
            if ' ' in kw and kw in query_lower:
                spec = 2 if 'parent' in info else 1
                if spec > best_specificity:
                    best_type = st
                    best_specificity = spec

    if not best_type:
        return (None, None)

    # Build effective types (include subcategories for parent types)
    effective = {best_type}
    info = SOURCE_TAXONOMY.get(best_type, {})
    for sub in info.get('subcategories', []):
        effective.add(sub)

    return (best_type, effective)


def get_source_label(source_type):
    """Get human-readable label with icon for a source type."""
    info = SOURCE_TAXONOMY.get(source_type, {})
    icon = info.get('icon', '')
    label = info.get('label', source_type)
    parent = info.get('parent')
    if parent:
        parent_label = SOURCE_TAXONOMY.get(parent, {}).get('label', parent)
        return f"{icon} {parent_label} > {label}"
    return f"{icon} {label}"


def extract_conversation_meta(text):
    """Extract metadata from konversation_*.txt format."""
    lines = text.split('\n')[:20]
    agent = ''
    date = ''
    user_messages = []
    first_sentence = ''

    for line in lines:
        line_s = line.strip()
        if line_s.startswith('Agent:'):
            agent = line_s.split(':', 1)[1].strip()
        elif line_s.startswith('Datum:') or line_s.startswith('Date:'):
            date = line_s.split(':', 1)[1].strip()
        elif line_s.startswith('Du:') or line_s.startswith('User:'):
            msg = line_s.split(':', 1)[1].strip()
            user_messages.append(msg)
        elif not first_sentence and len(line_s) > 10:
            first_sentence = line_s[:100]

    subject = user_messages[0] if user_messages else first_sentence
    return {'agent': agent, 'date': date, 'subject': subject, 'user_messages': user_messages}

# Time expression patterns
TIME_WORDS = {
    'gestern': -1, 'yesterday': -1, 'ontem': -1,
    'heute': 0, 'today': 0, 'hoje': 0,
}
TIME_PHRASES = {
    'letzte woche': 7, 'letzten woche': 7, 'letzter woche': 7,
    'last week': 7, 'semana passada': 7, 'ultima semana': 7,
    'letzten monat': 30, 'letztem monat': 30, 'letzter monat': 30,
    'last month': 30, 'mes passado': 30, 'ultimo mes': 30,
    'letzten 3 tage': 3, 'last 3 days': 3,
    'letzten 7 tage': 7, 'last 7 days': 7,
    'letzten 14 tage': 14, 'last 14 days': 14,
    'letzten 30 tage': 30, 'last 30 days': 30,
}

# Field indicators
FIELD_FROM = {'von', 'from', 'de', 'vom'}
FIELD_TO = {'an', 'to', 'para', 'fuer', 'für'}
FIELD_ABOUT = {'ueber', 'über', 'about', 'sobre', 'betreff', 'subject', 'thema', 'bezueglich'}


# ═══════════════════════════════════════════════════════════════════════════════
# TEIL 1 — SEARCH INDEX
# ═══════════════════════════════════════════════════════════════════════════════

class SearchIndex:
    """Builds and maintains a search index for an agent's memory directory."""

    def __init__(self, speicher_path):
        self.speicher = speicher_path
        self.memory_dir = os.path.join(speicher_path, 'memory')
        self.index_file = os.path.join(speicher_path, '.search_index.json')
        self.entries = {}  # filename -> entry dict
        self._lock = threading.Lock()

    def _detect_file_type(self, fname, preview=''):
        """Detect file type from filename using source taxonomy."""
        return detect_source_type(fname, preview)

    def _extract_email_headers(self, text):
        """Extract From/To/Subject from email text (first 10 lines)."""
        headers = {'from': '', 'to': '', 'subject': '', 'date': ''}
        for line in text.split('\n')[:15]:
            line_stripped = line.strip()
            ll = line_stripped.lower()
            if ll.startswith(('von:', 'from:')):
                headers['from'] = line_stripped.split(':', 1)[1].strip()
            elif ll.startswith(('an:', 'to:')):
                headers['to'] = line_stripped.split(':', 1)[1].strip()
            elif ll.startswith(('betreff:', 'subject:')):
                headers['subject'] = line_stripped.split(':', 1)[1].strip()
            elif ll.startswith(('datum:', 'date:')):
                headers['date'] = line_stripped.split(':', 1)[1].strip()
        return headers

    def _extract_keywords(self, text, max_keywords=10):
        """Extract top-N most frequent non-stopword keywords from text."""
        words = re.findall(r'[a-zA-ZäöüÄÖÜßàáâãéèêíìîóòôõúùûçñ]+', text.lower())
        freq = {}
        for w in words:
            if len(w) < 3 or w in STOPWORDS:
                continue
            w_norm = normalize_unicode(w)
            if len(w_norm) < 3:
                continue
            freq[w_norm] = freq.get(w_norm, 0) + 1
        sorted_words = sorted(freq.items(), key=lambda x: x[1], reverse=True)
        return [w for w, _ in sorted_words[:max_keywords]]

    def _read_file_text(self, fpath, max_chars=30000):
        """Read file content as text. Returns empty string for binaries/errors."""
        fname = os.path.basename(fpath)
        ext = os.path.splitext(fname.lower())[1]
        if ext in BINARY_EXTS:
            return ''
        try:
            with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                return f.read(max_chars)
        except Exception:
            return ''

    def _index_file(self, fname):
        """Create an index entry for a single file."""
        fpath = self._resolve_file_path(fname)
        try:
            stat = os.stat(fpath)
        except Exception:
            return None

        mtime = stat.st_mtime
        file_date = datetime.datetime.fromtimestamp(mtime).strftime('%Y-%m-%d')

        entry = {
            'path': fpath,
            'filename': fname,
            'type': 'file',  # will be updated below
            'source_type': 'file',
            'date': file_date,
            'mtime': mtime,
            'size': stat.st_size,
            'preview': '',
            'from': '',
            'to': '',
            'subject': '',
            'keywords': [],
        }

        # Read content for text-based files
        # Use basename for source type detection (files may be in subdirs like whatsapp/)
        fname_base = os.path.basename(fname)
        text = ''
        stype_initial = detect_source_type(fname_base)
        if stype_initial != 'screenshot':
            text = self._read_file_text(fpath)
            if text:
                entry['preview'] = text[:300].replace('\n', ' ').strip()
                entry['keywords'] = self._extract_keywords(text)

        # Detect source type with preview for notification detection
        stype = detect_source_type(fname_base, text[:500] if text else '')
        entry['type'] = stype
        entry['source_type'] = stype

        # Parse email headers
        if stype in ('email', 'notification') and text:
            headers = self._extract_email_headers(text)
            entry['from'] = headers['from']
            entry['to'] = headers['to']
            entry['subject'] = headers['subject']

        # Parse conversation metadata
        if stype == 'conversation' and text:
            conv_meta = extract_conversation_meta(text)
            entry['subject'] = conv_meta.get('subject', '')

        return entry

    @staticmethod
    def _is_duplicate_file(fname):
        """Check if a file is a macOS duplicate (e.g. 'file_2.txt', 'file_2 2.txt')."""
        base = os.path.basename(fname)
        # Pattern: name_2.ext, name_2 2.ext, name_2 3.ext, name 2.ext
        if re.search(r'_2( \d+)?\.[^.]+$', base):
            return True
        if re.search(r' 2( \d+)?\.[^.]+$', base):
            return True
        return False

    def _get_all_indexable_files(self):
        """Get all files to index: memory/ dir (incl. subdirs) + konversation_*.txt from agent root.
        Skips macOS duplicate files (_2.txt, _2 2.txt etc.)."""
        files = []
        if os.path.exists(self.memory_dir):
            for entry in os.listdir(self.memory_dir):
                entry_path = os.path.join(self.memory_dir, entry)
                if os.path.isfile(entry_path):
                    if not self._is_duplicate_file(entry):
                        files.append(entry)
                elif os.path.isdir(entry_path):
                    # Scan subdirectories (e.g. whatsapp/)
                    for sub_entry in os.listdir(entry_path):
                        if os.path.isfile(os.path.join(entry_path, sub_entry)):
                            if not self._is_duplicate_file(sub_entry):
                                files.append(os.path.join(entry, sub_entry))
        # Also include konversation_*.txt from agent root (not in memory/)
        if os.path.exists(self.speicher):
            for f in os.listdir(self.speicher):
                if f.startswith('konversation_') and f.endswith('.txt'):
                    if f not in files:
                        files.append(f)
        return files

    def _resolve_file_path(self, fname):
        """Resolve file path: check memory/ first (incl. subdirs), then agent root."""
        mem_path = os.path.join(self.memory_dir, fname)
        if os.path.exists(mem_path):
            return mem_path
        root_path = os.path.join(self.speicher, fname)
        if os.path.exists(root_path):
            return root_path
        # For subdir files like "whatsapp/chat.txt", basename might be used as key
        base = os.path.basename(fname)
        if base != fname:
            mem_base = os.path.join(self.memory_dir, base)
            if os.path.exists(mem_base):
                return mem_base
        return mem_path  # default fallback

    def build_index(self):
        """Build complete index from scratch."""
        if not os.path.exists(self.memory_dir) and not os.path.exists(self.speicher):
            return 0

        os.makedirs(self.memory_dir, exist_ok=True)
        start = time.time()
        files = self._get_all_indexable_files()
        entries = {}

        for fname in files:
            entry = self._index_file(fname)
            if entry:
                entries[fname] = entry

        with self._lock:
            self.entries = entries

        self._save_index()
        elapsed = time.time() - start
        agent_name = os.path.basename(self.speicher)
        count = len(entries)
        print(f"Index fuer {agent_name}: {count} Dateien in {elapsed:.1f}s")
        return count

    def update_index(self):
        """Incremental update: only re-index new/modified files, remove deleted."""
        if not os.path.exists(self.memory_dir) and not os.path.exists(self.speicher):
            return 0

        os.makedirs(self.memory_dir, exist_ok=True)
        self._load_index()
        current_files = set(self._get_all_indexable_files())
        indexed_files = set(self.entries.keys())

        # Remove deleted files
        for fname in indexed_files - current_files:
            del self.entries[fname]

        # Add/update new and modified files
        updated = 0
        for fname in current_files:
            fpath = self._resolve_file_path(fname)
            try:
                mtime = os.path.getmtime(fpath)
            except Exception:
                continue

            existing = self.entries.get(fname)
            if existing and existing.get('mtime', 0) >= mtime:
                continue  # Not modified

            entry = self._index_file(fname)
            if entry:
                self.entries[fname] = entry
                updated += 1

        if updated > 0:
            self._save_index()

        return updated

    def build_or_update(self):
        """Build index if none exists, otherwise update incrementally."""
        if os.path.exists(self.index_file):
            return self.update_index()
        else:
            return self.build_index()

    def add_file(self, fname):
        """Index a single newly added file. Also fires an async embedding
        update so that semantic retrieval is available immediately. Embedding
        failures are silent — they never block keyword indexing."""
        self._load_index()
        entry = self._index_file(fname)
        if entry:
            with self._lock:
                self.entries[fname] = entry
            self._save_index()
            # Fire-and-forget semantic indexing (network-bound, optional).
            try:
                t = threading.Thread(
                    target=index_file_with_embedding,
                    args=(self.speicher, fname),
                    daemon=True,
                )
                t.start()
            except Exception:
                pass

    def _save_index(self):
        """Save index to disk."""
        try:
            with open(self.index_file, 'w', encoding='utf-8') as f:
                json.dump(self.entries, f, ensure_ascii=False)
        except Exception as e:
            print(f"Index save error: {e}")

    def _load_index(self):
        """Load index from disk if not already loaded."""
        if self.entries:
            return
        if os.path.exists(self.index_file):
            try:
                with open(self.index_file, 'r', encoding='utf-8') as f:
                    self.entries = json.load(f)
            except Exception:
                self.entries = {}

    def get_entries(self):
        """Get all index entries."""
        self._load_index()
        return self.entries

    @property
    def file_count(self):
        self._load_index()
        return len(self.entries)


# ─── Recency helpers (used by HybridSearch and hybrid_rag_search) ────────────

# Standard email-watcher names start with YYYY-MM-DD_HH-MM-SS_…
# Older "email_…" exports start with email_YYYY-MM-DD_HH-MM-SS_…
_DATE_RX = re.compile(r'(\d{4}-\d{2}-\d{2})(?:[_T-](\d{2}-\d{2}-\d{2}|\d{2}:\d{2}:\d{2}))?')


def extract_date_from_name(fname):
    """Pull a date/time string out of a memory filename.

    Returns an ISO-ish string like '2026-04-14T12-50-42' (sortable, newest =
    largest) when one is found, or '' when nothing date-like is present.
    Both the new email_watcher v2 schema (YYYY-MM-DD_HH-MM-SS_*) and legacy
    forms (email_YYYY-MM-DD…, YYYY-MM-DD…) are recognised.
    """
    if not fname:
        return ''
    base = os.path.basename(fname)
    m = _DATE_RX.search(base)
    if not m:
        return ''
    date_part = m.group(1)
    time_part = m.group(2) or '00-00-00'
    # Normalise ':' -> '-' so the string sorts correctly.
    time_part = time_part.replace(':', '-')
    return f"{date_part}T{time_part}"


def _recency_key(name, entry=None):
    """Sort key for newest-first ordering. Falls back to mtime if no date in
    the filename, so files without a stamped date still sort sensibly."""
    iso = extract_date_from_name(name)
    if iso:
        return iso
    mt = (entry or {}).get('mtime', 0) if entry else 0
    if mt:
        # Convert to an ISO-shaped string so the comparison stays string-based
        # and we never mix int/str when sorting alongside ISO keys.
        try:
            return datetime.datetime.fromtimestamp(mt).strftime('%Y-%m-%dT%H-%M-%S')
        except Exception:
            return ''
    return ''


# ═══════════════════════════════════════════════════════════════════════════════
# TEIL 2 — QUERY PARSER
# ═══════════════════════════════════════════════════════════════════════════════

class QueryIntent:
    """Structured search intent parsed from natural language."""

    def __init__(self):
        self.keywords = []          # Cleaned search keywords
        self.date_filter = None     # (start_date, end_date) or None
        self.date_label = ''        # Human-readable: "gestern", "letzte Woche"
        self.file_type = None       # "email", "document", "image", etc. or None
        self.source_type = None     # Specific source type from taxonomy
        self.source_types_effective = None  # Set of types to match (includes subcategories)
        self.field_filter = None    # "from", "to", "about", or None
        self.field_label = ''       # "Von:", "An:", "Betreff:"
        self.person_names = []      # Detected proper nouns
        self.raw_query = ''         # Original query text
        self.is_search = False      # True if search intent detected
        self.max_results = None     # Optional result-count hint ("die letzten 10" -> 10)
        self.wants_global = False   # True when user explicitly asks for global / "everywhere"
        self.recency_first = False  # True -> sort results by date desc (score is tie-break)
        self.wants_deep = False     # True when user asks for detailed/thorough / /deep prefix

    def summary(self):
        """Human-readable summary of the search intent."""
        parts = []
        if self.keywords:
            parts.append(' '.join(self.keywords))
        if self.source_type:
            label = get_source_label(self.source_type)
            parts.append(f"Typ: {label}")
        elif self.file_type:
            type_labels = {'email': 'Email', 'document': 'Dokument', 'image': 'Bild',
                          'conversation': 'Konversation', 'salesforce': 'Salesforce', 'web': 'Web'}
            parts.append(f"Typ: {type_labels.get(self.file_type, self.file_type)}")
        if self.date_label:
            parts.append(f"Zeitraum: {self.date_label}")
        if self.field_label:
            parts.append(f"Feld: {self.field_label}")
        return ' | '.join(parts) if parts else self.raw_query


class QueryParser:
    """Parses natural language queries into structured QueryIntent."""

    @staticmethod
    def parse(msg, force_search=False):
        """Parse a message into a QueryIntent.
        force_search: if True, always treat as search (for /find, search dialog, etc.)
        """
        intent = QueryIntent()
        intent.raw_query = msg
        msg_lower = msg.lower()
        words = msg_lower.split()
        orig_words = msg.split()

        # ── 1. Detect search intent ─────────────────────────────────────
        has_action = any(w.rstrip('.,;:!?') in SEARCH_ACTIONS for w in words)
        if not has_action:
            has_action = any(phrase in msg_lower for phrase in SEARCH_PHRASES)

        has_object = any(w.rstrip('.,;:!?') in SEARCH_OBJECTS for w in words)
        has_question_intent = any(phrase in msg_lower for phrase in QUESTION_INTENT_PHRASES)
        if not has_question_intent:
            has_question_intent = any(
                re.search(pat, msg_lower) for pat in QUESTION_INTENT_REGEXES
            )

        # Legacy trigger
        if 'memory folder' in msg_lower or 'memory ordner' in msg_lower:
            has_action = True
            has_object = True

        # Detect proper nouns (capitalized words, not at sentence start)
        skip_set = SEARCH_ACTIONS | SEARCH_OBJECTS | STOPWORDS | FIELD_FROM | FIELD_TO | FIELD_ABOUT
        for i, w in enumerate(orig_words):
            w_clean = w.rstrip('.,;:!?')
            if w_clean.endswith("'s"):
                w_clean = w_clean[:-2]
            if not w_clean or len(w_clean) < 2 or i == 0:
                continue
            if w_clean[0].isupper() and w_clean.lower() not in skip_set:
                intent.person_names.append(w_clean)

        has_proper_noun = len(intent.person_names) > 0
        has_date = False
        # Simple date heuristic (full parser below refines this)
        if re.search(r'\b(\d{1,2}\.\d{1,2}\.\d{0,4}|\d{4}-\d{2}-\d{2})\b', msg_lower):
            has_date = True

        # Global / deep markers — evaluated EARLY so explicit user intent survives
        # even when the lexical trigger is weak ("Search extended memory for thomas").
        if any(t in msg_lower for t in GLOBAL_TRIGGERS):
            intent.wants_global = True
        DEEP_MARKERS_EARLY = (
            '/deep ', '/deep_', 'ausfuehrlich', 'ausführlich', 'detailliert',
            'in depth', 'in-depth', 'thorough', 'detailed', 'detalhado',
        )
        if any(m in msg_lower for m in DEEP_MARKERS_EARLY):
            intent.wants_deep = True

        # Conversational override: greetings / acknowledgements / meta should not trigger
        stripped = msg_lower.strip()
        is_conversational = False
        if any(stripped == p.strip() or stripped.startswith(p) for p in NO_SEARCH_OVERRIDES):
            is_conversational = True
        # Very short messages with NO substantive content (no name, no object, no
        # question-phrase) are treated as conversational. This keeps short topic
        # queries like "ExFlow Rechnung" or "Pitch Folien" alive.
        if (len(words) <= 2 and not has_proper_noun
                and not has_object and not has_question_intent):
            is_conversational = True

        if force_search:
            # Forced search mode (from /find, search dialog, etc.)
            # Treat all non-stopwords as keywords, all words as potential person names
            intent.is_search = True
            has_action = True
            has_object = True
            # In force mode: add all substantial words as person names too
            for w in orig_words:
                w_clean = w.rstrip('.,;:!?')
                if w_clean.endswith("'s"):
                    w_clean = w_clean[:-2]
                if w_clean and len(w_clean) >= 3 and w_clean.lower() not in STOPWORDS:
                    if w_clean not in intent.person_names:
                        intent.person_names.append(w_clean)
        else:
            if is_conversational:
                intent.is_search = False
                return intent

            # New, more permissive trigger logic:
            # A search is intended if there is any "query trigger" (action verb, question
            # phrase, search phrase) AND some concrete "subject" (search object, proper
            # noun, or explicit date). Also trigger if BOTH an object and a proper noun
            # are present (covers imperative-free sentences like "ExFlow Rechnung").
            has_trigger = has_action or has_question_intent
            has_subject = has_object or has_proper_noun or has_date

            trigger_with_subject = has_trigger and has_subject
            strong_subject_only = has_object and has_proper_noun
            # Short "topic-only" query: "ExFlow Rechnung", "Pitch Folien", "Thomas Mail"
            # → treat as search when there is a subject and no conversational flag.
            short_topic_query = (
                len(words) <= 4
                and has_subject
                and not is_conversational
            )
            # Explicit global/deep intent is enough on its own — the user clearly
            # wants a search even if the subject isn't a known SEARCH_OBJECT.
            explicit_search_mode = intent.wants_global or intent.wants_deep

            if not (trigger_with_subject or strong_subject_only
                    or short_topic_query or explicit_search_mode):
                intent.is_search = False
                return intent

        intent.is_search = True

        # ── Recency-first detection ─────────────────────────────────────
        # 1) Explicit recency words ("letzte mail", "neueste", "latest", ...)
        for trig in RECENCY_TRIGGER_PHRASES:
            # Whole-word/phrase containment so "letzte" doesn't fire inside "verletzte".
            if re.search(r'(?<![\w])' + re.escape(trig) + r'(?![\w])', msg_lower):
                intent.recency_first = True
                break

        # 2) Person-only short query: "Fabian Adam", "Thomas Smith" — assume
        #    the user wants the freshest item from that person.
        if not intent.recency_first and not force_search:
            substantial = [w for w in words if w.rstrip('.,;:!?') not in STOPWORDS]
            if (len(intent.person_names) >= 1
                    and len(substantial) <= 3
                    and not has_object
                    and not has_question_intent
                    and not has_action):
                intent.recency_first = True

        # 3) Conversational person-phrases ("hat sich gemeldet", "news von ...")
        #    imply recency too — the user is asking about the latest contact.
        if not intent.recency_first:
            CONVERSATIONAL_PERSON_HINTS = (
                'hat sich gemeldet', 'meldet sich', 'meldete sich',
                'gibt es was von', 'gibts was von', 'gibts neues von',
                'gibt es neues von', 'gibt es was neues von', 'gibts was neues von',
                'was gibt es von', 'was gibts von', 'was gibts neues von',
                'h\u00f6re ich von', 'hoere ich von', 'h\u00f6ren wir von', 'hoeren wir von',
                'news von', 'update zu', 'update von', 'neues von',
                'antwort von', 'antwort auf', 'feedback von',
                'r\u00fcckmeldung von', 'rueckmeldung von',
                'zuletzt von', 'zuletzt mit',
                'heard from', 'any news from', 'any update from', 'any update on',
                'update on', 'news from', 'reached out', 'replied',
                'response from', 'latest from', 'last from',
                'novidades de', 'resposta de', 'me escreveu', 'me mandou',
            )
            if any(p in msg_lower for p in CONVERSATIONAL_PERSON_HINTS):
                intent.recency_first = True

        # 4) Any QUESTION_INTENT_REGEXES match implies a conversational
        #    "what's the latest" question -> recency_first.
        if not intent.recency_first:
            if any(re.search(pat, msg_lower) for pat in QUESTION_INTENT_REGEXES):
                intent.recency_first = True

        # ── 2. Time filter ──────────────────────────────────────────────
        today = datetime.date.today()

        # Multi-word time phrases first
        for phrase, days in TIME_PHRASES.items():
            if phrase in msg_lower:
                intent.date_filter = (today - datetime.timedelta(days=days), today)
                intent.date_label = phrase
                msg_lower = msg_lower.replace(phrase, ' ')
                break

        # Single-word time expressions
        if not intent.date_filter:
            for word, offset in TIME_WORDS.items():
                if word in words:
                    d = today + datetime.timedelta(days=offset) if offset != 0 else today
                    if offset == 0:
                        intent.date_filter = (today, today)
                    else:
                        intent.date_filter = (d, d)
                    intent.date_label = word
                    break

        # Explicit date patterns: DD.MM. or DD.MM.YYYY
        if not intent.date_filter:
            m = re.search(r'(\d{1,2})\.(\d{1,2})\.(\d{2,4})?', msg_lower)
            if m:
                try:
                    day, month = int(m.group(1)), int(m.group(2))
                    year = int(m.group(3)) if m.group(3) else today.year
                    if year < 100:
                        year += 2000
                    d = datetime.date(year, month, day)
                    intent.date_filter = (d, d)
                    intent.date_label = d.strftime('%d.%m.%Y')
                except ValueError:
                    pass

        # ── 3. File type detection (source taxonomy) ──────────────────
        src_type, src_effective = detect_source_filter(msg)
        if src_type:
            intent.source_type = src_type
            intent.source_types_effective = src_effective
            # Map to legacy file_type for backward compat
            parent = SOURCE_TAXONOMY.get(src_type, {}).get('parent', src_type)
            intent.file_type = parent
        else:
            for ftype, kws in TYPE_KEYWORDS.items():
                if any(w.rstrip('.,;:!?') in kws for w in words):
                    intent.file_type = ftype
                    break

        # ── 4. Field filter ─────────────────────────────────────────────
        for w in words:
            w_clean = w.rstrip('.,;:!?')
            if w_clean in FIELD_FROM:
                intent.field_filter = 'from'
                intent.field_label = 'Von:'
                break
            elif w_clean in FIELD_TO:
                intent.field_filter = 'to'
                intent.field_label = 'An:'
                break
            elif w_clean in FIELD_ABOUT:
                intent.field_filter = 'about'
                intent.field_label = 'Betreff/Body:'
                break

        # ── 5. Extract keywords ─────────────────────────────────────────
        all_noise = (SEARCH_ACTIONS | SEARCH_OBJECTS | STOPWORDS |
                     FIELD_FROM | FIELD_TO | FIELD_ABOUT |
                     set(TIME_WORDS.keys()) |
                     {'memory', 'folder', 'ordner', 'suche', 'suchen'})
        query_text = msg_lower
        for phrase in SEARCH_PHRASES + list(TIME_PHRASES.keys()):
            query_text = query_text.replace(phrase, ' ')

        tokens = []
        for w in query_text.split():
            w_clean = w.strip('.,;:!?()[]{}"\'/').lower()
            if w_clean in all_noise or len(w_clean) < 3:
                continue
            tokens.append(w_clean)
            # Add hyphen-less variant
            no_h = w_clean.replace('-', '')
            if no_h != w_clean and len(no_h) >= 3:
                tokens.append(no_h)

        # Add proper nouns
        for pn in intent.person_names:
            pn_lower = pn.lower()
            if pn_lower not in tokens:
                tokens.append(pn_lower)

        # Fallback: use object words
        if not tokens:
            tokens = [w.rstrip('.,;:!?') for w in words if w.rstrip('.,;:!?') in SEARCH_OBJECTS]

        # Deduplicate preserving order
        seen = set()
        intent.keywords = []
        for t in tokens:
            if t not in seen:
                seen.add(t)
                intent.keywords.append(t)

        # ── 6. Result-count hint ("die letzten 10", "top 5", "the last 20") ──
        num_word_map = {
            'einen': 1, 'ein': 1, 'eine': 1, 'one': 1, 'zwei': 2, 'two': 2,
            'drei': 3, 'three': 3, 'vier': 4, 'four': 4, 'fuenf': 5, 'five': 5,
            'fünf': 5, 'sechs': 6, 'six': 6, 'sieben': 7, 'seven': 7,
            'acht': 8, 'eight': 8, 'neun': 9, 'nine': 9, 'zehn': 10, 'ten': 10,
        }
        count_triggers = {
            'letzte', 'letzten', 'letzter', 'letztes',
            'top', 'neueste', 'neuesten',
            'erste', 'ersten', 'erster', 'erstes',
            'last', 'latest', 'recent', 'first',
            'ultimo', 'ultima', 'ultimos', 'primeiros',
        }
        # Only pick up a count when it's adjacent (±1 token) to a count_trigger word.
        # Avoids false matches on dates like "20.03.2024".
        toks = [t.rstrip('.,;:!?') for t in msg_lower.split()]
        for i, tok in enumerate(toks):
            if tok in count_triggers:
                # Look at next and previous tokens
                candidates = []
                if i + 1 < len(toks):
                    candidates.append(toks[i + 1])
                if i - 1 >= 0:
                    candidates.append(toks[i - 1])
                for cand in candidates:
                    if cand.isdigit():
                        intent.max_results = max(1, min(200, int(cand)))
                        break
                    if cand in num_word_map:
                        intent.max_results = num_word_map[cand]
                        break
                if intent.max_results:
                    break

        # (wants_global / wants_deep were already evaluated earlier so they
        #  survive even when the lexical trigger would otherwise bail out.)
        return intent


# ═══════════════════════════════════════════════════════════════════════════════
# TEIL 3 — HYBRID SEARCH
# ═══════════════════════════════════════════════════════════════════════════════

class HybridSearch:
    """Multi-stage search combining index filtering, BM25 scoring, and full-text."""

    @staticmethod
    def search(intent, speicher_path, max_results=3, forced_type=None,
               recency_first=None):
        """Execute a structured search. Returns (results, feedback_info).
        results: list of {name, content, score}
        feedback_info: dict with query, date_label, field_label, found_count, index_count
        forced_type: if set, overrides intent's type filter (e.g. 'email', 'webclip', 'document', 'conversation', 'screenshot')
        recency_first: if True, sort by date desc and use score as tie-break.
            Defaults to intent.recency_first.
        """
        if recency_first is None:
            recency_first = bool(getattr(intent, 'recency_first', False))
        # Apply forced type filter
        if forced_type:
            st, eff = detect_source_filter(forced_type)
            if st:
                intent.source_type = st
                intent.source_types_effective = eff
            elif forced_type in SOURCE_TAXONOMY:
                intent.source_type = forced_type
                subs = set(SOURCE_TAXONOMY[forced_type].get('subcategories', []))
                intent.source_types_effective = {forced_type} | subs

        index = SearchIndex(speicher_path)
        index._load_index()

        if not index.entries:
            # Try building index first
            index.build_or_update()

        all_entries = index.entries
        feedback = {
            'query': ' '.join(intent.keywords) if intent.keywords else intent.raw_query,
            'date_label': intent.date_label,
            'field_label': intent.field_label,
            'file_type': intent.file_type,
            'found_count': 0,
            'index_count': len(all_entries),
        }

        if not all_entries:
            return ([], feedback)

        keywords_norm = [normalize_unicode(kw) for kw in intent.keywords]
        names_norm = [normalize_unicode(pn) for pn in intent.person_names]

        # ── SCHRITT 1: Filter by date and type (index only, no file reads) ──
        candidates = {}
        for fname, entry in all_entries.items():
            # Date filter
            if intent.date_filter:
                start_date, end_date = intent.date_filter
                start_ts = datetime.datetime.combine(start_date, datetime.time.min).timestamp()
                end_ts = datetime.datetime.combine(end_date, datetime.time.max).timestamp()
                mtime = entry.get('mtime', 0)
                if mtime < start_ts or mtime > end_ts:
                    continue

            # Type filter — use source_types_effective (includes subcategories)
            if intent.source_types_effective:
                entry_type = entry.get('source_type') or entry.get('type', 'file')
                if entry_type not in intent.source_types_effective:
                    continue
            elif intent.file_type:
                entry_type = entry.get('source_type') or entry.get('type', 'file')
                # Also match parent type
                entry_parent = SOURCE_TAXONOMY.get(entry_type, {}).get('parent', entry_type)
                if entry_type != intent.file_type and entry_parent != intent.file_type:
                    continue

            candidates[fname] = entry

        if not candidates:
            return ([], feedback)

        # ── SCHRITT 2: BM25 keyword scoring on index data ──────────────
        scored = {}

        for fname, entry in candidates.items():
            score = 0
            fname_norm = normalize_unicode(fname)
            fname_flat = fname_norm.replace('-', '').replace('_', '')

            # Date filter bonus
            if intent.date_filter:
                score += 10

            # Notification penalty (noise, lower priority)
            entry_stype = entry.get('source_type') or entry.get('type', 'file')
            if entry_stype == 'notification':
                score -= 15

            # Filename matching (exact + fuzzy)
            for kw in keywords_norm:
                if kw in fname_norm or kw in fname_flat:
                    score += 4
                else:
                    fm, fs = fuzzy_match(kw, fname_norm.replace('_', ' ').replace('-', ' '))
                    if fm:
                        score += min(fs, 3)

            # Email header matching
            entry_stype = entry.get('source_type') or entry.get('type', 'file')
            if entry_stype in ('email', 'notification'):
                from_norm = normalize_unicode(entry.get('from', ''))
                to_norm = normalize_unicode(entry.get('to', ''))
                subject_norm = normalize_unicode(entry.get('subject', ''))

                for kw in keywords_norm:
                    if intent.field_filter == 'from' and kw in from_norm:
                        score += 5
                    elif intent.field_filter == 'to' and kw in to_norm:
                        score += 5
                    elif intent.field_filter == 'about' and kw in subject_norm:
                        score += 3
                    elif intent.field_filter is None:
                        if kw in from_norm:
                            score += 5
                        if kw in subject_norm:
                            score += 3
                        if kw in to_norm:
                            score += 2

                # Person name matching in from/to/subject (exact + fuzzy)
                for pn in names_norm:
                    if pn in from_norm:
                        score += 8
                    elif pn in to_norm:
                        score += 6
                    elif pn in subject_norm:
                        score += 4
                    else:
                        # Fuzzy match against from, to, subject fields
                        fm_from, fs_from = fuzzy_match(pn, from_norm)
                        fm_to, fs_to = fuzzy_match(pn, to_norm)
                        fm_subj, fs_subj = fuzzy_match(pn, subject_norm)
                        if fm_from:
                            score += min(fs_from + 2, 6)
                        elif fm_to:
                            score += min(fs_to + 1, 4)
                        elif fm_subj:
                            score += min(fs_subj, 3)

            # Non-email: person names in filename (exact + fuzzy)
            for pn in names_norm:
                if pn in fname_norm or pn in fname_flat:
                    score += 5
                else:
                    fm, fs = fuzzy_match(pn, fname_norm.replace('_', ' ').replace('-', ' '))
                    if fm:
                        score += min(fs, 4)

            # Preview/keyword matching from index (exact + fuzzy)
            preview_norm = normalize_unicode(entry.get('preview', ''))
            entry_keywords = entry.get('keywords', [])
            for kw in keywords_norm:
                if kw in preview_norm:
                    score += 1
                elif len(kw) >= 4:
                    fm, fs = fuzzy_match(kw, preview_norm)
                    if fm:
                        score += 1
                if kw in entry_keywords:
                    score += 2

            if score > 0:
                scored[fname] = score

        # ── SCHRITT 3: Full-text search on top-20 candidates ───────────
        # Sort by score, take top 20 for full-text
        top_candidates = sorted(scored.items(), key=lambda x: x[1], reverse=True)[:20]

        # Also include unscored candidates if we have fewer than 3 hits
        if len(top_candidates) < max_results:
            unscored = [(fname, 0) for fname in candidates if fname not in scored]
            # Sort by recency
            unscored.sort(key=lambda x: candidates[x[0]].get('mtime', 0), reverse=True)
            top_candidates.extend(unscored[:20])

        memory_dir = os.path.join(speicher_path, 'memory')
        # Content cache: fname -> text (read once, used for both body scoring
        # and final result assembly).
        content_cache = {}
        for fname, base_score in top_candidates:
            entry = candidates.get(fname)
            if not entry or entry.get('type') == 'image':
                continue

            # Resolve path: check memory/ first, then agent root (for konversation_*.txt)
            fpath = entry.get('path') or os.path.join(memory_dir, fname)
            if not os.path.exists(fpath):
                fpath = os.path.join(speicher_path, fname)
            try:
                with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                    text = f.read(20000)
            except Exception:
                continue

            content_cache[fname] = text
            text_norm = normalize_unicode(text)
            body_score = 0

            for kw in keywords_norm:
                count = text_norm.count(kw)
                if count >= 3:
                    body_score += 3
                elif count >= 1:
                    body_score += 1

            for pn in names_norm:
                if pn in text_norm:
                    body_score += 5
                elif len(pn) >= 4:
                    # Fuzzy match person names in full text (sample first 2000 chars)
                    fm, fs = fuzzy_match(pn, text_norm[:2000])
                    if fm:
                        body_score += min(fs, 3)

            scored[fname] = scored.get(fname, 0) + body_score

        # ── SCHRITT 4: Ranking and return ──────────────────────────────
        # Default: sort by score desc, with date as tie-break (newest first).
        # When recency_first is set, we do the same Score sort first — but then
        # take a Top-N pool (still requires score > 0) and re-sort *that pool*
        # by date. When a person is mentioned in the query, the pool is further
        # restricted to hits that actually mention that person, so "Hat sich
        # Fabian gemeldet?" returns Fabian's mails sorted newest-first instead
        # of arbitrary newest mails from anyone.
        score_then_date = sorted(
            scored.items(),
            key=lambda kv: (kv[1], _recency_key(kv[0], candidates.get(kv[0], {}))),
            reverse=True,
        )
        if recency_first:
            relevant = [(n, s) for n, s in score_then_date if s > 0]
            # Person-aware narrowing: when the query names someone, only keep
            # hits where that name surfaces in the filename, sender, subject
            # or preview. Falls back to the unfiltered pool if nothing matches
            # so we never lose the user's hits entirely.
            if names_norm:
                filtered = []
                for n, s in relevant:
                    entry = candidates.get(n, {})
                    haystack = ' '.join([
                        normalize_unicode(n),
                        normalize_unicode(entry.get('from', '')),
                        normalize_unicode(entry.get('to', '')),
                        normalize_unicode(entry.get('subject', '')),
                        normalize_unicode(entry.get('preview', '')),
                    ]).lower()
                    if any(pn in haystack for pn in names_norm):
                        filtered.append((n, s))
                if filtered:
                    relevant = filtered
            pool_size = max(max_results + 3, 10)
            pool = relevant[:pool_size]
            final_sorted = sorted(
                pool,
                key=lambda kv: (_recency_key(kv[0], candidates.get(kv[0], {})), kv[1]),
                reverse=True,
            )
        else:
            final_sorted = score_then_date
        results = []
        for fname, score in final_sorted[:max_results]:
            if score <= 0:
                break
            if fname in content_cache:
                # Reuse cached content from full-text step.
                results.append({'name': fname, 'content': content_cache[fname], 'score': score})
                continue
            entry = candidates.get(fname, {})
            fpath = entry.get('path') or os.path.join(memory_dir, fname)
            if not os.path.exists(fpath):
                fpath = os.path.join(speicher_path, fname)
            try:
                with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                    content = f.read(20000)
                results.append({'name': fname, 'content': content, 'score': score})
            except Exception:
                continue

        feedback['found_count'] = len(results)
        return (results, feedback)


# ═══════════════════════════════════════════════════════════════════════════════
# CONVENIENCE FUNCTIONS (for web_server.py integration)
# ═══════════════════════════════════════════════════════════════════════════════

# Cache of indexes per agent
_index_cache = {}
_index_cache_lock = threading.Lock()


def get_or_build_index(speicher_path):
    """Get cached index or build a new one. Thread-safe."""
    with _index_cache_lock:
        if speicher_path in _index_cache:
            idx = _index_cache[speicher_path]
            # Update in background if stale
            return idx
        idx = SearchIndex(speicher_path)
        _index_cache[speicher_path] = idx

    # Build/update outside lock
    idx.build_or_update()
    return idx


def build_index_async(speicher_path):
    """Build/update index in a background thread."""
    def _build():
        get_or_build_index(speicher_path)
    t = threading.Thread(target=_build, daemon=True)
    t.start()


# ─── NLP KEYWORD EXTRACTION ──────────────────────────────────────────────────

_STOPWORDS_DE = {
    'ich', 'du', 'er', 'sie', 'es', 'wir', 'ihr', 'und', 'oder', 'aber', 'wenn',
    'als', 'dass', 'den', 'dem', 'des', 'die', 'der', 'das', 'ein', 'eine', 'einer',
    'einem', 'einen', 'eines', 'ist', 'sind', 'war', 'hat', 'habe', 'haben', 'hatte',
    'wird', 'werden', 'wurde', 'kann', 'koennen', 'soll', 'muss', 'darf', 'will',
    'mit', 'von', 'fuer', 'auf', 'aus', 'bei', 'nach', 'ueber', 'unter', 'vor',
    'zwischen', 'durch', 'ohne', 'bis', 'gegen', 'seit', 'waehrend', 'wegen',
    'nicht', 'kein', 'keine', 'keinen', 'auch', 'noch', 'schon', 'mal', 'nur',
    'sehr', 'dann', 'denn', 'weil', 'wie', 'was', 'wer', 'wo', 'wann', 'welche',
    'dieser', 'diese', 'dieses', 'jede', 'jeder', 'jedes', 'alle', 'einige',
    'mein', 'dein', 'sein', 'unser', 'euer', 'meine', 'seine', 'ihre', 'unsere',
    'an', 'in', 'im', 'am', 'zum', 'zur', 'vom', 'beim', 'ins', 'ans',
    'hier', 'dort', 'da', 'so', 'nun', 'ja', 'nein', 'etwa', 'also', 'doch',
    'habe', 'hast', 'bin', 'bist', 'gibt', 'gab', 'geben', 'machen',
}
_STOPWORDS_EN = {
    'i', 'you', 'he', 'she', 'it', 'we', 'they', 'me', 'him', 'her', 'us', 'them',
    'the', 'a', 'an', 'and', 'or', 'but', 'if', 'then', 'that', 'this', 'these',
    'is', 'are', 'was', 'were', 'has', 'have', 'had', 'will', 'would', 'can', 'could',
    'should', 'may', 'might', 'must', 'shall', 'do', 'does', 'did', 'been', 'being',
    'with', 'from', 'for', 'on', 'at', 'by', 'to', 'of', 'in', 'about', 'into',
    'through', 'during', 'before', 'after', 'above', 'below', 'between', 'under',
    'not', 'no', 'also', 'just', 'only', 'very', 'too', 'so', 'than', 'more',
    'some', 'any', 'all', 'each', 'every', 'both', 'few', 'many', 'much',
    'my', 'your', 'his', 'its', 'our', 'their', 'what', 'which', 'who', 'when',
    'where', 'how', 'why', 'here', 'there', 'now', 'then', 'again', 'once',
    'get', 'got', 'make', 'made', 'go', 'went', 'come', 'came',
}
_STOPWORDS_ALL = _STOPWORDS_DE | _STOPWORDS_EN | {
    'suche', 'finde', 'zeige', 'search', 'find', 'show', 'look',
    'letzte', 'letzten', 'letzter', 'last', 'recent', 'ago',
    'woche', 'monat', 'monate', 'tage', 'tagen', 'week', 'month', 'months', 'days',
    'drei', 'zwei', 'vier', 'fuenf', 'sechs', 'three', 'two', 'four', 'five', 'six',
}


def extract_search_keywords(freetext, search_type=None):
    """Extract significant search keywords from freetext.
    For short queries (<=5 words): return all words.
    For longer queries: filter stopwords, return max 6 keywords.
    """
    import re as _re
    words = _re.findall(r'[a-zA-ZäöüÄÖÜß@._-]{2,}', freetext)
    if len(words) <= 5:
        return [w for w in words if len(w) >= 2]

    # Filter stopwords for longer queries
    significant = []
    for w in words:
        if w.lower() not in _STOPWORDS_ALL and len(w) >= 3:
            significant.append(w)
        elif w[0].isupper() and len(w) >= 3:
            # Keep proper nouns even if in stopwords
            significant.append(w)

    # Deduplicate preserving order
    seen = set()
    deduped = []
    for w in significant:
        wl = w.lower()
        if wl not in seen:
            seen.add(wl)
            deduped.append(w)

    return deduped[:6]


def search_contacts(query, speicher_path):
    """Search contacts.json for matching contacts.
    Returns list of {name, content, score} compatible with HybridSearch results.
    """
    import json as _json
    contacts_path = os.path.join(speicher_path, 'memory', 'contacts.json')
    if not os.path.exists(contacts_path):
        return []

    try:
        with open(contacts_path, 'r') as f:
            data = _json.load(f)
    except Exception:
        return []

    contacts = data.get('contacts', [])
    keywords = [kw.lower() for kw in query.split() if len(kw) >= 2]
    if not keywords:
        return []

    results = []
    for c in contacts:
        searchable = ' '.join(filter(None, [
            c.get('name') or '', c.get('email') or '', c.get('company') or '',
            c.get('title') or '', c.get('phone') or '',
        ])).lower()

        score = sum(searchable.count(kw) for kw in keywords)
        # Bonus for name/email exact match
        score += sum(3 for kw in keywords if kw in ((c.get('name') or '').lower() + ' ' + (c.get('email') or '').lower()))

        if score > 0:
            # Format as readable content
            lines = [f"Name: {c.get('name', '?')}"]
            if c.get('email'):
                lines.append(f"E-Mail: {c['email']}")
            if c.get('company'):
                lines.append(f"Firma: {c['company']}")
            if c.get('title'):
                lines.append(f"Titel: {c['title']}")
            if c.get('phone'):
                lines.append(f"Tel: {c['phone']}")
            lines.append(f"Kontakte: {c.get('total_contacts', 0)} (gesendet: {c.get('sent', 0)}, empfangen: {c.get('received', 0)})")
            if c.get('last_contact'):
                lines.append(f"Letzter Kontakt: {c['last_contact']}")

            results.append({
                'name': f"contact_{(c.get('name') or c.get('email', 'unknown')).replace(' ', '_')[:30]}",
                'content': '\n'.join(lines),
                'score': score,
            })

    results.sort(key=lambda x: -x['score'])
    return results[:20]


def get_recent_files(speicher_path, category=None, limit=10, per_category=False):
    """Get most recent files from the search index, optionally filtered by category.
    If per_category=True and no category given, returns limit files per top-level category.
    Returns list of {name, path, source_type, date, mtime, preview, subject, from} sorted by mtime desc.
    """
    idx = SearchIndex(speicher_path)
    idx.build_or_update()
    entries = idx.entries

    # Top-level categories for grouping
    TOP_CATEGORIES = {
        'email': {'email', 'notification'},
        'webclip': {'webclip', 'webclip_salesforce', 'webclip_slack', 'webclip_linkedin', 'webclip_general'},
        'document': {'document', 'document_word', 'document_excel', 'document_pdf', 'document_pptx'},
        'conversation': {'conversation'},
        'screenshot': {'screenshot'},
        'whatsapp': {'whatsapp', 'whatsapp_direct', 'whatsapp_group'},
    }

    def matches_category(source_type, cat):
        if not cat:
            return True
        return source_type in TOP_CATEGORIES.get(cat, {cat})

    if per_category and not category:
        # Return limit items per category
        results = []
        for cat, type_set in TOP_CATEGORIES.items():
            cat_items = []
            for fname, entry in entries.items():
                st = entry.get('source_type') or entry.get('type', 'file')
                if st in type_set:
                    cat_items.append({
                        'name': fname,
                        'path': entry.get('path', ''),
                        'source_type': st,
                        'category': cat,
                        'date': entry.get('date', ''),
                        'mtime': entry.get('mtime', 0),
                        'preview': (entry.get('preview', '') or '')[:150],
                        'subject': entry.get('subject', ''),
                        'from': entry.get('from', ''),
                    })
            cat_items.sort(key=lambda x: -x['mtime'])
            results.extend(cat_items[:limit])
        return results

    # Single category or all
    items = []
    for fname, entry in entries.items():
        st = entry.get('source_type') or entry.get('type', 'file')
        if not matches_category(st, category):
            continue
        items.append({
            'name': fname,
            'path': entry.get('path', ''),
            'source_type': st,
            'category': category or st,
            'date': entry.get('date', ''),
            'mtime': entry.get('mtime', 0),
            'preview': (entry.get('preview', '') or '')[:150],
            'subject': entry.get('subject', ''),
            'from': entry.get('from', ''),
        })
    items.sort(key=lambda x: -x['mtime'])
    return items[:limit]


def auto_search(msg, speicher_path, max_results=None, use_rag=True,
                fast=True, enable_global=True):
    """Unified smart auto-search entry point.

    Behaviour:
      1. Parse intent. If not a search query -> return ([], None).
      2. If user explicitly asked for global search (intent.wants_global) and
         enable_global is set -> route to global_rag_search.
      3. Otherwise run hybrid_rag_search (BM25 + semantic + RRF fusion, with
         optional query expansion + contextual compression).
      4. Fallback to classic HybridSearch if the RAG pipeline returns empty
         or raises.

    Returns (results, feedback) where each result is {name, content, score}
    so the existing call sites in web_server.py keep working.
    """
    intent = QueryParser.parse(msg)
    if not intent.is_search:
        return ([], None)

    effective_max = max_results or intent.max_results or 5
    # Deep mode upgrade: /deep, "ausführlich" etc. in the message → pay for
    # query expansion + compression.
    effective_fast = fast and not intent.wants_deep

    # If the user is searching for emails (explicit "email"/"mail" keyword or
    # a file_type=email filter), default to recency_first so the freshest mail
    # always sits at position 1.
    msg_lower_for_rec = msg.lower()
    if not intent.recency_first:
        EMAIL_WORDS = ('email', 'emails', 'mail', 'mails', 'e-mail', 'e-mails',
                       'nachricht', 'nachrichten', 'message', 'messages',
                       'mensagem', 'mensagens')
        words_in_msg = msg_lower_for_rec.split()
        if (intent.file_type == 'email'
                or any(w.rstrip('.,;:!?') in EMAIL_WORDS for w in words_in_msg)):
            intent.recency_first = True

    # ── Route to global search when requested ──
    if enable_global and intent.wants_global:
        try:
            return global_rag_search(msg, max_results=effective_max,
                                     fast=effective_fast,
                                     compress=not effective_fast)
        except Exception as e:
            print(f"[auto_search] global_rag_search failed: {e}")
            # fall through to local

    # ── Primary path: Hybrid RAG ──
    rag_out = None
    if use_rag and speicher_path:
        try:
            n_variants = 0 if effective_fast else 2
            compress = not effective_fast
            rag_out = hybrid_rag_search(
                msg, speicher_path,
                max_results=effective_max,
                compress=compress,
                n_variants=n_variants,
                recency_first=intent.recency_first,
            )
        except Exception as e:
            print(f"[auto_search] hybrid_rag_search error: {e}")
            rag_out = None

    if rag_out and rag_out.get('results'):
        results = []
        for r in rag_out['results']:
            content = r.get('compressed') or r.get('snippet') or ''
            results.append({
                'name': r['name'],
                'content': content,
                'score': r.get('score', 0.0),
            })
        feedback = {
            'query': ' '.join(intent.keywords) if intent.keywords else msg,
            'keywords': intent.keywords,
            'date_label': intent.date_label,
            'field_label': intent.field_label,
            'file_type': intent.file_type,
            'found_count': len(results),
            'index_count': 0,
            'rag': True,
            'semantic': rag_out.get('semantic', False),
            'queries': rag_out.get('queries', [msg]),
            'fallback': rag_out.get('fallback'),
            'recency_first': intent.recency_first,
        }
        return (results, feedback)

    # ── Keyword-only fallback ──
    try:
        idx = get_or_build_index(speicher_path)
        idx.update_index()
    except Exception as e:
        print(f"[auto_search] index update error: {e}")

    results, feedback = HybridSearch.search(
        intent, speicher_path, max_results=effective_max,
        recency_first=intent.recency_first,
    )
    if feedback is not None:
        feedback['rag'] = False
        feedback['recency_first'] = intent.recency_first
    return (results, feedback)


def format_search_feedback(feedback, found_count):
    """Format search feedback as a user-visible string."""
    if not feedback:
        return ''

    parts = [feedback.get('query', '')]
    if feedback.get('file_type'):
        type_labels = {'email': 'Email', 'document': 'Dokument', 'image': 'Bild',
                      'conversation': 'Konversation', 'salesforce': 'Salesforce'}
        parts.append(f"Typ: {type_labels.get(feedback['file_type'], feedback['file_type'])}")
    if feedback.get('date_label'):
        parts.append(f"Zeitraum: {feedback['date_label']}")
    if feedback.get('field_label'):
        parts.append(f"Feld: {feedback['field_label']}")

    query_desc = ' | '.join(parts)
    idx_count = feedback.get('index_count', 0)

    mode_tags = []
    if feedback.get('global'):
        mode_tags.append('global')
    if feedback.get('rag'):
        mode_tags.append('RAG')
    if feedback.get('semantic'):
        mode_tags.append('semantic')
    mode_str = f" [{'/'.join(mode_tags)}]" if mode_tags else ''

    if found_count > 0:
        idx_part = f" | Index: {idx_count} Dateien" if idx_count else ''
        return f"\U0001f50d Suche{mode_str}: {query_desc} | Gefunden: {found_count} Datei(en){idx_part}"
    else:
        idx_part = f" Gesucht in: {idx_count} Dateien" if idx_count else ''
        return f"\U0001f50d Keine Treffer{mode_str} fuer '{query_desc}'.{idx_part}"


# ═══════════════════════════════════════════════════════════════════════════════
# TEIL 4 — GLOBAL SEARCH INDEX
# ═══════════════════════════════════════════════════════════════════════════════

DATALAKE_BASE = os.path.expanduser(
    "~/Library/Mobile Documents/com~apple~CloudDocs/Downloads shared/claude_datalake"
)
DOWNLOADS_SHARED = os.path.expanduser(
    "~/Library/Mobile Documents/com~apple~CloudDocs/Downloads shared"
)

SKIP_DIRS = {'backups', 'build', 'dist', '__pycache__', '.git', 'node_modules'}
SKIP_EXTS = {'.pyc', '.pyo', '.DS_Store', '.aifs'}

GLOBAL_TRIGGERS = [
    # Deutsch
    'erweitertes gedaechtnis', 'erweitertem gedaechtnis',
    'erweitertes gedächtnis', 'erweitertem gedächtnis',
    'ueberall suchen', 'überall suchen', 'ueberall', 'überall',
    'global suchen', 'alle agenten', 'allen agenten',
    'gesamtes memory', 'gesamten memory', 'alles durchsuchen',
    # Englisch
    'extended memory', 'global search', 'search everywhere',
    'all agents', 'everywhere', 'search all', 'across all',
    # Portugiesisch
    'memoria extendida', 'memória extendida', 'busca global',
    'procura tudo', 'em tudo', 'todos os agentes', 'pesquisa global',
]

NOTIFICATION_PATTERNS_SE = [
    'noreply', 'no-reply', 'no_reply', 'donotreply', 'do-not-reply',
    'mailer-daemon', 'postmaster', 'notifications@', 'notification@',
    'newsletter', 'digest', 'automated', 'auto-reply', 'autoreply',
    'bounce', 'teams.mail.microsoft', 'notify.microsoft',
    'posted a message', 'mentioned you', 'invited you',
    'automatic reply', 'automatische antwort', 'abwesenheit',
    'out of office', 'delivery failed', 'delivery status',
    'unsubscribe', 'abmelden',
]


def detect_global_trigger(msg):
    """Check if message contains a global search trigger. Returns True/False."""
    msg_lower = msg.lower()
    return any(t in msg_lower for t in GLOBAL_TRIGGERS)


def _detect_agent_from_path(fpath):
    """Derive agent name from file path."""
    fpath_lower = fpath.lower()
    for agent in ['signicat', 'privat', 'trustedcarrier']:
        if f'/{agent}/' in fpath_lower:
            return agent
    if '/email_inbox/' in fpath_lower:
        return 'inbox'
    return 'global'


def _detect_type_from_ext(fname):
    """Detect file type from extension."""
    ext = os.path.splitext(fname.lower())[1]
    if ext in {'.eml'}:
        return 'email'
    if ext in {'.pdf', '.docx', '.xlsx', '.pptx', '.doc', '.xls', '.ppt'}:
        return 'document'
    if ext in {'.png', '.jpg', '.jpeg', '.gif', '.webp'}:
        return 'image'
    if ext in {'.py', '.js', '.json', '.ts', '.css', '.html'}:
        return 'code'
    if ext in {'.txt', '.md'}:
        fname_lower = fname.lower()
        if fname_lower.startswith('email_'):
            return 'email'
        if fname_lower.startswith('konversation_'):
            return 'conversation'
        return 'file'
    return 'file'


def _is_notification_se(from_field, subject):
    """Check if email is a notification."""
    check = (from_field + ' ' + subject).lower()
    return any(p in check for p in NOTIFICATION_PATTERNS_SE)


def _extract_text_from_file(fpath, max_chars=30000):
    """Extract text content from various file types."""
    ext = os.path.splitext(fpath.lower())[1]

    # Binary / image — no text
    if ext in BINARY_EXTS:
        return ''

    # Plain text files
    if ext in {'.txt', '.eml', '.md', '.py', '.js', '.json', '.ts', '.css', '.html', '.csv', '.log'}:
        try:
            with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                raw = f.read(max_chars)
            # For JSON webclips: extract full_text + title for better searchability
            if ext == '.json':
                fname = os.path.basename(fpath).lower()
                if any(fname.startswith(p) for p in ('web_', 'slack_', 'salesforce_', 'linkedin_')):
                    try:
                        import json as _json
                        parsed = _json.loads(raw)
                        parts = []
                        if parsed.get('title'):
                            parts.append('Title: ' + parsed['title'])
                        if parsed.get('url'):
                            parts.append('URL: ' + parsed['url'])
                        if parsed.get('full_text'):
                            parts.append(parsed['full_text'])
                        if parts:
                            return '\n'.join(parts)[:max_chars]
                    except (ValueError, KeyError, TypeError):
                        pass
            return raw
        except Exception:
            return ''

    # PDF
    if ext == '.pdf':
        try:
            import PyPDF2
            with open(fpath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text_parts = []
                for page in reader.pages[:20]:
                    text_parts.append(page.extract_text() or '')
                return '\n'.join(text_parts)[:max_chars]
        except Exception:
            return ''

    # DOCX
    if ext == '.docx':
        try:
            from docx import Document
            doc = Document(fpath)
            return '\n'.join(p.text for p in doc.paragraphs)[:max_chars]
        except Exception:
            return ''

    # XLSX
    if ext == '.xlsx':
        try:
            from openpyxl import load_workbook
            wb = load_workbook(fpath, read_only=True, data_only=True)
            text_parts = []
            for ws in wb.worksheets[:5]:
                for row in ws.iter_rows(max_row=200, values_only=True):
                    vals = [str(c) for c in row if c is not None]
                    if vals:
                        text_parts.append(' | '.join(vals))
            wb.close()
            return '\n'.join(text_parts)[:max_chars]
        except Exception:
            return ''

    # PPTX
    if ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(fpath)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        text_parts.append(shape.text)
            return '\n'.join(text_parts)[:max_chars]
        except Exception:
            return ''

    # Fallback: try as text
    try:
        with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
            return f.read(max_chars)
    except Exception:
        return ''


class GlobalSearchIndex:
    """Builds and maintains a global search index across all agents and Downloads shared."""

    def __init__(self):
        self.index_file = os.path.join(DOWNLOADS_SHARED, '.global_search_index.json')
        self.entries = {}  # relative_path -> entry dict
        self._lock = threading.Lock()
        self._build_time = 0

    def _should_skip(self, name, is_dir=False):
        """Check if a file/dir should be skipped."""
        if is_dir:
            return name in SKIP_DIRS or name.startswith('.')
        ext = os.path.splitext(name)[1]
        return ext in SKIP_EXTS or name == '.DS_Store'

    def _scan_directory(self, root_dir, file_list):
        """Recursively scan a directory, appending (abs_path, rel_path) tuples."""
        try:
            entries = os.scandir(root_dir)
        except (PermissionError, OSError):
            return
        for entry in entries:
            try:
                if entry.is_dir(follow_symlinks=False):
                    if not self._should_skip(entry.name, is_dir=True):
                        self._scan_directory(entry.path, file_list)
                elif entry.is_file(follow_symlinks=False):
                    if not self._should_skip(entry.name):
                        file_list.append(entry.path)
            except (PermissionError, OSError):
                continue

    def _extract_email_headers(self, text):
        """Extract From/To/Subject from email text."""
        headers = {'from': '', 'to': '', 'subject': '', 'date': ''}
        for line in text.split('\n')[:15]:
            line_stripped = line.strip()
            ll = line_stripped.lower()
            if ll.startswith(('von:', 'from:')):
                headers['from'] = line_stripped.split(':', 1)[1].strip()
            elif ll.startswith(('an:', 'to:')):
                headers['to'] = line_stripped.split(':', 1)[1].strip()
            elif ll.startswith(('betreff:', 'subject:')):
                headers['subject'] = line_stripped.split(':', 1)[1].strip()
            elif ll.startswith(('datum:', 'date:')):
                headers['date'] = line_stripped.split(':', 1)[1].strip()
        return headers

    def _extract_keywords(self, text, max_keywords=10):
        """Extract top-N keywords."""
        words = re.findall(r'[a-zA-Z\u00e4\u00f6\u00fc\u00c4\u00d6\u00dc\u00df\u00e0-\u00ff]+', text.lower())
        freq = {}
        for w in words:
            if len(w) < 3 or w in STOPWORDS:
                continue
            w_norm = normalize_unicode(w)
            if len(w_norm) < 3:
                continue
            freq[w_norm] = freq.get(w_norm, 0) + 1
        sorted_words = sorted(freq.items(), key=lambda x: x[1], reverse=True)
        return [w for w, _ in sorted_words[:max_keywords]]

    def _index_file(self, fpath):
        """Create an index entry for a single file."""
        try:
            stat = os.stat(fpath)
        except Exception:
            return None

        fname = os.path.basename(fpath)
        ftype = detect_source_type(fname)
        mtime = stat.st_mtime
        file_date = datetime.datetime.fromtimestamp(mtime).strftime('%Y-%m-%d')
        agent = _detect_agent_from_path(fpath)

        # Compute relative folder
        try:
            folder = os.path.relpath(os.path.dirname(fpath), DOWNLOADS_SHARED)
        except ValueError:
            folder = os.path.dirname(fpath)

        entry = {
            'path': fpath,
            'filename': fname,
            'folder': folder,
            'agent': agent,
            'type': ftype,
            'date': file_date,
            'mtime': mtime,
            'size': stat.st_size,
            'from': '',
            'subject': '',
            'preview': '',
            'is_notification': False,
            'keywords': [],
        }

        # Extract text for non-image files
        if ftype != 'screenshot':
            text = _extract_text_from_file(fpath)
            if text:
                entry['preview'] = text[:200].replace('\n', ' ').strip()
                entry['keywords'] = self._extract_keywords(text)

                # Re-detect type with preview (for notification detection)
                ftype = detect_source_type(fname, text[:500])
                entry['type'] = ftype
                entry['source_type'] = ftype

                # Email headers
                if ftype in ('email', 'notification'):
                    headers = self._extract_email_headers(text)
                    entry['from'] = headers['from']
                    entry['subject'] = headers['subject']
                    entry['is_notification'] = (ftype == 'notification')

                # Conversation metadata
                if ftype == 'conversation':
                    conv_meta = extract_conversation_meta(text)
                    entry['subject'] = conv_meta.get('subject', '')

        entry['source_type'] = ftype
        return entry

    def build_global_index(self):
        """Build complete global index from both sources."""
        start = time.time()
        all_files = []

        # Source 1: All agent memory folders + email_inbox + claude_outputs
        if os.path.exists(DATALAKE_BASE):
            for item in os.listdir(DATALAKE_BASE):
                item_path = os.path.join(DATALAKE_BASE, item)
                if os.path.isdir(item_path) and not self._should_skip(item, is_dir=True):
                    memory_dir = os.path.join(item_path, 'memory')
                    if os.path.exists(memory_dir):
                        self._scan_directory(memory_dir, all_files)
            # email_inbox
            inbox = os.path.join(DATALAKE_BASE, 'email_inbox')
            if os.path.exists(inbox):
                self._scan_directory(inbox, all_files)
            # claude_outputs
            outputs = os.path.join(DATALAKE_BASE, 'claude_outputs')
            if os.path.exists(outputs):
                self._scan_directory(outputs, all_files)

        # Source 2: Entire Downloads shared folder
        if os.path.exists(DOWNLOADS_SHARED):
            self._scan_directory(DOWNLOADS_SHARED, all_files)

        # Deduplicate by absolute path
        all_files = list(set(all_files))

        print(f"Globaler Index: Scanne {len(all_files)} Dateien...")

        entries = {}
        for fpath in all_files:
            entry = self._index_file(fpath)
            if entry:
                entries[fpath] = entry

        with self._lock:
            self.entries = entries

        self._save_index()
        elapsed = time.time() - start
        self._build_time = elapsed

        # Count by agent
        agent_counts = {}
        for e in entries.values():
            a = e.get('agent', 'global')
            agent_counts[a] = agent_counts.get(a, 0) + 1
        agent_str = ', '.join(f"{k}: {v}" for k, v in sorted(agent_counts.items()))

        print(f"Globaler Index: {len(entries)} Dateien in {elapsed:.1f}s ({agent_str})")
        return len(entries)

    def update_index(self):
        """Incremental update: only re-index new/modified files."""
        self._load_index()
        if not self.entries:
            return self.build_global_index()

        start = time.time()

        # Rescan file list
        all_files = []
        if os.path.exists(DATALAKE_BASE):
            for item in os.listdir(DATALAKE_BASE):
                item_path = os.path.join(DATALAKE_BASE, item)
                if os.path.isdir(item_path) and not self._should_skip(item, is_dir=True):
                    memory_dir = os.path.join(item_path, 'memory')
                    if os.path.exists(memory_dir):
                        self._scan_directory(memory_dir, all_files)
            inbox = os.path.join(DATALAKE_BASE, 'email_inbox')
            if os.path.exists(inbox):
                self._scan_directory(inbox, all_files)
            outputs = os.path.join(DATALAKE_BASE, 'claude_outputs')
            if os.path.exists(outputs):
                self._scan_directory(outputs, all_files)
        if os.path.exists(DOWNLOADS_SHARED):
            self._scan_directory(DOWNLOADS_SHARED, all_files)
        all_files = list(set(all_files))

        current_files = set(all_files)
        indexed_files = set(self.entries.keys())

        # Remove deleted
        for fpath in indexed_files - current_files:
            del self.entries[fpath]

        # Add/update modified
        updated = 0
        for fpath in current_files:
            try:
                mtime = os.path.getmtime(fpath)
            except Exception:
                continue
            existing = self.entries.get(fpath)
            if existing and existing.get('mtime', 0) >= mtime:
                continue
            entry = self._index_file(fpath)
            if entry:
                self.entries[fpath] = entry
                updated += 1

        if updated > 0:
            self._save_index()

        elapsed = time.time() - start
        if updated > 0:
            print(f"Globaler Index: {updated} Dateien aktualisiert in {elapsed:.1f}s (gesamt: {len(self.entries)})")
        return updated

    def build_or_update(self):
        """Build if no index exists, otherwise update."""
        if os.path.exists(self.index_file):
            return self.update_index()
        return self.build_global_index()

    def _save_index(self):
        try:
            with open(self.index_file, 'w', encoding='utf-8') as f:
                json.dump(self.entries, f, ensure_ascii=False)
        except Exception as e:
            print(f"Global index save error: {e}")

    def _load_index(self):
        if self.entries:
            return
        if os.path.exists(self.index_file):
            try:
                with open(self.index_file, 'r', encoding='utf-8') as f:
                    self.entries = json.load(f)
            except Exception:
                self.entries = {}

    @property
    def file_count(self):
        self._load_index()
        return len(self.entries)


# Global index singleton + cache
_global_index = None
_global_index_lock = threading.Lock()


def get_global_index():
    """Get or build the global search index. Thread-safe."""
    global _global_index
    with _global_index_lock:
        if _global_index is not None:
            return _global_index
        _global_index = GlobalSearchIndex()
    _global_index.build_or_update()
    return _global_index


def build_global_index_async():
    """Build/update global index in a background thread."""
    def _build():
        get_global_index()
    t = threading.Thread(target=_build, daemon=True)
    t.start()


def _list_agent_speicher_paths():
    """Return all agent speicher directories with a memory/ subfolder."""
    base = DATALAKE_BASE
    if not os.path.exists(base):
        return []
    paths = []
    for item in os.listdir(base):
        p = os.path.join(base, item)
        if not os.path.isdir(p):
            continue
        if item.startswith('.') or item in ('config', 'email_inbox', 'claude_outputs'):
            continue
        if os.path.exists(os.path.join(p, 'memory')):
            paths.append(p)
    return paths


def update_all_indexes():
    """Incremental update of all agent indexes + global index.
    Called periodically by the web server background thread."""
    datalake = os.path.expanduser("~/Library/Mobile Documents/com~apple~CloudDocs/Downloads shared/claude_datalake")
    if not os.path.exists(datalake):
        return
    total_updated = 0
    for item in os.listdir(datalake):
        agent_path = os.path.join(datalake, item)
        if not os.path.isdir(agent_path):
            continue
        memory_dir = os.path.join(agent_path, 'memory')
        if not os.path.exists(memory_dir):
            continue
        # Skip config, email_inbox etc.
        if item.startswith('.') or item in ('config', 'email_inbox', 'claude_outputs'):
            continue
        try:
            idx = get_or_build_index(agent_path)
            updated = idx.update_index()
            total_updated += updated
        except Exception as e:
            print(f"Index update error for {item}: {e}")
    # Also update global index
    try:
        gi = get_global_index()
        gi.update_index()
    except Exception as e:
        print(f"Global index update error: {e}")
    if total_updated > 0:
        print(f"[INDEX] Periodischer Update: {total_updated} neue/geaenderte Dateien indexiert")
    return total_updated


def index_single_file(speicher_path, filename):
    """Index a single file immediately. Used by email_watcher after saving."""
    try:
        idx = get_or_build_index(speicher_path)
        idx.add_file(filename)
    except Exception as e:
        print(f"Single file index error ({filename}): {e}")


def global_search(msg, max_results=50):
    """Execute a global search across all agents.
    Returns (results, feedback) where results are sorted by score.
    """
    intent = QueryParser.parse(msg)
    # Force is_search for global trigger
    if not intent.is_search:
        intent.is_search = True

    gidx = get_global_index()
    gidx._load_index()

    if not gidx.entries:
        return ([], {'query': msg, 'found_count': 0, 'index_count': 0, 'global': True})

    keywords_norm = [normalize_unicode(kw) for kw in intent.keywords]
    names_norm = [normalize_unicode(pn) for pn in intent.person_names]

    feedback = {
        'query': ' '.join(intent.keywords) if intent.keywords else intent.raw_query,
        'date_label': intent.date_label,
        'field_label': intent.field_label,
        'file_type': intent.file_type,
        'found_count': 0,
        'index_count': len(gidx.entries),
        'global': True,
    }

    scored = {}

    for fpath, entry in gidx.entries.items():
        # Date filter
        if intent.date_filter:
            start_date, end_date = intent.date_filter
            start_ts = datetime.datetime.combine(start_date, datetime.time.min).timestamp()
            end_ts = datetime.datetime.combine(end_date, datetime.time.max).timestamp()
            mtime = entry.get('mtime', 0)
            if mtime < start_ts or mtime > end_ts:
                continue

        # Type filter — use source_types_effective when available so subcategories
        # (webclip_slack, document_pdf, ...) are matched against their parent.
        entry_type = entry.get('source_type') or entry.get('type', 'file')
        if intent.source_types_effective:
            if entry_type not in intent.source_types_effective:
                continue
        elif intent.file_type:
            entry_parent = SOURCE_TAXONOMY.get(entry_type, {}).get('parent', entry_type)
            if entry_type != intent.file_type and entry_parent != intent.file_type:
                continue

        score = 0
        fname_norm = normalize_unicode(entry.get('filename', ''))
        fname_flat = fname_norm.replace('-', '').replace('_', '')

        # Pre-ranking notification penalty (consistent with local HybridSearch).
        if entry_type == 'notification':
            score -= 15

        if intent.date_filter:
            score += 10

        # Filename matching (exact + fuzzy)
        for kw in keywords_norm:
            if kw in fname_norm or kw in fname_flat:
                score += 4
            else:
                fm, fs = fuzzy_match(kw, fname_norm.replace('_', ' ').replace('-', ' '))
                if fm:
                    score += min(fs, 3)

        # Email header matching (exact + fuzzy)
        if entry_type in ('email', 'notification'):
            from_norm = normalize_unicode(entry.get('from', ''))
            to_norm = normalize_unicode(entry.get('to', ''))
            subject_norm = normalize_unicode(entry.get('subject', ''))

            for kw in keywords_norm:
                if intent.field_filter == 'from' and kw in from_norm:
                    score += 5
                elif intent.field_filter == 'to' and kw in to_norm:
                    score += 5
                elif intent.field_filter == 'about' and kw in subject_norm:
                    score += 3
                elif intent.field_filter is None:
                    if kw in from_norm:
                        score += 5
                    if kw in subject_norm:
                        score += 3
                    if kw in to_norm:
                        score += 2

            for pn in names_norm:
                if pn in from_norm:
                    score += 8
                elif pn in to_norm:
                    score += 6
                elif pn in subject_norm:
                    score += 4
                else:
                    fm_from, fs_from = fuzzy_match(pn, from_norm)
                    fm_to, fs_to = fuzzy_match(pn, to_norm)
                    fm_subj, fs_subj = fuzzy_match(pn, subject_norm)
                    if fm_from:
                        score += min(fs_from + 2, 6)
                    elif fm_to:
                        score += min(fs_to + 1, 4)
                    elif fm_subj:
                        score += min(fs_subj, 3)

        # Non-email: person names in filename (exact + fuzzy)
        for pn in names_norm:
            if pn in fname_norm or pn in fname_flat:
                score += 5
            else:
                fm, fs = fuzzy_match(pn, fname_norm.replace('_', ' ').replace('-', ' '))
                if fm:
                    score += min(fs, 4)

        # Preview/keyword matching (exact + fuzzy)
        preview_norm = normalize_unicode(entry.get('preview', ''))
        entry_keywords = entry.get('keywords', [])
        for kw in keywords_norm:
            if kw in preview_norm:
                score += 1
            elif len(kw) >= 4:
                fm, fs = fuzzy_match(kw, preview_norm)
                if fm:
                    score += 1
            if kw in entry_keywords:
                score += 2

        if score > 0:
            scored[fpath] = score

    # Full-text top-20
    top_candidates = sorted(scored.items(), key=lambda x: x[1], reverse=True)[:20]
    for fpath, base_score in top_candidates:
        entry = gidx.entries.get(fpath)
        if not entry or entry.get('type') == 'image':
            continue
        try:
            text = _extract_text_from_file(fpath, max_chars=20000)
            if not text:
                continue
            text_norm = normalize_unicode(text)
            body_score = 0
            for kw in keywords_norm:
                count = text_norm.count(kw)
                if count >= 3:
                    body_score += 3
                elif count >= 1:
                    body_score += 1
            for pn in names_norm:
                if pn in text_norm:
                    body_score += 5
            scored[fpath] = scored.get(fpath, 0) + body_score
        except Exception:
            continue

    # Final ranking
    final_sorted = sorted(scored.items(), key=lambda x: x[1], reverse=True)
    results = []
    for fpath, score in final_sorted[:max_results]:
        if score <= 0:
            break
        entry = gidx.entries.get(fpath, {})
        from_field = entry.get('from', '')
        subject = entry.get('subject', '')
        from_norm = normalize_unicode(from_field) if from_field else ''
        from_person = any(pn in from_norm for pn in names_norm) if names_norm else False
        is_notif = entry.get('is_notification', False)
        if not is_notif and from_field:
            is_notif = _is_notification_se(from_field, subject)

        results.append({
            'name': entry.get('filename', os.path.basename(fpath)),
            'path': fpath,
            'type': entry.get('type', 'file'),
            'date': entry.get('date', ''),
            'from': from_field,
            'subject': subject,
            'preview': entry.get('preview', '')[:150],
            'score': score,
            'from_person': from_person,
            'is_notification': is_notif,
            'agent': entry.get('agent', 'global'),
            'folder': entry.get('folder', ''),
        })

    # Sort: from_person + not notification first, then score desc
    results.sort(key=lambda x: (not x['from_person'], x['is_notification'], -x['score']))

    feedback['found_count'] = len(results)
    return (results, feedback)


def global_rag_search(query, max_results=10, fast=True, compress=False):
    """Cross-agent hybrid RAG: keyword (existing global index) fused with
    semantic search over the union of per-agent embedding indexes.

    This is additive — we do NOT build a separate global embedding store.
    Instead we iterate the per-agent `.embedding_index.json` files, which
    are kept up-to-date by `index_file_with_embedding` at add-time.

    Returns (results, feedback) compatible with auto_search.
    """
    # ── 1. Keyword side: existing global_search ──
    try:
        kw_results, kw_feedback = global_search(query, max_results=50)
    except Exception as e:
        print(f"[global_rag_search] keyword side failed: {e}")
        kw_results = []
        kw_feedback = {'query': query, 'global': True, 'found_count': 0, 'index_count': 0}

    kw_by_name = {r['name']: r for r in kw_results}
    kw_ranked = [r['name'] for r in kw_results]

    # ── 2. Semantic side: union of per-agent embedding indexes ──
    semantic_lists = []
    chunk_by_key = {}
    first_speicher = None

    api_key = _get_openai_api_key()
    if api_key:
        q_vecs = _call_openai_embedding([query], api_key)
        qv = q_vecs[0] if q_vecs else None
        if qv:
            for speicher in _list_agent_speicher_paths():
                try:
                    emb_idx = get_embedding_index(speicher)
                    emb_idx._load()
                    if not emb_idx.entries:
                        continue
                    if first_speicher is None:
                        first_speicher = speicher
                    hits = emb_idx.search(qv, max_results=20)
                    sem_list = []
                    for h in hits:
                        name = h['name']
                        sem_list.append(name)
                        prev = chunk_by_key.get(name)
                        if prev is None or h['score'] > prev[1]:
                            chunk_by_key[name] = (h['chunk'], h['score'], speicher)
                    if sem_list:
                        semantic_lists.append(sem_list)
                except Exception as e:
                    print(f"[global_rag_search] embedding search for {speicher} failed: {e}")

    # ── 3. RRF fusion ──
    fused = rrf_fuse([kw_ranked] + semantic_lists, k=_RRF_K, top_n=max_results * 2)

    # ── 4. Build results ──
    results = []
    memory_dir_fallback = os.path.join(DATALAKE_BASE, '_unknown')
    for item in fused[:max_results]:
        fname = item['name']
        kw_entry = kw_by_name.get(fname)

        snippet = ''
        speicher_for_compress = None
        if fname in chunk_by_key:
            snippet = chunk_by_key[fname][0]
            speicher_for_compress = chunk_by_key[fname][2]
        elif kw_entry and kw_entry.get('path'):
            try:
                with open(kw_entry['path'], 'r', encoding='utf-8', errors='replace') as f:
                    snippet = f.read(8000)
            except Exception:
                snippet = kw_entry.get('preview', '')[:2000]

        if compress and snippet:
            speicher_for_compress = speicher_for_compress or first_speicher
            if speicher_for_compress:
                try:
                    snippet = compress_chunk(snippet, query, speicher_for_compress)
                except Exception:
                    pass

        results.append({
            'name': fname,
            'content': snippet,
            'score': item['rrf_score'],
            'path': (kw_entry or {}).get('path', ''),
            'agent': (kw_entry or {}).get('agent', 'global'),
        })

    feedback = {
        'query': query,
        'global': True,
        'rag': True,
        'semantic': bool(semantic_lists),
        'found_count': len(results),
        'index_count': kw_feedback.get('index_count', 0) if kw_feedback else 0,
    }
    return (results, feedback)


# ═══════════════════════════════════════════════════════════════════════════════
# TEIL 4 — SEMANTIC RAG (Embeddings + Query Expansion + RRF + Compression)
# ═══════════════════════════════════════════════════════════════════════════════
# Additive upgrade on top of the keyword pipeline. All calls are wrapped in
# try/except and fall back to the existing keyword search if the OpenAI or
# Anthropic/Mistral/Gemini APIs are unreachable, the key is missing, or quota
# is exceeded. Existing indexes (.search_index.json / .global_search_index.json)
# are never touched — embeddings live in their own .embedding_index.json files.

import math
import urllib.request
import urllib.error

try:
    import numpy as _np
    _HAS_NUMPY = True
except Exception:
    _HAS_NUMPY = False

_MODELS_JSON_PATH = os.path.expanduser(
    "~/Library/Mobile Documents/com~apple~CloudDocs/Downloads shared/claude_datalake/config/models.json"
)
_EMBEDDING_MODEL = "text-embedding-3-small"
_EMBEDDING_DIM = 1536
_EMBEDDING_HTTP_TIMEOUT = 10
_LLM_HTTP_TIMEOUT = 20
_RRF_K = 60
_CHUNK_TARGET_TOKENS = 500
_CHUNK_OVERLAP_TOKENS = 50
_TOKEN_CHAR_RATIO = 4          # grobe Heuristik: 1 token ≈ 4 chars
_MAX_CHUNK_CHARS = _CHUNK_TARGET_TOKENS * _TOKEN_CHAR_RATIO
_OVERLAP_CHARS = _CHUNK_OVERLAP_TOKENS * _TOKEN_CHAR_RATIO
_SHORT_DOC_CHAR_LIMIT = 1500   # darunter: ganzes Dokument als ein Chunk


def _load_models_config():
    """Read models.json once per call. Returns dict or None."""
    try:
        with open(_MODELS_JSON_PATH, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return None


def _get_openai_api_key():
    cfg = _load_models_config()
    if not cfg:
        return None
    return cfg.get('providers', {}).get('openai', {}).get('api_key') or None


def _get_agent_llm(speicher_path):
    """Return (provider, model_id, api_key) for the given agent's configured
    chat model. Falls back to the first Anthropic model in models.json."""
    cfg = _load_models_config() or {}
    providers = cfg.get('providers', {})

    # Try per-agent config file: [speicher]/memory/models.json
    agent_cfg_path = os.path.join(speicher_path, 'memory', 'models.json')
    try:
        if os.path.exists(agent_cfg_path):
            with open(agent_cfg_path, 'r', encoding='utf-8') as f:
                agent_cfg = json.load(f)
            prov = agent_cfg.get('provider')
            mid = agent_cfg.get('model')
            if prov and mid:
                key = providers.get(prov, {}).get('api_key')
                if key:
                    return (prov, mid, key)
    except Exception:
        pass

    # Fallback: anthropic haiku if available, otherwise first anthropic model
    anth = providers.get('anthropic', {})
    key = anth.get('api_key')
    if key:
        models = anth.get('models', [])
        for m in models:
            if 'haiku' in m.get('id', '').lower():
                return ('anthropic', m['id'], key)
        if models:
            return ('anthropic', models[0]['id'], key)
    return (None, None, None)


# ─── EMBEDDINGS (OpenAI text-embedding-3-small) ──────────────────────────────

# Circuit-breaker state for the embedding endpoint. When OpenAI returns 429
# or 5xx we stop calling for `_EMBEDDING_CB_COOLDOWN` seconds so the caller
# falls back to keyword-only retrieval without log-spam.
_EMBEDDING_CB_COOLDOWN = 600  # 10 minutes
_embedding_cb_block_until = 0.0
_embedding_cb_last_log = 0.0


def _call_openai_embedding(texts, api_key):
    """Batched embedding call. Returns list[list[float]] or None on failure.
    Includes a circuit breaker on HTTP 429/5xx to prevent log-spam & quota burn.
    """
    global _embedding_cb_block_until, _embedding_cb_last_log
    if not api_key or not texts:
        return None
    now = time.time()
    if now < _embedding_cb_block_until:
        # Throttle the cooldown message to once per minute
        if now - _embedding_cb_last_log > 60:
            _embedding_cb_last_log = now
            remaining = int(_embedding_cb_block_until - now)
            print(f"[embedding] circuit-breaker active ({remaining}s left); skipping call")
        return None

    batch_size = 64
    all_vecs = []
    url = "https://api.openai.com/v1/embeddings"
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        payload = json.dumps({
            'model': _EMBEDDING_MODEL,
            'input': batch,
        }).encode('utf-8')
        req = urllib.request.Request(
            url,
            data=payload,
            headers={
                'Authorization': f'Bearer {api_key}',
                'Content-Type': 'application/json',
            },
        )
        try:
            with urllib.request.urlopen(req, timeout=_EMBEDDING_HTTP_TIMEOUT) as resp:
                body = resp.read()
                data = json.loads(body.decode('utf-8'))
            for item in data.get('data', []):
                all_vecs.append(item.get('embedding', []))
        except urllib.error.HTTPError as e:
            # Quota / rate-limit / server error -> open circuit breaker.
            if e.code == 429 or 500 <= e.code < 600:
                _embedding_cb_block_until = time.time() + _EMBEDDING_CB_COOLDOWN
                print(f"[embedding] HTTP {e.code}: circuit-breaker opens for "
                      f"{_EMBEDDING_CB_COOLDOWN}s — fallback aktiv")
            else:
                print(f"[embedding] HTTP {e.code}: fallback aktiv")
            return None
        except Exception as e:
            print(f"[embedding] call failed: {e}")
            return None
    return all_vecs if all_vecs else None


def embed_text(text):
    """Embed a single string. Returns list[float] or None."""
    key = _get_openai_api_key()
    if not key:
        return None
    vecs = _call_openai_embedding([text], key)
    if not vecs:
        return None
    return vecs[0]


# ─── CHUNKING ────────────────────────────────────────────────────────────────

def _chunk_text(text, max_chars=None, overlap=None):
    """Split text into paragraph-aligned chunks of roughly ~500 tokens.

    - Short texts (< _SHORT_DOC_CHAR_LIMIT) return a single chunk.
    - Paragraphs separated by blank lines are accumulated greedily until the
      target chunk size is hit, then flushed with `overlap` chars of tail
      repeated into the next chunk.
    """
    if not text:
        return []
    max_chars = max_chars or _MAX_CHUNK_CHARS
    overlap = overlap or _OVERLAP_CHARS

    text = text.strip()
    if len(text) <= _SHORT_DOC_CHAR_LIMIT:
        return [text]

    paragraphs = re.split(r'\n\s*\n', text)
    chunks = []
    buf = ''
    for p in paragraphs:
        p = p.strip()
        if not p:
            continue
        if not buf:
            buf = p
            continue
        if len(buf) + 2 + len(p) <= max_chars:
            buf = buf + '\n\n' + p
        else:
            chunks.append(buf)
            tail = buf[-overlap:] if overlap and len(buf) > overlap else ''
            buf = (tail + '\n\n' + p) if tail else p
    if buf:
        chunks.append(buf)

    # Safety: if a single paragraph is larger than max_chars, hard-split it
    final = []
    for c in chunks:
        if len(c) <= max_chars:
            final.append(c)
            continue
        step = max_chars - overlap
        for i in range(0, len(c), max(step, 1)):
            final.append(c[i:i + max_chars])
    return final


# ─── EMBEDDING INDEX (per agent) ─────────────────────────────────────────────

class EmbeddingIndex:
    """Per-agent embedding store. Mirrors SearchIndex layout but keeps vectors
    in a separate .embedding_index.json file so that existing keyword logic is
    untouched. Lazy: only files added via `add_file` (or the reindex helper)
    ever get embeddings.
    """

    def __init__(self, speicher_path):
        self.speicher = speicher_path
        self.memory_dir = os.path.join(speicher_path, 'memory')
        self.index_file = os.path.join(speicher_path, '.embedding_index.json')
        self.entries = {}   # filename -> {mtime, chunks:[{text, embedding}]}
        self._lock = threading.Lock()
        self._loaded = False

    def _load(self):
        if self._loaded:
            return
        if os.path.exists(self.index_file):
            try:
                with open(self.index_file, 'r', encoding='utf-8') as f:
                    self.entries = json.load(f)
            except Exception:
                self.entries = {}
        self._loaded = True

    def _save(self):
        try:
            os.makedirs(os.path.dirname(self.index_file), exist_ok=True)
            with open(self.index_file, 'w', encoding='utf-8') as f:
                json.dump(self.entries, f, ensure_ascii=False)
        except Exception as e:
            print(f"[embedding_index] save error: {e}")

    def _resolve_path(self, fname):
        p = os.path.join(self.memory_dir, fname)
        if os.path.exists(p):
            return p
        p = os.path.join(self.speicher, fname)
        if os.path.exists(p):
            return p
        return None

    def add_file(self, fname, text=None, api_key=None):
        """Embed a single file (idempotent by mtime). Returns True on success."""
        self._load()
        fpath = self._resolve_path(fname)
        if not fpath:
            return False
        try:
            mtime = os.path.getmtime(fpath)
        except Exception:
            return False
        existing = self.entries.get(fname)
        if existing and existing.get('mtime', 0) >= mtime:
            return True  # already fresh

        if text is None:
            ext = os.path.splitext(fname.lower())[1]
            if ext in BINARY_EXTS:
                return False
            try:
                with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                    text = f.read(30000)
            except Exception:
                return False
        if not text or not text.strip():
            return False

        chunks = _chunk_text(text)
        if not chunks:
            return False

        key = api_key or _get_openai_api_key()
        vecs = _call_openai_embedding(chunks, key)
        if not vecs or len(vecs) != len(chunks):
            return False

        entry = {
            'mtime': mtime,
            'path': fpath,
            'chunks': [
                {'text': c, 'embedding': v}
                for c, v in zip(chunks, vecs)
            ],
        }
        with self._lock:
            self.entries[fname] = entry
            self._save()
        return True

    def search(self, query_vec, max_results=20):
        """Cosine similarity search. Returns list of
        {'name', 'chunk', 'score', 'chunk_index'}.
        """
        self._load()
        if not self.entries or not query_vec:
            return []
        qv = query_vec
        q_norm = math.sqrt(sum(x * x for x in qv)) or 1.0
        scored = []
        for fname, entry in self.entries.items():
            for i, chunk in enumerate(entry.get('chunks', [])):
                v = chunk.get('embedding')
                if not v:
                    continue
                dot = 0.0
                v_norm_sq = 0.0
                for a, b in zip(qv, v):
                    dot += a * b
                    v_norm_sq += b * b
                v_norm = math.sqrt(v_norm_sq) or 1.0
                sim = dot / (q_norm * v_norm)
                scored.append({
                    'name': fname,
                    'chunk': chunk.get('text', ''),
                    'score': sim,
                    'chunk_index': i,
                })
        scored.sort(key=lambda x: x['score'], reverse=True)
        return scored[:max_results]


_embedding_index_cache = {}
_embedding_index_lock = threading.Lock()


def get_embedding_index(speicher_path):
    with _embedding_index_lock:
        idx = _embedding_index_cache.get(speicher_path)
        if idx is None:
            idx = EmbeddingIndex(speicher_path)
            _embedding_index_cache[speicher_path] = idx
    return idx


def index_file_with_embedding(speicher_path, fname, text=None):
    """Convenience hook: called alongside SearchIndex.add_file so that newly
    arrived e-mails/web-clips/files get both a keyword and a semantic entry.
    Silent on failure (OpenAI down / no key / quota)."""
    try:
        idx = get_embedding_index(speicher_path)
        return idx.add_file(fname, text=text)
    except Exception as e:
        print(f"[embedding_index] add_file error: {e}")
        return False


# ─── QUERY EXPANSION ─────────────────────────────────────────────────────────

def _call_llm_for_text(provider, model_id, api_key, system_prompt, user_prompt,
                      max_tokens=300, timeout=None):
    """Minimal LLM call used for query expansion and compression. Uses urllib
    so we stay dependency-free. Returns the assistant text or None on failure.
    """
    timeout = timeout or _LLM_HTTP_TIMEOUT
    try:
        if provider == 'anthropic':
            payload = json.dumps({
                'model': model_id,
                'max_tokens': max_tokens,
                'system': system_prompt,
                'messages': [{'role': 'user', 'content': user_prompt}],
            }).encode('utf-8')
            req = urllib.request.Request(
                'https://api.anthropic.com/v1/messages',
                data=payload,
                headers={
                    'x-api-key': api_key,
                    'anthropic-version': '2023-06-01',
                    'content-type': 'application/json',
                },
            )
            with urllib.request.urlopen(req, timeout=timeout) as resp:
                data = json.loads(resp.read().decode('utf-8'))
            parts = data.get('content', [])
            return parts[0].get('text', '') if parts else None
        elif provider == 'openai':
            payload = json.dumps({
                'model': model_id,
                'max_tokens': max_tokens,
                'messages': [
                    {'role': 'system', 'content': system_prompt},
                    {'role': 'user', 'content': user_prompt},
                ],
            }).encode('utf-8')
            req = urllib.request.Request(
                'https://api.openai.com/v1/chat/completions',
                data=payload,
                headers={
                    'Authorization': f'Bearer {api_key}',
                    'Content-Type': 'application/json',
                },
            )
            with urllib.request.urlopen(req, timeout=timeout) as resp:
                data = json.loads(resp.read().decode('utf-8'))
            choices = data.get('choices', [])
            return choices[0].get('message', {}).get('content') if choices else None
        elif provider == 'mistral':
            payload = json.dumps({
                'model': model_id,
                'max_tokens': max_tokens,
                'messages': [
                    {'role': 'system', 'content': system_prompt},
                    {'role': 'user', 'content': user_prompt},
                ],
            }).encode('utf-8')
            req = urllib.request.Request(
                'https://api.mistral.ai/v1/chat/completions',
                data=payload,
                headers={
                    'Authorization': f'Bearer {api_key}',
                    'Content-Type': 'application/json',
                },
            )
            with urllib.request.urlopen(req, timeout=timeout) as resp:
                data = json.loads(resp.read().decode('utf-8'))
            choices = data.get('choices', [])
            return choices[0].get('message', {}).get('content') if choices else None
    except Exception as e:
        print(f"[llm] {provider} call failed: {e}")
        return None
    return None


def expand_query(query, speicher_path=None, n_variants=2):
    """Return a list of search queries: [original, variant1, variant2, ...].

    If the LLM is unreachable, returns [original] so downstream code still
    runs with the plain query.
    """
    original = (query or '').strip()
    if not original:
        return []
    provider, model_id, api_key = _get_agent_llm(speicher_path or '')
    if not provider:
        return [original]

    system_prompt = (
        "Du bist ein Retrieval-Assistent. Der Nutzer sucht in einer persoenlichen "
        f"Wissensbasis (E-Mails, Web-Clips, Dateien). Formuliere GENAU {n_variants} "
        "alternative, praegnante Such-Formulierungen zur Nutzer-Anfrage. Nutze "
        "Synonyme, Entitaeten, verwandte Begriffe. Pro Zeile eine Variante, keine "
        "Nummerierung, keine Erklaerungen."
    )
    user_prompt = f"Anfrage: {original}"
    out = _call_llm_for_text(provider, model_id, api_key, system_prompt,
                             user_prompt, max_tokens=200)
    if not out:
        return [original]
    variants = []
    for line in out.splitlines():
        line = line.strip(' -*•\t0123456789.)')
        if line and line.lower() != original.lower():
            variants.append(line)
        if len(variants) >= n_variants:
            break
    return [original] + variants


# ─── RRF FUSION ──────────────────────────────────────────────────────────────

def rrf_fuse(ranked_lists, k=_RRF_K, top_n=20):
    """Reciprocal Rank Fusion.

    ranked_lists: iterable of lists. Each list contains either strings
        (document IDs / filenames) or dicts with a 'name' key.
    Returns a ranked list of {'name', 'rrf_score', 'sources'} entries.
    """
    scores = {}
    sources = {}
    for lst in ranked_lists:
        if not lst:
            continue
        seen_in_list = set()
        for rank, item in enumerate(lst, start=1):
            key = item if isinstance(item, str) else item.get('name')
            if not key or key in seen_in_list:
                continue
            seen_in_list.add(key)
            scores[key] = scores.get(key, 0.0) + 1.0 / (k + rank)
            sources.setdefault(key, []).append(rank)
    fused = [
        {'name': name, 'rrf_score': score, 'sources': sources[name]}
        for name, score in scores.items()
    ]
    fused.sort(key=lambda x: x['rrf_score'], reverse=True)
    return fused[:top_n]


# ─── CONTEXTUAL COMPRESSION ──────────────────────────────────────────────────

def compress_chunk(chunk_text, query, speicher_path=None, max_tokens=220):
    """Ask the LLM to distil the chunk down to the sentences that actually
    answer the query. Falls back to the original chunk (truncated) on failure.
    """
    if not chunk_text:
        return ''
    provider, model_id, api_key = _get_agent_llm(speicher_path or '')
    if not provider:
        return chunk_text[:1200]
    system_prompt = (
        "Du erhaeltst einen Textabschnitt und eine Nutzer-Anfrage. Gib NUR die "
        "Saetze zurueck, die fuer die Anfrage wirklich relevant sind — wortgetreu, "
        "keine Zusammenfassung, keine Kommentare. Gibt es nichts Relevantes, "
        "antworte mit dem leeren String."
    )
    user_prompt = f"Anfrage: {query}\n\nAbschnitt:\n{chunk_text[:4000]}"
    out = _call_llm_for_text(provider, model_id, api_key, system_prompt,
                             user_prompt, max_tokens=max_tokens)
    if out is None:
        return chunk_text[:1200]
    out = out.strip()
    return out or chunk_text[:1200]


# ─── HYBRID RAG SEARCH (top-level orchestrator) ──────────────────────────────

def hybrid_rag_search(query, speicher_path, max_results=5,
                      fuse_top=20, compress=True, n_variants=2,
                      recency_first=None):
    """Full RAG pipeline:
        1) expand query into {original + variants} (skipped if n_variants=0)
        2) parallel keyword search (existing HybridSearch) + semantic search
        3) RRF fusion -> top `fuse_top`
        4) take first `max_results`, optionally contextually compress

    recency_first: if True, after RRF the top items are reordered so the newest
        date wins, with the RRF score acting as tie-break. If None (default),
        the parsed intent of `query` decides.

    Returns dict with:
        queries:    list of all expanded queries
        semantic:   bool, whether semantic search contributed
        results:    list of {name, snippet, compressed, score, sources}
        fallback:   reason string if the hybrid path degraded to keyword-only
    """
    out = {'queries': [query], 'semantic': False, 'results': [], 'fallback': None}
    if not query or not speicher_path:
        return out

    if n_variants and n_variants > 0:
        queries = expand_query(query, speicher_path=speicher_path, n_variants=n_variants)
    else:
        queries = [query]
    out['queries'] = queries

    # --- Keyword side: run existing HybridSearch per query variant ---
    # We use force_search=True so every word is treated as a candidate person
    # name — that's what makes "ExFlow Rechnung" find a PDF named
    # "Rechnung_*.pdf". The recency narrowing below uses the *normal* parse
    # (person_names only) so that step still distinguishes real people from
    # ordinary query terms.
    keyword_lists = []
    for q in queries:
        try:
            intent = QueryParser.parse(q, force_search=True)
        except TypeError:
            intent = QueryParser.parse(q)
            intent.is_search = True
        try:
            results, _ = HybridSearch.search(intent, speicher_path, max_results=30)
            keyword_lists.append([r['name'] for r in results])
        except Exception as e:
            print(f"[rag] keyword search failed for '{q}': {e}")
            keyword_lists.append([])

    # --- Semantic side: embed each query, search embedding index ---
    semantic_lists = []
    chunk_by_key = {}   # fname -> best chunk text seen
    emb_idx = get_embedding_index(speicher_path)
    emb_idx._load()
    if emb_idx.entries:
        api_key = _get_openai_api_key()
        q_vecs = _call_openai_embedding(queries, api_key) if api_key else None
        if q_vecs:
            out['semantic'] = True
            for qv in q_vecs:
                hits = emb_idx.search(qv, max_results=30)
                semantic_lists.append([h['name'] for h in hits])
                for h in hits:
                    prev = chunk_by_key.get(h['name'])
                    if prev is None or h['score'] > prev[1]:
                        chunk_by_key[h['name']] = (h['chunk'], h['score'])
        else:
            out['fallback'] = 'embedding_unavailable'
    else:
        out['fallback'] = 'no_embedding_index'

    # --- RRF fusion ---
    fused = rrf_fuse(keyword_lists + semantic_lists, k=_RRF_K, top_n=fuse_top)

    # Decide recency mode: explicit override > intent of original query
    if recency_first is None:
        try:
            recency_first = bool(getattr(QueryParser.parse(query), 'recency_first', False))
        except Exception:
            recency_first = False

    # Score-first sort with date tie-break.
    fused.sort(
        key=lambda it: (it['rrf_score'], extract_date_from_name(it['name'])),
        reverse=True,
    )
    if recency_first:
        # Take the top-N RRF candidates as a relevance pool, then re-sort that
        # pool by date so the freshest *relevant* item lands at position 1.
        # When the query mentions a person, narrow the pool to hits that
        # actually contain that person's name (in filename/snippet) — otherwise
        # the freshest *unrelated* mail of the day wins.
        master_intent_for_recency = QueryParser.parse(query)
        person_names = [
            normalize_unicode(p).lower() for p in master_intent_for_recency.person_names
        ]
        pool_size = max(max_results + 3, 10)
        pool = fused[:pool_size]
        if person_names:
            narrowed = []
            for it in pool:
                hay = normalize_unicode(it['name']).lower()
                if any(pn in hay for pn in person_names):
                    narrowed.append(it)
                    continue
                # Cheap snippet check — read the file head if we have nothing else.
                snippet_blob = chunk_by_key.get(it['name'], (None, 0))[0]
                if snippet_blob and any(pn in normalize_unicode(snippet_blob).lower()
                                        for pn in person_names):
                    narrowed.append(it)
            if narrowed:
                pool = narrowed
        pool.sort(
            key=lambda it: (extract_date_from_name(it['name']), it['rrf_score']),
            reverse=True,
        )
        top = pool[:max_results]
    else:
        top = fused[:max_results]

    # --- Build output: prefer semantic chunk if we have one, else file head ---
    results = []
    memory_dir = os.path.join(speicher_path, 'memory')
    for item in top:
        fname = item['name']
        snippet = ''
        if fname in chunk_by_key:
            snippet = chunk_by_key[fname][0]
        else:
            fpath = os.path.join(memory_dir, fname)
            if not os.path.exists(fpath):
                fpath = os.path.join(speicher_path, fname)
            try:
                with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                    snippet = f.read(4000)
            except Exception:
                snippet = ''
        compressed = compress_chunk(snippet, query, speicher_path) if compress else snippet
        results.append({
            'name': fname,
            'snippet': snippet,
            'compressed': compressed,
            'score': item['rrf_score'],
            'sources': item['sources'],
        })

    out['results'] = results
    return out


def reindex_embeddings_async(speicher_path, limit=None):
    """Run `reindex_embeddings` in a background thread.
    Safe to call on startup; no-op if OpenAI key missing.
    """
    def _run():
        try:
            res = reindex_embeddings(speicher_path, limit=limit)
            if res.get('ok') and res.get('indexed'):
                print(f"[embedding_index] {speicher_path}: indexed={res['indexed']} "
                      f"skipped={res.get('skipped', 0)} failed={res.get('failed', 0)}")
        except Exception as e:
            print(f"[embedding_index] reindex error for {speicher_path}: {e}")
    t = threading.Thread(target=_run, daemon=True)
    t.start()


def reindex_all_embeddings_async(limit_per_agent=None):
    """Backfill semantic embeddings for every known agent in the background.
    Each agent runs in its own thread so one slow agent doesn't block others."""
    for p in _list_agent_speicher_paths():
        reindex_embeddings_async(p, limit=limit_per_agent)


def reindex_embeddings(speicher_path, limit=None):
    """Generate embeddings for all files already present in the keyword index
    of this agent. Safe to re-run — existing fresh entries are skipped."""
    kw = SearchIndex(speicher_path)
    kw._load_index()
    if not kw.entries:
        kw.build_or_update()

    emb = get_embedding_index(speicher_path)
    emb._load()
    api_key = _get_openai_api_key()
    if not api_key:
        return {'ok': False, 'reason': 'no_openai_key', 'indexed': 0}

    done = 0
    skipped = 0
    failed = 0
    for fname in list(kw.entries.keys()):
        if limit and done >= limit:
            break
        existing = emb.entries.get(fname)
        fpath = emb._resolve_path(fname)
        if not fpath:
            continue
        try:
            mtime = os.path.getmtime(fpath)
        except Exception:
            continue
        if existing and existing.get('mtime', 0) >= mtime:
            skipped += 1
            continue
        ok = emb.add_file(fname, api_key=api_key)
        if ok:
            done += 1
        else:
            failed += 1
    return {'ok': True, 'indexed': done, 'skipped': skipped, 'failed': failed,
            'total_in_keyword_index': len(kw.entries)}


# ═══════════════════════════════════════════════════════════════════════════════
# CLI — Build indexes for all agents
# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == '__main__':
    BASE = os.path.expanduser(
        "~/Library/Mobile Documents/com~apple~CloudDocs/Downloads shared/claude_datalake"
    )
    agents_dir = os.path.join(BASE, 'config', 'agents')

    if not os.path.exists(agents_dir):
        print(f"Agents-Verzeichnis nicht gefunden: {agents_dir}")
        exit(1)

    print("Baue Search-Index fuer alle Agenten...\n")
    total_files = 0
    start = time.time()

    for fname in os.listdir(agents_dir):
        if not fname.endswith('.txt') or '_' in fname:
            continue  # Skip sub-agents
        agent = fname.replace('.txt', '')
        speicher_path = os.path.join(BASE, agent)
        memory_dir = os.path.join(speicher_path, 'memory')

        if os.path.exists(memory_dir):
            idx = SearchIndex(speicher_path)
            count = idx.build_index()
            total_files += count
        else:
            print(f"  {agent}: kein Memory-Ordner")

    elapsed = time.time() - start
    print(f"\nAgent-Indexes: {total_files} Dateien indexiert in {elapsed:.1f}s")

    # Build global index
    print("\nBaue globalen Index...")
    gidx = GlobalSearchIndex()
    global_count = gidx.build_global_index()
    print(f"Globaler Index: {global_count} Dateien")
