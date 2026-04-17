import os
import json
import datetime
import threading
import uuid
import copy
import re
import signal
import sys
try:
    import setproctitle
    setproctitle.setproctitle("AssistantDev WebServer")
except ImportError:
    pass

VALID_PROVIDERS = {'anthropic', 'openai', 'mistral', 'gemini', 'perplexity'}

PENDING_MARKER = '[ANTWORT AUSSTEHEND - Server-Neustart hat diese Antwort unterbrochen]'
RECOVERY_MARKER = '[Antwort verloren - Server wurde neu gestartet]'

_shutdown_event = threading.Event()
_active_requests = 0
_active_requests_lock = threading.Lock()

# Search engine
try:
    from search_engine import auto_search, format_search_feedback, build_index_async, build_global_index_async, detect_global_trigger, update_all_indexes, reindex_all_embeddings_async
except ImportError:
    auto_search = None
    format_search_feedback = None
    build_index_async = None
    build_global_index_async = None
    detect_global_trigger = None
    update_all_indexes = None
    reindex_all_embeddings_async = None
import requests
from flask import Flask, request, jsonify, render_template_string, make_response
from bs4 import BeautifulSoup

try:
    from capabilities_template import inject_capabilities_on_startup
except ImportError:
    inject_capabilities_on_startup = None

# ─── IMAGE HELPERS ───────────────────────────────────────────────────────────
# Anthropic API limitiert Bild-Dimensionen auf max. 8000 px pro Seite.
# Groessere Bilder muessen vor dem Senden herunterskaliert werden.
ANTHROPIC_MAX_IMAGE_DIM = 8000
_IMAGE_SAFE_DIM = 7900  # kleine Sicherheits-Reserve


def downscale_image_b64_if_needed(b64_data, media_type):
    """Nimmt base64-kodierte Bilddaten und skaliert herunter, falls eine
    Seite > ANTHROPIC_MAX_IMAGE_DIM. Gibt (neues_b64, neuer_media_type) zurueck.
    Bei nicht noetiger Skalierung: Original-Daten zurueck.
    Bei Decode-Fehlern: (None, None) — Caller verwirft das Bild, damit
    Anthropic die Anfrage nicht wegen Dimensions-Verstoss ablehnt.
    """
    if not b64_data:
        return b64_data, media_type
    try:
        import base64 as _b64
        from PIL import Image as _PILImage, ImageFile as _PILImageFile
        import io as _io

        _PILImage.MAX_IMAGE_PIXELS = None
        _PILImageFile.LOAD_TRUNCATED_IMAGES = True

        raw = _b64.b64decode(b64_data)
        img = _PILImage.open(_io.BytesIO(raw))
        w, h = img.size
        if w <= ANTHROPIC_MAX_IMAGE_DIM and h <= ANTHROPIC_MAX_IMAGE_DIM:
            return b64_data, media_type

        img.thumbnail((_IMAGE_SAFE_DIM, _IMAGE_SAFE_DIM), _PILImage.LANCZOS)

        fmt_map = {
            'image/png': ('PNG', 'image/png'),
            'image/jpeg': ('JPEG', 'image/jpeg'),
            'image/jpg': ('JPEG', 'image/jpeg'),
            'image/webp': ('WEBP', 'image/webp'),
            'image/gif': ('PNG', 'image/png'),  # GIF -> PNG nach Resize
        }
        out_fmt, out_mime = fmt_map.get((media_type or '').lower(), ('PNG', 'image/png'))

        if out_fmt == 'JPEG' and img.mode in ('RGBA', 'LA', 'P'):
            img = img.convert('RGB')

        buf = _io.BytesIO()
        save_kwargs = {}
        if out_fmt == 'JPEG':
            save_kwargs['quality'] = 90
            save_kwargs['optimize'] = True
        img.save(buf, format=out_fmt, **save_kwargs)
        new_b64 = _b64.b64encode(buf.getvalue()).decode('utf-8')
        print(f"[IMG-RESIZE] {w}x{h} -> {img.size[0]}x{img.size[1]} ({media_type} -> {out_mime})", flush=True)
        return new_b64, out_mime
    except Exception as _img_ex:
        print(f"[IMG-RESIZE] Fehler beim Downscalen, Bild wird verworfen: {_img_ex}", flush=True)
        return None, None


# ─── FILE CREATION HELPERS ───────────────────────────────────────────────────

OUTPUT_DIR = os.path.join(os.path.expanduser("~/Library/Mobile Documents/com~apple~CloudDocs/Downloads shared"), "claude_outputs")

def ensure_output_dir():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    return OUTPUT_DIR

def sanitize_llm_json(raw):
    """Sanitize JSON from LLM output that may have single quotes, trailing commas, etc."""
    import re as _sre
    s = raw.strip()
    # Remove markdown code fences
    if s.startswith('```'):
        s = _sre.sub(r'^```\w*\n?', '', s)
        s = _sre.sub(r'\n?```$', '', s)
        s = s.strip()
    # Try standard JSON first
    try:
        import json as _sjson
        return _sjson.loads(s)
    except Exception:
        pass
    # Fix single quotes → double quotes (careful with apostrophes in text)
    # Strategy: replace single-quoted keys and values
    try:
        import ast
        parsed = ast.literal_eval(s)
        if isinstance(parsed, dict):
            return parsed
        # ast.literal_eval can produce sets/tuples/lists — only dicts are valid here
    except Exception:
        pass
    # Manual fixes: trailing commas before } or ]
    s2 = _sre.sub(r',\s*([}\]])', r'\1', s)
    # Replace single quotes with double quotes (simple approach)
    s2 = s2.replace("'", '"')
    try:
        import json as _sjson
        return _sjson.loads(s2)
    except Exception:
        pass
    # Last resort: raise the original error
    import json as _sjson
    return _sjson.loads(raw)


def create_docx_from_spec(spec):
    from docx import Document
    from docx.shared import Pt, RGBColor
    doc = Document()
    title = spec.get('title', 'Dokument')
    doc.add_heading(title, 0)
    for item in spec.get('content', []):
        t = item.get('type', 'paragraph')
        text = item.get('text', '')
        if t == 'heading':
            doc.add_heading(text, level=item.get('level', 1))
        elif t == 'paragraph':
            doc.add_paragraph(text)
        elif t == 'bullet':
            doc.add_paragraph(text, style='List Bullet')
        elif t == 'table':
            rows = item.get('rows', [])
            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                table.style = 'Table Grid'
                for i, row in enumerate(rows):
                    for j, cell_text in enumerate(row):
                        table.cell(i, j).text = str(cell_text)
    fname = title.replace(' ', '_').replace('/', '_')[:50] + '.docx'
    fpath = os.path.join(ensure_output_dir(), fname)
    doc.save(fpath)
    return fname, fpath

def create_xlsx_from_spec(spec):
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = openpyxl.Workbook()
    title = spec.get('title', 'Tabelle')
    sheets = spec.get('sheets', [{'name': 'Sheet1', 'rows': spec.get('rows', [])}])
    wb.remove(wb.active)
    for sheet_spec in sheets:
        ws = wb.create_sheet(sheet_spec.get('name', 'Sheet'))
        rows = sheet_spec.get('rows', [])
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = ws.cell(row=i+1, column=j+1, value=val)
                if i == 0:  # Header row bold
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill(start_color='2D3A2D', end_color='2D3A2D', fill_type='solid')
                    cell.font = Font(bold=True, color='F0F0F0')
    fname = title.replace(' ', '_').replace('/', '_')[:50] + '.xlsx'
    fpath = os.path.join(ensure_output_dir(), fname)
    wb.save(fpath)
    return fname, fpath

def create_pdf_from_spec(spec):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors
    title = spec.get('title', 'Dokument')
    fname = title.replace(' ', '_').replace('/', '_')[:50] + '.pdf'
    fpath = os.path.join(ensure_output_dir(), fname)
    doc = SimpleDocTemplate(fpath, pagesize=A4,
                            rightMargin=2*cm, leftMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph(title, styles['Title']))
    story.append(Spacer(1, 0.5*cm))
    for item in spec.get('content', []):
        t = item.get('type', 'paragraph')
        text = item.get('text', '')
        if t == 'heading':
            level = item.get('level', 1)
            style = styles['Heading' + str(min(level, 3))]
            story.append(Paragraph(text, style))
        elif t == 'paragraph':
            story.append(Paragraph(text, styles['Normal']))
            story.append(Spacer(1, 0.3*cm))
        elif t == 'bullet':
            story.append(Paragraph('• ' + text, styles['Normal']))
        elif t == 'table':
            rows = item.get('rows', [])
            if rows:
                t_obj = Table(rows)
                t_obj.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#2D3A2D')),
                    ('TEXTCOLOR', (0,0), (-1,0), colors.white),
                    ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.grey),
                    ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#F5F5F5')]),
                ]))
                story.append(t_obj)
                story.append(Spacer(1, 0.3*cm))
    doc.build(story)
    return fname, fpath

def create_pptx_from_spec(spec):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    title_str = spec.get('title', 'Praesentation')

    # Dark theme colors
    BG_COLOR    = RGBColor(0x1a, 0x1a, 0x1a)
    GOLD_COLOR  = RGBColor(0xf0, 0xc0, 0x60)
    WHITE_COLOR = RGBColor(0xf0, 0xf0, 0xf0)
    GRAY_COLOR  = RGBColor(0x88, 0x88, 0x88)

    slide_width  = prs.slide_width
    slide_height = prs.slide_height

    def set_bg(slide):
        from pptx.oxml.ns import qn
        from lxml import etree
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_textbox(slide, text, left, top, width, height,
                    size=18, bold=False, color=WHITE_COLOR, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox

    slides_spec = spec.get('slides', [])

    for i, slide_spec in enumerate(slides_spec):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        set_bg(slide)

        stype   = slide_spec.get('type', 'content')
        heading = slide_spec.get('heading', '')
        body    = slide_spec.get('body', '')
        bullets = slide_spec.get('bullets', [])
        footer  = slide_spec.get('footer', '')

        margin  = Inches(0.5)
        w       = slide_width - 2 * margin

        if stype == 'title':
            # Big title slide
            add_textbox(slide, title_str,
                        margin, Inches(1.8), w, Inches(1.2),
                        size=36, bold=True, color=WHITE_COLOR, align=PP_ALIGN.CENTER)
            if heading:
                add_textbox(slide, heading,
                            margin, Inches(3.2), w, Inches(0.8),
                            size=22, color=GOLD_COLOR, align=PP_ALIGN.CENTER)
            # Gold rule
            from pptx.util import Pt as PtU
            line = slide.shapes.add_shape(
                1, margin, Inches(3.0), w, Emu(40000))
            line.fill.solid()
            line.fill.fore_color.rgb = GOLD_COLOR
            line.line.fill.background()
        else:
            # Heading bar
            bar = slide.shapes.add_shape(
                1, 0, 0, slide_width, Inches(0.95))
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(0x2d, 0x3a, 0x2d)
            bar.line.fill.background()

            if heading:
                add_textbox(slide, heading,
                            margin, Emu(120000), w, Inches(0.75),
                            size=24, bold=True, color=WHITE_COLOR)

            y = Inches(1.1)

            if body:
                add_textbox(slide, body,
                            margin, y, w, Inches(1.2),
                            size=16, color=WHITE_COLOR)
                y += Inches(1.3)

            for bullet in bullets:
                add_textbox(slide, '– ' + bullet,
                            margin + Inches(0.2), y, w - Inches(0.2), Inches(0.5),
                            size=15, color=RGBColor(0xd0, 0xd0, 0xd0))
                y += Inches(0.52)

        # Slide number
        add_textbox(slide, str(i + 1),
                    slide_width - Inches(0.8), slide_height - Inches(0.4),
                    Inches(0.6), Inches(0.3),
                    size=10, color=GRAY_COLOR, align=PP_ALIGN.RIGHT)

        if footer:
            add_textbox(slide, footer,
                        margin, slide_height - Inches(0.4), w, Inches(0.3),
                        size=10, color=GRAY_COLOR)

    fname = title_str.replace(' ', '_').replace('/', '_')[:50] + '.pptx'
    fpath = os.path.join(ensure_output_dir(), fname)
    prs.save(fpath)
    return fname, fpath

def send_email_draft(spec):
    import subprocess
    to      = spec.get('to', '')
    subject = spec.get('subject', '')
    body    = spec.get('body', '')
    cc      = spec.get('cc', '')
    sender  = spec.get('from', '')

    # Escape for AppleScript - handle all special chars
    def esc(s):
        s = s.replace('\\', '\\\\')
        s = s.replace('"', '\\"')
        s = s.replace('\n', '\\n')
        s = s.replace('\r', '')
        return s

    cc_line = f'\n        make new cc recipient at end of cc recipients with properties {{address:\"{esc(cc)}\"}}' if cc else ''
    sender_line = f'\n    set sender of newMessage to "{esc(sender)}"' if sender else ''

    script = f'''tell application "Mail"
    set newMessage to make new outgoing message with properties {{subject:"{esc(subject)}", content:"{esc(body)}", visible:true}}
    tell newMessage
        make new to recipient at end of to recipients with properties {{address:"{esc(to)}"}}{cc_line}
    end tell{sender_line}
    activate
end tell'''

    subprocess.Popen(['osascript', '-e', script],
                     stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return True

BASE = os.path.expanduser("~/Library/Mobile Documents/com~apple~CloudDocs/Downloads shared/claude_datalake")
AGENTS_DIR = os.path.join(BASE, "config/agents")
MODELS_FILE = os.path.join(BASE, "config/models.json")

app = Flask(__name__)

def cleanup_agent_files():
    """Clean old memory content from agent txt files on startup."""
    STRIP_MARKERS = [
        "--- GEDAECHTNIS:", "--- DATEI-ERSTELLUNG ---",
        "--- WEITERE FAEHIGKEITEN ---", "VERGANGENE KONVERSATIONEN:",
        "GESPEICHERTE DATEIEN IM MEMORY", "Von: ", "An: ", "Betreff: ",
    ]
    if not os.path.exists(AGENTS_DIR):
        return
    for fname in os.listdir(AGENTS_DIR):
        if not fname.endswith('.txt'):
            continue
        fpath = os.path.join(AGENTS_DIR, fname)
        try:
            with open(fpath, 'r', errors='ignore') as f:
                raw = f.read()
            cut = len(raw)
            for marker in STRIP_MARKERS:
                pos = raw.find(marker)
                if 0 < pos < cut:
                    cut = pos
            clean = raw[:cut].strip()
            if len(clean) < len(raw) - 10:
                with open(fpath, 'w') as f:
                    f.write(clean)
                print(f"Cleaned: {fname}")
        except Exception as e:
            print(f"Cleanup error {fname}: {e}")

cleanup_agent_files()

# ─── DYNAMIC CAPABILITIES: Agent-Prompts beim Start aktualisieren ─────────────
def _inject_capabilities():
    if inject_capabilities_on_startup is None:
        return
    try:
        inject_capabilities_on_startup(
            agents_dir=AGENTS_DIR,
            models_file=MODELS_FILE,
            datalake_base=BASE,
        )
    except Exception as e:
        print(f"[CAPABILITIES] Startup-Injection fehlgeschlagen: {e}")

_inject_capabilities()

# ─── RECOVERY: Pending-Marker aus abgebrochenen Sessions ersetzen ─────────────
def _recover_pending_markers():
    recovered = 0
    try:
        for agent_dir in os.listdir(BASE):
            agent_path = os.path.join(BASE, agent_dir)
            if not os.path.isdir(agent_path):
                continue
            for fname in os.listdir(agent_path):
                if not fname.startswith('konversation_') or not fname.endswith('.txt'):
                    continue
                fpath = os.path.join(agent_path, fname)
                try:
                    with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                    if PENDING_MARKER in content:
                        content = content.replace(PENDING_MARKER, RECOVERY_MARKER)
                        tmp = fpath + '.tmp'
                        with open(tmp, 'w', encoding='utf-8') as f:
                            f.write(content)
                        os.replace(tmp, fpath)
                        recovered += 1
                        print(f'[RECOVERY] Pending-Marker ersetzt in {agent_dir}/{fname}')
                except Exception as e:
                    print(f'[RECOVERY] Fehler bei {agent_dir}/{fname}: {e}')
    except Exception as e:
        print(f'[RECOVERY] Scan-Fehler: {e}')
    if recovered:
        print(f'[RECOVERY] {recovered} Konversation(en) mit Pending-Marker repariert')

_recover_pending_markers()

# ─── PROVIDER ADAPTERS ────────────────────────────────────────────────────────

def _sanitize_anthropic_images(messages):
    """Walk message list und downscale alle base64-Bilder, die die
    Anthropic-Dimensions-Grenze (8000 px) ueberschreiten. Mutiert die
    Strukturen in-place, damit gecachte Session-Verlaeufe ebenfalls
    dauerhaft korrigiert sind. Bilder, die nicht herunterskaliert werden
    koennen (Decode-Fehler / Decompression-Bomb), werden entfernt, damit
    die gesamte Anfrage nicht 400't.
    """
    for msg in messages:
        content = msg.get('content')
        if not isinstance(content, list):
            continue
        kept = []
        for part in content:
            if not isinstance(part, dict) or part.get('type') != 'image':
                kept.append(part)
                continue
            source = part.get('source') or {}
            if source.get('type') != 'base64':
                kept.append(part)
                continue
            new_b64, new_mime = downscale_image_b64_if_needed(
                source.get('data'), source.get('media_type')
            )
            if not new_b64:
                print("[IMG-SANITIZE] Oversize-Bild entfernt (Anthropic-Limit 8000 px)", flush=True)
                continue
            source['data'] = new_b64
            source['media_type'] = new_mime
            kept.append(part)
        msg['content'] = kept


def call_anthropic(api_key, model_id, system_prompt, messages):
    from anthropic import Anthropic
    _sanitize_anthropic_images(messages)
    client = Anthropic(api_key=api_key)
    r = client.messages.create(model=model_id, max_tokens=4096, system=system_prompt, messages=messages)
    return r.content[0].text

def call_openai(api_key, model_id, system_prompt, messages):
    import openai
    client = openai.OpenAI(api_key=api_key)
    oai_messages = [{"role": "system", "content": system_prompt}] + messages
    r = client.chat.completions.create(model=model_id, messages=oai_messages, max_tokens=4096)
    return r.choices[0].message.content

def call_perplexity(api_key, model_id, system_prompt, messages):
    import requests as _pplx_req
    from requests.exceptions import ReadTimeout, ConnectionError as ReqConnectionError

    # Perplexity requires strictly alternating user/assistant messages after system.
    # Merge consecutive same-role messages to enforce this.
    raw_msgs = list(messages)
    # Remove empty messages
    raw_msgs = [m for m in raw_msgs if m.get('content')]
    # Flatten list-type content (vision messages) to text
    for m in raw_msgs:
        if isinstance(m.get('content'), list):
            m = dict(m)
            m['content'] = ' '.join(
                p.get('text', '') for p in m['content']
                if isinstance(p, dict) and p.get('type') == 'text'
            )
    # Merge consecutive same-role messages
    merged = []
    for m in raw_msgs:
        content_text = m.get('content', '')
        if isinstance(content_text, list):
            content_text = ' '.join(
                p.get('text', '') for p in content_text
                if isinstance(p, dict) and p.get('type') == 'text'
            )
        if not content_text or not content_text.strip():
            continue
        if merged and merged[-1]['role'] == m['role']:
            merged[-1]['content'] += '\n\n' + content_text
        else:
            merged.append({'role': m['role'], 'content': content_text})
    # Ensure first non-system message is 'user' (Perplexity requirement)
    if merged and merged[0]['role'] != 'user':
        merged.insert(0, {'role': 'user', 'content': '.'})
    # Ensure last message is 'user' (Perplexity sends to model for completion)
    if merged and merged[-1]['role'] != 'user':
        merged.append({'role': 'user', 'content': '.'})
    pplx_messages = [{"role": "system", "content": system_prompt}] + merged
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {"model": model_id, "messages": pplx_messages, "max_tokens": 8000}

    # Model-specific timeouts (connect_timeout, read_timeout)
    PERPLEXITY_TIMEOUTS = {
        'sonar-deep-research': (10, 300),   # Deep Research: bis zu 5 Min
        'sonar-reasoning-pro': (10, 180),   # Reasoning Pro: bis zu 3 Min
        'sonar-reasoning': (10, 180),       # Reasoning: bis zu 3 Min
        'sonar-pro': (10, 120),             # Sonar Pro: 2 Min
        'sonar': (10, 120),                 # Sonar: 2 Min
    }
    timeout = PERPLEXITY_TIMEOUTS.get(model_id, (10, 120))

    MODEL_LABELS = {
        'sonar-deep-research': 'Sonar Deep Research',
        'sonar-reasoning-pro': 'Sonar Reasoning Pro',
        'sonar-reasoning': 'Sonar Reasoning',
        'sonar-pro': 'Sonar Pro',
        'sonar': 'Sonar',
    }
    model_label = MODEL_LABELS.get(model_id, model_id)

    try:
        resp = _pplx_req.post(
            "https://api.perplexity.ai/chat/completions",
            headers=headers, json=payload, timeout=timeout
        )
    except ReadTimeout:
        raise Exception(
            f"Perplexity {model_label} hat zu lange gebraucht (Timeout nach {timeout[1]}s). "
            f"Versuche es erneut oder waehle ein schnelleres Modell."
        )
    except ReqConnectionError:
        raise Exception(
            f"Verbindung zu Perplexity fehlgeschlagen. Pruefe deine Internetverbindung."
        )

    data = resp.json()
    if resp.status_code != 200:
        raise Exception(f"Perplexity API Fehler: {data.get('error', {}).get('message', str(data))}")
    text = data["choices"][0]["message"]["content"]
    # Append citations as clickable Markdown links if present
    citations = data.get("citations")
    if citations and isinstance(citations, list) and len(citations) > 0:
        text += "\n\n**Quellen:**\n"
        for i, url in enumerate(citations, 1):
            try:
                from urllib.parse import urlparse
                domain = urlparse(url).netloc.replace("www.", "")
            except Exception:
                domain = url[:60]
            text += f"[{i}] [{domain}]({url})\n"
    return text

def call_mistral(api_key, model_id, system_prompt, messages):
    headers = {"Authorization": "Bearer " + api_key, "Content-Type": "application/json"}
    mistral_messages = [{"role": "system", "content": system_prompt}] + messages
    payload = {"model": model_id, "messages": mistral_messages, "max_tokens": 4096}
    r = requests.post("https://api.mistral.ai/v1/chat/completions", headers=headers, json=payload, timeout=60)
    return r.json()["choices"][0]["message"]["content"]

def call_gemini(api_key, model_id, system_prompt, messages):
    """Gemini via REST API (kein SDK noetig)."""
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_id}:generateContent?key={api_key}"
    contents = []
    for m in messages:
        role = 'user' if m['role'] == 'user' else 'model'
        c = m['content']
        if isinstance(c, list):
            parts = []
            for item in c:
                if item.get('type') == 'text':
                    parts.append({'text': item['text']})
                elif item.get('type') == 'image':
                    parts.append({'inline_data': {'mime_type': item['source']['media_type'], 'data': item['source']['data']}})
            contents.append({'role': role, 'parts': parts})
        else:
            contents.append({'role': role, 'parts': [{'text': str(c)}]})
    payload = {'contents': contents, 'generationConfig': {'maxOutputTokens': 4096}}
    if system_prompt:
        payload['systemInstruction'] = {'parts': [{'text': system_prompt}]}
    try:
        r = requests.post(url, json=payload, timeout=120)
        data = r.json()
        if r.status_code != 200:
            err_msg = data.get('error', {}).get('message', str(data))
            if '429' in str(r.status_code) or 'quota' in err_msg.lower() or 'rate' in err_msg.lower():
                return ("Gemini Rate Limit erreicht. Entweder kurz warten "
                        "oder auf bezahlten API Plan upgraden: "
                        "https://ai.google.dev/pricing")
            raise Exception(f"Gemini API Fehler ({r.status_code}): " + err_msg)
        parts = data['candidates'][0]['content'].get('parts', [])
        return ''.join(p.get('text', '') for p in parts)
    except requests.exceptions.Timeout:
        raise Exception("Gemini API Timeout (120s)")

# Human-readable provider names
PROVIDER_DISPLAY = {
    'anthropic': 'Anthropic',
    'openai': 'OpenAI',
    'mistral': 'Mistral',
    'gemini': 'Google',
    'perplexity': 'Perplexity',
}

# Human-readable model names
MODEL_DISPLAY = {
    'claude-sonnet-4-6': 'Claude Sonnet 4.6',
    'claude-opus-4-6': 'Claude Opus 4.6',
    'claude-haiku-4-5': 'Claude Haiku 4.5',
    'claude-sonnet-4-20250514': 'Claude Sonnet 4',
    'claude-3-5-sonnet-20241022': 'Claude 3.5 Sonnet',
    'gpt-4o': 'GPT-4o',
    'gpt-4o-mini': 'GPT-4o Mini',
    'o1': 'o1',
    'o1-mini': 'o1 Mini',
    'mistral-large-latest': 'Mistral Large',
    'mistral-small-latest': 'Mistral Small',
    'open-mistral-nemo': 'Mistral Nemo',
    'gemini-2.0-flash': 'Gemini 2.0 Flash',
    'gemini-2.5-pro': 'Gemini 2.5 Pro',
    'gemini-2.5-flash': 'Gemini 2.5 Flash',
    'gemini-3-flash-preview': 'Gemini 3 Flash',
    'gemini-3-pro-preview': 'Gemini 3 Pro',
    'gemini-3.1-pro-preview': 'Gemini 3.1 Pro',
    'gemini-2.5-flash-lite': 'Gemini 2.5 Flash Lite',
    'gemini-1.5-pro': 'Gemini 1.5 Pro',
    'sonar': 'Sonar',
    'sonar-pro': 'Sonar Pro',
    'sonar-reasoning': 'Sonar Reasoning',
    'sonar-reasoning-pro': 'Sonar Reasoning Pro',
    'sonar-deep-research': 'Sonar Deep Research',
}

ADAPTERS = {"anthropic": call_anthropic, "openai": call_openai, "perplexity": call_perplexity, "mistral": call_mistral, "gemini": call_gemini}

# ─── STATE (see session-based state below) ────────────────────────────────────

def load_models():
    if os.path.exists(MODELS_FILE):
        with open(MODELS_FILE) as f:
            return json.load(f)
    return {"providers": {}}

# ─── FILE EXTRACTION ──────────────────────────────────────────────────────────

def extract_file_content(raw, filename):
    fname = filename.lower()
    if fname.endswith('.pdf'):
        try:
            import PyPDF2, io
            reader = PyPDF2.PdfReader(io.BytesIO(raw))
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        except Exception as e:
            return "[PDF Fehler: " + str(e) + "]"
    if fname.endswith('.docx'):
        try:
            import zipfile, io
            from xml.etree import ElementTree as ET
            z = zipfile.ZipFile(io.BytesIO(raw))
            xml = z.read('word/document.xml')
            tree = ET.fromstring(xml)
            ns = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
            return ' '.join(node.text for node in tree.iter(ns+'t') if node.text)
        except Exception as e:
            return "[DOCX Fehler: " + str(e) + "]"
    if fname.endswith('.xlsx'):
        try:
            import zipfile, io
            from xml.etree import ElementTree as ET
            z = zipfile.ZipFile(io.BytesIO(raw))
            strings = []
            if 'xl/sharedStrings.xml' in z.namelist():
                tree = ET.fromstring(z.read('xl/sharedStrings.xml'))
                for si in tree.iter('{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t'):
                    strings.append(si.text or '')
            result = []
            if 'xl/worksheets/sheet1.xml' in z.namelist():
                tree = ET.fromstring(z.read('xl/worksheets/sheet1.xml'))
                ns = '{http://schemas.openxmlformats.org/spreadsheetml/2006/main}'
                for row in tree.iter(ns+'row'):
                    vals = []
                    for cell in row.iter(ns+'c'):
                        v = cell.find(ns+'v')
                        if v is not None and v.text:
                            if cell.get('t','') == 's':
                                vals.append(strings[int(v.text)] if int(v.text) < len(strings) else '')
                            else:
                                vals.append(v.text)
                    if vals: result.append('\t'.join(vals))
            return '\n'.join(result)
        except Exception as e:
            return "[XLSX Fehler: " + str(e) + "]"
    if fname.endswith('.eml'):
        try:
            import email
            msg = email.message_from_bytes(raw)
            parts = ["Von: " + str(msg.get('From','')), "An: " + str(msg.get('To','')),
                     "Betreff: " + str(msg.get('Subject','')), "Datum: " + str(msg.get('Date','')), ""]
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        parts.append(part.get_payload(decode=True).decode('utf-8', errors='ignore'))
            else:
                parts.append(msg.get_payload(decode=True).decode('utf-8', errors='ignore'))
            return '\n'.join(parts)
        except Exception as e:
            return "[EML Fehler: " + str(e) + "]"
    if any(fname.endswith(ext) for ext in ['.mp4','.mov','.avi','.mkv','.mp3','.wav','.m4a']):
        return "[Mediendatei: " + filename + " — Inhalt kann nicht als Text gelesen werden.]"
    return raw.decode('utf-8', errors='ignore')

# ─── MEMORY SYSTEM ────────────────────────────────────────────────────────────

def load_index(speicher):
    """Load session index for an agent.

    Robust gegenueber korrupten _index.json (iCloud-Sync-Konflikte,
    Race-Conditions beim Schreiben). Bei JSONDecodeError / unerwartetem
    Typ: die defekte Datei wird als `.corrupt-<ts>.bak` umbenannt und
    ein leerer Index zurueckgegeben. `migrate_old_conversations()` baut
    den Index danach aus den vorhandenen `konversation_*.txt` neu auf.
    """
    index_file = os.path.join(speicher, '_index.json')
    if not os.path.exists(index_file):
        return []
    try:
        with open(index_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
        if not isinstance(data, list):
            raise ValueError(f'_index.json ist keine Liste, sondern {type(data).__name__}')
        return data
    except (json.JSONDecodeError, ValueError, OSError) as e:
        import datetime as _dt
        ts = _dt.datetime.now().strftime('%Y%m%d_%H%M%S')
        bak = index_file + f'.corrupt-{ts}.bak'
        try:
            os.rename(index_file, bak)
            print(f'[INDEX] KORRUPT: {index_file} ({e}). Umbenannt in {bak}. Rebuild leer.')
        except OSError as _mv_err:
            print(f'[INDEX] KORRUPT: {index_file} ({e}). Backup fehlgeschlagen: {_mv_err}')
        return []

def save_index(speicher, index):
    """Save session index."""
    index_file = os.path.join(speicher, '_index.json')
    with open(index_file, 'w') as f:
        json.dump(index, f, ensure_ascii=False, indent=2)

def migrate_old_conversations(speicher):
    """Import existing conversation files that aren't yet in the index."""
    index = load_index(speicher)
    indexed_files = {e.get('file','') for e in index}

    conv_files = sorted([
        f for f in os.listdir(speicher)
        if f.startswith('konversation_') and f.endswith('.txt')
    ])

    new_entries = []
    for fname in conv_files:
        if fname in indexed_files:
            continue
        fpath = os.path.join(speicher, fname)
        try:
            with open(fpath) as f:
                content = f.read()
            # Extract date from filename: konversation_2026-03-31_15-27.txt
            parts = fname.replace('konversation_','').replace('.txt','').split('_')
            if len(parts) >= 2:
                date_str = parts[0] + ' ' + parts[1].replace('-',':')
            else:
                date_str = fname
            # Build a short summary from content (first 800 chars)
            summary = "[Importiert] " + content[:400].replace('\n',' ').strip()
            new_entries.append({
                "date": date_str,
                "file": fname,
                "summary": summary,
                "referenced_files": []
            })
        except:
            pass

    if new_entries:
        # Insert old entries at beginning, sorted by date
        combined = new_entries + index
        save_index(speicher, combined)
        return len(new_entries)
    return 0

def summarize_conversation(verlauf, agent_prompt):
    """Ask Claude to summarize the current conversation in 2-3 sentences."""
    if not verlauf or len(verlauf) < 2:
        return None
    try:
        config = load_models()
        api_key = config['providers']['anthropic']['api_key']
        from anthropic import Anthropic
        client = Anthropic(api_key=api_key)
        conv_text = "\n".join(
            ("User: " if m['role'] == 'user' else "Assistant: ") + str(m['content'])[:500]
            for m in verlauf[-10:]  # last 10 messages max
        )
        r = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=200,
            system="Fasse die folgende Konversation in 2-3 praezisen Saetzen zusammen. Nur die Zusammenfassung, kein Kommentar.",
            messages=[{"role": "user", "content": conv_text}]
        )
        return r.content[0].text
    except:
        return None

def close_current_session(state=None):
    """Summarize and index the current session before switching agents."""
    if state is None:
        state = get_session()
    if not state['agent'] or not state['verlauf']:
        return
    summary = summarize_conversation(state['verlauf'], state['system_prompt'])
    if not summary:
        return
    speicher = state['speicher']
    index = load_index(speicher)
    entry = {
        "date": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
        "file": os.path.basename(state['dateiname']) if state['dateiname'] else "",
        "summary": summary,
        "referenced_files": list(state['session_files'])
    }
    index.append(entry)
    # Keep max 50 entries
    if len(index) > 50:
        index = index[-50:]
    save_index(speicher, index)

def build_memory_context(speicher, agent_name):
    """Compact memory index - no file contents, no full conversations."""
    index = load_index(speicher)
    parts = [f"\n\n--- GEDAECHTNIS: {agent_name.upper()} ---"]

    # Instructions for memory access
    parts.append("""
Du hast Zugriff auf ein persistentes Gedaechtnis mit vergangenen Konversationen und Dateien.
Um Inhalte zu laden: schreibe "Memory Folder Suche [Stichwort]" in deiner Nachricht.
""")

    # Compact session index - titles only, no content
    if index:
        parts.append("VERGANGENE SESSIONS (neueste zuerst):")
        for entry in reversed(index[-10:]):
            date = entry.get('date', '')
            summary = entry.get('summary', '')
            # Only first line, max 80 chars
            title = summary.strip().split('\n')[0][:80] if summary.strip() else entry.get('file', '')
            parts.append(f"  • {date}: {title}")

    # List memory files - names only, no content
    memory_dir = os.path.join(speicher, 'memory')
    if os.path.exists(memory_dir):
        # Exclude binary/image files from listing
        text_exts = {'.txt', '.docx', '.pdf', '.xlsx', '.csv', '.json', '.md', '.eml'}
        image_exts = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'}
        
        all_files = sorted(
            os.listdir(memory_dir),
            key=lambda f: os.path.getmtime(os.path.join(memory_dir, f)),
            reverse=True
        )
        
        doc_files = [f for f in all_files if any(f.lower().endswith(e) for e in text_exts)][:10]
        img_files = [f for f in all_files if any(f.lower().endswith(e) for e in image_exts)][:5]
        
        if doc_files:
            parts.append(f"\nDOKUMENTE IM MEMORY ({len(doc_files)} neueste):")
            for f in doc_files:
                parts.append(f"  - {f}")
        
        if img_files:
            parts.append(f"\nBILDER IM MEMORY ({len(img_files)}):")
            for f in img_files:
                parts.append(f"  - {f}")

    # Memory-Suche Instruktion
    parts.append("""
## Memory-Suche
Du kannst dein Memory gezielt durchsuchen. Schreibe dazu in deiner Antwort:
MEMORY_SEARCH: {"query": "Suchbegriff", "date_from": "YYYY-MM-DD", "date_to": "YYYY-MM-DD", "direction": "IN/OUT", "contact": "email-oder-name"}
Alle Felder ausser query sind optional.
Wenn der User nach E-Mails, Dokumenten oder vergangenen Informationen fragt, nutze MEMORY_SEARCH bevor du antwortest.
Expliziter Befehl vom User: /search [suchbegriff]

Beispiele:
- MEMORY_SEARCH: {"query": "ExFlow Rechnungen", "date_from": "2026-01-01"}
- MEMORY_SEARCH: {"query": "Onboarding", "direction": "IN", "contact": "thomas"}
""")

    parts.append("\n--- ENDE GEDAECHTNIS ---")
    return '\n'.join(parts)


def _get_wm_dir(agent_name):
    """Return the working_memory directory for this agent.
    Parent agents use <speicher>/working_memory/.
    Sub-agents get an isolated subdir <speicher>/working_memory/_<subname>/ so their
    working memory is NOT shared with the parent or other sub-agents.
    The broader agent memory (speicher itself) stays shared by design."""
    speicher = get_agent_speicher(agent_name) if agent_name else BASE
    parent = get_parent_agent(agent_name) if agent_name else None
    if parent:
        subname = agent_name.split('_', 1)[1]
        return os.path.join(speicher, 'working_memory', '_' + subname)
    return os.path.join(speicher, 'working_memory')


def load_working_memory(agent_name):
    """Load persistent working memory files for an agent from working_memory/ directory."""
    wm_dir = _get_wm_dir(agent_name)
    manifest_path = os.path.join(wm_dir, '_manifest.json')
    if not os.path.exists(manifest_path):
        return ''
    try:
        with open(manifest_path, 'r', encoding='utf-8') as f:
            manifest = json.load(f)
    except Exception:
        return ''
    max_tokens = manifest.get('max_tokens', 8000)
    files = manifest.get('files', [])
    if not files:
        return ''
    files_sorted = sorted(files, key=lambda x: (-x.get('priority', 5), x.get('added', '')))
    parts = []
    total_text = ''
    cleanup_warnings = []
    for entry in files_sorted:
        fname = entry.get('filename', '')
        fpath = os.path.join(wm_dir, fname)
        if not os.path.exists(fpath):
            continue
        try:
            with open(fpath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        except Exception:
            continue
        candidate = total_text + content
        if len(candidate) / 4 > max_tokens:
            cleanup_warnings.append(fname)
            continue
        total_text = candidate
        prio = entry.get('priority', 5)
        parts.append(f"[{fname}] (Prioritaet: {prio}):\n{content}")
    if not parts and not cleanup_warnings:
        return ''
    result = '\n\n--- WORKING MEMORY ---\n'
    result += '\n\n'.join(parts)
    for w in cleanup_warnings:
        result += f'\n[Working Memory Auto-Cleanup: {w} entfernt (Token-Limit erreicht)]'
    result += '\n--- ENDE WORKING MEMORY ---\n'
    if cleanup_warnings:
        remaining = [e for e in files if e.get('filename') not in cleanup_warnings]
        manifest['files'] = remaining
        try:
            tmp = manifest_path + '.tmp'
            with open(tmp, 'w', encoding='utf-8') as f:
                json.dump(manifest, f, ensure_ascii=False, indent=2)
            os.replace(tmp, manifest_path)
        except Exception:
            pass
    return result


def working_memory_add(agent_name, filename, content, priority=5, description=''):
    """Add a file to agent's working memory."""
    wm_dir = _get_wm_dir(agent_name)
    os.makedirs(wm_dir, exist_ok=True)
    manifest_path = os.path.join(wm_dir, '_manifest.json')
    if os.path.exists(manifest_path):
        with open(manifest_path, 'r', encoding='utf-8') as f:
            manifest = json.load(f)
    else:
        manifest = {'max_tokens': 8000, 'auto_cleanup': True, 'files': []}
    fpath = os.path.join(wm_dir, filename)
    tmp = fpath + '.tmp'
    with open(tmp, 'w', encoding='utf-8') as f:
        f.write(content)
    os.replace(tmp, fpath)
    manifest['files'] = [e for e in manifest['files'] if e.get('filename') != filename]
    manifest['files'].append({
        'filename': filename,
        'added': datetime.datetime.now().strftime('%Y-%m-%d'),
        'added_by': 'agent',
        'priority': priority,
        'description': description,
    })
    tmp_m = manifest_path + '.tmp'
    with open(tmp_m, 'w', encoding='utf-8') as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    os.replace(tmp_m, manifest_path)
    return manifest


def working_memory_remove(agent_name, filename):
    """Remove a file from agent's working memory."""
    wm_dir = _get_wm_dir(agent_name)
    manifest_path = os.path.join(wm_dir, '_manifest.json')
    if not os.path.exists(manifest_path):
        return None
    with open(manifest_path, 'r', encoding='utf-8') as f:
        manifest = json.load(f)
    manifest['files'] = [e for e in manifest['files'] if e.get('filename') != filename]
    fpath = os.path.join(wm_dir, filename)
    if os.path.exists(fpath):
        os.remove(fpath)
    tmp_m = manifest_path + '.tmp'
    with open(tmp_m, 'w', encoding='utf-8') as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    os.replace(tmp_m, manifest_path)
    return manifest


def working_memory_list(agent_name):
    """List all files in agent's working memory."""
    wm_dir = _get_wm_dir(agent_name)
    manifest_path = os.path.join(wm_dir, '_manifest.json')
    if not os.path.exists(manifest_path):
        return {'max_tokens': 8000, 'auto_cleanup': True, 'files': []}
    with open(manifest_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def ensure_output_dir():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    return OUTPUT_DIR

def create_docx_from_spec(spec):
    from docx import Document
    from docx.shared import Pt, RGBColor
    doc = Document()
    title = spec.get('title', 'Dokument')
    doc.add_heading(title, 0)
    for item in spec.get('content', []):
        t = item.get('type', 'paragraph')
        text = item.get('text', '')
        if t == 'heading':
            doc.add_heading(text, level=item.get('level', 1))
        elif t == 'paragraph':
            doc.add_paragraph(text)
        elif t == 'bullet':
            doc.add_paragraph(text, style='List Bullet')
        elif t == 'table':
            rows = item.get('rows', [])
            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                table.style = 'Table Grid'
                for i, row in enumerate(rows):
                    for j, cell_text in enumerate(row):
                        table.cell(i, j).text = str(cell_text)
    fname = title.replace(' ', '_').replace('/', '_')[:50] + '.docx'
    fpath = os.path.join(ensure_output_dir(), fname)
    doc.save(fpath)
    return fname, fpath

def create_xlsx_from_spec(spec):
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = openpyxl.Workbook()
    title = spec.get('title', 'Tabelle')
    sheets = spec.get('sheets', [{'name': 'Sheet1', 'rows': spec.get('rows', [])}])
    wb.remove(wb.active)
    for sheet_spec in sheets:
        ws = wb.create_sheet(sheet_spec.get('name', 'Sheet'))
        rows = sheet_spec.get('rows', [])
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = ws.cell(row=i+1, column=j+1, value=val)
                if i == 0:  # Header row bold
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill(start_color='2D3A2D', end_color='2D3A2D', fill_type='solid')
                    cell.font = Font(bold=True, color='F0F0F0')
    fname = title.replace(' ', '_').replace('/', '_')[:50] + '.xlsx'
    fpath = os.path.join(ensure_output_dir(), fname)
    wb.save(fpath)
    return fname, fpath

def create_pdf_from_spec(spec):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors
    title = spec.get('title', 'Dokument')
    fname = title.replace(' ', '_').replace('/', '_')[:50] + '.pdf'
    fpath = os.path.join(ensure_output_dir(), fname)
    doc = SimpleDocTemplate(fpath, pagesize=A4,
                            rightMargin=2*cm, leftMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph(title, styles['Title']))
    story.append(Spacer(1, 0.5*cm))
    for item in spec.get('content', []):
        t = item.get('type', 'paragraph')
        text = item.get('text', '')
        if t == 'heading':
            level = item.get('level', 1)
            style = styles['Heading' + str(min(level, 3))]
            story.append(Paragraph(text, style))
        elif t == 'paragraph':
            story.append(Paragraph(text, styles['Normal']))
            story.append(Spacer(1, 0.3*cm))
        elif t == 'bullet':
            story.append(Paragraph('• ' + text, styles['Normal']))
        elif t == 'table':
            rows = item.get('rows', [])
            if rows:
                t_obj = Table(rows)
                t_obj.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#2D3A2D')),
                    ('TEXTCOLOR', (0,0), (-1,0), colors.white),
                    ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.grey),
                    ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#F5F5F5')]),
                ]))
                story.append(t_obj)
                story.append(Spacer(1, 0.3*cm))
    doc.build(story)
    return fname, fpath

def create_pptx_from_spec(spec):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    title_str = spec.get('title', 'Praesentation')

    # Dark theme colors
    BG_COLOR    = RGBColor(0x1a, 0x1a, 0x1a)
    GOLD_COLOR  = RGBColor(0xf0, 0xc0, 0x60)
    WHITE_COLOR = RGBColor(0xf0, 0xf0, 0xf0)
    GRAY_COLOR  = RGBColor(0x88, 0x88, 0x88)

    slide_width  = prs.slide_width
    slide_height = prs.slide_height

    def set_bg(slide):
        from pptx.oxml.ns import qn
        from lxml import etree
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_textbox(slide, text, left, top, width, height,
                    size=18, bold=False, color=WHITE_COLOR, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox

    slides_spec = spec.get('slides', [])

    for i, slide_spec in enumerate(slides_spec):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        set_bg(slide)

        stype   = slide_spec.get('type', 'content')
        heading = slide_spec.get('heading', '')
        body    = slide_spec.get('body', '')
        bullets = slide_spec.get('bullets', [])
        footer  = slide_spec.get('footer', '')

        margin  = Inches(0.5)
        w       = slide_width - 2 * margin

        if stype == 'title':
            # Big title slide
            add_textbox(slide, title_str,
                        margin, Inches(1.8), w, Inches(1.2),
                        size=36, bold=True, color=WHITE_COLOR, align=PP_ALIGN.CENTER)
            if heading:
                add_textbox(slide, heading,
                            margin, Inches(3.2), w, Inches(0.8),
                            size=22, color=GOLD_COLOR, align=PP_ALIGN.CENTER)
            # Gold rule
            from pptx.util import Pt as PtU
            line = slide.shapes.add_shape(
                1, margin, Inches(3.0), w, Emu(40000))
            line.fill.solid()
            line.fill.fore_color.rgb = GOLD_COLOR
            line.line.fill.background()
        else:
            # Heading bar
            bar = slide.shapes.add_shape(
                1, 0, 0, slide_width, Inches(0.95))
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(0x2d, 0x3a, 0x2d)
            bar.line.fill.background()

            if heading:
                add_textbox(slide, heading,
                            margin, Emu(120000), w, Inches(0.75),
                            size=24, bold=True, color=WHITE_COLOR)

            y = Inches(1.1)

            if body:
                add_textbox(slide, body,
                            margin, y, w, Inches(1.2),
                            size=16, color=WHITE_COLOR)
                y += Inches(1.3)

            for bullet in bullets:
                add_textbox(slide, '– ' + bullet,
                            margin + Inches(0.2), y, w - Inches(0.2), Inches(0.5),
                            size=15, color=RGBColor(0xd0, 0xd0, 0xd0))
                y += Inches(0.52)

        # Slide number
        add_textbox(slide, str(i + 1),
                    slide_width - Inches(0.8), slide_height - Inches(0.4),
                    Inches(0.6), Inches(0.3),
                    size=10, color=GRAY_COLOR, align=PP_ALIGN.RIGHT)

        if footer:
            add_textbox(slide, footer,
                        margin, slide_height - Inches(0.4), w, Inches(0.3),
                        size=10, color=GRAY_COLOR)

    fname = title_str.replace(' ', '_').replace('/', '_')[:50] + '.pptx'
    fpath = os.path.join(ensure_output_dir(), fname)
    prs.save(fpath)
    return fname, fpath

def send_email_draft(spec):
    import subprocess
    to      = spec.get('to', '')
    subject = spec.get('subject', '')
    body    = spec.get('body', '')
    cc      = spec.get('cc', '')
    sender  = spec.get('from', '')

    # Escape for AppleScript - handle all special chars
    def esc(s):
        s = s.replace('\\', '\\\\')
        s = s.replace('"', '\\"')
        s = s.replace('\n', '\\n')
        s = s.replace('\r', '')
        return s

    cc_line = f'\n        make new cc recipient at end of cc recipients with properties {{address:\"{esc(cc)}\"}}' if cc else ''
    sender_line = f'\n    set sender of newMessage to "{esc(sender)}"' if sender else ''

    script = f'''tell application "Mail"
    set newMessage to make new outgoing message with properties {{subject:"{esc(subject)}", content:"{esc(body)}", visible:true}}
    tell newMessage
        make new to recipient at end of to recipients with properties {{address:"{esc(to)}"}}{cc_line}
    end tell{sender_line}
    activate
end tell'''

    subprocess.Popen(['osascript', '-e', script],
                     stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return True

def send_email_reply(spec):
    """Opens Apple Mail reply to an existing email identified by message_id."""
    import subprocess
    to      = spec.get('to', '')
    subject = spec.get('subject', '')
    body    = spec.get('body', '')
    cc      = spec.get('cc', '')
    sender  = spec.get('from', '')
    message_id = spec.get('message_id', '')
    quote_original = spec.get('quote_original', True)

    def esc(s):
        s = s.replace('\\', '\\\\')
        s = s.replace('"', '\\"')
        s = s.replace('\n', '\\n')
        s = s.replace('\r', '')
        return s

    # Build CC recipients AppleScript lines
    cc_lines = ''
    if cc:
        for addr in [a.strip() for a in cc.split(',') if a.strip()]:
            cc_lines += f'\n            make new cc recipient at end of cc recipients with properties {{address:"{esc(addr)}"}}'

    # AppleScript: try to find message by message-id and reply, fallback to new email
    if message_id:
        script = f'''tell application "Mail"
    set foundMsg to missing value
    set msgId to "{esc(message_id)}"
    repeat with acct in accounts
        repeat with mbox in mailboxes of acct
            try
                set msgs to (messages of mbox whose message id is msgId)
                if (count of msgs) > 0 then
                    set foundMsg to item 1 of msgs
                    exit repeat
                end if
            end try
        end repeat
        if foundMsg is not missing value then exit repeat
    end repeat
    if foundMsg is not missing value then
        set replyMsg to reply foundMsg with opening window
        delay 0.5
        tell replyMsg
            set subject to "{esc(subject)}"
            set content to "{esc(body)}"{cc_lines}
        end tell
        {f'set sender of replyMsg to "{esc(sender)}"' if sender else ''}
    else
        set newMessage to make new outgoing message with properties {{subject:"{esc(subject)}", content:"{esc(body)}", visible:true}}
        tell newMessage
            make new to recipient at end of to recipients with properties {{address:"{esc(to)}"}}{cc_lines}
        end tell
        {f'set sender of newMessage to "{esc(sender)}"' if sender else ''}
    end if
    activate
end tell'''
    else:
        # No message_id: fallback to regular new email
        sender_line = f'\n    set sender of newMessage to "{esc(sender)}"' if sender else ''
        script = f'''tell application "Mail"
    set newMessage to make new outgoing message with properties {{subject:"{esc(subject)}", content:"{esc(body)}", visible:true}}
    tell newMessage
        make new to recipient at end of to recipients with properties {{address:"{esc(to)}"}}{cc_lines}
    end tell{sender_line}
    activate
end tell'''

    subprocess.Popen(['osascript', '-e', script],
                     stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return True

def send_email_reply(spec):
    """Opens Apple Mail reply to an existing email identified by message_id."""
    import subprocess
    to      = spec.get('to', '')
    subject = spec.get('subject', '')
    body    = spec.get('body', '')
    cc      = spec.get('cc', '')
    message_id = spec.get('message_id', '')
    quote_original = spec.get('quote_original', True)

    def esc(s):
        s = s.replace('\\', '\\\\')
        s = s.replace('"', '\\"')
        s = s.replace('\n', '\\n')
        s = s.replace('\r', '')
        return s

    # Build CC recipients AppleScript lines
    cc_lines = ''
    if cc:
        for addr in [a.strip() for a in cc.split(',') if a.strip()]:
            cc_lines += f'\n            make new cc recipient at end of cc recipients with properties {{address:"{esc(addr)}"}}'

    # AppleScript: try to find message by message-id and reply, fallback to new email
    if message_id:
        script = f'''tell application "Mail"
    set foundMsg to missing value
    set msgId to "{esc(message_id)}"
    repeat with acct in accounts
        repeat with mbox in mailboxes of acct
            try
                set msgs to (messages of mbox whose message id is msgId)
                if (count of msgs) > 0 then
                    set foundMsg to item 1 of msgs
                    exit repeat
                end if
            end try
        end repeat
        if foundMsg is not missing value then exit repeat
    end repeat
    if foundMsg is not missing value then
        set replyMsg to reply foundMsg with opening window
        delay 0.5
        tell replyMsg
            set subject to "{esc(subject)}"
            set content to "{esc(body)}"{cc_lines}
        end tell
    else
        set newMessage to make new outgoing message with properties {{subject:"{esc(subject)}", content:"{esc(body)}", visible:true}}
        tell newMessage
            make new to recipient at end of to recipients with properties {{address:"{esc(to)}"}}{cc_lines}
        end tell
    end if
    activate
end tell'''
    else:
        # No message_id: fallback to regular new email
        script = f'''tell application "Mail"
    set newMessage to make new outgoing message with properties {{subject:"{esc(subject)}", content:"{esc(body)}", visible:true}}
    tell newMessage
        make new to recipient at end of to recipients with properties {{address:"{esc(to)}"}}{cc_lines}
    end tell
    activate
end tell'''

    result = subprocess.run(['osascript', '-e', script],
                          capture_output=True, text=True, timeout=30)
    if result.returncode != 0:
        raise Exception(f"AppleScript Fehler: {result.stderr.strip()}")
    return True

def send_whatsapp_draft(spec, agent_name=None):
    """Opens WhatsApp with a pre-filled message. Looks up phone in contacts.json, then macOS Contacts."""
    import subprocess
    import urllib.parse

    to_name = spec.get('to', '')
    message = spec.get('message', '')
    phone = spec.get('phone', '')  # Optional direct phone

    # Step 1: Look up in agent's contacts.json
    if not phone and to_name and agent_name:
        parent = agent_name.split('_')[0] if '_' in agent_name else agent_name
        contacts_path = os.path.join(BASE, parent, "memory", "contacts.json")
        if os.path.exists(contacts_path):
            try:
                with open(contacts_path) as f:
                    cdata = json.load(f)
                to_lower = to_name.lower()
                for c in cdata.get('contacts', []):
                    cname = (c.get('name') or '').lower()
                    if to_lower in cname or cname in to_lower:
                        if c.get('phone'):
                            phone = c['phone']
                            break
            except Exception:
                pass

    # Step 2: Look up in ALL agents' contacts.json files (cross-agent)
    if not phone and to_name:
        try:
            to_lower = to_name.lower()
            for agent_dir in os.listdir(BASE):
                cpath = os.path.join(BASE, agent_dir, "memory", "contacts.json")
                if os.path.isfile(cpath):
                    try:
                        with open(cpath) as f:
                            cdata = json.load(f)
                        for c in cdata.get('contacts', []):
                            cname = (c.get('name') or '').lower()
                            if to_lower in cname or cname in to_lower:
                                if c.get('phone'):
                                    phone = c['phone']
                                    break
                    except Exception:
                        continue
                if phone:
                    break
        except Exception:
            pass

    # Step 3: Look up in macOS Contacts app via AppleScript
    if not phone and to_name:
        try:
            script = f"""
            tell application "Contacts"
                set matchedPeople to every person whose name contains "{to_name}"
                if (count of matchedPeople) > 0 then
                    set thePerson to item 1 of matchedPeople
                    set thePhones to phones of thePerson
                    if (count of thePhones) > 0 then
                        set theNumber to value of item 1 of thePhones
                        return theNumber
                    end if
                end if
                return ""
            end tell
            """
            result = subprocess.run(
                ['osascript', '-e', script],
                capture_output=True, text=True, timeout=10
            )
            found_phone = result.stdout.strip()
            if found_phone:
                phone = found_phone
        except Exception:
            pass

    if not phone:
        # Fallback: copy message to clipboard and just open WhatsApp
        if message:
            subprocess.run(['pbcopy'], input=message.encode('utf-8'), timeout=5)
        subprocess.run(['open', '-a', 'WhatsApp'], capture_output=True, text=True, timeout=10)
        return to_name, None  # None signals clipboard fallback

    # Normalize phone: remove spaces, dashes, dots, leading +
    phone_clean = re.sub(r'[\s./-]', '', phone).lstrip('+')

    encoded_msg = urllib.parse.quote(message)
    whatsapp_url = f"whatsapp://send?phone={phone_clean}&text={encoded_msg}"

    result = subprocess.run(['open', whatsapp_url], capture_output=True, text=True, timeout=10)
    if result.returncode != 0:
        raise Exception(f"WhatsApp konnte nicht geoeffnet werden: {result.stderr.strip()}")
    return to_name, phone


def send_slack_draft(spec):
    """SLACK_API_V1: Sendet Slack-Nachricht via API wenn Bot-Token vorhanden,
    sonst Fallback auf Desktop-App + Clipboard.
    """
    import subprocess

    channel = spec.get('channel', '')
    to = spec.get('to', '')
    message = spec.get('message', '')
    if not message:
        raise Exception("Slack: Kein Nachrichtentext angegeben")
    if not channel and not to:
        raise Exception("Slack: Weder 'channel' noch 'to' angegeben")

    target = channel or to

    # Versuch 1: Slack API (wenn Bot-Token konfiguriert)
    sc = _get_slack_config() if '_get_slack_config' in dir() or True else None
    try:
        sc = _get_slack_config()
    except Exception:
        sc = None
    if sc:
        # Channel/User-ID aufloesen
        resolved = target
        if target.startswith('#'):
            ch_id = slack_find_channel_id(target)
            if ch_id:
                resolved = ch_id
        elif not target.startswith(('C', 'U', 'D', 'G')):
            # Kein Channel-Prefix → vermutlich Personenname
            uid = slack_find_user_id(target)
            if uid:
                # DM Channel oeffnen
                ok, dm_data = _slack_api('conversations.open', json_body={'users': uid})
                if ok and dm_data.get('channel', {}).get('id'):
                    resolved = dm_data['channel']['id']
        ok, resp = slack_send_message(resolved, message)
        if ok:
            print(f"[SLACK API] Nachricht gesendet an {target}", flush=True)
            return target, False  # clipboard_only=False, erfolgreich via API
        else:
            print(f"[SLACK API] Fehler: {resp.get('error', '?')} — Fallback auf Desktop", flush=True)

    # Fallback: Desktop-App + Clipboard (alter Mechanismus)
    subprocess.run(['pbcopy'], input=message.encode('utf-8'), timeout=5)
    subprocess.run(['open', '-a', 'Slack'], capture_output=True, text=True, timeout=10)
    script = '''
delay 1.5
tell application "System Events"
    tell process "Slack"
        set frontmost to true
        delay 0.3
        keystroke "v" using command down
    end tell
end tell'''
    result = subprocess.run(['osascript', '-e', script],
                          capture_output=True, text=True, timeout=10)
    if result.returncode != 0:
        return target, True
    return target, False


BASE = os.path.expanduser("~/Library/Mobile Documents/com~apple~CloudDocs/Downloads shared/claude_datalake")
AGENTS_DIR = os.path.join(BASE, "config/agents")
MODELS_FILE = os.path.join(BASE, "config/models.json")
SUBAGENT_KEYWORDS_FILE = os.path.join(BASE, "config/subagent_keywords.json")

# ─── SUB-AGENT DELEGATION ────────────────────────────────────────────────────

_DELEGATION_ACTIONS = {
    'nutze', 'benutze', 'verwende', 'delegiere', 'frag',
    'use', 'delegate', 'ask', 'using',
    'usa', 'delega', 'pergunta',
}
_DELEGATION_PHRASES = [
    'übergib an', 'uebergib an', 'schick das an', 'das soll', 'mit dem',
    'lass den', 'lass die',
    'send to', 'switch to', 'hand off', 'let the', 'with the',
    'manda para', 'deixa o', 'deixa a',
]
_DELEGATION_OBJECTS = {
    'sub-agent', 'subagent', 'sub agent', 'agent', 'spezialist',
    'specialist', 'spezialisten', 'assistente', 'assistenten',
}


def _load_subagent_keywords():
    """Load keyword mapping from config file."""
    if os.path.exists(SUBAGENT_KEYWORDS_FILE):
        try:
            with open(SUBAGENT_KEYWORDS_FILE) as f:
                return json.load(f)
        except Exception:
            pass
    return {}


def _levenshtein(s1, s2):
    """Simple Levenshtein distance."""
    if len(s1) < len(s2):
        return _levenshtein(s2, s1)
    if len(s2) == 0:
        return len(s1)
    prev = range(len(s2) + 1)
    for i, c1 in enumerate(s1):
        curr = [i + 1]
        for j, c2 in enumerate(s2):
            curr.append(min(prev[j + 1] + 1, curr[j] + 1, prev[j] + (c1 != c2)))
        prev = curr
    return prev[-1]


def _get_available_subagents(parent_name):
    """Get list of sub-agent names for a parent."""
    if not os.path.exists(AGENTS_DIR):
        return []
    result = []
    prefix = parent_name + '_'
    for fname in os.listdir(AGENTS_DIR):
        if fname.endswith('.txt') and fname.startswith(prefix):
            full_name = fname[:-4]
            sub_label = full_name[len(prefix):]
            result.append({'full_name': full_name, 'sub_label': sub_label})
    return result


def detect_delegation(msg, current_agent):
    """Detect if user wants to delegate to a sub-agent.
    Returns dict with full_name, score, matched_keywords — or None."""
    msg_lower = msg.lower()
    parent = get_parent_agent(current_agent) or current_agent

    words = msg_lower.split()
    has_action = any(w.rstrip('.,;:!?') in _DELEGATION_ACTIONS for w in words)
    if not has_action:
        has_action = any(phrase in msg_lower for phrase in _DELEGATION_PHRASES)
    if not has_action:
        return None

    subs = _get_available_subagents(parent)
    if not subs:
        return None

    kw_map = _load_subagent_keywords()
    best_match = None
    best_score = 0
    best_keywords = []

    for sub in subs:
        sub_label = sub['sub_label']
        score = 0
        matched_kws = []

        # 1. Exact match
        if sub_label in msg_lower:
            score = 100
            matched_kws.append(sub_label)

        # 2. Keyword match from config
        if score == 0 and sub_label in kw_map:
            for kw in kw_map[sub_label]:
                if kw.lower() in msg_lower:
                    score = max(score, 50)
                    matched_kws.append(kw)

        # 3. Partial / fuzzy match
        if score == 0:
            for w in words:
                w_clean = w.rstrip('.,;:!?')
                if len(w_clean) < 3:
                    continue
                if w_clean in sub_label or sub_label in w_clean:
                    score = max(score, 30)
                    matched_kws.append(w_clean)
                dist = _levenshtein(w_clean, sub_label)
                if dist <= 2 and len(sub_label) > 3:
                    score = max(score, 40 - dist * 10)
                    matched_kws.append(w_clean)

        if score > best_score:
            best_score = score
            best_match = sub['full_name']
            best_keywords = matched_kws[:]

    if best_score > 0:
        return {
            'full_name': best_match,
            'score': best_score,
            'matched_keywords': list(set(best_keywords)),
            'display_name': get_agent_display_name(best_match),
        }
    return None


def execute_delegation(sub_agent_name, original_msg, kontext_items, state=None):
    if state is None:
        state = get_session()
    """Execute a delegated call to a sub-agent. Returns result dict."""
    prompt_file = os.path.join(AGENTS_DIR, sub_agent_name + '.txt')
    if not os.path.exists(prompt_file):
        return {'error': f'Sub-Agent "{sub_agent_name}" nicht gefunden'}

    with open(prompt_file) as f:
        sub_prompt = f.read()
    for marker in ['\n\n--- GEDAECHTNIS:', '\n--- GEDAECHTNIS:',
                   '\n\n--- DATEI-ERSTELLUNG ---', '\n--- DATEI-ERSTELLUNG ---',
                   '\n\n--- WEITERE FAEHIGKEITEN ---']:
        sub_prompt = sub_prompt.split(marker)[0]
    sub_prompt = sub_prompt.strip()

    speicher = get_agent_speicher(sub_agent_name)
    display_name = get_agent_display_name(sub_agent_name)
    memory_ctx = build_memory_context(speicher, display_name)

    # Reuse file capability from current system prompt
    current_sp = state.get('system_prompt', '')
    cap_start = current_sp.find('\n--- DATEI-ERSTELLUNG ---')
    cap_end = current_sp.find('--- ENDE DATEI-ERSTELLUNG ---')
    if cap_start > -1 and cap_end > -1:
        file_cap_block = current_sp[cap_start:cap_end + len('--- ENDE DATEI-ERSTELLUNG ---')]
        full_sub_prompt = sub_prompt + memory_ctx + file_cap_block
    else:
        full_sub_prompt = sub_prompt + memory_ctx

    # Last 5 messages as context
    recent_history = state['verlauf'][-5:] if state['verlauf'] else []

    # Build context from kontext_items
    text_ctx = ''
    if kontext_items:
        text_ctx = '\n\n--- KONTEXT ---\n'
        for item in kontext_items:
            if not item.get('image_b64'):
                text_ctx += '\n[' + item['name'] + ']:\n' + item['content'][:10000] + '\n'
        text_ctx += '\n--- ENDE KONTEXT ---\n'

    full_msg = original_msg + text_ctx if text_ctx else original_msg
    sub_messages = list(recent_history) + [{'role': 'user', 'content': full_msg}]

    config = load_models()
    provider_key = state.get('provider', 'anthropic')
    model_id = state.get('model_id', 'claude-sonnet-4-6')
    provider_cfg = config['providers'].get(provider_key, {})
    api_key = provider_cfg.get('api_key', '')
    model_name = next((m['name'] for m in provider_cfg.get('models', []) if m['id'] == model_id), model_id)
    adapter = ADAPTERS.get(provider_key)
    if not adapter:
        return {'error': 'Unbekannter Anbieter: ' + provider_key}

    text = adapter(api_key, model_id, full_sub_prompt, sub_messages)
    response_text = '\U0001f916 ' + display_name + ' uebernimmt:\n\n' + text

    return {
        'response': response_text,
        'model_name': model_name,
        'provider_display': PROVIDER_DISPLAY.get(provider_key, provider_key),
        'model_display': MODEL_DISPLAY.get(model_id, model_name),
        'delegated_to': sub_agent_name,
        'delegated_display': display_name,
    }


app = Flask(__name__)

def cleanup_agent_files():
    """Clean old memory content from agent txt files on startup."""
    STRIP_MARKERS = [
        "--- GEDAECHTNIS:", "--- DATEI-ERSTELLUNG ---",
        "--- WEITERE FAEHIGKEITEN ---", "VERGANGENE KONVERSATIONEN:",
        "GESPEICHERTE DATEIEN IM MEMORY", "Von: ", "An: ", "Betreff: ",
    ]
    if not os.path.exists(AGENTS_DIR):
        return
    for fname in os.listdir(AGENTS_DIR):
        if not fname.endswith('.txt'):
            continue
        fpath = os.path.join(AGENTS_DIR, fname)
        try:
            with open(fpath, 'r', errors='ignore') as f:
                raw = f.read()
            cut = len(raw)
            for marker in STRIP_MARKERS:
                pos = raw.find(marker)
                if 0 < pos < cut:
                    cut = pos
            clean = raw[:cut].strip()
            if len(clean) < len(raw) - 10:
                with open(fpath, 'w') as f:
                    f.write(clean)
                print(f"Cleaned: {fname}")
        except Exception as e:
            print(f"Cleanup error {fname}: {e}")

cleanup_agent_files()

# ─── DYNAMIC CAPABILITIES: Agent-Prompts beim Start aktualisieren ─────────────
if inject_capabilities_on_startup is not None:
    try:
        inject_capabilities_on_startup(
            agents_dir=AGENTS_DIR,
            models_file=MODELS_FILE,
            datalake_base=BASE,
        )
    except Exception as e:
        print(f"[CAPABILITIES] Startup-Injection fehlgeschlagen: {e}")

# ─── PROVIDER ADAPTERS ────────────────────────────────────────────────────────

def _sanitize_anthropic_images(messages):
    """Walk message list und downscale alle base64-Bilder, die die
    Anthropic-Dimensions-Grenze (8000 px) ueberschreiten. Mutiert die
    Strukturen in-place, damit gecachte Session-Verlaeufe ebenfalls
    dauerhaft korrigiert sind. Bilder, die nicht herunterskaliert werden
    koennen (Decode-Fehler / Decompression-Bomb), werden entfernt, damit
    die gesamte Anfrage nicht 400't.
    """
    for msg in messages:
        content = msg.get('content')
        if not isinstance(content, list):
            continue
        kept = []
        for part in content:
            if not isinstance(part, dict) or part.get('type') != 'image':
                kept.append(part)
                continue
            source = part.get('source') or {}
            if source.get('type') != 'base64':
                kept.append(part)
                continue
            new_b64, new_mime = downscale_image_b64_if_needed(
                source.get('data'), source.get('media_type')
            )
            if not new_b64:
                print("[IMG-SANITIZE] Oversize-Bild entfernt (Anthropic-Limit 8000 px)", flush=True)
                continue
            source['data'] = new_b64
            source['media_type'] = new_mime
            kept.append(part)
        msg['content'] = kept


def call_anthropic(api_key, model_id, system_prompt, messages):
    from anthropic import Anthropic
    _sanitize_anthropic_images(messages)
    client = Anthropic(api_key=api_key)
    r = client.messages.create(model=model_id, max_tokens=4096, system=system_prompt, messages=messages)
    return r.content[0].text

def call_openai(api_key, model_id, system_prompt, messages):
    import openai
    client = openai.OpenAI(api_key=api_key)
    oai_messages = [{"role": "system", "content": system_prompt}] + messages
    r = client.chat.completions.create(model=model_id, messages=oai_messages, max_tokens=4096)
    return r.choices[0].message.content

def call_perplexity(api_key, model_id, system_prompt, messages):
    import requests as _pplx_req
    from requests.exceptions import ReadTimeout, ConnectionError as ReqConnectionError

    # Perplexity requires strictly alternating user/assistant messages after system.
    # Merge consecutive same-role messages to enforce this.
    raw_msgs = list(messages)
    # Remove empty messages
    raw_msgs = [m for m in raw_msgs if m.get('content')]
    # Flatten list-type content (vision messages) to text
    for m in raw_msgs:
        if isinstance(m.get('content'), list):
            m = dict(m)
            m['content'] = ' '.join(
                p.get('text', '') for p in m['content']
                if isinstance(p, dict) and p.get('type') == 'text'
            )
    # Merge consecutive same-role messages
    merged = []
    for m in raw_msgs:
        content_text = m.get('content', '')
        if isinstance(content_text, list):
            content_text = ' '.join(
                p.get('text', '') for p in content_text
                if isinstance(p, dict) and p.get('type') == 'text'
            )
        if not content_text or not content_text.strip():
            continue
        if merged and merged[-1]['role'] == m['role']:
            merged[-1]['content'] += '\n\n' + content_text
        else:
            merged.append({'role': m['role'], 'content': content_text})
    # Ensure first non-system message is 'user' (Perplexity requirement)
    if merged and merged[0]['role'] != 'user':
        merged.insert(0, {'role': 'user', 'content': '.'})
    # Ensure last message is 'user' (Perplexity sends to model for completion)
    if merged and merged[-1]['role'] != 'user':
        merged.append({'role': 'user', 'content': '.'})
    pplx_messages = [{"role": "system", "content": system_prompt}] + merged
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {"model": model_id, "messages": pplx_messages, "max_tokens": 8000}

    # Model-specific timeouts (connect_timeout, read_timeout)
    PERPLEXITY_TIMEOUTS = {
        'sonar-deep-research': (10, 300),   # Deep Research: bis zu 5 Min
        'sonar-reasoning-pro': (10, 180),   # Reasoning Pro: bis zu 3 Min
        'sonar-reasoning': (10, 180),       # Reasoning: bis zu 3 Min
        'sonar-pro': (10, 120),             # Sonar Pro: 2 Min
        'sonar': (10, 120),                 # Sonar: 2 Min
    }
    timeout = PERPLEXITY_TIMEOUTS.get(model_id, (10, 120))

    MODEL_LABELS = {
        'sonar-deep-research': 'Sonar Deep Research',
        'sonar-reasoning-pro': 'Sonar Reasoning Pro',
        'sonar-reasoning': 'Sonar Reasoning',
        'sonar-pro': 'Sonar Pro',
        'sonar': 'Sonar',
    }
    model_label = MODEL_LABELS.get(model_id, model_id)

    try:
        resp = _pplx_req.post(
            "https://api.perplexity.ai/chat/completions",
            headers=headers, json=payload, timeout=timeout
        )
    except ReadTimeout:
        raise Exception(
            f"Perplexity {model_label} hat zu lange gebraucht (Timeout nach {timeout[1]}s). "
            f"Versuche es erneut oder waehle ein schnelleres Modell."
        )
    except ReqConnectionError:
        raise Exception(
            f"Verbindung zu Perplexity fehlgeschlagen. Pruefe deine Internetverbindung."
        )

    data = resp.json()
    if resp.status_code != 200:
        raise Exception(f"Perplexity API Fehler: {data.get('error', {}).get('message', str(data))}")
    text = data["choices"][0]["message"]["content"]
    # Append citations as clickable Markdown links if present
    citations = data.get("citations")
    if citations and isinstance(citations, list) and len(citations) > 0:
        text += "\n\n**Quellen:**\n"
        for i, url in enumerate(citations, 1):
            try:
                from urllib.parse import urlparse
                domain = urlparse(url).netloc.replace("www.", "")
            except Exception:
                domain = url[:60]
            text += f"[{i}] [{domain}]({url})\n"
    return text

def call_mistral(api_key, model_id, system_prompt, messages):
    headers = {"Authorization": "Bearer " + api_key, "Content-Type": "application/json"}
    mistral_messages = [{"role": "system", "content": system_prompt}] + messages
    payload = {"model": model_id, "messages": mistral_messages, "max_tokens": 4096}
    r = requests.post("https://api.mistral.ai/v1/chat/completions", headers=headers, json=payload, timeout=60)
    return r.json()["choices"][0]["message"]["content"]

def call_gemini(api_key, model_id, system_prompt, messages):
    """Gemini via REST API (kein SDK noetig)."""
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_id}:generateContent?key={api_key}"
    contents = []
    for m in messages:
        role = 'user' if m['role'] == 'user' else 'model'
        c = m['content']
        if isinstance(c, list):
            parts = []
            for item in c:
                if item.get('type') == 'text':
                    parts.append({'text': item['text']})
                elif item.get('type') == 'image':
                    parts.append({'inline_data': {'mime_type': item['source']['media_type'], 'data': item['source']['data']}})
            contents.append({'role': role, 'parts': parts})
        else:
            contents.append({'role': role, 'parts': [{'text': str(c)}]})
    payload = {'contents': contents, 'generationConfig': {'maxOutputTokens': 4096}}
    if system_prompt:
        payload['systemInstruction'] = {'parts': [{'text': system_prompt}]}
    try:
        r = requests.post(url, json=payload, timeout=120)
        data = r.json()
        if r.status_code != 200:
            err_msg = data.get('error', {}).get('message', str(data))
            if '429' in str(r.status_code) or 'quota' in err_msg.lower() or 'rate' in err_msg.lower():
                return ("Gemini Rate Limit erreicht. Entweder kurz warten "
                        "oder auf bezahlten API Plan upgraden: "
                        "https://ai.google.dev/pricing")
            raise Exception(f"Gemini API Fehler ({r.status_code}): " + err_msg)
        parts = data['candidates'][0]['content'].get('parts', [])
        return ''.join(p.get('text', '') for p in parts)
    except requests.exceptions.Timeout:
        raise Exception("Gemini API Timeout (120s)")

ADAPTERS = {"anthropic": call_anthropic, "openai": call_openai, "perplexity": call_perplexity, "mistral": call_mistral, "gemini": call_gemini}

# ─── IMAGE & VIDEO GENERATION ────────────────────────────────────────────────

OUTPUT_DIR = os.path.expanduser("~/Library/Mobile Documents/com~apple~CloudDocs/Downloads shared/claude_outputs")

# Provider -> Image support mapping
IMAGE_PROVIDERS = {
    "openai": "gpt-image-1",
    "gemini": "imagen-4.0-generate-001",
}
VIDEO_PROVIDERS = {
    "gemini": "veo-3.1-generate-preview",
}

# Capability-Tags fuer Model-Dropdown
MODEL_CAPABILITIES = {
    "gemini-2.5-flash": ["video", "image"],
    "gemini-2.5-pro": ["reasoning", "video", "image"],
    "gemini-3-flash-preview": ["video", "image"],
    "gemini-3-pro-preview": ["reasoning", "video", "image"],
    "gemini-3.1-pro-preview": ["reasoning", "video", "image"],
    "gemini-2.0-flash": ["video", "image"],
    "gpt-4o": ["image"],
    "gpt-image-1": ["image"],
    "o1": ["reasoning"],
    "sonar-deep-research": ["reasoning"],
    "sonar-reasoning-pro": ["reasoning"],
    "sonar-reasoning": ["reasoning"],
}
CAPABILITY_EMOJI = {"video": "\U0001f3ac", "image": "\U0001f5bc\ufe0f", "reasoning": "\U0001f9e0"}


def _generate_image_single(prompt, fpath, provider_key, api_key):
    """Try generating an image with a single provider. Returns True on success."""
    import base64 as _b64

    if provider_key == 'openai':
        r = requests.post("https://api.openai.com/v1/images/generations",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": "gpt-image-1", "prompt": prompt, "size": "1024x1024",
                  "quality": "medium", "output_format": "png"},
            timeout=120)
        data = r.json()
        if r.status_code != 200:
            raise Exception(f"OpenAI: {data.get('error', {}).get('message', str(data))}")
        with open(fpath, 'wb') as f:
            f.write(_b64.b64decode(data['data'][0]['b64_json']))
        return True

    elif provider_key == 'gemini':
        # Gemini Bildgenerierung: Imagen 4 zuerst (predict API), dann Gemini-Native Fallbacks
        import base64 as _img_b64

        # Versuch 1: Imagen 4 (beste Qualitaet, predict API)
        imagen_models = ["imagen-4.0-generate-001", "imagen-4.0-fast-generate-001"]
        for imodel in imagen_models:
            try:
                url = f"https://generativelanguage.googleapis.com/v1beta/models/{imodel}:predict?key={api_key}"
                r = requests.post(url,
                    headers={"Content-Type": "application/json"},
                    json={
                        "instances": [{"prompt": prompt}],
                        "parameters": {"sampleCount": 1, "aspectRatio": "1:1"}
                    },
                    timeout=120)
                data = r.json()
                if r.status_code == 200:
                    predictions = data.get('predictions', [])
                    if predictions and predictions[0].get('bytesBase64Encoded'):
                        with open(fpath, 'wb') as f:
                            f.write(_img_b64.b64decode(predictions[0]['bytesBase64Encoded']))
                        print(f"  Imagen {imodel}: Bild generiert")
                        return True
                print(f"  Imagen {imodel}: {data.get('error', {}).get('message', 'Kein Bild')}")
            except Exception as ie:
                print(f"  Imagen {imodel} Fehler: {ie}")

        # Versuch 2: Gemini-Native Image-Modelle (generateContent API mit responseModalities)
        gemini_models = ["gemini-2.5-flash-image", "gemini-3.1-flash-image-preview", "gemini-3-pro-image-preview"]
        last_err = None
        for gmodel in gemini_models:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{gmodel}:generateContent?key={api_key}"
            r = requests.post(url,
                headers={"Content-Type": "application/json"},
                json={
                    "contents": [{"parts": [{"text": "Generate an image: " + prompt}]}],
                    "generationConfig": {"responseModalities": ["TEXT", "IMAGE"]}
                },
                timeout=120)
            data = r.json()
            if r.status_code != 200:
                last_err = f"Gemini ({gmodel}): {data.get('error', {}).get('message', str(data))}"
                print(f"  Gemini Image {gmodel} fehlgeschlagen: {last_err}")
                continue
            parts = data.get('candidates', [{}])[0].get('content', {}).get('parts', [])
            for p in parts:
                if 'inlineData' in p:
                    with open(fpath, 'wb') as f:
                        f.write(_b64.b64decode(p['inlineData']['data']))
                    return True
            last_err = f"Gemini ({gmodel}): Kein Bild generiert (moeglicherweise Content-Filter)"
            print(f"  {last_err}")
            continue
        raise Exception(last_err or "Gemini: Bildgenerierung fehlgeschlagen (alle Modelle versucht)")

    raise Exception(f"Unbekannter Image-Provider: {provider_key}")


def generate_image(prompt, agent_name, provider_key=None, task_id=None):
    """Generate an image with the selected provider. No fallback to other providers.
    If task_id is provided, progress is written to TASK_STATUS so the frontend
    can show a spinner / progress bar until the request completes.
    """
    config = load_models()
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    ts = datetime.datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
    fname = f"{agent_name}_image_{ts}.png"
    fpath = os.path.join(OUTPUT_DIR, fname)

    # Strikte Provider-Pruefung – kein Fallback
    if not provider_key or provider_key not in IMAGE_PROVIDERS:
        supported = ', '.join(IMAGE_PROVIDERS.keys())
        provider_name = provider_key or 'keiner'
        raise Exception(
            f"Bildgenerierung nicht verfuegbar mit {provider_name}. "
            f"Wechsle zu Google Gemini oder OpenAI fuer Bildgenerierung"
        )

    api_key = config['providers'].get(provider_key, {}).get('api_key', '')
    if not api_key:
        raise Exception(f"Kein API-Key fuer {provider_key} konfiguriert.")

    task_update(task_id, progress=15, message='Sende Anfrage an Bildmodell...')
    try:
        _generate_image_single(prompt, fpath, provider_key, api_key)
    except Exception as e:
        raise Exception(f"Fehler bei Bildgenerierung mit {provider_key}: {e}")

    task_done(task_id, message='Bild fertig')
    return fname, fpath, ""


def generate_video(prompt, agent_name, provider_key=None, task_id=None):
    """Generate a video using Gemini Veo 3.1. Returns (filename, filepath).
    If task_id is provided, progress is written to TASK_STATUS during the
    long-running poll loop so the frontend can display a live progress bar.
    """
    config = load_models()
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    ts = datetime.datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
    fname = f"{agent_name}_video_{ts}.mp4"
    fpath = os.path.join(OUTPUT_DIR, fname)

    # Strikte Provider-Pruefung – kein Fallback auf Gemini
    if not provider_key or provider_key not in VIDEO_PROVIDERS:
        provider_name = provider_key or 'keiner'
        raise Exception(
            f"Videogenerierung ist nur mit Google Gemini verfuegbar. "
            f"Bitte wechsle zu Gemini fuer Videogenerierung"
        )
    api_key = config['providers'].get(provider_key, {}).get('api_key')
    if not api_key:
        raise Exception("Video-Generierung benoetigt einen Gemini API-Key")

    # VEO_RETRY_V3: Retry-Schleife fuer transiente Fehler (code 13/14/429)
    # Retry 0: default (16:9, 8s)
    # Retry 1: durationSeconds=5 (kuerzer)
    # Retry 2: aspectRatio flip + stable fallback-model (veo-2.0-generate-001)
    default_model = VIDEO_PROVIDERS.get(provider_key, 'veo-3.1-generate-preview')
    STABLE_FALLBACK_MODEL = 'veo-2.0-generate-001'
    MAX_RETRIES = 3
    RETRYABLE_CODES = {13, 14, 429}
    BACKOFF = [0, 10, 20]  # seconds before retry 0/1/2
    import time as _vt

    op_name = None
    last_error_text = ''
    for _retry in range(MAX_RETRIES):
        # Retry-spezifische Parameter
        video_model = default_model
        aspect = "16:9"
        duration = 8
        retry_label = ''
        if _retry == 1:
            duration = 5
            retry_label = ' (kuerzere Dauer)'
        elif _retry == 2:
            aspect = "9:16"  # flipped
            video_model = STABLE_FALLBACK_MODEL
            retry_label = f' (Fallback {STABLE_FALLBACK_MODEL})'

        if _retry == 0:
            task_update(task_id, progress=2, message='Sende Anfrage an Gemini Veo...')
        else:
            wait_s = BACKOFF[_retry]
            task_update(
                task_id,
                progress=2,
                message=f'Retry [{_retry+1}/{MAX_RETRIES}]{retry_label}... warte {wait_s}s',
            )
            print(f"[VEO] Retry {_retry+1}/{MAX_RETRIES}{retry_label}: backoff {wait_s}s", flush=True)
            _vt.sleep(wait_s)

        url = f"https://generativelanguage.googleapis.com/v1beta/models/{video_model}:predictLongRunning?key={api_key}"
        payload = {
            "instances": [{"prompt": prompt}],
            "parameters": {"aspectRatio": aspect, "durationSeconds": duration},
        }
        print(f"[VEO] POST model={video_model} aspect={aspect} dur={duration}s", flush=True)
        try:
            r = requests.post(url,
                headers={"Content-Type": "application/json"},
                json=payload,
                timeout=30)
        except Exception as post_ex:
            last_error_text = f"Netzwerk-Fehler: {post_ex}"
            print(f"[VEO] POST Exception: {post_ex}", flush=True)
            task_update(task_id, progress=3, message=f'Netzwerk-Fehler - wird erneut versucht ({_retry+1}/{MAX_RETRIES})...')
            continue

        try:
            data = r.json()
        except Exception:
            data = {}

        if r.status_code == 200 and data.get('name'):
            op_name = data['name']
            print(f"[VEO] Operation gestartet: {op_name}", flush=True)
            break

        # Fehler analysieren
        err_obj = data.get('error', {}) or {}
        err_code = err_obj.get('code', r.status_code)
        err_msg = err_obj.get('message', str(data) if data else f'HTTP {r.status_code}')
        last_error_text = f"code {err_code}: {err_msg}"
        print(f"[VEO] POST-Fehler code={err_code}: {err_msg}", flush=True)

        # User-facing Status je Error-Code
        try:
            err_code_int = int(err_code)
        except Exception:
            err_code_int = 0

        if err_code_int == 13:
            status_msg = f'Gemini Veo Server-Fehler - wird erneut versucht ({_retry+1}/{MAX_RETRIES})...'
        elif err_code_int == 429:
            status_msg = f'Gemini Veo Rate Limit - warte {BACKOFF[min(_retry+1, MAX_RETRIES-1)]}s und versuche erneut ({_retry+1}/{MAX_RETRIES})...'
        elif err_code_int == 14:
            status_msg = f'Gemini Veo kurz nicht verfuegbar - wird erneut versucht ({_retry+1}/{MAX_RETRIES})...'
        else:
            status_msg = f'Unbekannter Fehler [{err_code}] - wird erneut versucht ({_retry+1}/{MAX_RETRIES})...'
        task_update(task_id, progress=3, message=status_msg)

        # Nicht-retryable Fehler bubbeln sofort hoch
        if err_code_int not in RETRYABLE_CODES and r.status_code not in (500, 502, 503, 504):
            raise Exception(f"Gemini Veo API-Fehler (code {err_code}): {err_msg}")
        # Sonst weiter in die naechste Retry-Runde

    if not op_name:
        raise Exception(
            f"Video-Generierung nach {MAX_RETRIES} Versuchen fehlgeschlagen. "
            f"Bitte spaeter erneut versuchen. Letzter Fehler: {last_error_text}"
        )

    task_update(task_id, progress=5, message='Video wird generiert...')

    # Poll for completion (max 6 minutes: 72 Versuche x 5 Sekunden = 360 s)
    # VEO_PATCH_V2: 6-Minuten-Timeout, detailliertes Logging, Content-Filter-Erkennung
    import base64
    import json as _vjson
    poll_url = f"https://generativelanguage.googleapis.com/v1beta/{op_name}?key={api_key}"
    MAX_ATTEMPTS = 72
    POLL_INTERVAL = 5  # seconds → 72 * 5 = 360 s (6 min)
    TOTAL_SECS = MAX_ATTEMPTS * POLL_INTERVAL
    print(f"[VEO] Poll gestartet: op={op_name} max_attempts={MAX_ATTEMPTS} interval={POLL_INTERVAL}s total={TOTAL_SECS}s", flush=True)

    last_api_snapshot = ''
    for _attempt in range(MAX_ATTEMPTS):
        import time as _t
        _t.sleep(POLL_INTERVAL)
        # Simulated progress (5 -> 95) based on attempt number, clamped
        pct = min(95, 5 + int((_attempt + 1) * 90 / MAX_ATTEMPTS))
        if _attempt < 8:
            phase = 'Video wird initialisiert...'
        elif _attempt < 28:
            phase = 'Video wird gerendert...'
        elif _attempt < 56:
            phase = 'Rendering laeuft, fast fertig...'
        else:
            phase = 'Letzte Schritte...'
        task_update(task_id, progress=pct, message=phase)

        try:
            pr = requests.get(poll_url, timeout=30)
            pdata = pr.json()
        except Exception as poll_ex:
            # Transient network error — log and keep trying
            print(f"[VEO] Poll #{_attempt+1}/{MAX_ATTEMPTS} Netzwerk-Fehler: {poll_ex}", flush=True)
            continue

        # Compact log for every attempt with current state
        _state_keys = sorted(list(pdata.keys()))
        _done_flag = pdata.get('done')
        _has_err = bool(pdata.get('error'))
        print(f"[VEO] Poll #{_attempt+1}/{MAX_ATTEMPTS} done={_done_flag} err={_has_err} keys={_state_keys}", flush=True)
        last_api_snapshot = _vjson.dumps(pdata)[:600]

        if pdata.get('error'):
            err_obj = pdata['error']
            err_msg = err_obj.get('message', str(err_obj))
            err_code = err_obj.get('code', '?')
            print(f"[VEO] API-Fehler code={err_code}: {err_msg}", flush=True)
            raise Exception(f"Gemini Veo API-Fehler (code {err_code}): {err_msg}")

        if not pdata.get('done'):
            continue

        # done=True — response auswerten
        response = pdata.get('response', {}) or {}
        gvr = response.get('generateVideoResponse', {}) or {}

        # Content-Filter: raiMediaFilteredCount > 0 → Video wurde von Safety-Filtern blockiert
        rai_count = gvr.get('raiMediaFilteredCount', 0)
        rai_reasons = gvr.get('raiMediaFilteredReasons', [])
        if rai_count and not gvr.get('generatedSamples'):
            reason_txt = '; '.join(str(r) for r in rai_reasons) if rai_reasons else 'Keine Details'
            print(f"[VEO] CONTENT FILTER: rai_count={rai_count} reasons={reason_txt}", flush=True)
            raise Exception(f"Gemini Veo: Video wurde vom Content-Filter blockiert ({reason_txt}). Bitte Prompt anpassen.")

        # Samples in allen bekannten Formaten suchen
        samples = gvr.get('generatedSamples', [])
        if not samples:
            samples = response.get('predictions', response.get('generatedVideos', []))
        if not samples:
            # Log full response for debugging
            resp_preview = _vjson.dumps(response)[:800]
            print(f"[VEO] done=true aber keine samples gefunden. Response: {resp_preview}", flush=True)
            raise Exception(
                f"Gemini Veo: Kein Video generiert (done=true, leere Response). "
                f"Details: {resp_preview[:200]}"
            )

        # Try URI download first (current API returns download URL)
        vid_uri = samples[0].get('video', {}).get('uri', '')
        if vid_uri:
            print(f"[VEO] Erfolgreich: Download-URI erhalten, lade Video...", flush=True)
            task_update(task_id, progress=97, message='Lade Video herunter...')
            dl_url = vid_uri + ('&' if '?' in vid_uri else '?') + f'key={api_key}'
            vr = requests.get(dl_url, timeout=180)
            if vr.status_code != 200:
                raise Exception(f"Gemini Veo: Video-Download fehlgeschlagen (HTTP {vr.status_code})")
            with open(fpath, 'wb') as f:
                f.write(vr.content)
            print(f"[VEO] Video gespeichert: {fpath} ({len(vr.content)} bytes)", flush=True)
            task_done(task_id, message='Video fertig')
            return fname, fpath, aspect

        # Fallback: base64 encoded video
        vid_b64 = samples[0].get('bytesBase64Encoded', samples[0].get('video', {}).get('bytesBase64Encoded', ''))
        if vid_b64:
            print(f"[VEO] Erfolgreich: base64-Video erhalten, dekodiere...", flush=True)
            task_update(task_id, progress=97, message='Dekodiere Video...')
            with open(fpath, 'wb') as f:
                f.write(base64.b64decode(vid_b64))
            task_done(task_id, message='Video fertig')
            return fname, fpath, aspect

        # Bekanntes Format nicht erkannt → volle Debug-Ausgabe
        sample_preview = _vjson.dumps(samples[0])[:500]
        print(f"[VEO] Sample-Format nicht erkannt. Erstes Sample: {sample_preview}", flush=True)
        raise Exception(f"Gemini Veo: Video-Format nicht erkannt. Sample: {sample_preview[:200]}")

    print(f"[VEO] TIMEOUT nach {TOTAL_SECS}s (letzte API-Snapshot: {last_api_snapshot[:300]})", flush=True)
    raise Exception(f"Gemini Veo: Timeout nach {TOTAL_SECS//60} Minuten ({MAX_ATTEMPTS} Versuche * {POLL_INTERVAL}s). Letzter Status: {last_api_snapshot[:150]}")


# ─── SLACK API INTEGRATION (SLACK_API_V1) ──────────────────────────────────────
# Echter Slack-API-Client via Bot-Token. Ersetzt den Clipboard-Paste-Hack fuer
# Outbound und fuegt Lese-Faehigkeiten hinzu (Channels, History, Users).
# Docs: https://docs.slack.dev/reference/methods/chat.postMessage/

def _get_slack_config():
    """Laedt Slack-Config aus models.json. Gibt dict oder None zurueck."""
    config = load_models()
    sc = config.get('slack')
    if not sc or not sc.get('bot_token'):
        return None
    return sc

def _slack_api(method, params=None, json_body=None):
    """Generischer Slack Web API Aufruf. Returns (ok, data_dict)."""
    sc = _get_slack_config()
    if not sc:
        return False, {'error': 'Kein Slack bot_token in models.json konfiguriert'}
    token = sc['bot_token']
    url = f"https://slack.com/api/{method}"
    headers = {'Authorization': f'Bearer {token}'}
    try:
        if json_body is not None:
            headers['Content-Type'] = 'application/json'
            r = requests.post(url, headers=headers, json=json_body, timeout=15)
        elif params:
            r = requests.get(url, headers=headers, params=params, timeout=15)
        else:
            r = requests.get(url, headers=headers, timeout=15)
        data = r.json()
        return data.get('ok', False), data
    except Exception as e:
        return False, {'error': str(e)}

def slack_send_message(channel, text, thread_ts=None):
    """Sendet eine Nachricht an einen Slack-Channel oder User.
    channel: '#channel-name' oder Channel-ID oder User-ID
    Returns: (ok, response_dict)
    """
    body = {'channel': channel.lstrip('#'), 'text': text}
    if thread_ts:
        body['thread_ts'] = thread_ts
    return _slack_api('chat.postMessage', json_body=body)

def slack_list_channels(limit=100):
    """Listet alle Channels auf die der Bot Zugriff hat."""
    return _slack_api('conversations.list', params={
        'types': 'public_channel,private_channel',
        'limit': limit, 'exclude_archived': True,
    })

def slack_list_users(limit=200):
    """Listet alle User im Workspace."""
    return _slack_api('users.list', params={'limit': limit})

def slack_channel_history(channel_id, limit=20):
    """Liest die letzten N Nachrichten aus einem Channel."""
    return _slack_api('conversations.history', params={
        'channel': channel_id, 'limit': limit,
    })

def slack_find_channel_id(name):
    """Sucht Channel-ID anhand des Namens."""
    ok, data = slack_list_channels(limit=500)
    if not ok:
        return None
    for ch in data.get('channels', []):
        if ch.get('name') == name.lstrip('#') or ch.get('name_normalized') == name.lstrip('#'):
            return ch['id']
    return None

def slack_find_user_id(name):
    """Sucht User-ID anhand des Display-Namens oder Real-Namens."""
    ok, data = slack_list_users(limit=500)
    if not ok:
        return None
    name_low = name.lower()
    for u in data.get('members', []):
        rn = (u.get('real_name') or '').lower()
        dn = (u.get('profile', {}).get('display_name') or '').lower()
        un = (u.get('name') or '').lower()
        if name_low in (rn, dn, un) or name_low in rn:
            return u['id']
    return None


# ─── CANVA API INTEGRATION (CANVA_API_V1) ─────────────────────────────────────
# Canva Connect REST API fuer Design-Operationen.
# Docs: https://www.canva.dev/docs/connect/

def _get_canva_config():
    """Laedt Canva-Config aus models.json. Gibt dict oder None zurueck."""
    config = load_models()
    cc = config.get('canva')
    if not cc or not cc.get('access_token'):
        return None
    return cc

def _canva_refresh_token():
    """CANVA_TOKEN_REFRESH: Erneuert den Access Token via Refresh Token."""
    import base64 as _cb64
    config = load_models()
    cc = config.get('canva', {})
    rt = cc.get('refresh_token', '')
    cid = cc.get('client_id', '')
    csec = cc.get('client_secret', '')
    if not rt or not cid or not csec:
        return False
    creds = _cb64.b64encode(f"{cid}:{csec}".encode()).decode()
    try:
        r = requests.post(
            'https://api.canva.com/rest/v1/oauth/token',
            headers={'Authorization': f'Basic {creds}', 'Content-Type': 'application/x-www-form-urlencoded'},
            data={'grant_type': 'refresh_token', 'refresh_token': rt},
            timeout=15,
        )
        if r.status_code != 200:
            print(f"[CANVA] Token-Refresh fehlgeschlagen: HTTP {r.status_code}", flush=True)
            return False
        td = r.json()
        # In models.json speichern
        config['canva']['access_token'] = td['access_token']
        if td.get('refresh_token'):
            config['canva']['refresh_token'] = td['refresh_token']
        config['canva']['expires_in'] = td.get('expires_in', 0)
        config['canva']['scope'] = td.get('scope', cc.get('scope', ''))
        with open(MODELS_FILE, 'w') as f:
            json.dump(config, f, indent=4)
        print(f"[CANVA] Token erfolgreich erneuert (expires_in={td.get('expires_in')}s)", flush=True)
        return True
    except Exception as e:
        print(f"[CANVA] Token-Refresh Exception: {e}", flush=True)
        return False


def _canva_api(method, path, json_body=None, params=None):
    """Generischer Canva Connect API Aufruf mit automatischem Token-Refresh."""
    cc = _get_canva_config()
    if not cc:
        return False, {'error': 'Kein Canva access_token in models.json konfiguriert'}
    token = cc['access_token']
    base = cc.get('api_base', 'https://api.canva.com/rest/v1')
    url = f"{base}{path}"
    headers = {'Authorization': f'Bearer {token}'}
    try:
        if method == 'GET':
            r = requests.get(url, headers=headers, params=params or {}, timeout=30)
        elif method == 'POST':
            headers['Content-Type'] = 'application/json'
            r = requests.post(url, headers=headers, json=json_body or {}, timeout=30)
        else:
            return False, {'error': f'Unbekannte Methode: {method}'}
        # CANVA_TOKEN_REFRESH: Bei 401 automatisch Token erneuern und Retry
        if r.status_code == 401:
            print("[CANVA] 401 — versuche Token-Refresh...", flush=True)
            if _canva_refresh_token():
                cc2 = _get_canva_config()
                headers['Authorization'] = f'Bearer {cc2["access_token"]}'
                if method == 'GET':
                    r = requests.get(url, headers=headers, params=params or {}, timeout=30)
                else:
                    r = requests.post(url, headers=headers, json=json_body or {}, timeout=30)
        data = r.json() if r.content else {}
        if r.status_code >= 400:
            return False, {'error': data.get('message', f'HTTP {r.status_code}'), 'status': r.status_code}
        return True, data
    except Exception as e:
        return False, {'error': str(e)}

def canva_list_designs(query=None, count=20):
    """Sucht Designs im Canva-Account."""
    params = {'count': count}
    if query:
        params['query'] = query
    return _canva_api('GET', '/designs', params=params)

def canva_get_design(design_id):
    """Gibt Details zu einem Design zurueck."""
    return _canva_api('GET', f'/designs/{design_id}')

def canva_create_design(title, design_type='doc', width=None, height=None):
    """Erstellt ein neues leeres Canva-Design."""
    body = {'title': title}
    if design_type:
        body['design_type'] = {'type': design_type}
    if width and height:
        body['design_type'] = {'type': 'custom', 'width': width, 'height': height}
    return _canva_api('POST', '/designs', json_body=body)

def canva_export_design(design_id, format_type='pdf'):
    """Exportiert ein Design als PDF/PNG/JPG."""
    body = {'design_id': design_id, 'format': {'type': format_type}}
    return _canva_api('POST', '/exports', json_body=body)

def canva_list_folders(count=50):
    """Listet Ordner im Canva-Account."""
    return _canva_api('GET', '/folders', params={'count': count})


# CANVA_CAMPAIGNS_V1: Brand Templates + Autofill fuer Ad-Kampagnen
# Docs: https://www.canva.dev/docs/connect/api-reference/autofills/

def canva_list_brand_templates(query=None, count=50):
    """Listet alle Brand Templates des Users (erfordert Canva Enterprise oder Team)."""
    params = {'count': count}
    if query:
        params['query'] = query
    return _canva_api('GET', '/brand-templates', params=params)

def canva_get_brand_template(template_id):
    """Gibt Metadaten eines Brand Templates zurueck."""
    return _canva_api('GET', f'/brand-templates/{template_id}')

def canva_get_template_dataset(template_id):
    """Gibt die ausfuellbaren Felder (Dataset) eines Brand Templates zurueck.
    Zeigt welche Platzhalter (Text, Bild, Chart) befuellt werden koennen."""
    return _canva_api('GET', f'/brand-templates/{template_id}/dataset')

def canva_autofill(template_id, data, title=None):
    """Erstellt ein neues Design aus einem Brand Template mit automatisch
    befuellten Feldern (Texte, Bilder, Charts).

    template_id: Brand Template ID
    data: dict mit Feld-Mappings, z.B.:
      {
        "headline": {"type": "text", "text": "Summer Sale 50% Off"},
        "body": {"type": "text", "text": "Shop now at example.com"},
        "hero_image": {"type": "image", "asset_id": "abc123"}
      }
    title: Name des generierten Designs (optional)

    Returns: (ok, job_data) — Job ist async, Status via canva_get_autofill_job()
    """
    body = {'brand_template_id': template_id, 'data': data}
    if title:
        body['title'] = title[:255]
    return _canva_api('POST', '/autofills', json_body=body)

def canva_get_autofill_job(job_id):
    """Prueft den Status eines Autofill-Jobs (in_progress/success/failed)."""
    return _canva_api('GET', f'/autofills/{job_id}')

def canva_upload_asset(url, name='uploaded_image'):
    """Laedt ein Bild/Asset von einer URL in Canva hoch. Gibt asset_id zurueck."""
    body = {
        'asset_upload': {
            'type': 'external_url',
            'url': url,
            'name': name,
        }
    }
    return _canva_api('POST', '/assets/upload', json_body=body)

def canva_batch_campaign(template_id, rows, title_prefix='Campaign'):
    """Generiert mehrere Designs aus einem Template mit verschiedenen Daten.
    rows: list of dicts, jeder dict = ein Design mit Feld-Mappings
    Gibt Liste von Job-IDs zurueck.

    Beispiel:
      rows = [
        {"headline": {"type":"text","text":"Ad Variant A"}, "cta": {"type":"text","text":"Buy Now"}},
        {"headline": {"type":"text","text":"Ad Variant B"}, "cta": {"type":"text","text":"Shop Now"}},
      ]
    """
    jobs = []
    for i, row in enumerate(rows):
        title = f"{title_prefix} {i+1}"
        ok, data = canva_autofill(template_id, row, title=title)
        if ok:
            job_id = data.get('job', {}).get('id', data.get('id', ''))
            jobs.append({'index': i, 'ok': True, 'job_id': job_id, 'title': title})
        else:
            jobs.append({'index': i, 'ok': False, 'error': data.get('error', '?'), 'title': title})
    return jobs


# ─── CALENDAR INTEGRATION (CALENDAR_INTEGRATION_V1) ────────────────────────────
# AppleScript-basiertes Auslesen von Fantastical/Apple Calendar Events.
# Fantastical teilt den macOS CalendarStore, deshalb funktioniert Apple Calendar
# AppleScript direkt mit Fantastical-Daten.

_CALENDAR_TARGETS = ["Arbeit", "Privat", "Familie"]
_cal_cache = {'events': [], 'cals': set(), 'ts': 0, 'key': ''}

_CALENDAR_INTENT_DE = [
    "kalender", "termin", "termine", "meeting", "meetings", "heute", "morgen",
    "diese woche", "naechste woche", "wann", "agenda", "tagesplan", "zeitplan",
    "verfuegbar", "verfügbar", "frei", "besetzt", "schedule",
]
_CALENDAR_INTENT_EN = [
    "calendar", "schedule", "today", "tomorrow", "appointment", "appointments",
    "meeting", "meetings", "when", "agenda", "free", "busy", "available",
    "this week", "next week",
]
_CALENDAR_KEYWORDS = set(_CALENDAR_INTENT_DE + _CALENDAR_INTENT_EN)


def _has_calendar_intent(msg):
    """Prueft ob eine User-Nachricht nach Kalender-Daten fragt."""
    low = msg.lower()
    return any(kw in low for kw in _CALENDAR_KEYWORDS)


def get_calendar_events(days_back=0, days_ahead=7, calendars=None, search=None):
    """Liest Events aus Apple Calendar via AppleScript.
    Returns: (events_list, calendars_found_set, error_str_or_None)

    Jedes Event: {title, start, end, location, calendar_name, notes, all_day}
    Events sind chronologisch nach start sortiert.
    """
    target_cals = calendars if calendars else _CALENDAR_TARGETS

    # Cache pruefen (120s Gueltigkeitsdauer)
    import time as _cal_time
    cache_key = f"{days_back}:{days_ahead}:{','.join(target_cals)}:{search or ''}"
    if _cal_cache['key'] == cache_key and (_cal_time.time() - _cal_cache['ts']) < 120:
        return list(_cal_cache['events']), set(_cal_cache['cals']), None

    # AppleScript bauen: pro Kalender einen try-Block
    cal_blocks = []
    for cname in target_cals:
        safe = cname.replace('"', '\\"')
        cal_blocks.append(f"""
        try
            set c to calendar "{safe}"
            set es to (every event of c whose start date >= startD and start date <= endD)
            repeat with e in es
                set t to summary of e
                set s to start date of e
                set eEnd to end date of e
                set loc to ""
                try
                    set loc to location of e
                end try
                set nt to ""
                try
                    set nt to description of e
                end try
                set ad to allday event of e
                set out to out & t & "|||" & (s as string) & "|||" & (eEnd as string) & "|||" & loc & "|||" & "{safe}" & "|||" & (ad as string) & "|||" & nt & "\n"
            end repeat
        end try""")

    script = f"""
set today to current date
set startD to today - ({days_back} * days)
set endD to today + ({days_ahead} * days)
tell application "Calendar"
    set out to ""
    {"".join(cal_blocks)}
    return out
end tell
"""

    try:
        import subprocess as _cal_sp
        r = _cal_sp.run(
            ['osascript', '-e', script],
            capture_output=True, text=True, timeout=45,
        )
        if r.returncode != 0:
            err = (r.stderr or '').strip()
            if 'not allowed' in err.lower() or 'permission' in err.lower():
                return [], set(), "Keine Kalender-Berechtigung. Bitte in Systemeinstellungen > Datenschutz > Kalender erlauben."
            return [], set(), f"AppleScript-Fehler: {err[:200]}"
        raw = r.stdout.strip()
    except Exception as e:
        return [], set(), f"Kalender-Zugriff fehlgeschlagen: {e}"

    if not raw:
        return [], set(), None  # keine Events, kein Fehler

    events = []
    cals_found = set()
    for line in raw.split("\n"):
        line = line.strip()
        if not line or '|||' not in line:
            continue
        parts = line.split('|||')
        if len(parts) < 6:
            continue
        title = parts[0].strip()
        start_str = parts[1].strip()
        end_str = parts[2].strip()
        location = parts[3].strip()
        cal_name = parts[4].strip()
        all_day = parts[5].strip().lower() == 'true'
        notes = parts[6].strip() if len(parts) > 6 else ''
        # AppleScript gibt 'missing value' statt leer zurueck
        if location == 'missing value':
            location = ''
        if notes == 'missing value':
            notes = ''

        # Freitext-Suche
        if search:
            s = search.lower()
            if s not in title.lower() and s not in notes.lower() and s not in location.lower():
                continue

        # Datum parsen (macOS AppleScript Format: "Monday, 14 April 2026 at 07:30:00")
        dt_start = _parse_applescript_date(start_str)
        dt_end = _parse_applescript_date(end_str)

        events.append({
            'title': title,
            'start': dt_start.isoformat() if dt_start else start_str,
            'end': dt_end.isoformat() if dt_end else end_str,
            'location': location,
            'calendar_name': cal_name,
            'notes': notes[:500],
            'all_day': all_day,
            '_sort_key': dt_start.timestamp() if dt_start else 0,
        })
        cals_found.add(cal_name)

    events.sort(key=lambda e: e['_sort_key'])
    for e in events:
        e.pop('_sort_key', None)
    # Cache aktualisieren
    _cal_cache['events'] = list(events)
    _cal_cache['cals'] = set(cals_found)
    _cal_cache['ts'] = _cal_time.time()
    _cal_cache['key'] = cache_key
    return events, cals_found, None


def _parse_applescript_date(s):
    """Parst ein macOS AppleScript Datum wie 'Monday, 14 April 2026 at 07:30:00'."""
    if not s:
        return None
    import re as _cal_re
    # Format: "Weekday, DD Month YYYY at HH:MM:SS"
    m = _cal_re.search(r'(\d{1,2})\s+(\w+)\s+(\d{4})\s+(?:at|um)?\s*(\d{1,2}):(\d{2}):(\d{2})', s)
    if m:
        day, month_str, year, h, mi, sec = m.group(1), m.group(2), m.group(3), m.group(4), m.group(5), m.group(6)
        months = {
            'january': 1, 'february': 2, 'march': 3, 'april': 4, 'may': 5, 'june': 6,
            'july': 7, 'august': 8, 'september': 9, 'october': 10, 'november': 11, 'december': 12,
            'januar': 1, 'februar': 2, 'maerz': 3, 'märz': 3, 'mai': 5, 'juni': 6,
            'juli': 7, 'oktober': 10, 'dezember': 12,
        }
        month = months.get(month_str.lower())
        if month:
            try:
                return datetime.datetime(int(year), month, int(day), int(h), int(mi), int(sec))
            except Exception:
                pass
    return None


def format_calendar_context(events, max_items=15):
    """Formatiert Events als Kontext-Block fuer den System-Prompt."""
    if not events:
        return ""
    lines = ["--- KALENDER (kommende Termine) ---"]
    for e in events[:max_items]:
        start = e.get('start', '')
        title = e.get('title', '')
        cal = e.get('calendar_name', '')
        loc = e.get('location', '')
        ad = ' (ganztaegig)' if e.get('all_day') else ''
        line = f"[{start}] {title}{ad}"
        if cal:
            line += f" | Kalender: {cal}"
        if loc:
            line += f" | Ort: {loc}"
        lines.append(line)
    lines.append("--- ENDE KALENDER ---")
    return "\n".join(lines)


# ─── STATE (SESSION-BASED) ─────────────────────────────────────────────────────

import time as _time

queue_lock = threading.Lock()
sessions = {}  # session_id -> state dict

def _new_state():
    return {
        "agent": None, "system_prompt": None, "speicher": None,
        "verlauf": [], "dateiname": None, "kontext_items": [],
        "provider": "anthropic", "model_id": "claude-sonnet-4-6",
        "session_files": [],
        "queue": [], "processing": False, "stop_requested": False,
        "completed_responses": [], "current_prompt": "",
        "last_active": _time.time(),
    }

def get_session(session_id=None):
    if not session_id:
        session_id = 'default'
    if session_id not in sessions:
        sessions[session_id] = _new_state()
    sessions[session_id]['last_active'] = _time.time()
    return sessions[session_id]

def cleanup_old_sessions():
    cutoff = _time.time() - (24 * 60 * 60)
    to_delete = [sid for sid, s in sessions.items() if s.get('last_active', 0) < cutoff]
    for sid in to_delete:
        try:
            auto_save_session(sid)
        except Exception:
            pass
        del sessions[sid]


# ─── TASK STATUS (PROGRESS TRACKING FOR LONG-RUNNING OPERATIONS) ───────────────
# In-memory registry for tracking the progress of long-running tasks
# like video and image generation. Frontend polls /task_status/<task_id>.

TASK_STATUS = {}
task_lock = threading.Lock()


def task_create(kind, session_id, estimated_total=180, initial_message=None):
    """Create a new task entry and return its task_id.
    kind: 'video' | 'image' | generic label
    estimated_total: rough expected duration in seconds (used for ETA)
    """
    task_id = str(uuid.uuid4())
    now = _time.time()
    default_msg = {
        'video': 'Video-Generierung gestartet...',
        'image': 'Bild-Generierung gestartet...',
    }.get(kind, 'Aufgabe gestartet...')
    with task_lock:
        TASK_STATUS[task_id] = {
            'task_id': task_id,
            'kind': kind,
            'session_id': session_id or 'default',
            'status': 'running',
            'progress': 0,
            'message': initial_message or default_msg,
            'started_at': now,
            'updated_at': now,
            'finished_at': None,
            'estimated_total_seconds': estimated_total,
            'error': None,
        }
    return task_id


def task_update(task_id, progress=None, message=None):
    """Update a task's progress and message. Silently ignored if task_id unknown."""
    if not task_id:
        return
    with task_lock:
        t = TASK_STATUS.get(task_id)
        if not t or t['status'] in ('done', 'error'):
            return
        if progress is not None:
            try:
                t['progress'] = max(0, min(99, int(progress)))
            except Exception:
                pass
        if message is not None:
            t['message'] = str(message)
        t['updated_at'] = _time.time()


def task_done(task_id, message='Fertig'):
    if not task_id:
        return
    with task_lock:
        t = TASK_STATUS.get(task_id)
        if not t:
            return
        t['status'] = 'done'
        t['progress'] = 100
        t['message'] = message
        t['finished_at'] = _time.time()
        t['updated_at'] = t['finished_at']


def task_error(task_id, err):
    if not task_id:
        return
    with task_lock:
        t = TASK_STATUS.get(task_id)
        if not t:
            return
        t['status'] = 'error'
        t['message'] = f"Fehler: {err}"
        t['error'] = str(err)
        t['finished_at'] = _time.time()
        t['updated_at'] = t['finished_at']


def task_get(task_id):
    """Return a JSON-serializable snapshot of a task, or None."""
    with task_lock:
        t = TASK_STATUS.get(task_id)
        if not t:
            return None
        out = dict(t)
    now = _time.time()
    end_t = out.get('finished_at') or now
    out['elapsed_seconds'] = round(end_t - out['started_at'], 1)
    est = out.get('estimated_total_seconds') or 0
    if est and out['status'] == 'running':
        out['eta_seconds'] = max(0, round(est - out['elapsed_seconds'], 1))
    elif out['status'] == 'done':
        out['eta_seconds'] = 0
    else:
        out['eta_seconds'] = None
    return out


def tasks_for_session(session_id, max_done_age=15):
    """Return all active tasks for a session, plus recently-finished ones
    (so the frontend can display completion state briefly).
    """
    now = _time.time()
    ids = []
    with task_lock:
        for tid, t in TASK_STATUS.items():
            if t['session_id'] != session_id:
                continue
            if t['status'] in ('done', 'error'):
                fin = t.get('finished_at') or 0
                if (now - fin) > max_done_age:
                    continue
            ids.append(tid)
    return [task_get(tid) for tid in ids if task_get(tid)]


def tasks_cleanup(max_age=3600):
    """Drop tasks finished longer than max_age seconds ago."""
    now = _time.time()
    with task_lock:
        stale = [tid for tid, t in TASK_STATUS.items()
                 if t.get('finished_at') and (now - t['finished_at']) > max_age]
        for tid in stale:
            TASK_STATUS.pop(tid, None)


def auto_save_session(session_id):
    """Auto-save the full conversation of a session to its konversation file.
    Overwrites the file with the complete conversation (header + all messages).
    Silent fail — logs errors but never raises.

    Neu: Wenn die Session noch keine Konversationsdatei hat (dateiname is None),
    wird sie hier beim ersten Aufruf angelegt — mit Sekunden-Genauigkeit im
    Zeitstempel, damit zwei Tabs mit demselben Agenten innerhalb einer Minute
    nicht dieselbe Datei bekommen. So sieht der Nutzer im History-Sidebar genau
    dann einen neuen Eintrag, wenn er den ersten Prompt gesendet hat."""
    try:
        if session_id not in sessions:
            return
        st = sessions[session_id]
        if not st.get('agent') or not st.get('verlauf'):
            return
        # Lazy-Create: Konversationsdatei erst beim ersten Prompt anlegen
        if not st.get('dateiname'):
            speicher = st.get('speicher')
            agent_name = st.get('agent')
            if not speicher or not agent_name:
                return
            os.makedirs(speicher, exist_ok=True)
            datum = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
            parent = get_parent_agent(agent_name)
            if parent:
                sub_label = agent_name.split('_', 1)[1]
                dateiname = os.path.join(speicher, 'konversation_' + datum + '_' + sub_label + '.txt')
            else:
                dateiname = os.path.join(speicher, 'konversation_' + datum + '.txt')
            st['dateiname'] = dateiname
            print(f'[AUTO-SAVE] Neue Konversation angelegt: {os.path.basename(dateiname)}')
        dateiname = st['dateiname']
        agent = st['agent']
        provider_key = st.get('provider', 'anthropic')
        model_id = st.get('model_id', 'claude-sonnet-4-6')

        # Build full file content: header + all messages
        # Extract date from dateiname (konversation_2026-04-06_08-45.txt)
        import re
        basename = os.path.basename(dateiname).replace('.txt', '').replace('konversation_', '')
        lines = ['Agent: ' + agent, 'Datum: ' + basename, '']

        for i, m in enumerate(st['verlauf']):
            role = m.get('role', '')
            text = m.get('content', '')
            if isinstance(text, list):
                # Vision content — extract text parts
                text = ' '.join(p.get('text', '') for p in text if isinstance(p, dict) and p.get('type') == 'text')
            if role == 'user':
                lines.append('[' + provider_key + '/' + model_id + ']')
                lines.append('Du: ' + text)
            elif role == 'assistant':
                lines.append('Assistant: ' + text)
                lines.append('')  # blank line after each exchange

        # Append context files block if any are loaded
        ctx_items = st.get('kontext_items', [])
        if ctx_items:
            ctx_entries = []
            for ci in ctx_items:
                entry = {'name': ci.get('name', ''), 'type': 'file'}
                if ci.get('image_b64'):
                    entry['type'] = 'image'
                ctx_entries.append(entry)
            lines.append('')
            lines.append('[KONTEXT_DATEIEN:' + json.dumps(ctx_entries, ensure_ascii=False) + ']')

        # Atomic write: write to .tmp then rename
        tmp_path = dateiname + '.tmp'
        with open(tmp_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        os.replace(tmp_path, dateiname)
        print(f'[AUTO-SAVE] Session {session_id[:12]} gesichert -> {os.path.basename(dateiname)}')
    except Exception as e:
        print(f'[AUTO-SAVE] Fehler bei {session_id[:12]}: {e}')

def load_models():
    if os.path.exists(MODELS_FILE):
        with open(MODELS_FILE) as f:
            return json.load(f)
    return {"providers": {}}

# ─── FILE EXTRACTION ──────────────────────────────────────────────────────────

def extract_file_content(raw, filename):
    fname = filename.lower()
    if fname.endswith('.pdf'):
        try:
            import PyPDF2, io
            reader = PyPDF2.PdfReader(io.BytesIO(raw))
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        except Exception as e:
            return "[PDF Fehler: " + str(e) + "]"
    if fname.endswith('.docx'):
        try:
            import zipfile, io
            from xml.etree import ElementTree as ET
            z = zipfile.ZipFile(io.BytesIO(raw))
            xml = z.read('word/document.xml')
            tree = ET.fromstring(xml)
            ns = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
            return ' '.join(node.text for node in tree.iter(ns+'t') if node.text)
        except Exception as e:
            return "[DOCX Fehler: " + str(e) + "]"
    if fname.endswith('.xlsx'):
        try:
            import zipfile, io
            from xml.etree import ElementTree as ET
            z = zipfile.ZipFile(io.BytesIO(raw))
            strings = []
            if 'xl/sharedStrings.xml' in z.namelist():
                tree = ET.fromstring(z.read('xl/sharedStrings.xml'))
                for si in tree.iter('{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t'):
                    strings.append(si.text or '')
            result = []
            if 'xl/worksheets/sheet1.xml' in z.namelist():
                tree = ET.fromstring(z.read('xl/worksheets/sheet1.xml'))
                ns = '{http://schemas.openxmlformats.org/spreadsheetml/2006/main}'
                for row in tree.iter(ns+'row'):
                    vals = []
                    for cell in row.iter(ns+'c'):
                        v = cell.find(ns+'v')
                        if v is not None and v.text:
                            if cell.get('t','') == 's':
                                vals.append(strings[int(v.text)] if int(v.text) < len(strings) else '')
                            else:
                                vals.append(v.text)
                    if vals: result.append('\t'.join(vals))
            return '\n'.join(result)
        except Exception as e:
            return "[XLSX Fehler: " + str(e) + "]"
    if fname.endswith('.eml'):
        try:
            import email
            msg = email.message_from_bytes(raw)
            parts = ["Von: " + str(msg.get('From','')), "An: " + str(msg.get('To','')),
                     "Betreff: " + str(msg.get('Subject','')), "Datum: " + str(msg.get('Date','')), ""]
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        parts.append(part.get_payload(decode=True).decode('utf-8', errors='ignore'))
            else:
                parts.append(msg.get_payload(decode=True).decode('utf-8', errors='ignore'))
            return '\n'.join(parts)
        except Exception as e:
            return "[EML Fehler: " + str(e) + "]"
    if any(fname.endswith(ext) for ext in ['.mp4','.mov','.avi','.mkv','.mp3','.wav','.m4a']):
        return "[Mediendatei: " + filename + " — Inhalt kann nicht als Text gelesen werden.]"
    return raw.decode('utf-8', errors='ignore')

# ─── MEMORY SYSTEM ────────────────────────────────────────────────────────────

def load_index(speicher):
    """Load session index for an agent.

    Robust gegenueber korrupten _index.json (iCloud-Sync-Konflikte,
    Race-Conditions beim Schreiben). Bei JSONDecodeError / unerwartetem
    Typ: die defekte Datei wird als `.corrupt-<ts>.bak` umbenannt und
    ein leerer Index zurueckgegeben. `migrate_old_conversations()` baut
    den Index danach aus den vorhandenen `konversation_*.txt` neu auf.
    """
    index_file = os.path.join(speicher, '_index.json')
    if not os.path.exists(index_file):
        return []
    try:
        with open(index_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
        if not isinstance(data, list):
            raise ValueError(f'_index.json ist keine Liste, sondern {type(data).__name__}')
        return data
    except (json.JSONDecodeError, ValueError, OSError) as e:
        import datetime as _dt
        ts = _dt.datetime.now().strftime('%Y%m%d_%H%M%S')
        bak = index_file + f'.corrupt-{ts}.bak'
        try:
            os.rename(index_file, bak)
            print(f'[INDEX] KORRUPT: {index_file} ({e}). Umbenannt in {bak}. Rebuild leer.')
        except OSError as _mv_err:
            print(f'[INDEX] KORRUPT: {index_file} ({e}). Backup fehlgeschlagen: {_mv_err}')
        return []

def save_index(speicher, index):
    """Save session index."""
    index_file = os.path.join(speicher, '_index.json')
    with open(index_file, 'w') as f:
        json.dump(index, f, ensure_ascii=False, indent=2)

def migrate_old_conversations(speicher):
    """Import existing conversation files that aren't yet in the index."""
    index = load_index(speicher)
    indexed_files = {e.get('file','') for e in index}

    conv_files = sorted([
        f for f in os.listdir(speicher)
        if f.startswith('konversation_') and f.endswith('.txt')
    ])

    new_entries = []
    for fname in conv_files:
        if fname in indexed_files:
            continue
        fpath = os.path.join(speicher, fname)
        try:
            with open(fpath) as f:
                content = f.read()
            # Extract date from filename: konversation_2026-03-31_15-27.txt
            parts = fname.replace('konversation_','').replace('.txt','').split('_')
            if len(parts) >= 2:
                date_str = parts[0] + ' ' + parts[1].replace('-',':')
            else:
                date_str = fname
            # Build a short summary from content (first 800 chars)
            summary = "[Importiert] " + content[:400].replace('\n',' ').strip()
            new_entries.append({
                "date": date_str,
                "file": fname,
                "summary": summary,
                "referenced_files": []
            })
        except:
            pass

    if new_entries:
        # Insert old entries at beginning, sorted by date
        combined = new_entries + index
        save_index(speicher, combined)
        return len(new_entries)
    return 0

def summarize_conversation(verlauf, agent_prompt):
    """Ask Claude to summarize the current conversation in 2-3 sentences."""
    if not verlauf or len(verlauf) < 2:
        return None
    try:
        config = load_models()
        api_key = config['providers']['anthropic']['api_key']
        from anthropic import Anthropic
        client = Anthropic(api_key=api_key)
        conv_text = "\n".join(
            ("User: " if m['role'] == 'user' else "Assistant: ") + str(m['content'])[:500]
            for m in verlauf[-10:]  # last 10 messages max
        )
        r = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=200,
            system="Fasse die folgende Konversation in 2-3 praezisen Saetzen zusammen. Nur die Zusammenfassung, kein Kommentar.",
            messages=[{"role": "user", "content": conv_text}]
        )
        return r.content[0].text
    except:
        return None

def close_current_session(state=None):
    """Summarize and index the current session before switching agents."""
    if state is None:
        state = get_session()
    if not state['agent'] or not state['verlauf']:
        return
    summary = summarize_conversation(state['verlauf'], state['system_prompt'])
    if not summary:
        return
    speicher = state['speicher']
    index = load_index(speicher)
    entry = {
        "date": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
        "file": os.path.basename(state['dateiname']) if state['dateiname'] else "",
        "summary": summary,
        "referenced_files": list(state['session_files'])
    }
    index.append(entry)
    # Keep max 50 entries
    if len(index) > 50:
        index = index[-50:]
    save_index(speicher, index)

HTML = """<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<title>Assistant</title>
<link rel="icon" href="data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 32 32'%3E%3Crect width='32' height='32' rx='6' fill='%23111'/%3E%3Ctext x='16' y='24' text-anchor='middle' font-family='system-ui' font-weight='700' font-size='22' fill='%23f0c060'%3EA%3C/text%3E%3C/svg%3E">
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap" rel="stylesheet">
<style>
  * { margin:0; padding:0; box-sizing:border-box; }
  body { font-family:'Inter',sans-serif; background:#1a1a1a; color:#e0e0e0; height:100vh; display:flex; flex-direction:column; overflow:hidden; }
  #header { background:#111; border-bottom:1px solid #333; padding:10px 20px; display:flex; align-items:center; gap:12px; flex-shrink:0; z-index:10; }
  #header h1 { font-size:11px; font-weight:700; color:#aaa; letter-spacing:2.5px; text-transform:uppercase; margin-right:8px; }
  .hdr-btn { background:none; border:1px solid #444; color:#aaa; padding:5px 12px; cursor:pointer; font-size:12px; border-radius:6px; font-family:Inter,sans-serif; transition:all 0.15s; }
  .hdr-btn:hover { border-color:#666; color:#eee; }
  #agent-label { font-size:13px; color:#f0c060; font-weight:600; }
  #prompt-btn { background:none; border:1px solid #444; color:#aaa; padding:5px 12px; cursor:pointer; font-size:12px; border-radius:6px; font-family:Inter,sans-serif; }
  #prompt-btn:hover { border-color:#f0c060; color:#f0c060; }
  select.hdr-select { background:#1a1a1a; border:1px solid #444; color:#aaa; padding:5px 10px; font-size:12px; border-radius:6px; font-family:Inter,sans-serif; cursor:pointer; }
  select.hdr-select:hover { border-color:#666; }
  select.hdr-select option { background:#1a1a1a; color:#e0e0e0; padding:4px 8px; }
  select.hdr-select { -webkit-appearance:menulist; appearance:menulist; min-width:120px; }
  /* Nav Menu */
  .nav-wrap { position:relative; display:inline-block; }
  .nav-btn { background:none; border:1px solid #444; color:#ccc; padding:5px 10px; cursor:pointer; font-size:16px; border-radius:6px; font-family:Inter,sans-serif; transition:all 0.15s; line-height:1; }
  .nav-btn:hover, .nav-btn.active { border-color:#f0c060; color:#f0c060; background:#1a1a10; }
  .nav-menu { display:none; position:absolute; right:0; top:calc(100% + 6px); background:#1a1a2e; border:1px solid #444; border-radius:10px; min-width:280px; z-index:200; box-shadow:0 8px 24px rgba(0,0,0,0.6); padding:6px 0; max-height:80vh; overflow-y:auto; }
  .nav-menu.open { display:block; }
  .nav-menu-section { padding:6px 14px 4px; font-size:10px; color:#888; font-weight:700; text-transform:uppercase; letter-spacing:1.2px; }
  .nav-menu-item { display:flex; align-items:center; gap:10px; padding:9px 14px; color:#d0d0e0; font-size:13px; cursor:pointer; text-decoration:none; transition:background 0.12s; }
  .nav-menu-item:hover { background:#22224a; color:#fff; }
  .nav-menu-item .nav-icon { font-size:16px; width:22px; text-align:center; flex-shrink:0; }
  .nav-menu-item .nav-label { flex:1; }
  .nav-menu-item .nav-hint { font-size:10px; color:#666; }
  .nav-menu-divider { height:1px; background:#333; margin:4px 0; }
  /* Chat Tabs */
  #tab-bar { display:flex; align-items:center; background:#0d0d0d; border-bottom:1px solid #333; padding:0 8px; flex-shrink:0; overflow-x:auto; gap:2px; min-height:32px; }
  .chat-tab { display:flex; align-items:center; gap:6px; padding:5px 12px; background:none; border:none; border-bottom:2px solid transparent; color:#888; font-size:11px; font-family:Inter,sans-serif; cursor:pointer; white-space:nowrap; transition:all .15s; max-width:180px; }
  .chat-tab:hover { color:#ccc; background:#1a1a1a; }
  .chat-tab.active { color:#f0c060; border-bottom-color:#f0c060; background:#1a1a1a; }
  .chat-tab .tab-name { overflow:hidden; text-overflow:ellipsis; }
  .chat-tab .tab-close { font-size:13px; color:#555; padding:0 2px; line-height:1; cursor:pointer; border:none; background:none; }
  .chat-tab .tab-close:hover { color:#f87171; }
  #tab-add { background:none; border:1px solid #333; color:#888; font-size:14px; padding:2px 8px; border-radius:4px; cursor:pointer; margin-left:4px; line-height:1; flex-shrink:0; }
  #tab-add:hover { border-color:#f0c060; color:#f0c060; }
  /* Services status in nav */
  .svc-row { display:flex; align-items:center; gap:8px; padding:6px 14px; font-size:12px; color:#ccc; }
  .svc-dot { width:8px; height:8px; border-radius:50%; flex-shrink:0; }
  .svc-dot.online { background:#4ade80; box-shadow:0 0 4px #4ade80; }
  .svc-dot.offline { background:#f87171; box-shadow:0 0 4px #f87171; }
  .svc-name { flex:1; }
  .svc-restart { background:none; border:1px solid #444; color:#888; font-size:10px; padding:2px 8px; border-radius:4px; cursor:pointer; }
  .svc-restart:hover { border-color:#f0c060; color:#f0c060; }
  #main { display:flex; flex:1; overflow:hidden; }
  #sidebar { width:30%; min-width:280px; background:#141414; border-right:1px solid #2a2a2a; display:flex; flex-direction:column; overflow:hidden; transition:width 0.2s,min-width 0.2s; flex-shrink:0; }
  #sidebar-header { padding:12px 16px; border-bottom:1px solid #333; display:flex; align-items:center; justify-content:space-between; flex-shrink:0; }
  #prompt-editor { flex:1; background:#111; color:#e0e0e0; border:none; padding:16px; font-size:12px; font-family:Inter,sans-serif; resize:none; outline:none; line-height:1.6; min-height:0; }
  #sidebar-btns { padding:10px 16px; border-top:1px solid #333; display:flex; gap:8px; flex-shrink:0; }
  #history-section { border-top:2px solid #333; flex:1; min-height:0; display:flex; flex-direction:column; padding:10px 16px 6px; }
  #history-section span { font-size:11px; color:#aaa; font-weight:600; letter-spacing:1.5px; text-transform:uppercase; }
  #history-list { overflow-y:auto; flex:1; min-height:0; padding:0 8px 8px; }
  .history-item { display:block; width:100%; background:none; border:none; border-radius:6px; padding:7px 10px; cursor:pointer; text-align:left; font-family:Inter,sans-serif; margin-bottom:3px; }
  .history-item:hover { background:#222; }
  .history-item.active { background:#2a3a2a; }
  .h-date { font-size:10px; color:#666; display:block; margin-bottom:2px; }
  .h-summary { font-size:11px; color:#bbb; display:block; line-height:1.4; }
  #chat-area { flex:1; display:flex; flex-direction:column; overflow:hidden; }
  #ctx-bar { background:#161616; border-bottom:1px solid #252525; padding:6px 16px; display:flex; align-items:center; gap:8px; flex-wrap:wrap; flex-shrink:0; min-height:36px; }
  #ctx-bar .ctx-label { font-size:10px; color:#555; text-transform:uppercase; letter-spacing:1px; font-weight:600; }
  .ctx-item { background:#202020; border:1px solid #333; border-radius:5px; padding:3px 8px; font-size:11px; color:#888; display:flex; align-items:center; gap:5px; }
  .ctx-item.auto-loaded { border-color:#4a8a4a; background:#1a2a1a; color:#8aba8a; }
  .ctx-item button { background:none; border:none; color:#555; cursor:pointer; font-size:11px; padding:0; line-height:1; }
  .ctx-item button:hover { color:#f0c060; }
  #search-overlay { display:none; position:fixed; top:0; left:0; width:100%; height:100%; background:rgba(0,0,0,0.7); z-index:1000; justify-content:center; align-items:center; }
  #search-overlay.show { display:flex; }
  #search-panel { background:#1a1a1a; border:1px solid #444; border-radius:12px; width:min(600px,90%); max-height:80vh; display:flex; flex-direction:column; box-shadow:0 8px 32px rgba(0,0,0,0.5); }
  #search-panel-header { padding:14px 18px; border-bottom:1px solid #333; display:flex; align-items:center; justify-content:space-between; }
  #search-panel-header h3 { font-size:13px; color:#e0e0e0; font-weight:600; margin:0; }
  #search-panel-header span { font-size:11px; color:#888; }
  #search-results-list { flex:1; overflow-y:auto; padding:6px 0; max-height:50vh; }
  .search-result-item { padding:10px 18px; border-bottom:1px solid #222; display:flex; gap:10px; align-items:flex-start; cursor:pointer; transition:background 0.1s; }
  .search-result-item:hover { background:#252525; }
  .search-result-item input[type=checkbox] { margin-top:3px; accent-color:#f0c060; flex-shrink:0; }
  .search-result-info { flex:1; min-width:0; }
  .search-result-name { font-size:12px; color:#ccc; font-weight:500; white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
  .search-result-meta { font-size:11px; color:#777; margin-top:2px; }
  .search-result-meta .from-person { color:#6aba6a; font-weight:500; }
  .search-result-meta .from-auto { color:#666; }
  .search-result-preview { font-size:11px; color:#555; margin-top:3px; line-height:1.4; display:-webkit-box; -webkit-line-clamp:2; -webkit-box-orient:vertical; overflow:hidden; }
  #search-panel-footer { padding:12px 18px; border-top:1px solid #333; display:flex; gap:8px; justify-content:flex-end; }
  #search-panel-footer button { padding:7px 16px; border-radius:6px; font-size:12px; font-weight:600; cursor:pointer; font-family:Inter,sans-serif; border:none; }
  .search-result-item.notification .search-result-name { color:#777; font-style:italic; }
  .search-result-item.notification .search-result-meta { color:#555; }
  .search-result-item.notification { opacity:0.7; }
  .search-section-divider { padding:6px 18px; font-size:11px; color:#888; font-weight:600; letter-spacing:0.5px; border-bottom:1px solid #333; background:#141414; }
  #search-selection-counter { font-size:11px; color:#f0c060; font-weight:600; margin-right:auto; align-self:center; }
  #search-filter-bar { padding:8px 18px; display:flex; gap:6px; flex-wrap:wrap; border-bottom:1px solid #333; }
  .search-filter-btn { padding:4px 10px; border-radius:4px; font-size:11px; cursor:pointer; border:1px solid #333; background:#1a1a1a; color:#888; font-family:Inter,sans-serif; transition:all 0.15s; }
  .search-filter-btn:hover { border-color:#666; color:#ccc; }
  .search-filter-btn.active { border-color:#f0c060; color:#f0c060; background:#2a2200; }
  #search-subfilter-bar { padding:4px 18px 6px; display:none; gap:4px; flex-wrap:wrap; border-bottom:1px solid #2a2a2a; }
  #search-subfilter-bar.show { display:flex; }
  .search-subfilter-btn { padding:3px 8px; border-radius:3px; font-size:10px; cursor:pointer; border:1px solid #2a2a2a; background:#111; color:#666; font-family:Inter,sans-serif; }
  .search-subfilter-btn:hover { border-color:#555; color:#aaa; }
  .search-subfilter-btn.active { border-color:#f0c060; color:#f0c060; }
  .search-btn-primary { background:#f0c060; color:#111; }
  .search-btn-primary:hover { background:#f5d080; }
  .search-btn-secondary { background:#2a2a2a; color:#aaa; border:1px solid #444 !important; }
  .search-btn-secondary:hover { color:#eee; border-color:#666 !important; }
  .search-btn-cancel { background:none; color:#666; }
  .search-btn-cancel:hover { color:#aaa; }
  #messages { flex:1; overflow-y:auto; padding:20px 24px; display:flex; flex-direction:column; gap:12px; max-width:900px; width:100%; margin:0 auto; }
  .msg { max-width:72%; }
  .msg.user { align-self:flex-end; }
  .msg.assistant { align-self:flex-start; }
  .bubble { padding:10px 16px; border-radius:10px; font-size:14px; line-height:1.65; }
  .msg.user .bubble { background:#2d4a2d; color:#d8ecd8; border-radius:10px 10px 3px 10px; }
  .msg.assistant .bubble { background:#1e1e1e; border:1px solid #2a2a2a; color:#d8d8d8; border-radius:10px 10px 10px 3px; }
  .msg .meta { font-size:10px; color:#444; margin-top:4px; padding:0 4px; }
  .msg.user .meta { text-align:right; }
  .status-msg { text-align:center; color:#555; font-size:12px; font-style:italic; }
  .memory-msg { text-align:center; color:#6a9a5a; font-size:12px; }
  #typing-indicator { display:none; background:#141414; border-top:1px solid #2d3a2d; padding:8px 20px; align-items:center; gap:10px; flex-shrink:0; }
  .typing-dots { display:flex; gap:4px; }
  .typing-dots div { width:6px; height:6px; background:#f0c060; border-radius:50%; animation:bounce 1.2s infinite; }
  .typing-dots div:nth-child(2) { animation-delay:0.2s; }
  .typing-dots div:nth-child(3) { animation-delay:0.4s; }
  @keyframes bounce { 0%,60%,100%{transform:translateY(0)} 30%{transform:translateY(-6px)} }
  #typing-text { font-size:12px; color:#888; font-family:Inter,sans-serif; }
  #send-btn.loading { position:relative; color:transparent; }
  #send-btn.loading::after { content:''; position:absolute; width:16px; height:16px; top:50%; left:50%; margin:-8px 0 0 -8px; border:2px solid #111; border-top-color:transparent; border-radius:50%; animation:spin 0.8s linear infinite; }
  @keyframes spin { to { transform:rotate(360deg); } }
  /* Progress Bar fuer Video/Bild-Generierung */
  .task-progress { background:#1a1a1a; border:1px solid #2a2a2a; border-radius:10px; padding:12px 16px; margin:10px auto; max-width:600px; font-family:Inter,sans-serif; }
  .task-progress .tp-head { display:flex; align-items:center; gap:8px; margin-bottom:8px; font-size:12px; color:#e0e0e0; }
  .task-progress .tp-pulse { width:8px; height:8px; border-radius:50%; background:#f0c060; animation:tpPulse 1.4s ease-in-out infinite; flex-shrink:0; }
  .task-progress .tp-label { flex:1; color:#c0c0c0; }
  .task-progress .tp-times { font-size:11px; color:#777; font-variant-numeric:tabular-nums; }
  .task-progress .tp-bar-outer { background:#0d0d0d; border:1px solid #2a2a2a; border-radius:6px; height:10px; overflow:hidden; }
  .task-progress .tp-bar-inner { height:100%; background:linear-gradient(90deg,#f0c060,#f5cc70); border-radius:5px; width:0%; transition:width 0.6s ease; position:relative; }
  .task-progress .tp-bar-inner::after { content:''; position:absolute; top:0; left:0; right:0; bottom:0; background:linear-gradient(90deg,transparent,rgba(255,255,255,0.18),transparent); animation:tpShimmer 1.8s linear infinite; }
  .task-progress.done .tp-bar-inner { background:#6a9a5a; }
  .task-progress.done .tp-bar-inner::after { display:none; }
  .task-progress.done .tp-pulse { background:#6a9a5a; animation:none; }
  .task-progress.error .tp-bar-inner { background:#a05050; }
  .task-progress.error .tp-bar-inner::after { display:none; }
  .task-progress.error .tp-pulse { background:#a05050; animation:none; }
  @keyframes tpPulse { 0%,100%{opacity:1;transform:scale(1)} 50%{opacity:0.45;transform:scale(0.82)} }
  @keyframes tpShimmer { 0%{transform:translateX(-100%)} 100%{transform:translateX(100%)} }
  #input-area { background:#111; border-top:1px solid #2a2a2a; padding:12px 16px; flex-shrink:0; }
  #input-row { display:flex; gap:10px; align-items:flex-end; margin-bottom:8px; }
  #msg-input { flex:1; background:#1e1e1e; border:1px solid #333; color:#e0e0e0; padding:10px 14px; font-size:14px; font-family:Inter,sans-serif; resize:none; border-radius:8px; min-height:44px; max-height:160px; outline:none; line-height:1.5; }
  #msg-input:focus { border-color:#555; }
  #send-btn { background:#f0c060; color:#111; border:none; padding:10px 20px; cursor:pointer; font-size:13px; border-radius:8px; height:44px; font-weight:700; font-family:Inter,sans-serif; flex-shrink:0; }
  #send-btn:hover { background:#f5cc70; }
  #send-btn:disabled { background:#333; color:#666; cursor:not-allowed; }
  #stop-btn { background:#cc3333; color:#fff; border:none; padding:10px 16px; cursor:pointer; font-size:13px; border-radius:8px; height:44px; font-weight:700; font-family:Inter,sans-serif; flex-shrink:0; display:none; }
  #stop-btn:hover { background:#dd4444; }
  #queue-display { display:none; padding:6px 16px; background:#161616; border-top:1px solid #252525; font-size:12px; color:#888; font-family:Inter,sans-serif; }
  .msg.queued .bubble { background:#1a1a1a; border:1px dashed #333; color:#666; font-size:12px; font-style:italic; }
  #tools-row { display:flex; gap:8px; align-items:center; }
  #url-input { width:160px; background:#1a1a1a; border:1px solid #2a2a2a; color:#888; padding:7px 12px; font-size:12px; border-radius:6px; outline:none; font-family:Inter,sans-serif; }
  #url-input::placeholder { color:#444; }
  #url-input:focus { border-color:#444; color:#ccc; }
  #file-ac-wrap { position:relative; flex:1; }
  #file-ac-input { width:100%; background:#1a1a1a; border:1px solid #2a2a2a; color:#888; padding:7px 12px; font-size:12px; border-radius:6px; outline:none; font-family:Inter,sans-serif; box-sizing:border-box; }
  #file-ac-input::placeholder { color:#444; }
  #file-ac-input:focus { border-color:#444; color:#ccc; }
  #slash-ac-dropdown {
  display:none; position:absolute; bottom:100%; left:0; right:0;
  .slash-ac-group { padding:6px 12px 3px; font-size:10px; color:#f0c060; font-weight:700; letter-spacing:1px; text-transform:uppercase; border-top:1px solid #2a2a2a; margin-top:2px; }
  .slash-ac-group:first-child { border-top:none; margin-top:0; }
  background:#1a1a1a; border:1px solid #333; border-radius:8px;
  max-height:400px; overflow-y:auto; z-index:200; margin-bottom:4px;
  box-shadow: 0 -4px 12px rgba(0,0,0,0.4);
}
.slash-ac-item {
  padding:10px 14px; cursor:pointer; font-size:13px; color:#ccc;
  font-family:Inter,system-ui,sans-serif; border-bottom:1px solid #252525;
}
.slash-ac-item:last-child { border-bottom:none; }
.slash-ac-item:hover, .slash-ac-item.active { background:#252525; color:#fff; }
#type-ac-dropdown { display:none; position:absolute; bottom:100%; left:0; right:0; background:#1a1a1a; border:1px solid #333; border-radius:8px; margin-bottom:4px; max-height:220px; overflow-y:auto; z-index:201; font-family:Inter,sans-serif; font-size:13px; }
.type-ac-item { padding:8px 14px; cursor:pointer; display:flex; justify-content:space-between; align-items:center; color:#ccc; transition:background 0.1s; }
.type-ac-item:hover, .type-ac-item.active { background:#252525; color:#fff; }
.type-ac-item .type-shortcut { color:#555; font-size:11px; font-weight:600; min-width:18px; text-align:center; }

#file-ac-dropdown { display:none; position:absolute; bottom:100%; left:0; right:0; background:#1a1a1a; border:1px solid #444; border-radius:6px; margin-bottom:4px; max-height:280px; overflow-y:auto; z-index:150; box-shadow:0 -4px 12px rgba(0,0,0,0.5); }
  .file-ac-item { padding:8px 12px; cursor:pointer; font-size:12px; color:#ccc; border-bottom:1px solid #252525; }
  .file-ac-item:hover, .file-ac-item.selected { background:#2a2a2a; color:#f0c060; }
  .file-ac-item .ac-filename { font-weight:500; }
  .file-ac-item .ac-snippet { color:#666; font-size:10px; margin-top:2px; display:block; white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
  .tool-btn { background:none; border:1px solid #333; color:#777; padding:6px 12px; cursor:pointer; font-size:12px; border-radius:6px; font-family:Inter,sans-serif; white-space:nowrap; }
  .tool-btn:hover { border-color:#555; color:#bbb; }
  #memory-search { flex:1; background:#1a1a1a; border:1px solid #2a2a2a; color:#888; padding:7px 12px; font-size:12px; border-radius:6px; outline:none; font-family:Inter,sans-serif; }
  #memory-search::placeholder { color:#444; }
  #file-input { display:none; }
  #search-toggle { background:none; border:1px solid #333; padding:6px 10px; cursor:pointer; font-size:16px; border-radius:6px; line-height:1; }
  #search-toggle.active { border-color:#f0c060; background:#2a2a1a; }
  #agent-modal { display:none; position:fixed; inset:0; background:rgba(0,0,0,0.85); z-index:200; align-items:center; justify-content:center; }
  #agent-modal.show { display:flex; }
  #agent-box { background:#1a1a1a; border:1px solid #333; padding:28px; min-width:320px; border-radius:12px; max-height:80vh; overflow-y:auto; }
  #agent-box h2 { font-size:11px; color:#666; text-transform:uppercase; letter-spacing:2px; margin-bottom:20px; font-weight:600; }
  .agent-opt { display:block; width:100%; background:#222; border:1px solid #333; color:#ccc; padding:10px 14px; margin-bottom:8px; cursor:pointer; font-size:13px; text-align:left; border-radius:8px; font-family:Inter,sans-serif; }
  .agent-opt:hover { border-color:#f0c060; color:#f0c060; }
  .agent-parent-row { display:flex; align-items:center; gap:0; margin-bottom:8px; }
  .agent-parent-row .agent-opt { flex:1; margin-bottom:0; border-radius:8px 0 0 8px; }
  .agent-expand-btn { background:#222; border:1px solid #333; border-left:none; color:#666; padding:10px 12px; cursor:pointer; font-size:13px; border-radius:0 8px 8px 0; transition:transform 0.2s; font-family:Inter,sans-serif; }
  .agent-expand-btn:hover { color:#f0c060; border-color:#f0c060; }
  .agent-expand-btn.expanded { color:#f0c060; }
  .agent-expand-btn .arrow { display:inline-block; transition:transform 0.2s; }
  .agent-expand-btn.expanded .arrow { transform:rotate(90deg); }
  .agent-subs { overflow:hidden; max-height:0; transition:max-height 0.25s ease-out; }
  .agent-subs.open { max-height:300px; transition:max-height 0.3s ease-in; }
  .agent-sub-opt { display:block; width:calc(100% - 16px); margin-left:16px; background:#1a1a1a; border:1px solid #2a2a2a; color:#999; padding:8px 14px; margin-bottom:4px; cursor:pointer; font-size:11px; text-align:left; border-radius:6px; font-family:Inter,sans-serif; }
  .agent-sub-opt:hover { border-color:#f0c060; color:#f0c060; }
  .agent-sub-opt::before { content:'\21b3 '; color:#555; }
  .new-agent-row { display:flex; gap:8px; margin-top:12px; }
  .new-agent-input { flex:1; background:#1a1a1a; border:1px solid #333; color:#ccc; padding:8px 12px; font-size:13px; border-radius:8px; outline:none; font-family:Inter,sans-serif; }
  .new-agent-btn { background:#f0c060; color:#111; border:none; padding:8px 14px; cursor:pointer; font-size:12px; border-radius:8px; font-weight:700; font-family:Inter,sans-serif; }
  #drop-overlay { display:none; position:fixed; inset:0; background:rgba(240,192,96,0.12); border:3px dashed #f0c060; z-index:500; pointer-events:none; align-items:center; justify-content:center; flex-direction:column; gap:12px; }
  #drop-overlay.active { display:flex; }
  pre { white-space:pre-wrap; word-wrap:break-word; font-family:Inter,sans-serif; font-size:14px; }
  .history-item { display:block; width:100%; background:none; border:none; border-radius:6px; padding:7px 10px; cursor:pointer; text-align:left; font-family:Inter,sans-serif; margin-bottom:3px; transition:background 0.15s; }
  .bubble { position:relative; }
  .bubble.markdown-rendered h1, .bubble.markdown-rendered h2, .bubble.markdown-rendered h3, .bubble.markdown-rendered h4 { color:#e8e8e8; margin:14px 0 6px 0; font-family:Inter,sans-serif; }
  .bubble.markdown-rendered h1 { font-size:18px; font-weight:700; }
  .bubble.markdown-rendered h2 { font-size:16px; font-weight:700; }
  .bubble.markdown-rendered h3 { font-size:15px; font-weight:600; }
  .bubble.markdown-rendered h4 { font-size:14px; font-weight:600; }
  .bubble.markdown-rendered p { margin:6px 0; }
  .bubble.markdown-rendered ul, .bubble.markdown-rendered ol { margin:6px 0 6px 20px; padding:0; }
  .bubble.markdown-rendered li { margin:3px 0; }
  .bubble.markdown-rendered a { color:#6aafff; text-decoration:underline; }
  .bubble.markdown-rendered a:hover { color:#8ac4ff; }
  .bubble.markdown-rendered strong { color:#f0f0f0; }
  .bubble.markdown-rendered code { background:#0d0d0d; padding:1px 5px; border-radius:3px; font-size:13px; font-family:monospace; }
  .bubble.markdown-rendered pre { background:#0d0d0d; border:1px solid #333; border-radius:6px; padding:10px 12px; margin:8px 0; overflow-x:auto; }
  .bubble.markdown-rendered pre code { background:none; padding:0; font-size:12px; }
  .bubble.markdown-rendered table { border-collapse:collapse; margin:8px 0; width:100%; font-size:13px; }
  .bubble.markdown-rendered th, .bubble.markdown-rendered td { border:1px solid #333; padding:6px 10px; text-align:left; }
  .bubble.markdown-rendered th { background:#1a1a1a; color:#ccc; font-weight:600; }
  .bubble.markdown-rendered blockquote { border-left:3px solid #444; padding-left:12px; margin:8px 0; color:#999; }
  .bubble.markdown-rendered hr { border:none; border-top:1px solid #333; margin:12px 0; }
  .bubble.markdown-rendered img { max-width:100%; border-radius:4px; }
  .code-block-wrapper { position:relative; background:#0d0d0d; border:1px solid #333; border-radius:6px; margin:8px 0; padding:0; }
  .code-block-wrapper pre { margin:0; padding:10px 12px; white-space:pre-wrap; word-wrap:break-word; font-size:12px; }
  .code-block-lang { position:absolute; top:4px; left:10px; font-size:9px; color:#555; text-transform:uppercase; font-family:Inter,sans-serif; }
  /* Email Search Modal */
  #email-search-modal { display:none; position:fixed; inset:0; background:rgba(0,0,0,0.85); z-index:200; align-items:center; justify-content:center; }
  #email-search-modal.show { display:flex; }
  #email-search-box { background:#1a1a1a; border:1px solid #333; border-radius:12px; padding:20px; width:600px; max-width:92vw; max-height:85vh; display:flex; flex-direction:column; }
  #email-search-box h2 { margin:0 0 12px; font-size:16px; color:#e0e0e0; }
  .esm-filters { display:flex; gap:8px; margin-bottom:10px; flex-wrap:wrap; }
  .esm-filter { flex:1; min-width:120px; }
  .esm-filter label { display:block; font-size:10px; color:#666; margin-bottom:3px; text-transform:uppercase; letter-spacing:0.5px; }
  .esm-filter input { width:100%; background:#111; border:1px solid #333; border-radius:5px; color:#e0e0e0; padding:6px 8px; font-size:13px; font-family:Inter,sans-serif; box-sizing:border-box; }
  .esm-filter input:focus { border-color:#4a8aca; outline:none; }
  .esm-filter input::placeholder { color:#555; }
  #esm-results { flex:1; overflow-y:auto; max-height:50vh; margin-top:8px; }
  .esm-result { padding:10px 12px; border:1px solid #2a2a3e; border-radius:8px; margin:5px 0; cursor:pointer; transition:background 0.15s, border-color 0.15s; background:#111; }
  .esm-result:hover { background:#1a1a3a; border-color:#4a8aca; }
  .esm-result-from { font-size:13px; color:#e0e0e0; }
  .esm-result-email { font-size:11px; color:#777; margin-left:4px; }
  .esm-result-subject { font-size:12px; color:#aaa; margin-top:3px; white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
  .esm-result-date { font-size:10px; color:#555; float:right; }
  .esm-result-meta { font-size:10px; color:#555; margin-top:2px; }
  .esm-hint { text-align:center; color:#555; font-size:12px; padding:20px 0; }
  .esm-loading { text-align:center; color:#888; font-size:12px; padding:15px 0; }
  .esm-footer { display:flex; justify-content:flex-end; margin-top:10px; gap:8px; }
  .esm-btn { padding:6px 16px; border-radius:6px; border:none; cursor:pointer; font-size:13px; font-family:Inter,sans-serif; }
  .esm-btn-close { background:#333; color:#aaa; }
  .esm-btn-close:hover { background:#444; color:#fff; }
  /* Email Card in Chat */
  .email-card { background:#1a1a2e; border:1px solid #334; border-radius:10px; margin:8px 0; overflow:hidden; font-family:Inter,sans-serif; }
  .email-card-header { padding:12px 16px 8px; border-bottom:1px solid #2a2a3e; }
  .email-card-label { font-size:10px; color:#6a8aca; text-transform:uppercase; letter-spacing:1px; font-weight:600; margin-bottom:8px; }
  .email-card-row { font-size:12px; color:#bbb; margin:3px 0; line-height:1.5; }
  .email-card-row strong { color:#e0e0e0; font-weight:500; min-width:60px; display:inline-block; }
  .email-card-subject { font-size:14px; color:#e8e8e8; font-weight:600; margin:6px 0 2px; }
  .email-card-date { font-size:11px; color:#666; }
  .email-card-body { padding:12px 16px; max-height:400px; overflow-y:auto; font-size:13px; color:#ccc; line-height:1.6; white-space:pre-wrap; word-wrap:break-word; border-top:1px solid #2a2a3e; }
  .email-card-body::-webkit-scrollbar { width:6px; }
  .email-card-body::-webkit-scrollbar-thumb { background:#444; border-radius:3px; }
  .email-card-msgid { padding:4px 16px 8px; font-size:10px; color:#444; word-break:break-all; }
  .email-card-actions { padding:8px 16px 12px; display:flex; gap:8px; border-top:1px solid #2a2a3e; }
  .email-card-btn { padding:6px 16px; border-radius:6px; border:none; cursor:pointer; font-size:12px; font-family:Inter,sans-serif; transition:background 0.15s; }
  .email-card-btn-reply { background:#4a8aca; color:#fff; }
  .email-card-btn-reply:hover { background:#5a9ada; }
  .email-card-btn-close { background:#333; color:#aaa; }
  .email-card-btn-close:hover { background:#444; color:#fff; }
  .email-search-card { background:#1a1a2e; border:1px solid #334; border-radius:8px; padding:10px 14px; margin:6px 0; cursor:pointer; transition:background 0.15s, border-color 0.15s; }
  .email-search-card:hover { background:#22224a; border-color:#4a8aca; }
  .email-search-card-from { font-size:13px; color:#e0e0e0; }
  .email-search-card-email { font-size:11px; color:#888; margin-left:4px; }
  .email-search-card-subject { font-size:12px; color:#aaa; margin-top:3px; white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
  .email-search-card-date { font-size:10px; color:#555; float:right; }
    .code-copy-btn { position:absolute; top:4px; right:6px; background:rgba(255,255,255,0.06); border:1px solid rgba(255,255,255,0.12); border-radius:5px; color:#888; font-size:11px; padding:2px 8px; cursor:pointer; transition:background 0.15s, color 0.15s; z-index:10; font-family:Inter,sans-serif; }
  .code-copy-btn:hover { background:rgba(255,255,255,0.14); color:#fff; }
  .code-copy-btn.copied { color:#4caf50; border-color:#4caf50; }
  .output-block { position:relative; background:#0f1a0f; border-left:3px solid #4a8a4a; border-radius:4px; padding:12px 16px; margin:8px 0; }
  .output-block pre { margin:0; white-space:pre-wrap; word-wrap:break-word; }
  .output-copy-btn { position:absolute; top:6px; right:6px; background:rgba(74,138,74,0.15); border:1px solid rgba(74,138,74,0.3); border-radius:5px; color:#6a9a5a; font-size:11px; padding:2px 8px; cursor:pointer; transition:background 0.15s, color 0.15s; z-index:10; font-family:Inter,sans-serif; }
  .output-copy-btn:hover { background:rgba(74,138,74,0.3); color:#8fc87f; }
  .output-copy-btn.copied { color:#4caf50; border-color:#4caf50; }
  .snippet-copy-btn { position:absolute; top:6px; right:6px; background:rgba(255,255,255,0.04); border:1px solid rgba(255,255,255,0.08); border-radius:5px; color:#555; font-size:10px; padding:2px 8px; cursor:pointer; transition:background 0.15s, color 0.15s; z-index:10; font-family:Inter,sans-serif; }
  .snippet-copy-btn:hover { background:rgba(255,255,255,0.14); color:#fff; }
  .snippet-copy-btn.copied { color:#4caf50; border-color:#4caf50; }
  .section-copy-marker { position:relative; height:0; overflow:visible; pointer-events:none; }
  .section-copy-btn { position:absolute; right:0; top:2px; background:transparent; border:1px solid rgba(255,255,255,0.08); border-radius:4px; color:#555; font-size:10px; padding:1px 7px; cursor:pointer; transition:opacity 0.15s, color 0.15s, border-color 0.15s; z-index:5; font-family:Inter,sans-serif; opacity:0.4; pointer-events:auto; }
  .section-copy-btn:hover { opacity:1; color:#f0c060; border-color:#f0c060; }
  .shortcut-label { opacity:0.35; font-size:9px; font-weight:400; margin-left:2px; }
  #find-chips-bar { display:none; padding:6px 0 4px; gap:6px; flex-wrap:wrap; align-items:center; }
  #find-chips-bar.visible { display:flex; }
  .find-chip { padding:4px 10px; background:#222; border:1px solid #444; border-radius:14px; font-size:12px; color:#aaa; cursor:pointer; font-family:Inter,sans-serif; transition:all 0.15s; white-space:nowrap; }
  .find-chip:hover { border-color:#f0c060; color:#f0c060; }
  .find-chip.active { border-color:#f0c060; color:#f0c060; background:#2a2a1a; }
  .find-chip .chip-shortcut { opacity:0.4; font-size:9px; margin-left:3px; }
  #find-live-dropdown { display:none; position:absolute; bottom:100%; left:0; right:0; background:#1a1a1a; border:1px solid #333; border-radius:8px; margin-bottom:4px; max-height:280px; overflow-y:auto; z-index:202; font-family:Inter,sans-serif; }
  .find-live-item { padding:8px 14px; cursor:pointer; color:#ccc; font-size:13px; border-bottom:1px solid #222; display:flex; justify-content:space-between; }
  .find-live-item:hover, .find-live-item.active { background:#252525; color:#fff; }
  .find-live-item .flr-name { flex:1; overflow:hidden; text-overflow:ellipsis; white-space:nowrap; }
  .find-live-item .flr-type { font-size:10px; color:#666; margin-left:8px; white-space:nowrap; }
  .section-copy-btn.copied { color:#4caf50; border-color:#4caf50; opacity:1; }
  /* TOOLTIPS_V1 — dezente Hover-Tooltips fuer Provider/Modell/Agent */
  #tt-box { position:fixed; z-index:9999; background:#1a1a1a; color:#eaeaea; border:1px solid #3a3a3a; border-radius:8px; padding:8px 12px; font-size:12px; font-family:Inter,sans-serif; max-width:320px; line-height:1.4; box-shadow:0 6px 18px rgba(0,0,0,0.6); pointer-events:none; opacity:0; transition:opacity 0.12s ease; white-space:normal; }
  #tt-box.show { opacity:1; }
  #tt-box .tt-title { color:#f0c060; font-weight:700; margin-bottom:3px; font-size:12px; }
  #tt-box .tt-body { color:#cfcfcf; font-size:11px; }
</style>
<script src="https://cdn.jsdelivr.net/npm/marked/marked.min.js"></script>
</head>
<body>
<div id="tt-box"><div class="tt-title"></div><div class="tt-body"></div></div><!-- TOOLTIPS_V1 -->
<div id="drop-overlay">
  <div style="font-size:48px;">📂</div>
  <div style="font-size:18px;font-weight:600;color:#f0c060;font-family:Inter,sans-serif;">Dateien hier ablegen</div>
</div>

<div id="header">
  <h1>ASSISTANT</h1>
  <button id="prompt-btn" onclick="toggleSidebar()">☰ Prompt <span class="shortcut-label">[P]</span></button>
  <div id="header-spacer" style="flex:1;"></div><!-- AGENT_BTN_V1 -->
  <button class="hdr-btn" onclick="newSession()" style="background:#2a3a2a;border-color:#4a6a4a;color:#a0d090;">+ Neu <span class="shortcut-label">[N]</span></button>
  <button id="agent-btn" class="hdr-btn" data-tooltip-kind="agent" onclick="showAgentModal()"><span id="agent-label">Kein Agent</span> <span class="shortcut-label">[A]</span></button>
  <select id="provider-select" class="hdr-select" onchange="onProviderChange()">
    <option>Anthropic</option>
  </select>
  <select id="model-select" class="hdr-select" onchange="onModelChange()">
    <option>Claude Sonnet 4.6</option>
  </select>
  <div id="token-indicator" style="position:relative;display:inline-block;">
    <button id="token-btn" class="hdr-btn" onclick="toggleTokenPanel()" style="font-size:11px;padding:2px 8px;min-width:80px;">~0k tokens</button>
    <div id="token-panel" style="display:none;position:absolute;right:0;top:100%;background:#1a1a1a;border:1px solid #333;border-radius:8px;padding:12px;min-width:280px;z-index:100;box-shadow:0 4px 12px rgba(0,0,0,0.5);margin-top:4px;">
      <div style="font-size:11px;color:#888;margin-bottom:8px;font-weight:600;">KONTEXT-DETAILS</div>
      <div id="token-details" style="font-size:12px;color:#ccc;"></div>
      <button onclick="slimMode()" style="margin-top:8px;background:#2a2020;border:1px solid #6a3a3a;color:#d09090;border-radius:4px;padding:4px 10px;font-size:11px;cursor:pointer;width:100%;">Slim Mode — alle Memory-Dateien entfernen</button>
    </div>
  </div>
  <div class="nav-wrap" id="nav-wrap">
    <button class="nav-btn" onclick="toggleNavMenu()" title="Navigation">&#9776;</button>
    <div class="nav-menu" id="nav-menu">
      <div class="nav-menu-section">Services</div>
      <div id="svc-list"><div class="svc-row" style="color:#666;">Lade...</div></div>
      <div class="nav-menu-divider"></div>
      <div class="nav-menu-section">Administration</div>
      <a class="nav-menu-item" onclick="navigateTo('/admin')"><span class="nav-icon">&#9881;</span><span class="nav-label">Admin Panel</span><span class="nav-hint">Status &amp; Uebersicht</span></a>
      <a class="nav-menu-item" onclick="navigateTo('/admin/access-control')"><span class="nav-icon">&#128274;</span><span class="nav-label">Access Control</span><span class="nav-hint">Matrix</span></a>
      <a class="nav-menu-item" onclick="navigateTo('/admin/permissions')"><span class="nav-icon">&#128272;</span><span class="nav-label">Berechtigungen</span><span class="nav-hint">Memory-Zugriff</span></a>
      <div class="nav-menu-divider"></div>
      <div class="nav-menu-section">Dokumentation</div>
      <a class="nav-menu-item" onclick="navigateTo('/admin/docs')"><span class="nav-icon">&#128214;</span><span class="nav-label">Technische Docs</span><span class="nav-hint">API &amp; Architektur</span></a>
      <a class="nav-menu-item" onclick="navigateTo('/admin/changelog')"><span class="nav-icon">&#128203;</span><span class="nav-label">Changelog</span><span class="nav-hint">Aenderungshistorie</span></a>
    </div>
  </div>
</div>

<div id="tab-bar">
  <button id="tab-add" onclick="addChatTab()" title="Neuer Chat-Tab">+</button>
</div>

<div id="main">
  <div id="sidebar">
    <div id="sidebar-header">
      <span style="font-size:11px;color:#aaa;font-weight:600;letter-spacing:1.5px;text-transform:uppercase;">System Prompt</span>
      <div style="display:flex;align-items:center;gap:8px;">
        <span id="memory-indicator" style="font-size:10px;color:#6a9a5a;display:none;">● inkl. Memory</span>
        <span id="sidebar-agent-name" style="font-size:11px;color:#f0c060;font-style:italic;"></span>
      </div>
    </div>
    <div id="memory-note" style="padding:6px 16px;background:#161616;border-bottom:1px solid #222;flex-shrink:0;display:none;">
      <span style="font-size:10px;color:#666;">Hier wird der vollstaendige aktive System Prompt angezeigt. 'Basis speichern' speichert nur den oberen Teil (bis --- GEDAECHTNIS).</span>
    </div>
    <textarea id="prompt-editor" placeholder="Kein Agent aktiv..."></textarea>
    <div id="sidebar-btns">
      <button onclick="savePrompt()" style="flex:1;background:#f0c060;color:#111;border:none;padding:7px;border-radius:6px;font-size:12px;font-weight:700;cursor:pointer;font-family:Inter,sans-serif;">Basis speichern</button>
      <button onclick="reloadPrompt()" style="background:#2a2a2a;color:#aaa;border:1px solid #444;padding:7px 12px;border-radius:6px;font-size:12px;cursor:pointer;font-family:Inter,sans-serif;">Neu laden</button>
    </div>
    <div id="history-section">
      <span>Konversationen</span>
      <div id="history-list">
        <p style="font-size:11px;color:#555;padding:8px;font-style:italic;">Kein Agent aktiv</p>
      </div>
    </div>
  </div>

  <div id="chat-area">
    <div id="ctx-bar">
      <span class="ctx-label">Kontext</span>
      <div id="ctx-items" style="display:flex;gap:6px;flex-wrap:wrap;"></div>
    </div>
    <div id="messages"></div>
    <div id="typing-indicator">
      <div class="typing-dots"><div></div><div></div><div></div></div>
      <span id="typing-text">Denkt nach...</span>
    </div>
    <div id="queue-display"><span id="queue-text"></span></div>
    <div id="input-area">
      <div id="find-chips-bar">
        <span style="font-size:10px;color:#555;margin-right:4px;">Typ:</span>
        <span class="find-chip" data-cat="email" onclick="toggleFindChip('email')">✉ E-Mail<span class="chip-shortcut">[1]</span></span>
        <span class="find-chip" data-cat="webclip" onclick="toggleFindChip('webclip')">🌐 Web Clip<span class="chip-shortcut">[2]</span></span>
        <span class="find-chip" data-cat="document" onclick="toggleFindChip('document')">📄 Dokument<span class="chip-shortcut">[3]</span></span>
        <span class="find-chip" data-cat="conversation" onclick="toggleFindChip('conversation')">💬 Konversation<span class="chip-shortcut">[4]</span></span>
        <span class="find-chip" data-cat="screenshot" onclick="toggleFindChip('screenshot')">📸 Screenshot<span class="chip-shortcut">[5]</span></span>
      </div>
      <div id="input-row" style="position:relative;">
        <div id="find-live-dropdown"></div>
        <textarea id="msg-input" placeholder="Nachricht..." rows="1" disabled onkeydown="onKey(event)" oninput="autoResize(this); onInputHandler(this);"></textarea>
        <button id="send-btn" title="Ctrl+Enter" onclick="sendMessage()" disabled>Senden</button>
        <button id="stop-btn" onclick="stopQueue()">&#9209; Stop</button>
      </div>
      <div id="tools-row">
        <input type="text" id="url-input" placeholder="URL..." onkeydown="if(event.key==='Enter')addUrl()" />
        <button class="tool-btn" onclick="addUrl()">+ URL</button>

        <button class="tool-btn" onclick="document.getElementById('file-input').click()">+ Datei <span class="shortcut-label">[U]</span></button>
        <input type="file" id="file-input" multiple onchange="addFileFromInput()" />
      </div>
    </div>
  </div>
</div>

<div id="agent-modal">
  <div id="agent-box">
    <h2>Agent auswaehlen</h2>
    <div id="agent-list"></div>
    <div class="new-agent-row">
      <input type="text" class="new-agent-input" id="new-agent-name" placeholder="Neuer Agent..." />
      <button class="new-agent-btn" onclick="createAgent()">+ Neu</button>
    </div>
  </div>
</div>

<div id="email-search-modal">
  <div id="email-search-box">
    <h2>\u2709 E-Mail suchen</h2>
    <div class="esm-filters">
      <div class="esm-filter"><label>Von</label><input type="text" id="esm-from" placeholder="Absender..." /></div>
      <div class="esm-filter"><label>Betreff</label><input type="text" id="esm-subject" placeholder="Betreff..." /></div>
    </div>
    <div class="esm-filters">
      <div class="esm-filter"><label>An / CC</label><input type="text" id="esm-to" placeholder="Empfaenger..." /></div>
      <div class="esm-filter"><label>Freitext</label><input type="text" id="esm-body" placeholder="Inhalt..." /></div>
    </div>
    <div id="esm-results"><div class="esm-hint">Mindestens 2 Zeichen in ein Feld eingeben...</div></div>
    <div class="esm-footer">
      <button class="esm-btn esm-btn-close" onclick="closeEmailSearchModal()">Schliessen</button>
    </div>
  </div>
</div>

<div id="search-overlay">
  <div id="search-panel">
    <div id="search-panel-header">
      <h3 id="search-panel-title">🔍 Suche</h3>
      <span id="search-panel-count"></span>
    </div>
    <div id="search-filter-bar">
      <button class="search-filter-btn active" data-filter="all" onclick="applySearchFilter('all')">Alle</button>
      <button class="search-filter-btn" data-filter="email" onclick="applySearchFilter('email')">&#9993; E-Mail</button>
      <button class="search-filter-btn" data-filter="webclip" onclick="applySearchFilter('webclip')">&#127760; Web Clip</button>
      <button class="search-filter-btn" data-filter="document" onclick="applySearchFilter('document')">&#128196; Dokument</button>
      <button class="search-filter-btn" data-filter="conversation" onclick="applySearchFilter('conversation')">&#128172; Konversation</button>
      <button class="search-filter-btn" data-filter="screenshot" onclick="applySearchFilter('screenshot')">&#128248; Screenshot</button>
      <button class="search-filter-btn" data-filter="whatsapp" onclick="applySearchFilter('whatsapp')">&#128172; WhatsApp</button>
    </div>
    <div id="search-subfilter-bar"></div>
    <div id="search-info-bar" style="padding:4px 18px;display:flex;align-items:center;gap:12px;font-size:13px;color:#aaa;">
      <span id="search-hit-count" style="font-weight:600;color:#e0e0e0;"></span>
      <button id="search-toggle-all-btn" onclick="toggleAllSearchCheckboxes()" style="background:none;border:1px solid #555;color:#ccc;padding:3px 10px;border-radius:4px;font-size:12px;cursor:pointer;font-family:Inter,sans-serif;">Alle markieren</button>
    </div>
    <div id="search-results-list"></div>
    <div id="search-panel-footer">
      <span id="search-selection-counter">0 / 50 ausgewaehlt</span>
      <button class="search-btn-cancel" onclick="closeSearchDialog(true)">Ohne Dateien senden</button>
      <button class="search-btn-secondary" onclick="loadSelectedResults()">Auswahl laden</button>
      <button class="search-btn-primary" onclick="loadAllResults()">Alle laden (max 50)</button>
    </div>
  </div>
</div>

<script>
// ─── SESSION ID + TABS ──────────────────────────────────────────────────────
function makeSessionId() {
  return 'sess_' + Date.now() + '_' + Math.random().toString(36).substr(2, 9);
}
// Tab state
var _tabs = [];  // [{id, sessionId, agentName, label, messagesHtml, ctxHtml}]
var _activeTabId = null;
var SESSION_ID = '';

function addChatTab(agentName) {
  var tabId = 'tab_' + Date.now();
  var sessId = makeSessionId();
  var tab = {id: tabId, sessionId: sessId, agentName: agentName || '', label: agentName || 'Neuer Chat', messagesHtml: '', ctxHtml: ''};
  _tabs.push(tab);
  switchToTab(tabId);
  renderTabs();
  if (!agentName) showAgentModal();
  return tabId;
}

function switchToTab(tabId) {
  // Save current tab state
  if (_activeTabId) {
    var cur = _tabs.find(function(t){ return t.id === _activeTabId; });
    if (cur) {
      cur.messagesHtml = document.getElementById('messages').innerHTML;
      cur.ctxHtml = document.getElementById('ctx-items').innerHTML;
      cur.agentName = getAgentName();
      var lbl = document.getElementById('agent-label');
      cur.label = lbl ? lbl.textContent : cur.agentName;
      // Aktuell gewaehlten Provider/Model merken, damit Rueckwechsel den
      // Dropdown auf den tatsaechlich aktiven Stand dieses Tabs zurueckholt.
      var psCur = document.getElementById('provider-select');
      var msCur = document.getElementById('model-select');
      if (psCur && msCur && psCur.value) {
        cur.provider = psCur.value;
        cur.model_id = msCur.value;
        cur.model_name = (msCur.selectedOptions[0] && msCur.selectedOptions[0].textContent) || msCur.value;
      }
    }
  }
  // Activate new tab
  _activeTabId = tabId;
  var tab = _tabs.find(function(t){ return t.id === tabId; });
  if (!tab) return;
  SESSION_ID = tab.sessionId;
  // Restore DOM
  document.getElementById('messages').innerHTML = tab.messagesHtml || '';
  document.getElementById('ctx-items').innerHTML = tab.ctxHtml || '';
  if (tab.agentName) {
    var lbl = document.getElementById('agent-label');
    if (lbl) {
      lbl.textContent = tab.label || tab.agentName;
      lbl.dataset.agentName = tab.agentName;
    }
    document.getElementById('msg-input').disabled = false;
    document.getElementById('send-btn').disabled = false;
    document.getElementById('msg-input').placeholder = 'Nachricht an ' + (tab.label || tab.agentName) + '...';
    loadHistory(tab.agentName);
    // Provider/Model-Dropdown auf den Stand des Tabs zurueckholen + Backend
    // neu synchronisieren, damit UI und state['provider'] nicht divergieren.
    if (tab.provider && tab.model_id) {
      restoreTabProvider(tab);
    }
  }
  // Per-Tab Processing-State ins DOM spiegeln (typing-indicator, stop-btn,
  // queue-display) und im Hintergrund gesammelte Responses einfliessen lassen
  renderActiveTabState();
  renderTabs();
}

function closeTab(tabId, evt) {
  if (evt) { evt.stopPropagation(); evt.preventDefault(); }
  if (_tabs.length <= 1) return; // mindestens 1 Tab
  var idx = _tabs.findIndex(function(t){ return t.id === tabId; });
  if (idx < 0) return;
  var closingTab = _tabs[idx];
  // Polling und Typing-Animation fuer diese Session beenden, sonst laufen
  // sie im Hintergrund weiter und lecken Speicher.
  if (closingTab && closingTab.sessionId && _tabStates[closingTab.sessionId]) {
    var cs = _tabStates[closingTab.sessionId];
    if (cs.pollIntervalId) { clearInterval(cs.pollIntervalId); cs.pollIntervalId = null; }
    if (cs.typingIntervalId) { clearInterval(cs.typingIntervalId); cs.typingIntervalId = null; }
    delete _tabStates[closingTab.sessionId];
  }
  _tabs.splice(idx, 1);
  if (_activeTabId === tabId) {
    var newIdx = Math.min(idx, _tabs.length - 1);
    switchToTab(_tabs[newIdx].id);
  }
  renderTabs();
}

function renderTabs() {
  var bar = document.getElementById('tab-bar');
  var addBtn = document.getElementById('tab-add');
  // Remove old tabs (keep + button)
  var old = bar.querySelectorAll('.chat-tab');
  old.forEach(function(el){ el.remove(); });
  // Add tabs before + button
  _tabs.forEach(function(tab) {
    var el = document.createElement('button');
    el.className = 'chat-tab' + (tab.id === _activeTabId ? ' active' : '');
    el.onclick = function(){ switchToTab(tab.id); };
    var nameSpan = '<span class="tab-name">' + escHtml(tab.label || 'Neuer Chat') + '</span>';
    var closeSpan = _tabs.length > 1 ? '<span class="tab-close" onclick="closeTab(\\'' + tab.id + '\\', event)">&times;</span>' : '';
    el.innerHTML = nameSpan + closeSpan;
    bar.insertBefore(el, addBtn);
  });
}

function updateActiveTabLabel(name, displayName) {
  var tab = _tabs.find(function(t){ return t.id === _activeTabId; });
  if (tab) {
    tab.agentName = name;
    tab.label = displayName || name;
    renderTabs();
  }
}

// Init first tab
(function() {
  // WICHTIG: Session-ID wird bewusst NICHT mehr in sessionStorage persistiert.
  // In pywebview kann sessionStorage zwischen Fenstern desselben Origins geteilt
  // werden — zwei Fenster wuerden dann dieselbe Server-Session teilen, und
  // Responses eines Fensters taucht im anderen auf (cross-window pollution).
  // Jede Seiten-Ladung bekommt daher eine komplett frische Session-ID.
  var sessId = makeSessionId();
  var savedAgent = localStorage.getItem('last_active_agent');
  var tab = {id: 'tab_init', sessionId: sessId, agentName: savedAgent || '', label: savedAgent || 'Neuer Chat', messagesHtml: '', ctxHtml: ''};
  _tabs.push(tab);
  _activeTabId = 'tab_init';
  SESSION_ID = sessId;
  // Alten, ggf. geteilten Wert in sessionStorage aufraeumen
  try { sessionStorage.removeItem('assistant_session_id'); } catch(e) {}
  renderTabs();
})();

let sidebarOpen = true;
let currentModel = {provider:'anthropic', model_id:'claude-sonnet-4-6', model_name:'Claude Sonnet 4.6'};
let searchVisible = false;
function getAgentName() { return document.getElementById('agent-label').dataset.agentName || document.getElementById('agent-label').textContent; }

// ─── TOOLTIPS_V1 ─────────────────────────────────────────────────────────────
var PROVIDER_TOOLTIPS = {
  'Anthropic': 'Anthropic Claude — Stark in Analyse, Schreiben, Coding und komplexem Reasoning',
  'OpenAI': 'OpenAI GPT — Vielseitig, stark in Coding, Bildgenerierung (gpt-image-1) und strukturierten Outputs',
  'Google Gemini': 'Google Gemini — Stark in Multimodal, langer Kontext, Video- & Bildgenerierung (Veo, Imagen)',
  'Gemini': 'Google Gemini — Stark in Multimodal, langer Kontext, Video- & Bildgenerierung (Veo, Imagen)',
  'Mistral': 'Mistral — Schnell und effizient, stark in europaeischen Sprachen und Code',
  'Perplexity': 'Perplexity — Spezialisiert auf Web-Suche und aktuelle Informationen'
};
var MODEL_TOOLTIPS = {
  'Claude Sonnet 4.6': 'Bestes Preis-Leistungs-Verhaeltnis — schnell, intelligent, fuer taegliche Aufgaben',
  'Claude Opus 4.6': 'Staerkstes Claude-Modell — fuer komplexe, anspruchsvolle Aufgaben',
  'Claude Haiku 4.5': 'Schnellstes Claude-Modell — fuer einfache, schnelle Antworten',
  'GPT-4o': 'OpenAIs Flagship — stark in Multimodal, Coding, strukturierten Outputs',
  'GPT-4o Mini': 'Schnell und guenstig — fuer einfache Aufgaben und hohe Volumen',
  'o1': 'OpenAI Reasoning-Modell — fuer komplexe logische und mathematische Probleme',
  'Mistral Large': 'Mistral Flagship — stark in Mehrsprachigkeit und komplexen Aufgaben',
  'Mistral Small': 'Mistral kompakt — effizient fuer einfachere Aufgaben',
  'Mistral Nemo': 'Mistral Nemo — leichtgewichtig, sehr schnell',
  'Gemini 2.0 Flash': 'Schnell und guenstig — fuer alltaegliche Aufgaben mit Multimodal-Unterstuetzung',
  'Gemini 2.5 Pro': 'Googles staerkstes Modell — langer Kontext (1M tokens), komplexes Reasoning',
  'Gemini 2.5 Flash': 'Googles schnellstes Flagship — ideal fuer Video/Bild-Generierung und schnelle Antworten',
  'Gemini 3 Flash (Preview)': 'Gemini 3 Flash Preview — neue Generation, sehr schnell',
  'Gemini 3 Pro (Preview)': 'Gemini 3 Pro Preview — neue Generation, fuer anspruchsvolle Aufgaben',
  'Gemini 3.1 Pro (Preview)': 'Gemini 3.1 Pro Preview — neueste Generation mit Reasoning',
  'Sonar': 'Perplexity Web-Suche — schnell, aktuelle Informationen',
  'Sonar Pro': 'Perplexity Pro-Suche — tiefere Recherche, mehr Quellen',
  'Sonar Reasoning': 'Perplexity mit Reasoning — Suche + logische Schlussfolgerung',
  'Sonar Reasoning Pro': 'Perplexity Pro Reasoning — tiefste Analyse mit aktuellen Daten',
  'Sonar Deep Research': 'Perplexity Deep Research — ausfuehrliche Recherche fuer komplexe Themen'
};
var AGENT_DESCRIPTIONS = {};  // befuellt durch loadAgents()
var _ttTimer = null;

function ttGetContent(el) {
  var kind = el.getAttribute('data-tooltip-kind') || '';
  if (kind === 'provider') {
    var ps = document.getElementById('provider-select');
    var name = ps ? (ps.selectedOptions[0] ? ps.selectedOptions[0].textContent : ps.value) : '';
    var body = PROVIDER_TOOLTIPS[name] || ('KI-Provider: ' + name);
    return {title: name, body: body};
  }
  if (kind === 'model') {
    var ms = document.getElementById('model-select');
    if (!ms || !ms.selectedOptions[0]) return {title: 'Modell', body: ''};
    var txt = ms.selectedOptions[0].textContent;
    // Capabilities-Emojis am Ende abtrennen
    var name = txt.replace(/[^a-zA-Z0-9().\s-]+$/g, '').trim();
    var body = MODEL_TOOLTIPS[name];
    if (!body) {
      var pv = document.getElementById('provider-select');
      var pname = pv && pv.selectedOptions[0] ? pv.selectedOptions[0].textContent : '';
      body = 'KI-Modell von ' + (pname || 'unbekannt');
    }
    return {title: name, body: body};
  }
  if (kind === 'agent') {
    var lbl = document.getElementById('agent-label');
    var name = lbl ? (lbl.dataset.agentName || lbl.textContent) : '';
    if (!name || name === 'Kein Agent') {
      return {title: 'Kein Agent aktiv', body: 'Waehle oben rechts einen Agenten aus.'};
    }
    var desc = AGENT_DESCRIPTIONS[name] || '';
    return {title: 'Agent: ' + name, body: desc || 'Aktiver Agent - klicken zum Wechseln.'};
  }
  // Generic fallback
  var t = el.getAttribute('data-tooltip') || '';
  return {title: '', body: t};
}

function ttShow(el) {
  var box = document.getElementById('tt-box');
  if (!box) return;
  var c = ttGetContent(el);
  if (!c.body && !c.title) return;
  box.querySelector('.tt-title').textContent = c.title || '';
  box.querySelector('.tt-title').style.display = c.title ? 'block' : 'none';
  box.querySelector('.tt-body').textContent = c.body || '';
  var r = el.getBoundingClientRect();
  // Zuerst sichtbar machen um Groesse zu messen
  box.style.left = '-9999px'; box.style.top = '-9999px';
  box.classList.add('show');
  var bw = box.offsetWidth;
  var bh = box.offsetHeight;
  var left = r.left + r.width/2 - bw/2;
  if (left < 8) left = 8;
  if (left + bw > window.innerWidth - 8) left = window.innerWidth - bw - 8;
  var top = r.bottom + 8;
  if (top + bh > window.innerHeight - 8) top = r.top - bh - 8; // ueber dem Element
  box.style.left = left + 'px';
  box.style.top = top + 'px';
}

function ttHide() {
  if (_ttTimer) { clearTimeout(_ttTimer); _ttTimer = null; }
  var box = document.getElementById('tt-box');
  if (box) box.classList.remove('show');
}

function ttAttach(el) {
  if (!el || el.dataset.ttBound === '1') return;
  el.dataset.ttBound = '1';
  el.addEventListener('mouseenter', function(){
    if (_ttTimer) clearTimeout(_ttTimer);
    _ttTimer = setTimeout(function(){ ttShow(el); }, 300);
  });
  el.addEventListener('mouseleave', ttHide);
  el.addEventListener('mousedown', ttHide);
}

function ttAttachAll() {
  var ps = document.getElementById('provider-select');
  if (ps) { ps.setAttribute('data-tooltip-kind','provider'); ttAttach(ps); }
  var ms = document.getElementById('model-select');
  if (ms) { ms.setAttribute('data-tooltip-kind','model'); ttAttach(ms); }
  var ab = document.getElementById('agent-btn');
  if (ab) ttAttach(ab);
}

// ─── PROVIDERS ────────────────────────────────────────────────────────────────
async function loadProviders() {
  const r = await fetch('/models');
  const data = await r.json();
  const ps = document.getElementById('provider-select');
  const ms = document.getElementById('model-select');
  ps.innerHTML = '';
  data.forEach(p => {
    const o = document.createElement('option');
    o.value = p.provider; o.textContent = p.name;
    ps.appendChild(o);
  });
  if (data.length) populateModels(data[0]);
  const saved = localStorage.getItem('claude_model');
  if (saved) {
    try {
      const m = JSON.parse(saved);
      currentModel = m;
      for (let o of ps.options) { if (o.value === m.provider) { ps.value = m.provider; break; } }
      const pd = data.find(p => p.provider === m.provider);
      if (pd) { populateModels(pd); ms.value = m.model_id; }
    } catch(e) {}
  }
  ttAttachAll(); // TOOLTIPS_V1
  // Initiales Laden der Agent-Descriptions fuer Tooltip auf Agent-Button
  try { fetch('/agents').then(function(r){return r.json();}).then(function(ags){
    ags.forEach(function(a){
      if (a.description) AGENT_DESCRIPTIONS[a.name] = a.description;
      if (a.subagents) a.subagents.forEach(function(s){ if (s.description) AGENT_DESCRIPTIONS[s.name] = s.description; });
    });
  }); } catch(e) {}
}

function populateModels(providerData) {
  const ms = document.getElementById('model-select');
  ms.innerHTML = '';
  (providerData.models || []).forEach(m => {
    const o = document.createElement('option');
    o.value = m.id; o.textContent = m.name + (m.capabilities && m.capabilities.length ? ' ' + m.capabilities.join('') : '');
    ms.appendChild(o);
  });
}

function onProviderChange() {
  fetch('/models').then(r=>r.json()).then(data => {
    const pv = document.getElementById('provider-select').value;
    const pd = data.find(p => p.provider === pv);
    if (pd) populateModels(pd);
    onModelChange();
  });
}

async function onModelChange() {
  const pv = document.getElementById('provider-select').value;
  const mv = document.getElementById('model-select').value;
  const mt = document.getElementById('model-select').selectedOptions[0]?.textContent || mv;
  const r = await fetch('/select_model', {method:'POST', headers:{'Content-Type':'application/json'},
    body: JSON.stringify({provider: pv, model_id: mv, model_name: mt, session_id: SESSION_ID})});
  const data = await r.json();
  if (data.ok) {
    currentModel = {provider: pv, model_id: mv, model_name: mt};
    localStorage.setItem('claude_model', JSON.stringify(currentModel));
    // Aktuellen Tab aktualisieren — so wird die Auswahl beim Tab-Switch
    // wiederhergestellt und bleibt in sync mit dem Backend.
    if (typeof _activeTabId !== 'undefined' && _activeTabId && typeof _tabs !== 'undefined') {
      var curTab = _tabs.find(function(t){ return t.id === _activeTabId; });
      if (curTab) { curTab.provider = pv; curTab.model_id = mv; curTab.model_name = mt; }
    }
    var agName = document.getElementById('agent-label').dataset.agentName; if (agName) fetch('/api/agent-model-preference', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({agent:agName, provider:pv, model:mv})});
  }
}

// Stellt Provider/Model-Dropdown auf den Stand eines Tabs + synchronisiert
// den Backend-State der zugehoerigen Session. Verhindert, dass UI-Auswahl
// und state['provider'] divergieren.
async function restoreTabProvider(tab) {
  if (!tab || !tab.provider || !tab.model_id) return;
  try {
    var ps = document.getElementById('provider-select');
    var ms = document.getElementById('model-select');
    if (!ps || !ms) return;
    ps.value = tab.provider;
    var mdata = await fetch('/models').then(function(r){ return r.json(); });
    var pd = mdata.find(function(p){ return p.provider === tab.provider; });
    if (pd) {
      populateModels(pd);
      ms.value = tab.model_id;
    }
    // Backend fuer diese Session auf denselben Stand bringen, damit
    // der naechste /chat-Request tatsaechlich diesen Provider nutzt.
    await fetch('/select_model', {method:'POST', headers:{'Content-Type':'application/json'},
      body: JSON.stringify({provider: tab.provider, model_id: tab.model_id,
                             model_name: tab.model_name || tab.model_id,
                             session_id: SESSION_ID})});
    currentModel = {provider: tab.provider, model_id: tab.model_id, model_name: tab.model_name || tab.model_id};
  } catch(e) { console.log('restoreTabProvider error:', e); }
}

// ─── AGENTS ────────────────────────────────────────────────────────────────────
async function loadAgents() {
  const r = await fetch('/agents');
  const agents = await r.json();
  // TOOLTIPS_V1: Agent-Beschreibungen in globaler Map cachen
  try {
    agents.forEach(function(a){
      if (a.description) AGENT_DESCRIPTIONS[a.name] = a.description;
      if (a.subagents) a.subagents.forEach(function(s){ if (s.description) AGENT_DESCRIPTIONS[s.name] = s.description; });
    });
  } catch(e) {}
  const list = document.getElementById('agent-list');
  list.innerHTML = '';
  const expandedAgent = localStorage.getItem('agent_expanded') || '';
  agents.forEach(a => {
    if (a.has_subagents) {
      // Parent row with expand button
      const row = document.createElement('div');
      row.className = 'agent-parent-row';
      const btn = document.createElement('button');
      btn.className = 'agent-opt';
      btn.textContent = a.label;
      btn.onclick = () => selectAgent(a.name);
      row.appendChild(btn);
      const expBtn = document.createElement('button');
      expBtn.className = 'agent-expand-btn' + (expandedAgent === a.name ? ' expanded' : '');
      expBtn.innerHTML = '<span class="arrow">\u25B6</span>';
      const subsDiv = document.createElement('div');
      subsDiv.className = 'agent-subs' + (expandedAgent === a.name ? ' open' : '');
      expBtn.onclick = (e) => {
        e.stopPropagation();
        const isOpen = subsDiv.classList.contains('open');
        // Close all others
        document.querySelectorAll('.agent-subs.open').forEach(s => s.classList.remove('open'));
        document.querySelectorAll('.agent-expand-btn.expanded').forEach(b => b.classList.remove('expanded'));
        if (!isOpen) {
          subsDiv.classList.add('open');
          expBtn.classList.add('expanded');
          localStorage.setItem('agent_expanded', a.name);
        } else {
          localStorage.removeItem('agent_expanded');
        }
      };
      row.appendChild(expBtn);
      list.appendChild(row);
      // Sub-agents
      a.subagents.forEach(sub => {
        const subBtn = document.createElement('button');
        subBtn.className = 'agent-sub-opt';
        subBtn.textContent = sub.label;
        subBtn.onclick = () => selectAgent(sub.name);
        subsDiv.appendChild(subBtn);
      });
      list.appendChild(subsDiv);
    } else {
      // Simple agent without sub-agents
      const btn = document.createElement('button');
      btn.className = 'agent-opt';
      btn.textContent = a.label;
      btn.onclick = () => selectAgent(a.name);
      btn.style.marginBottom = '8px';
      list.appendChild(btn);
    }
  });
}

function showAgentModal() { loadAgents(); const m = document.getElementById('agent-modal'); m.classList.add('show'); m.style.display = 'flex'; }

async function selectAgent(name) {
  const m = document.getElementById('agent-modal'); m.classList.remove('show'); m.style.display = 'none';
  // Display name: "signicat > outbound" for sub-agents
  const displayName = name.includes('_') ? name.split('_')[0]+' \u203a '+name.split('_').slice(1).join('_') : name;
  addStatusMsg('Lade Gedaechtnis fuer ' + displayName + '...');
  let r, data;
  try {
    r = await fetch('/select_agent', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({agent:name, session_id:SESSION_ID})});
    data = await r.json();
  } catch(e) {
    localStorage.removeItem('last_active_agent');
    addStatusMsg('Agent-Laden fehlgeschlagen: ' + e.message);
    showAgentModal();
    return;
  }
  if (!data || data.ok === false) {
    localStorage.removeItem('last_active_agent');
    addStatusMsg('Agent "' + name + '" nicht verfuegbar' + (data && data.error ? ' ('+data.error+')' : ''));
    showAgentModal();
    return;
  }
  // Persist for auto-restore on next page load (survives server restart)
  localStorage.setItem('last_active_agent', name);
  document.getElementById('agent-label').textContent = displayName;
  document.getElementById('agent-label').dataset.agentName = name;
  updateActiveTabLabel(name, displayName);
  document.getElementById('msg-input').disabled = false;
  document.getElementById('send-btn').disabled = false;
  document.getElementById('msg-input').placeholder = 'Nachricht an ' + displayName + '...';
  document.getElementById('messages').innerHTML = '';
  document.getElementById('ctx-items').innerHTML = '';
  // Restore draft text for this agent
  var savedDraft = localStorage.getItem('draft_' + name);
  if (savedDraft) { document.getElementById('msg-input').value = savedDraft; autoResize(document.getElementById('msg-input')); }
  // Render recovered messages from today's session
  if (data.recovered_messages && data.recovered_messages.length > 0) {
    data.recovered_messages.forEach(function(m) {
      if (m.role === 'user') { addUserMsg(m.content); }
      else if (m.role === 'assistant') { addBotMsg(m.content); }
    });
  }
  // Restore saved provider/model preference for this agent (synchronous from response)
  if (data.pref_provider && data.pref_model) {
    try {
      var ps = document.getElementById('provider-select');
      ps.value = data.pref_provider;
      var mdata = await fetch('/models').then(r=>r.json());
      var pd = mdata.find(p => p.provider === data.pref_provider);
      if (pd) {
        populateModels(pd);
        var ms = document.getElementById('model-select');
        ms.value = data.pref_model;
        await onModelChange();
      }
    } catch(e) { console.log('Preference restore error:', e); }
  }
  addStatusMsg('Agent "' + displayName + '" \u2014 ' + currentModel.model_name);
  if (data.memory_info) addMemoryMsg(data.memory_info);
  if (data.base_prompt) updateSidebar(name, data.base_prompt);
  loadHistory(name);
}

async function createAgent() {
  const name = document.getElementById('new-agent-name').value.trim();
  if (!name) return;
  await fetch('/create_agent', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({name, session_id:SESSION_ID})});
  document.getElementById('new-agent-name').value = '';
  selectAgent(name);
}

// ─── TOKEN/KONTEXT MANAGER ──────────────────────────────────────────────────
async function updateTokenCount() {
  try {
    const r = await fetch('/api/context-info?session_id=' + SESSION_ID);
    const d = await r.json();
    var total = d.total_tokens || 0;
    var btn = document.getElementById('token-btn');
    if (!btn) return;
    var label = total >= 1000 ? '~' + (total / 1000).toFixed(1) + 'k tokens' : '~' + total + ' tokens';
    btn.textContent = label;
    if (d.limit_warning) {
      btn.style.borderColor = '#d09050';
      btn.style.color = '#d09050';
    } else {
      btn.style.borderColor = '';
      btn.style.color = '';
    }
    // Update details panel
    var det = document.getElementById('token-details');
    if (det) {
      var html = '<div style="margin-bottom:6px;">System-Prompt: <b>' + d.system_prompt_tokens + '</b></div>';
      html += '<div style="margin-bottom:6px;">Konversation: <b>' + d.conversation_tokens + '</b></div>';
      if (d.memory_files_loaded && d.memory_files_loaded.length > 0) {
        html += '<div style="margin-bottom:4px;color:#888;">Memory-Dateien:</div>';
        for (var i = 0; i < d.memory_files_loaded.length; i++) {
          var mf = d.memory_files_loaded[i];
          html += '<div style="display:flex;justify-content:space-between;padding:2px 0;font-size:11px;">';
          html += '<span style="max-width:180px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;">' + escHtml(mf.name) + '</span>';
          html += '<span style="color:#888;">' + mf.tokens + '</span>';
          html += '</div>';
        }
      }
      html += '<div style="margin-top:8px;border-top:1px solid #333;padding-top:6px;font-weight:600;">Gesamt: ' + total + ' tokens</div>';
      det.innerHTML = html;
    }
  } catch(e) {}
}
function toggleTokenPanel() {
  var panel = document.getElementById('token-panel');
  if (!panel) return;
  var visible = panel.style.display !== 'none';
  panel.style.display = visible ? 'none' : 'block';
  if (!visible) updateTokenCount();
}
async function slimMode() {
  try {
    await fetch('/remove_all_ctx', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({session_id:SESSION_ID})});
    document.getElementById('ctx-items').innerHTML = '';
    updateTokenCount();
    addStatusMsg('Alle Memory-Dateien aus dem Kontext entfernt.');
  } catch(e) {}
}
// Update token count after each message
setInterval(function() {
  var btn = document.getElementById('token-btn');
  if (btn) updateTokenCount();
}, 30000);

async function newSession() {
  await fetch('/close_session', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({session_id:SESSION_ID})});
  document.getElementById('messages').innerHTML = '';
  document.getElementById('ctx-items').innerHTML = '';
  addStatusMsg('Neue Session gestartet.');
  const name = getAgentName();
  if (name && name !== 'Kein Agent') {
    const r = await fetch('/select_agent', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({agent:name, session_id:SESSION_ID})});
    const data = await r.json();
    if (data.memory_info) addMemoryMsg(data.memory_info);
    loadHistory(name);
  }
}

// ─── NAV MENU ────────────────────────────────────────────────────────────────────
function toggleNavMenu() {
  var menu = document.getElementById('nav-menu');
  var btn = document.querySelector('.nav-btn');
  var isOpen = menu.classList.contains('open');
  menu.classList.toggle('open');
  btn.classList.toggle('active');
  if (!isOpen) {
    loadServices();
    setTimeout(function() {
      document.addEventListener('click', closeNavOnClickOutside, {once: true});
    }, 10);
  }
}
function closeNavOnClickOutside(e) {
  const wrap = document.getElementById('nav-wrap');
  if (!wrap.contains(e.target)) {
    document.getElementById('nav-menu').classList.remove('open');
    document.querySelector('.nav-btn').classList.remove('active');
  } else {
    setTimeout(function() {
      document.addEventListener('click', closeNavOnClickOutside, {once: true});
    }, 10);
  }
}
function navigateTo(path) {
  document.getElementById('nav-menu').classList.remove('open');
  document.querySelector('.nav-btn').classList.remove('active');
  window.location.href = path;
}
async function loadServices() {
  try {
    var r = await fetch('/api/services');
    var d = await r.json();
    var el = document.getElementById('svc-list');
    if (!el) return;
    var h = '';
    (d.services || []).forEach(function(s) {
      var dot = s.online ? 'online' : 'offline';
      var label = s.online ? 'online' : 'offline';
      var port = s.port ? ' :' + s.port : '';
      var extra = '';
      if (s.last_run) extra = '<span style="font-size:9px;color:#666;display:block;">Letzter Lauf: ' + s.last_run + '</span>';
      if (s.periodic) label = (s.online ? 'aktiv' : 'inaktiv') + ' (' + s.periodic + ')';
      h += '<div class="svc-row" style="flex-wrap:wrap;">';
      h += '<span class="svc-dot ' + dot + '"></span>';
      h += '<span class="svc-name">' + s.name + port + '</span>';
      h += '<span style="font-size:10px;color:' + (s.online ? '#4ade80' : '#f87171') + ';">' + label + '</span>';
      h += '<button class="svc-restart" onclick="restartService(\\''+s.id+'\\')">&#x21bb;</button>';
      if (extra) h += extra;
      h += '</div>';
    });
    el.innerHTML = h;
  } catch(e) { /* silent */ }
}
async function restartService(id) {
  var btn = event.target;
  btn.textContent = '...';
  btn.disabled = true;
  try {
    await fetch('/api/services/restart', {
      method: 'POST', headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({service: id})
    });
    setTimeout(loadServices, 3000);
  } catch(e) { /* silent */ }
}
async function openNewWindow(path) {
  document.getElementById('nav-menu').classList.remove('open');
  document.querySelector('.nav-btn').classList.remove('active');
  try {
    await fetch('/api/open-window', {
      method: 'POST', headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({path: path})
    });
  } catch(e) {
    window.open(path, '_blank');
  }
}

// ─── SIDEBAR ────────────────────────────────────────────────────────────────────
let currentBasePrompt = '';
async function toggleSidebar() {
  sidebarOpen = !sidebarOpen;
  const sb = document.getElementById('sidebar');
  sb.style.width = sidebarOpen ? '30%' : '0';
  sb.style.minWidth = sidebarOpen ? '280px' : '0';
  if (sidebarOpen) {
    const name = getAgentName();
    if (name && name !== 'Kein Agent') {
      const r = await fetch('/get_prompt?agent=' + encodeURIComponent(name) + '&session_id=' + SESSION_ID);
      const data = await r.json();
      if (data.ok) updateSidebar(name, data.prompt);
    }
  }
}

function updateSidebar(agentName, prompt) {
  document.getElementById('sidebar-agent-name').textContent = agentName;
  document.getElementById('prompt-editor').value = prompt;
  currentBasePrompt = prompt;
  const hasMemory = prompt.includes('--- GEDAECHTNIS:');
  document.getElementById('memory-indicator').style.display = hasMemory ? 'inline' : 'none';
  document.getElementById('memory-note').style.display = hasMemory ? 'block' : 'none';
}

async function savePrompt() {
  const fullText = document.getElementById('prompt-editor').value;
  const name = getAgentName();
  if (!name || name === 'Kein Agent') return;
  const r = await fetch('/save_prompt', {method:'POST', headers:{'Content-Type':'application/json'},
    body: JSON.stringify({agent: name, prompt: fullText, session_id: SESSION_ID})});
  const data = await r.json();
  if (data.ok) { addStatusMsg('System Prompt gespeichert.'); }
}

async function reloadPrompt() {
  const name = getAgentName();
  if (!name || name === 'Kein Agent') return;
  const r = await fetch('/get_prompt?agent=' + encodeURIComponent(name) + '&session_id=' + SESSION_ID);
  const data = await r.json();
  if (data.ok) { document.getElementById('prompt-editor').value = data.prompt; addStatusMsg('Neu geladen.'); }
}

// ─── HISTORY ────────────────────────────────────────────────────────────────────
var _histSessions = {};
async function loadHistory(agentName) {
  var r = await fetch('/get_history?agent=' + encodeURIComponent(agentName) + '&session_id=' + SESSION_ID);
  var data = await r.json();
  var list = document.getElementById('history-list');
  if (!data.sessions || !data.sessions.length) {
    list.innerHTML = '<p style="font-size:11px;color:#555;padding:8px;font-style:italic;">Keine Konversationen</p>';
    return;
  }
  _histSessions = {};
  var h = '';
  data.sessions.forEach(function(s, i) {
    _histSessions[s.file] = s;
    h += '<button class="history-item' + (i===0?' active':'') + '" data-file="' + escHtml(s.file) + '" onclick="onHistoryClick(this)">';
    h += '<span class="h-date">' + escHtml(s.date) + '</span>';
    h += '<span class="h-summary">' + escHtml(s.title || s.date) + '</span>';
    h += '</button>';
  });
  list.innerHTML = h;
}
function onHistoryClick(btn) {
  var file = btn.getAttribute('data-file');
  var session = _histSessions[file];
  if (session) loadConversation(session, btn);
}

async function loadConversation(session, btn) {
  document.querySelectorAll('.history-item').forEach(function(b) { b.classList.remove('active'); });
  btn.classList.add('active');
  const name = getAgentName();
  const r = await fetch('/load_conversation', {method:'POST', headers:{'Content-Type':'application/json'},
    body: JSON.stringify({agent: name, file: session.file, session_id: SESSION_ID, resume: true})});
  const data = await r.json();
  if (data.ok) {
    document.getElementById('messages').innerHTML = '';
    if (data.messages) data.messages.forEach(m => addMessage(m.role, m.content, m.role==='assistant'?name:null));
    // Restore model selector from conversation
    if (data.provider && data.model_id) {
      var ps = document.getElementById('provider-select');
      if (ps) {
        ps.value = data.provider;
        fetch('/models').then(r=>r.json()).then(mdata => {
          var pd = mdata.find(p => p.provider === data.provider);
          if (pd) {
            populateModels(pd);
            var ms = document.getElementById('model-select');
            if (ms) ms.value = data.model_id;
          }
        });
      }
    }
    // Restore context files from saved conversation
    document.getElementById('ctx-items').innerHTML = '';
    if (data.restored_ctx && data.restored_ctx.length) {
      data.restored_ctx.forEach(n => addCtxItem(n, 'file', true));
      addStatusMsg('Konversation geladen — ' + data.restored_ctx.length + ' Kontext-Datei(en) wiederhergestellt.');
    } else {
      addStatusMsg('Konversation geladen — deine naechste Nachricht wird hier fortgesetzt.');
    }
    if (data.missing_ctx && data.missing_ctx.length) {
      addStatusMsg('\u26A0 ' + data.missing_ctx.length + ' Kontext-Datei(en) nicht mehr verfuegbar: ' + data.missing_ctx.join(', '));
    }
  }
}

// ─── MESSAGES ────────────────────────────────────────────────────────────────────
function addStatusMsg(text) {
  const div = document.createElement('p');
  div.className = 'status-msg'; div.textContent = text;
  document.getElementById('messages').appendChild(div);
  scrollDown();
}

function addMemoryMsg(text) {
  const div = document.createElement('p');
  div.className = 'memory-msg'; div.textContent = '🧠 ' + text;
  document.getElementById('messages').appendChild(div);
  scrollDown();
}

function renderCodeBlocks(text) {
  // Split text into code blocks (```) and regular text
  var parts = text.split(/```/);
  if (parts.length < 3) {
    return '<pre>' + escHtml(text) + '</pre>';
  }
  var html = '';
  for (var i = 0; i < parts.length; i++) {
    if (i % 2 === 0) {
      var t = parts[i].trim();
      if (t) html += '<pre>' + escHtml(t) + '</pre>';
    } else {
      var lines = parts[i].split('\\n');
      var lang = '';
      var code = parts[i];
      if (lines[0] && /^[a-zA-Z0-9_+-]+$/.test(lines[0].trim())) {
        lang = lines[0].trim();
        code = lines.slice(1).join('\\n');
      }
      var langLabel = lang ? '<span class="code-block-lang">' + escHtml(lang) + '</span>' : '';
      html += '<div class="code-block-wrapper" data-code="' + escHtml(code).replace(/"/g, '&quot;') + '">' + langLabel + '<pre>' + escHtml(code) + '</pre></div>';
    }
  }
  return html;
}

function renderMessageContent(text) {
  // Parse <output>...</output> blocks first
  var outputMatch = text.match(/<output>([\s\S]*?)<\/output>/);
  if (outputMatch) {
    var before = text.substring(0, outputMatch.index).trim();
    var outputContent = outputMatch[1].trim();
    var after = text.substring(outputMatch.index + outputMatch[0].length).trim();
    var html = '';
    if (before) html += renderMarkdown(before);
    html += '<div class="output-block" data-output="' + escHtml(outputContent).replace(/"/g, '&quot;') + '">' + renderCodeBlocks(outputContent) + '</div>';
    if (after) html += renderMarkdown(after);
    return html;
  }
  // No output block — render with Markdown
  return renderMarkdown(text);
}

function renderMarkdown(text) {
  // Use marked.js if available, otherwise fall back to renderCodeBlocks
  if (typeof marked !== 'undefined') {
    try {
      marked.setOptions({
        breaks: true,
        gfm: true,
        headerIds: false,
        mangle: false,
      });
      // Let marked handle everything including code blocks
      var html = marked.parse(text);
      // Make links open in new tab
      html = html.replace(/<a /g, '<a target="_blank" rel="noopener" ');
      return html;
    } catch(e) {
      console.warn('Markdown parse error, falling back:', e);
    }
  }
  return renderCodeBlocks(text);
}

function addMessage(role, text, modelName, providerDisplay, modelDisplay) {
  var msgs = document.getElementById('messages');
  var div = document.createElement('div');
  div.className = 'msg ' + role;
  var time = new Date().toLocaleTimeString('de-DE', {hour:'2-digit', minute:'2-digit'});
  var meta = time;
  if (providerDisplay && modelDisplay) {
    meta += ' · <span style="color:#888;font-style:italic">' + escHtml(providerDisplay) + ' / ' + escHtml(modelDisplay) + '</span>';
  } else if (modelName) {
    meta += ' · ' + escHtml(modelName);
  }
  if (role === 'assistant') {
    div.innerHTML = '<div class="bubble markdown-rendered">' + renderMessageContent(text) + '</div><div class="meta">' + meta + '</div>';
    addCodeCopyButtons(div);
    addCopyButton(div, text);
    addSectionCopyButtons(div, text);
  } else {
    div.innerHTML = '<div class="bubble"><pre>' + escHtml(text) + '</pre></div><div class="meta">' + meta + '</div>';
    addCodeCopyButtons(div);
  }
  msgs.appendChild(div);
  scrollDown();
}

// Robuste Copy-Funktion mit Fallback auf document.execCommand('copy').
// Problem vorher: In pywebview (native macOS-App) ist navigator.clipboard
// unter Umstaenden undefined — dann wirft navigator.clipboard.writeText
// synchron eine TypeError und das .catch() wurde nicht erreicht, sodass
// gar kein Feedback und kein Fallback lief. Jetzt erst Feature-Detection,
// dann try/catch um beide Pfade, und der execCommand-Rueckgabewert wird
// geprueft, damit "Kopiert" nur bei echtem Erfolg erscheint.
function copyToClipboard(text, btn, label) {
  var origLabel = (btn && btn.textContent) || label || 'Kopieren';
  var effectiveLabel = label || origLabel;
  function showSuccess() {
    if (!btn) return;
    btn.textContent = 'Kopiert';
    btn.classList.add('copied');
    setTimeout(function(){ btn.textContent = effectiveLabel; btn.classList.remove('copied'); }, 2000);
  }
  function showFail(err) {
    if (!btn) return;
    btn.textContent = 'Fehler';
    setTimeout(function(){ btn.textContent = effectiveLabel; }, 2000);
    try { if (window.console && err) console.warn('Clipboard-Fehler:', err); } catch(e) {}
  }
  function fallback() {
    try {
      var ta = document.createElement('textarea');
      ta.value = text;
      ta.setAttribute('readonly', '');
      ta.style.position = 'fixed';
      ta.style.top = '0';
      ta.style.left = '-9999px';
      ta.style.opacity = '0';
      document.body.appendChild(ta);
      ta.focus(); ta.select();
      ta.setSelectionRange(0, text.length);
      var ok = false;
      try { ok = document.execCommand('copy'); } catch(e) { ok = false; }
      document.body.removeChild(ta);
      if (ok) showSuccess(); else showFail('execCommand returned false');
    } catch(e) { showFail(e); }
  }
  try {
    if (navigator && navigator.clipboard && window.isSecureContext && typeof navigator.clipboard.writeText === 'function') {
      navigator.clipboard.writeText(text).then(showSuccess, fallback);
    } else {
      fallback();
    }
  } catch(e) { fallback(); }
}

function addCodeCopyButtons(msgEl) {
  // Handle custom code-block-wrapper (renderCodeBlocks fallback)
  var blocks = msgEl.querySelectorAll('.code-block-wrapper');
  blocks.forEach(function(wrapper) {
    if (wrapper.querySelector('.code-copy-btn')) return;
    var code = wrapper.getAttribute('data-code');
    if (!code) return;
    var btn = document.createElement('button');
    btn.className = 'code-copy-btn';
    btn.textContent = 'Kopieren';
    btn.onclick = function() { copyToClipboard(code, btn, 'Kopieren'); };
    wrapper.appendChild(btn);
  });
  // Handle marked.js rendered <pre><code> blocks
  var pres = msgEl.querySelectorAll('pre');
  pres.forEach(function(pre) {
    if (pre.closest('.code-block-wrapper') || pre.closest('.output-block')) return;
    if (pre.parentNode.classList && pre.parentNode.classList.contains('code-block-wrapper')) return;
    if (pre.querySelector('.code-copy-btn')) return;
    var codeEl = pre.querySelector('code');
    var codeText = codeEl ? codeEl.textContent : pre.textContent;
    if (!codeText || codeText.trim().length < 2) return;
    // Wrap in code-block-wrapper for consistent styling
    var wrapper = document.createElement('div');
    wrapper.className = 'code-block-wrapper';
    pre.parentNode.insertBefore(wrapper, pre);
    wrapper.appendChild(pre);
    // Detect language from class (marked adds language-xxx)
    if (codeEl) {
      var langClass = (codeEl.className || '').match(/language-(\w+)/);
      if (langClass) {
        var langLabel = document.createElement('span');
        langLabel.className = 'code-block-lang';
        langLabel.textContent = langClass[1];
        wrapper.insertBefore(langLabel, pre);
      }
    }
    var btn = document.createElement('button');
    btn.className = 'code-copy-btn';
    btn.textContent = 'Kopieren';
    btn.onclick = function() { copyToClipboard(codeText, btn, 'Kopieren'); };
    wrapper.appendChild(btn);
  });
}

function addCopyButton(msgEl, rawText) {
  if (!rawText || rawText.length < 80) return;
  var bubble = msgEl.querySelector('.bubble');
  if (!bubble) return;
  // Output block gets its own prominent copy button
  var outputBlock = bubble.querySelector('.output-block');
  if (outputBlock && !outputBlock.querySelector('.output-copy-btn')) {
    var outputText = outputBlock.getAttribute('data-output') || outputBlock.innerText;
    var obtn = document.createElement('button');
    obtn.className = 'output-copy-btn';
    obtn.textContent = 'Kopieren';
    obtn.onclick = function() { copyToClipboard(outputText, obtn, 'Kopieren'); };
    outputBlock.appendChild(obtn);
  }
  // Small "Alles" button on bubble level (for copying everything)
  if (!bubble.querySelector('.snippet-copy-btn')) {
    var hasBlocks = bubble.querySelectorAll('.code-block-wrapper, .output-block').length > 0;
    if (hasBlocks) {
      var btn = document.createElement('button');
      btn.className = 'snippet-copy-btn';
      btn.textContent = 'Alles';
      btn.onclick = function() { copyToClipboard(rawText, btn, 'Alles'); };
      bubble.appendChild(btn);
    } else {
      var btn = document.createElement('button');
      btn.className = 'snippet-copy-btn';
      btn.textContent = 'Kopieren';
      btn.onclick = function() { copyToClipboard(rawText, btn, 'Kopieren'); };
      bubble.appendChild(btn);
    }
  }
}

function addSectionCopyButtons(msgEl, rawText) {
  if (!rawText || rawText.length < 40) return;
  var bubble = msgEl.querySelector('.bubble');
  if (!bubble) return;
  // Find ## headings in rawText, but only outside code blocks
  var codeParts = rawText.split('```');
  var headingPositions = [];
  var charPos = 0;
  for (var ci = 0; ci < codeParts.length; ci++) {
    if (ci % 2 === 0) {
      var re = /^## .+/gm;
      var m;
      while ((m = re.exec(codeParts[ci])) !== null) {
        headingPositions.push(charPos + m.index);
      }
    }
    charPos += codeParts[ci].length + 3;
  }
  if (headingPositions.length === 0) return;
  // Build rest-texts for each heading position
  var restTexts = headingPositions.map(function(p) { return rawText.substring(p).trim(); });
  // Find content <pre> elements (not in code-block-wrapper or output-block)
  var allPres = Array.from(bubble.querySelectorAll('pre'));
  var contentPres = allPres.filter(function(pre) {
    return !pre.closest('.code-block-wrapper') && !pre.closest('.output-block');
  });
  var restIdx = 0;
  contentPres.forEach(function(pre) {
    var preText = pre.textContent;
    var preLines = preText.split('\\n');
    var headingLineIdxs = [];
    for (var li = 0; li < preLines.length; li++) {
      if (/^## /.test(preLines[li]) && restIdx < restTexts.length) {
        headingLineIdxs.push({ lineIdx: li, restText: restTexts[restIdx] });
        restIdx++;
      }
    }
    if (headingLineIdxs.length === 0) return;
    // Split <pre> at heading boundaries and insert buttons
    var frag = document.createDocumentFragment();
    var lastCut = 0;
    headingLineIdxs.forEach(function(h) {
      if (h.lineIdx > lastCut) {
        var beforeText = preLines.slice(lastCut, h.lineIdx).join('\\n').trim();
        if (beforeText) {
          var p = document.createElement('pre');
          p.textContent = beforeText;
          frag.appendChild(p);
        }
      }
      var marker = document.createElement('div');
      marker.className = 'section-copy-marker';
      var btn = document.createElement('button');
      btn.className = 'section-copy-btn';
      btn.textContent = '\u2193 Kopieren';
      (function(rt, b) {
        b.onclick = function() { copyToClipboard(rt, b, '\u2193 Kopieren'); };
      })(h.restText, btn);
      marker.appendChild(btn);
      frag.appendChild(marker);
      lastCut = h.lineIdx;
    });
    if (lastCut < preLines.length) {
      var rest = preLines.slice(lastCut).join('\\n').trim();
      if (rest) {
        var p = document.createElement('pre');
        p.textContent = rest;
        frag.appendChild(p);
      }
    }
    pre.parentNode.replaceChild(frag, pre);
  });
}

function addDownloadButton(file) {
  const msgs = document.getElementById('messages');
  const div = document.createElement('div');
  div.style.cssText = 'text-align:center;padding:8px 0;';
  const icon = file.type==='xlsx'?'📊':file.type==='pdf'?'📕':file.type==='pptx'?'📊':'📄';
  const agentName = getAgentName();
  div.innerHTML = icon + ' <strong style="color:#f0f0f0;font-size:13px;">' + file.filename + '</strong><br>' +
    '<a href="/download_file?path=' + encodeURIComponent(file.path) + '" download="' + file.filename + '" ' +
    'style="display:inline-block;margin-top:6px;background:#f0c060;color:#111;padding:6px 20px;border-radius:6px;font-size:12px;font-weight:700;text-decoration:none;font-family:Inter,sans-serif;">⬇ Herunterladen</a>' +
    ' <button onclick="openInFinder(' + JSON.stringify(agentName) + ',' + JSON.stringify(file.filename) + ')" ' +
    'style="background:#2a2a2a;border:1px solid #444;color:#888;padding:6px 14px;border-radius:6px;font-size:12px;cursor:pointer;font-family:Inter,sans-serif;margin-left:6px;">Im Finder zeigen</button>';
  msgs.appendChild(div);
  scrollDown();
}

function addImagePreview(img) {
  const msgs = document.getElementById('messages');
  const div = document.createElement('div');
  div.style.cssText = 'text-align:center;padding:12px 0;';
  div.innerHTML = '<img src="/download_file?path=' + encodeURIComponent(img.path) + '" style="max-width:600px;max-height:500px;border-radius:8px;border:1px solid #333;" /><br>' +
    '<span style="font-size:11px;color:#888;font-family:Inter,sans-serif;">' + img.filename + '</span><br>' +
    '<a href="/download_file?path=' + encodeURIComponent(img.path) + '" download="' + img.filename + '" ' +
    'style="display:inline-block;margin-top:6px;background:#f0c060;color:#111;padding:6px 20px;border-radius:6px;font-size:12px;font-weight:700;text-decoration:none;font-family:Inter,sans-serif;">\u2b07 Bild herunterladen</a>';
  msgs.appendChild(div);
  scrollDown();
}

// ── Progress Bar fuer Video-/Bild-Generierung ────────────────────────────────
// progressBars: task_id -> { el, pollTimer, kind, finished }
var progressBars = {};

function fmtDuration(sec) {
  if (sec == null || isNaN(sec)) return '--:--';
  sec = Math.max(0, Math.round(sec));
  var m = Math.floor(sec / 60);
  var s = sec % 60;
  return m + ':' + (s < 10 ? '0' + s : s);
}

function createProgressBar(taskId, kind) {
  if (progressBars[taskId]) return progressBars[taskId].el;
  const msgs = document.getElementById('messages');
  const wrap = document.createElement('div');
  wrap.className = 'task-progress';
  wrap.dataset.taskId = taskId;
  const label = kind === 'video' ? 'Video wird generiert' : 'Bild wird erstellt';
  wrap.innerHTML =
    '<div class="tp-head">' +
      '<span class="tp-pulse"></span>' +
      '<span class="tp-label">' + label + '...</span>' +
      '<span class="tp-times">0:00</span>' +
    '</div>' +
    '<div class="tp-bar-outer"><div class="tp-bar-inner"></div></div>';
  msgs.appendChild(wrap);
  progressBars[taskId] = { el: wrap, pollTimer: null, kind: kind, finished: false };
  scrollDown();
  // Start dedicated per-task polling (2 s) so progress keeps updating even
  // while the /chat request is still pending.
  pollTaskStatus(taskId);
  progressBars[taskId].pollTimer = setInterval(function(){ pollTaskStatus(taskId); }, 2000);
  return wrap;
}

function updateProgressBar(taskId, data) {
  const entry = progressBars[taskId];
  if (!entry || !data) return;
  const el = entry.el;
  const bar = el.querySelector('.tp-bar-inner');
  const lbl = el.querySelector('.tp-label');
  const tms = el.querySelector('.tp-times');
  const pct = Math.max(0, Math.min(100, data.progress || 0));
  bar.style.width = pct + '%';
  if (data.message) lbl.textContent = data.message;
  var elapsed = fmtDuration(data.elapsed_seconds);
  var timeStr = elapsed;
  if (data.status === 'running' && data.eta_seconds != null && data.eta_seconds > 0) {
    timeStr = elapsed + ' · ~' + fmtDuration(data.eta_seconds) + ' verbleibend';
  }
  tms.textContent = timeStr;
  if (data.status === 'done') {
    el.classList.add('done');
    entry.finished = true;
    if (entry.pollTimer) { clearInterval(entry.pollTimer); entry.pollTimer = null; }
    // Remove the bar shortly after completion — the actual media preview
    // (addImagePreview / addVideoPreview) is inserted by the chat response.
    setTimeout(function(){ removeProgressBar(taskId); }, 800);
  } else if (data.status === 'error') {
    el.classList.add('error');
    entry.finished = true;
    if (entry.pollTimer) { clearInterval(entry.pollTimer); entry.pollTimer = null; }
    setTimeout(function(){ removeProgressBar(taskId); }, 5000);
  }
}

function removeProgressBar(taskId) {
  const entry = progressBars[taskId];
  if (!entry) return;
  if (entry.pollTimer) { clearInterval(entry.pollTimer); }
  if (entry.el && entry.el.parentNode) entry.el.parentNode.removeChild(entry.el);
  delete progressBars[taskId];
}

async function pollTaskStatus(taskId) {
  try {
    const r = await fetch('/task_status/' + encodeURIComponent(taskId));
    if (!r.ok) return;
    const data = await r.json();
    if (data && !data.error) updateProgressBar(taskId, data);
  } catch(e) {}
}

function addVideoPreview(vid) {
  const msgs = document.getElementById('messages');
  const div = document.createElement('div');
  div.style.cssText = 'text-align:center;padding:12px 0;';
  var videoStyle = vid.is_portrait
    ? 'max-width:280px;max-height:500px;aspect-ratio:9/16;border-radius:8px;border:1px solid #333;display:block;margin:0 auto;'
    : 'max-width:600px;border-radius:8px;border:1px solid #333;';
  div.innerHTML = '<video controls style="' + videoStyle + '"><source src="/download_file?path=' + encodeURIComponent(vid.path) + '" type="video/mp4"></video><br>' +
    '<span style="font-size:11px;color:#888;font-family:Inter,sans-serif;">' + vid.filename + (vid.is_portrait ? ' (Portrait 9:16)' : '') + '</span><br>' +
    '<a href="/download_file?path=' + encodeURIComponent(vid.path) + '" download="' + vid.filename + '" ' +
    'style="display:inline-block;margin-top:6px;background:#f0c060;color:#111;padding:6px 20px;border-radius:6px;font-size:12px;font-weight:700;text-decoration:none;font-family:Inter,sans-serif;">\u2b07 Video herunterladen</a>';
  msgs.appendChild(div);
  scrollDown();
}

function addMailtoFallback(e) {
  const msgs = document.getElementById('messages');
  const div = document.createElement('div');
  div.style.cssText = 'text-align:center;padding:10px 0;';
  const to = e.to||'', subject = e.subject||'', body = e.body||'';
  const mailto = 'mailto:' + encodeURIComponent(to) + '?subject=' + encodeURIComponent(subject) + '&body=' + encodeURIComponent(body);
  div.innerHTML = '<span style="font-size:12px;color:#aaa;font-family:Inter,sans-serif;">✉ ' + subject + '</span><br>' +
    '<a href="' + mailto + '" style="display:inline-block;margin-top:6px;background:#f0c060;color:#111;padding:7px 20px;border-radius:6px;font-size:12px;font-weight:700;text-decoration:none;font-family:Inter,sans-serif;">✉ In Apple Mail oeffnen</a>';
  msgs.appendChild(div);
  scrollDown();
}

function addFinderLink(agent, filename) {
  const msgs = document.getElementById('messages');
  const div = document.createElement('div');
  div.style.cssText = 'text-align:center;padding:4px 0;';
  div.innerHTML = '📄 <span style="font-size:12px;color:#aaa;font-family:Inter,sans-serif;">' + filename + '</span> ' +
    '<button onclick="openInFinder(' + JSON.stringify(agent) + ',' + JSON.stringify(filename) + ')" ' +
    'style="background:#2a3a2a;border:1px solid #4a6a4a;color:#a0d090;padding:3px 10px;border-radius:5px;font-size:11px;cursor:pointer;font-family:Inter,sans-serif;margin-left:6px;">Im Finder zeigen</button> ' +
    '<button onclick="openMemoryFolder(' + JSON.stringify(agent) + ')" ' +
    'style="background:#2a2a2a;border:1px solid #444;color:#888;padding:3px 10px;border-radius:5px;font-size:11px;cursor:pointer;font-family:Inter,sans-serif;margin-left:4px;">Memory Ordner oeffnen</button>';
  msgs.appendChild(div);
  scrollDown();
}

async function openInFinder(agent, filename) {
  await fetch('/open_in_finder', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({agent, filename, session_id:SESSION_ID})});
}
async function openMemoryFolder(agent) {
  await fetch('/open_in_finder', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({agent, filename:'', session_id:SESSION_ID})});
}

function handleChatResponse(data) {
  if (data.error) { addStatusMsg('Fehler: '+data.error); return; }
  if (data.type === 'subagent_confirmation_required') {
    showSubagentConfirmation(data);
    return;
  }
  if (data.delegated_to) {
    addStatusMsg('\u2192 Delegiert an ' + (data.delegated_display || data.delegated_to));
  }
  if (data.auto_search_info) addStatusMsg(data.auto_search_info);
  if (data.auto_loaded && data.auto_loaded.length) {
    const agent = data.agent || getAgentName();
    data.auto_loaded.forEach(n => { addCtxItem(n,'file',true); addFinderLink(agent,n); });
  }
  addMessage('assistant', data.response, data.model_name, data.provider_display, data.model_display);
  if (data.created_files && data.created_files.length) data.created_files.forEach(f => addDownloadButton(f));
  if (data.created_images && data.created_images.length) data.created_images.forEach(img => addImagePreview(img));
  if (data.created_videos && data.created_videos.length) data.created_videos.forEach(vid => addVideoPreview(vid));
  if (data.created_emails && data.created_emails.length) {
    data.created_emails.forEach(e => {
      if (e && e.ok) addStatusMsg('\u2709 Apple Mail ' + (e.reply ? 'Reply' : 'Draft') + ' geoeffnet: '+e.subject);
      else addMailtoFallback(e||{});
    });
  }
  if (data.created_whatsapps && data.created_whatsapps.length) {
    data.created_whatsapps.forEach(wa => {
      if (wa && wa.ok && wa.clipboard_fallback) addStatusMsg('\u260E WhatsApp geoeffnet \u2014 Keine Nummer fuer \u201c'+wa.to+'\u201d gefunden. Nachricht in Zwischenablage \u2014 bitte manuell einfuegen.');
      else if (wa && wa.ok) addStatusMsg('\u260E WhatsApp geoeffnet \u2014 Chat mit '+wa.to+' wird geoeffnet. Bitte auf Senden klicken.');
      else addStatusMsg('\u26A0 WhatsApp-Fehler: '+(wa.error||'Unbekannt'));
    });
  }
  if (data.created_slacks && data.created_slacks.length) {
    data.created_slacks.forEach(sl => {
      if (sl && sl.ok) addStatusMsg('\U0001f4ac Slack geoeffnet: ' + (sl.target || ''));
      else addStatusMsg('\u26A0 Slack-Fehler: '+(sl.error||'Unbekannt'));
    });
  }
}

function showSubagentConfirmation(data) {
  var msgs = document.getElementById('messages');
  var div = document.createElement('div');
  div.className = 'msg assistant';
  var kws = data.matched_keywords ? data.matched_keywords.join(', ') : '';
  var inner = '<div class="bubble" style="background:#1a2a1a;border:1px solid #4a8a4a;padding:12px;">';
  inner += '<div style="margin-bottom:8px;">&#128256; <b>Sub-Agent erkannt:</b> ' + escHtml(data.subagent_display || data.suggested_subagent) + '</div>';
  if (kws) inner += '<div style="margin-bottom:8px;color:#888;">Keywords: ' + escHtml(kws) + '</div>';
  inner += '<div style="margin-bottom:8px;color:#888;font-size:12px;">Nachricht: \u201c' + escHtml(data.original_message || '') + '\u201d</div>';
  inner += '<div id="subagent-btns-' + data.confirmation_id + '">';
  inner += '<button onclick="confirmSubagent(&apos;' + data.confirmation_id + '&apos;,true)" style="background:#4a8a4a;color:#fff;border:none;padding:6px 16px;border-radius:4px;cursor:pointer;margin-right:8px;">\u2713 Ja, weiterleiten</button>';
  inner += '<button onclick="confirmSubagent(&apos;' + data.confirmation_id + '&apos;,false)" style="background:#555;color:#fff;border:none;padding:6px 16px;border-radius:4px;cursor:pointer;">\u2717 Nein, selbst antworten</button>';
  inner += '</div></div>';
  var time = new Date().toLocaleTimeString('de-DE', {hour:'2-digit', minute:'2-digit'});
  div.innerHTML = inner + '<div class="meta">' + time + '</div>';
  msgs.appendChild(div);
  scrollDown();
  // Auto-decline after 5 minutes
  setTimeout(function() {
    var btns = document.getElementById('subagent-btns-' + data.confirmation_id);
    if (btns && btns.querySelector('button')) {
      confirmSubagent(data.confirmation_id, false);
    }
  }, 300000);
}

async function confirmSubagent(confirmationId, confirmed) {
  var btnsDiv = document.getElementById('subagent-btns-' + confirmationId);
  if (btnsDiv) {
    btnsDiv.innerHTML = confirmed
      ? '<span style="color:#8aba8a;">\u2713 Weiterleitung bestaetigt</span>'
      : '<span style="color:#888;">\u2717 Aktueller Agent antwortet</span>';
  }
  addStatusMsg(confirmed ? 'Wird an Sub-Agent weitergeleitet...' : 'Aktueller Agent verarbeitet die Anfrage...');
  try {
    var r = await fetch('/api/subagent_confirm', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({session_id: SESSION_ID, confirmation_id: confirmationId, confirmed: confirmed})
    });
    var data = await r.json();
    if (data.error) {
      addStatusMsg('Fehler: ' + data.error);
      return;
    }
    handleChatResponse(data);
  } catch(e) {
    addStatusMsg('Fehler bei Sub-Agent Verarbeitung: ' + e.message);
  }
}

function addCtxItem(name, type, autoLoaded) {
  const bar = document.getElementById('ctx-items');
  if ([...bar.children].some(c => c.dataset.name === name)) return;
  const div = document.createElement('div');
  div.className = 'ctx-item' + (autoLoaded ? ' auto-loaded' : ''); div.dataset.name = name;
  const icon = autoLoaded ? '🔍' : (type==='url'?'🔗':'📄');
  const span = document.createElement('span');
  span.textContent = icon + ' ' + name.substring(0, 40) + (name.length > 40 ? '...' : '');
  div.title = name;
  const btn = document.createElement('button');
  btn.textContent = '\u2715';
  btn.onclick = () => removeCtx(name, div);
  div.appendChild(span);
  div.appendChild(btn);
  bar.appendChild(div);
}

async function removeCtx(name, el) {
  el.remove();
  await fetch('/remove_ctx', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({name, session_id:SESSION_ID})});
}

function escHtml(t) { return (t||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;'); }
function scrollDown() { const m = document.getElementById('messages'); m.scrollTop = m.scrollHeight; }
function autoResize(el) { el.style.height='auto'; el.style.height=Math.min(el.scrollHeight,160)+'px'; }
function onKey(e) {
  const input = document.getElementById('msg-input');
  if (onFindLiveKey(e)) return;
  if (onTypeAcKey(e, input)) return;
  if (onSlashAcKey(e, input)) return;
  if (e.key==='Enter' && !e.shiftKey) { e.preventDefault(); sendMessage(); }
}

// ─── KEYBOARD SHORTCUTS (GLOBAL) ─────────────────────────────────────────────
let _findActiveChip = null;
let _findDebounceTimer = null;
let _findLiveIdx = -1;
let _findLiveResults = [];

document.addEventListener('keydown', function(e) {
  // Skip if inside input fields (except for our specific shortcuts)
  const tag = document.activeElement.tagName;
  const inInput = (tag === 'INPUT' || tag === 'TEXTAREA' || tag === 'SELECT');

  if (e.altKey && !e.ctrlKey && !e.metaKey) {
    switch(e.key.toLowerCase()) {
      case 'p': e.preventDefault(); toggleSidebar(); return;
      case 'n': e.preventDefault(); newSession(); return;
      case 'a': e.preventDefault(); showAgentModal(); return;
      case 'm': e.preventDefault(); document.getElementById('model-select').focus(); return;
      case 'f': e.preventDefault(); var inp=document.getElementById('msg-input'); inp.focus(); inp.value='/find '; onInputHandler(inp); return;
      case 'u': e.preventDefault(); document.getElementById('file-input').click(); return;
      case 'c': e.preventDefault(); copyLastAssistantMessage(); return;
      case 's': e.preventDefault(); saveConversation(); return;
      case '1': if (isFindChipsVisible()) { e.preventDefault(); toggleFindChip('email'); } return;
      case '2': if (isFindChipsVisible()) { e.preventDefault(); toggleFindChip('webclip'); } return;
      case '3': if (isFindChipsVisible()) { e.preventDefault(); toggleFindChip('document'); } return;
      case '4': if (isFindChipsVisible()) { e.preventDefault(); toggleFindChip('conversation'); } return;
      case '5': if (isFindChipsVisible()) { e.preventDefault(); toggleFindChip('screenshot'); } return;
    }
  }
  // Ctrl+Enter to send from anywhere
  if (e.ctrlKey && e.key === 'Enter' && !e.altKey && !e.metaKey) {
    e.preventDefault();
    sendMessage();
  }
});

function copyLastAssistantMessage() {
  var msgs = document.querySelectorAll('.msg.assistant .bubble');
  if (!msgs.length) return;
  var last = msgs[msgs.length-1];
  var text = last.innerText || last.textContent || '';
  function ok()   { addStatusMsg('Letzte Antwort kopiert'); }
  function fail() { addStatusMsg('Kopieren fehlgeschlagen'); }
  function fallback() {
    try {
      var ta = document.createElement('textarea');
      ta.value = text;
      ta.setAttribute('readonly', '');
      ta.style.position = 'fixed';
      ta.style.top = '0';
      ta.style.left = '-9999px';
      ta.style.opacity = '0';
      document.body.appendChild(ta);
      ta.focus(); ta.select();
      ta.setSelectionRange(0, text.length);
      var res = false;
      try { res = document.execCommand('copy'); } catch(e) { res = false; }
      document.body.removeChild(ta);
      if (res) ok(); else fail();
    } catch(e) { fail(); }
  }
  try {
    if (navigator && navigator.clipboard && window.isSecureContext && typeof navigator.clipboard.writeText === 'function') {
      navigator.clipboard.writeText(text).then(ok, fallback);
    } else { fallback(); }
  } catch(e) { fallback(); }
}

function isFindChipsVisible() {
  var bar = document.getElementById('find-chips-bar');
  return bar && bar.classList.contains('visible');
}

function showFindChips() {
  var bar = document.getElementById('find-chips-bar');
  if (bar) bar.classList.add('visible');
}

function hideFindChips() {
  var bar = document.getElementById('find-chips-bar');
  if (bar) bar.classList.remove('visible');
  _findActiveChip = null;
  document.querySelectorAll('.find-chip').forEach(c => c.classList.remove('active'));
  hideFindLiveDropdown();
}

function toggleFindChip(cat) {
  if (_findActiveChip === cat) {
    _findActiveChip = null;
    document.querySelectorAll('.find-chip').forEach(c => c.classList.remove('active'));
  } else {
    _findActiveChip = cat;
    document.querySelectorAll('.find-chip').forEach(c => {
      c.classList.toggle('active', c.dataset.cat === cat);
    });
  }
  // Update the input if in /find mode
  var inp = document.getElementById('msg-input');
  var val = inp.value;
  var m = val.match(/^\/find(?:_global)?/);
  if (m) {
    var prefix = m[0];
    var knownTypes = ['email','webclip','screenshot','contact','document','conversation'];
    var rest = val.substring(prefix.length).trim();
    var firstWord = rest.split(/\s+/)[0]||'';
    if (knownTypes.includes(firstWord.toLowerCase())) rest = rest.substring(firstWord.length).trim();
    inp.value = prefix + (_findActiveChip ? ' ' + _findActiveChip : '') + (rest ? ' ' + rest : ' ');
    inp.focus();
    triggerFindLiveSearch();
  }
}

function hideFindLiveDropdown() {
  var dd = document.getElementById('find-live-dropdown');
  if (dd) dd.style.display = 'none';
  _findLiveIdx = -1;
  _findLiveResults = [];
}

function triggerFindLiveSearch() {
  clearTimeout(_findDebounceTimer);
  _findDebounceTimer = setTimeout(doFindLiveSearch, 300);
}

async function doFindLiveSearch() {
  var inp = document.getElementById('msg-input');
  var val = inp.value;
  var m = val.match(/^\/find(_global)?(?:-(email|webclip|screenshot|document|conversation))?\s+(.*)/i);
  if (!m) { hideFindLiveDropdown(); return; }
  var isGlobal = !!m[1];
  var typePart = (m[2]||'').trim();
  var queryPart = (m[3]||'').trim();
  // Strip old-style type prefix from query if present
  var knTypes = ['email','webclip','screenshot','contact','document','conversation'];
  var qFirst = queryPart.split(/\s+/)[0];
  if (knTypes.includes(qFirst.toLowerCase())) queryPart = queryPart.substring(qFirst.length).trim();
  if (!queryPart || queryPart.length < 2) { hideFindLiveDropdown(); return; }

  try {
    var agent = getAgentName();
    var url = '/api/memory-files-search?q=' + encodeURIComponent(queryPart) + '&agent=' + encodeURIComponent(agent);
    var r = await fetch(url);
    var data = await r.json();
    if (!data || !data.length) { hideFindLiveDropdown(); return; }

    _findLiveResults = data.slice(0, 8);
    _findLiveIdx = -1;
    var dd = document.getElementById('find-live-dropdown');
    dd.innerHTML = '<div style="padding:4px 14px;font-size:10px;color:#555;">' + _findLiveResults.length + ' Treffer | ↑↓ navigieren, Enter auswaehlen, Esc schliessen</div>';
    _findLiveResults.forEach(function(item, i) {
      var div = document.createElement('div');
      div.className = 'find-live-item';
      div.innerHTML = '<span class="flr-name">' + escHtml(item.filename || item.name || '') + '</span><span class="flr-type">' + escHtml(item.snippet || '').substring(0,40) + '</span>';
      div.onmousedown = function(e) { e.preventDefault(); selectFindLiveItem(i); };
      dd.appendChild(div);
    });
    dd.style.display = 'block';
  } catch(e) {
    hideFindLiveDropdown();
  }
}

function selectFindLiveItem(idx) {
  if (idx < 0 || idx >= _findLiveResults.length) return;
  var item = _findLiveResults[idx];
  var inp = document.getElementById('msg-input');
  var val = inp.value;
  var m = val.match(/^\/find(?:_global)?/);
  var prefix = m ? m[0] : '/find';
  var typePart = _findActiveChip ? ' ' + _findActiveChip : '';
  inp.value = prefix + typePart + ' ' + (item.filename || item.name || '');
  hideFindLiveDropdown();
  inp.focus();
}

function onFindLiveKey(e) {
  var dd = document.getElementById('find-live-dropdown');
  if (!dd || dd.style.display === 'none' || !_findLiveResults.length) return false;
  var items = dd.querySelectorAll('.find-live-item');
  if (e.key === 'ArrowDown') { e.preventDefault(); _findLiveIdx = Math.min(_findLiveIdx+1, items.length-1); items.forEach(function(it,i){it.classList.toggle('active', i===_findLiveIdx);}); return true; }
  if (e.key === 'ArrowUp') { e.preventDefault(); _findLiveIdx = Math.max(_findLiveIdx-1, 0); items.forEach(function(it,i){it.classList.toggle('active', i===_findLiveIdx);}); return true; }
  if (e.key === 'Tab' && _findLiveResults.length > 0) { e.preventDefault(); selectFindLiveItem(_findLiveIdx >= 0 ? _findLiveIdx : 0); return true; }
  if (e.key === 'Enter' && _findLiveIdx >= 0) { e.preventDefault(); selectFindLiveItem(_findLiveIdx); return true; }
  if (e.key === 'Escape') { hideFindLiveDropdown(); return true; }
  return false;
}

// ─── SEND ────────────────────────────────────────────────────────────────────────
// ─── QUEUE SYSTEM (per-session / per-tab isolated state) ──────────────────
// Jeder Browser-Tab hat eine eigene SESSION_ID. Processing-State (laeuft/idle,
// aktueller Prompt, Stop-Button-Sichtbarkeit, Queue-Laenge, Pending-Responses,
// Placeholder-DOM-Refs und setInterval-IDs) wird pro Session in _tabStates
// gefuehrt. DOM-Updates passieren nur, wenn die betroffene Session aktuell
// aktiv ist; sonst wird der State silent aktualisiert und Responses fuer den
// inaktiven Tab werden gepuffert, bis dieser Tab wieder aktiv wird.
var _tabStates = {};  // sessionId -> state

function _tabState(sid) {
  sid = sid || SESSION_ID;
  if (!_tabStates[sid]) {
    _tabStates[sid] = {
      processing: false,
      currentPrompt: '',
      stopBtnVisible: false,
      queueLength: 0,
      pollIntervalId: null,
      typingIntervalId: null,
      queuedPlaceholders: {},  // queue_id -> DOM element (only valid while tab active)
      pendingResponses: []     // responses received while tab inactive
    };
  }
  return _tabStates[sid];
}

function _isActiveSession(sid) {
  return (sid || SESSION_ID) === SESSION_ID;
}

// Backwards-compat shim: legacy code reads/writes queuedPlaceholders directly.
// Via Proxy aliasen wir das auf den aktuellen Tab-State.
var queuedPlaceholders = new Proxy({}, {
  get: function(_, key) {
    var s = _tabState(SESSION_ID);
    if (key === Symbol.iterator) return s.queuedPlaceholders[Symbol.iterator];
    return s.queuedPlaceholders[key];
  },
  set: function(_, key, value) {
    _tabState(SESSION_ID).queuedPlaceholders[key] = value;
    return true;
  },
  deleteProperty: function(_, key) {
    delete _tabState(SESSION_ID).queuedPlaceholders[key];
    return true;
  },
  ownKeys: function() { return Object.keys(_tabState(SESSION_ID).queuedPlaceholders); },
  getOwnPropertyDescriptor: function(_, key) {
    var v = _tabState(SESSION_ID).queuedPlaceholders[key];
    return v === undefined ? undefined : {enumerable: true, configurable: true, value: v};
  }
});

// Rendert den gespeicherten State des aktiven Tabs ins globale DOM.
// Wird von switchToTab() aufgerufen, wenn der Nutzer Tabs wechselt.
function renderActiveTabState() {
  var s = _tabState(SESSION_ID);
  // Typing indicator
  var indicator = document.getElementById('typing-indicator');
  var typingText = document.getElementById('typing-text');
  if (s.processing) {
    indicator.style.display = 'flex';
    typingText.textContent = s.currentPrompt ? 'Verarbeite: '+s.currentPrompt+'...' : 'Verarbeitet...';
  } else {
    indicator.style.display = 'none';
  }
  // Stop btn
  document.getElementById('stop-btn').style.display = s.stopBtnVisible ? 'inline-block' : 'none';
  // Queue badge
  var qel = document.getElementById('queue-display');
  var qtxt = document.getElementById('queue-text');
  if (s.queueLength > 0) {
    qel.style.display = 'block';
    qtxt.textContent = '\u23f3 '+s.queueLength+' Prompt'+(s.queueLength>1?'s':'')+' in der Warteschlange';
  } else {
    qel.style.display = 'none';
  }
  // Flush buffered responses collected while this tab was inactive
  if (s.pendingResponses && s.pendingResponses.length) {
    var pending = s.pendingResponses;
    s.pendingResponses = [];
    pending.forEach(function(resp){
      // Placeholder aus dem jetzt aktiven DOM entfernen (data-queue-id Lookup)
      if (resp.queue_id) {
        var ph = document.querySelector('.msg.queued[data-queue-id="'+resp.queue_id+'"]');
        if (ph) ph.remove();
        if (s.queuedPlaceholders[resp.queue_id]) delete s.queuedPlaceholders[resp.queue_id];
      }
      handleResponse(resp);
    });
  }
}

function handleResponse(data) {
  if (data.error) { addStatusMsg('Fehler: '+data.error); return; }
  if (data.auto_search_info) {
    addStatusMsg(data.auto_search_info);
  }
  if (data.auto_loaded && data.auto_loaded.length) {
    const agent = data.agent || getAgentName();
    data.auto_loaded.forEach(n => { addCtxItem(n,'file',true); addFinderLink(agent,n); });
  }
  if (data.type === 'subagent_confirmation_required') {
    showSubagentConfirmation(data);
    return;
  }
  if (data.delegated_to) {
    addStatusMsg('\u2192 Delegiert an ' + (data.delegated_display || data.delegated_to));
  }
  addMessage('assistant', data.response, data.model_name, data.provider_display, data.model_display);
  if (data.created_files && data.created_files.length) data.created_files.forEach(f => addDownloadButton(f));
  if (data.created_images && data.created_images.length) data.created_images.forEach(img => addImagePreview(img));
  if (data.created_videos && data.created_videos.length) data.created_videos.forEach(vid => addVideoPreview(vid));
  if (data.created_emails && data.created_emails.length) {
    data.created_emails.forEach(e => {
      if (e && e.ok) addStatusMsg('\u2709 Apple Mail ' + (e.reply ? 'Reply' : 'Draft') + ' geoeffnet: '+e.subject);
      else addMailtoFallback(e||{});
    });
  }
  if (data.created_whatsapps && data.created_whatsapps.length) {
    data.created_whatsapps.forEach(wa => {
      if (wa && wa.ok && wa.clipboard_fallback) addStatusMsg('\u260E WhatsApp geoeffnet \u2014 Keine Nummer fuer \u201c'+wa.to+'\u201d gefunden. Nachricht in Zwischenablage \u2014 bitte manuell einfuegen.');
      else if (wa && wa.ok) addStatusMsg('\u260E WhatsApp geoeffnet \u2014 Chat mit '+wa.to+' wird geoeffnet. Bitte auf Senden klicken.');
      else addStatusMsg('\u26A0 WhatsApp-Fehler: '+(wa.error||'Unbekannt'));
    });
  }
  if (data.created_slacks && data.created_slacks.length) {
    data.created_slacks.forEach(sl => {
      if (sl && sl.ok && sl.clipboard_only) addStatusMsg('\U0001f4ac Slack geoeffnet \u2014 Nachricht in Zwischenablage. Bitte mit Cmd+V einfuegen und absenden.');
      else if (sl && sl.ok) addStatusMsg('\U0001f4ac Slack geoeffnet \u2014 Nachricht an '+sl.target+' eingefuegt. Bitte manuell absenden.');
      else addStatusMsg('\u26A0 Slack-Fehler: '+(sl.error||'Unbekannt'));
    });
  }
  updateTokenCount();
}

function addQueuedMessage(text, position, queueId) {
  const msgs = document.getElementById('messages');
  const div = document.createElement('div');
  div.className = 'msg queued';
  div.dataset.queueId = queueId;
  div.innerHTML = '<div class="bubble">\u23f3 In Warteschlange (Position '+position+'): '+escHtml(text.substring(0,80))+'</div>';
  msgs.appendChild(div);
  scrollDown();
  return div;
}

function updateQueueDisplay(count, sid) {
  sid = sid || SESSION_ID;
  _tabState(sid).queueLength = count || 0;
  if (!_isActiveSession(sid)) return;  // Nur DOM aktualisieren wenn dieser Tab aktiv ist
  const el = document.getElementById('queue-display');
  const txt = document.getElementById('queue-text');
  if (count > 0) {
    el.style.display = 'block';
    txt.textContent = '\u23f3 '+count+' Prompt'+(count>1?'s':'')+' in der Warteschlange';
  } else {
    el.style.display = 'none';
  }
}

function showStopBtn(show, sid) {
  sid = sid || SESSION_ID;
  _tabState(sid).stopBtnVisible = !!show;
  if (!_isActiveSession(sid)) return;
  document.getElementById('stop-btn').style.display = show ? 'inline-block' : 'none';
}

function startTyping(prompt, sid) {
  sid = sid || SESSION_ID;
  var s = _tabState(sid);
  s.processing = true;
  s.currentPrompt = prompt || '';
  showStopBtn(true, sid);
  if (s.typingIntervalId) { clearInterval(s.typingIntervalId); s.typingIntervalId = null; }
  const typingMsgs = ['Denkt nach...','Verarbeitet...','Analysiert...','Arbeitet...','Formuliert Antwort...','Noch einen Moment...'];
  var mi = 0;
  // Rotations-Animation laeuft pro Tab. Wenn inaktiv: State wird gesetzt, DOM nicht.
  s.typingIntervalId = setInterval(function(){
    mi = (mi+1) % typingMsgs.length;
    if (!_isActiveSession(sid)) return;
    var st = _tabState(sid);
    var ttx = document.getElementById('typing-text');
    if (ttx) ttx.textContent = st.currentPrompt ? 'Verarbeite: '+st.currentPrompt+'...' : typingMsgs[mi];
  }, 2500);
  if (!_isActiveSession(sid)) return;
  const indicator = document.getElementById('typing-indicator');
  const typingText = document.getElementById('typing-text');
  indicator.style.display = 'flex';
  typingText.textContent = prompt ? 'Verarbeite: '+prompt+'...' : typingMsgs[0];
}

function stopTyping(sid) {
  sid = sid || SESSION_ID;
  var s = _tabState(sid);
  s.processing = false;
  s.currentPrompt = '';
  if (s.typingIntervalId) { clearInterval(s.typingIntervalId); s.typingIntervalId = null; }
  if (!_isActiveSession(sid)) return;
  document.getElementById('typing-indicator').style.display = 'none';
}

function startPolling(sid) {
  sid = sid || SESSION_ID;
  var s = _tabState(sid);
  if (s.pollIntervalId) return;
  s.pollIntervalId = setInterval(async function(){
    try {
      const r = await fetch('/poll_responses?session_id=' + sid);
      const data = await r.json();
      if (data.responses && data.responses.length) {
        data.responses.forEach(function(resp){
          if (_isActiveSession(sid)) {
            if (s.queuedPlaceholders[resp.queue_id]) {
              s.queuedPlaceholders[resp.queue_id].remove();
              delete s.queuedPlaceholders[resp.queue_id];
            }
            handleResponse(resp);
          } else {
            // Tab inaktiv: Response puffern und auf Tab-Switch flushen
            s.pendingResponses.push(resp);
          }
        });
      }
    } catch(e) {}
    try {
      const r2 = await fetch('/queue_status?session_id=' + sid);
      const st = await r2.json();
      updateQueueDisplay(st.queue_length, sid);
      if (st.processing) {
        s.processing = true;
        if (st.current_prompt) s.currentPrompt = st.current_prompt;
        if (_isActiveSession(sid) && st.current_prompt) {
          const ttx = document.getElementById('typing-text');
          if (ttx) ttx.textContent = 'Verarbeite: '+st.current_prompt+'...';
        }
      }
      // Progress bars nur fuer aktiven Tab
      if (_isActiveSession(sid) && st && Array.isArray(st.active_tasks)) {
        st.active_tasks.forEach(function(t){
          if (!progressBars[t.task_id]) createProgressBar(t.task_id, t.kind || 'video');
          updateProgressBar(t.task_id, t);
        });
      }
      if (!st.processing && st.queue_length === 0) {
        stopPolling(sid);
      }
    } catch(e) {}
  }, 2000);
}

function stopPolling(sid) {
  sid = sid || SESSION_ID;
  var s = _tabState(sid);
  if (s.pollIntervalId) { clearInterval(s.pollIntervalId); s.pollIntervalId = null; }
  stopTyping(sid);
  showStopBtn(false, sid);
  updateQueueDisplay(0, sid);
}

async function stopQueue() {
  var sid = SESSION_ID;
  var s = _tabState(sid);
  try {
    const r = await fetch('/stop_queue', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({session_id: sid})});
    const data = await r.json();
    if (data.cancelled > 0) addStatusMsg(data.cancelled+' Prompt(s) abgebrochen.');
    // Entfernt Placeholders nur im DOM des aktiven Tabs (ist ohnehin dieser Tab)
    Object.values(s.queuedPlaceholders).forEach(function(el){ try { el.remove(); } catch(e) {} });
    s.queuedPlaceholders = {};
    stopPolling(sid);
  } catch(e) { addStatusMsg('Stop Fehler: '+e.message); }
}

// ─── SLASH COMMAND AUTOCOMPLETE ──────────────────────────────────────────────
const _SLASH_COMMANDS = [
  // SLASH_CLUSTER_V1: Gruppierte Slash Commands
  // ─── Kommunikation ───
  {cmd: '/create-email', label: '/create-email', desc: 'E-Mail Draft erstellen', template: 'Erstelle eine E-Mail an [Empfaenger] zum Thema: ', group: 'Kommunikation'},
  {cmd: '/create-email-reply', label: '/create-email-reply', desc: 'E-Mail Antwort erstellen', template: 'Antworte auf die E-Mail von [Absender] zum Thema [Betreff]: ', group: 'Kommunikation'},
  {cmd: '/reply', label: '/reply [suche]', desc: 'E-Mail suchen und im Chat oeffnen', template: '/reply ', group: 'Kommunikation'},
  {cmd: '/create-whatsapp', label: '/create-whatsapp', desc: 'WhatsApp-Nachricht', template: 'Schreibe eine WhatsApp-Nachricht an [Name]: ', group: 'Kommunikation'},
  {cmd: '/create-slack', label: '/create-slack', desc: 'Slack-Nachricht', template: 'Schreibe eine Slack-Nachricht an [#channel oder Name]: ', group: 'Kommunikation'},
  // ─── Kalender ───
  {cmd: '/calendar-today', label: '/calendar-today', desc: 'Heutige Termine', template: '', group: 'Kalender'},
  {cmd: '/calendar-tomorrow', label: '/calendar-tomorrow', desc: 'Termine morgen', template: '', group: 'Kalender'},
  {cmd: '/calendar-week', label: '/calendar-week', desc: 'Naechste 7 Tage', template: '', group: 'Kalender'},
  {cmd: '/calendar-search', label: '/calendar-search [query]', desc: 'Termine durchsuchen', template: 'Suche Termin: ', group: 'Kalender'},
  // ─── Medien ───
  {cmd: '/create-image', label: '/create-image', desc: 'Bild generieren (Gemini/OpenAI)', template: 'Erstelle ein Bild: ', group: 'Medien'},
  {cmd: '/create-video', label: '/create-video', desc: 'Video generieren (Gemini Veo)', template: 'Erstelle ein Video: ', group: 'Medien'},
  // ─── Dokumente ───
  {cmd: '/create-file-docx', label: '/create-file-docx', desc: 'Word-Dokument', template: 'Erstelle ein Word-Dokument: ', group: 'Dokumente'},
  {cmd: '/create-file-xlsx', label: '/create-file-xlsx', desc: 'Excel-Tabelle', template: 'Erstelle eine Excel-Tabelle: ', group: 'Dokumente'},
  {cmd: '/create-file-pdf', label: '/create-file-pdf', desc: 'PDF erstellen', template: 'Erstelle ein PDF: ', group: 'Dokumente'},
  {cmd: '/create-file-pptx', label: '/create-file-pptx', desc: 'PowerPoint', template: 'Erstelle eine PowerPoint-Praesentation: ', group: 'Dokumente'},
  // ─── Canva ───
  {cmd: '/canva-search', label: '/canva-search [query]', desc: 'Designs durchsuchen', template: 'Suche in Canva nach: ', group: 'Canva'},
  {cmd: '/canva-create', label: '/canva-create [titel]', desc: 'Neues Design', template: 'Erstelle ein Canva Design: ', group: 'Canva'},
  {cmd: '/canva-templates', label: '/canva-templates', desc: 'Brand Templates', template: '', group: 'Canva'},
  {cmd: '/canva-campaign', label: '/canva-campaign', desc: 'Ad-Kampagne generieren', template: 'Erstelle eine Canva Kampagne mit Template: ', group: 'Canva'},
  {cmd: '/canva-export', label: '/canva-export [id] [format]', desc: 'Design exportieren', template: 'Exportiere Canva Design: ', group: 'Canva'},
  // ─── Suche (Agent) ───
  {cmd: '/find', label: '/find [query]', desc: 'Alle Dateien durchsuchen', group: 'Suche'},
  {cmd: '/find-email', label: '/find-email [query]', desc: 'Nur E-Mails', group: 'Suche'},
  {cmd: '/find-whatsapp', label: '/find-whatsapp [query]', desc: 'Nur WhatsApp', group: 'Suche'},
  {cmd: '/find-webclip', label: '/find-webclip [query]', desc: 'Web Clips', group: 'Suche'},
  {cmd: '/find-slack', label: '/find-slack [query]', desc: 'Slack Nachrichten', group: 'Suche'},
  {cmd: '/find-salesforce', label: '/find-salesforce [query]', desc: 'Salesforce Records', group: 'Suche'},
  {cmd: '/find-document', label: '/find-document [query]', desc: 'Dokumente', group: 'Suche'},
  {cmd: '/find-conversation', label: '/find-conversation [query]', desc: 'Konversationen', group: 'Suche'},
  {cmd: '/find-screenshot', label: '/find-screenshot [query]', desc: 'Screenshots', group: 'Suche'},
  // ─── Suche (Global) ───
  {cmd: '/find_global', label: '/find_global [query]', desc: 'Alle Agenten', group: 'Globale Suche'},
  {cmd: '/find_global-email', label: '/find_global-email [query]', desc: 'E-Mails global', group: 'Globale Suche'},
  {cmd: '/find_global-webclip', label: '/find_global-webclip [query]', desc: 'Web Clips global', group: 'Globale Suche'},
  {cmd: '/find_global-document', label: '/find_global-document [query]', desc: 'Dokumente global', group: 'Globale Suche'},
  {cmd: '/find_global-conversation', label: '/find_global-conversation [query]', desc: 'Konversationen global', group: 'Globale Suche'},
  {cmd: '/find_global-screenshot', label: '/find_global-screenshot [query]', desc: 'Screenshots global', group: 'Globale Suche'},
];
let _slashAcIdx = -1;

function onInputHandler(el) {
  autoResize(el);
  const val = el.value;
  // Draft-save: persist typed text per agent
  var _draftAgent = getAgentName();
  if (_draftAgent && _draftAgent !== 'Kein Agent') { localStorage.setItem('draft_' + (document.getElementById('agent-label').dataset.agentName || _draftAgent), val); }
  // Show/filter slash dropdown when typing "/"
  if (val.startsWith('/') && !val.includes(' ')) {
    showSlashAutocomplete(el, val);
    hideTypeAutocomplete();
    hideFindChips();
    return;
  }
  // Show find chips after any /find variant (with trailing space)
  const isFindMode = /^\/find(?:_global)?(?:-\w+)?(?:\s|$)/i.test(val);
  if (isFindMode) {
    hideSlashAutocomplete();
    hideTypeAutocomplete();
    showFindChips();
    // Trigger live search if there's query text
    const queryMatch = val.match(/^\/find(?:_global)?(?:-\w+)?\s+(.*)/i);
    if (queryMatch && queryMatch[1] && queryMatch[1].trim().length >= 2) {
      triggerFindLiveSearch();
    } else {
      hideFindLiveDropdown();
    }
    return;
  }
  // Hide everything when not in slash/find context
  if (!val.startsWith('/')) { hideSlashAutocomplete(); hideTypeAutocomplete(); hideFindChips(); }
}

function showSlashAutocomplete(inputEl, filterText) {
  let dd = document.getElementById('slash-ac-dropdown');
  if (!dd) {
    dd = document.createElement('div');
    dd.id = 'slash-ac-dropdown';
    inputEl.parentElement.style.position = 'relative';
    inputEl.parentElement.appendChild(dd);
  }
  dd.innerHTML = '';
  _slashAcIdx = -1;
  const filter = (filterText || '/').toLowerCase();
  const filtered = _SLASH_COMMANDS.filter(c => c.cmd.toLowerCase().startsWith(filter));
  if (!filtered.length) { dd.style.display = 'none'; return; }
  var lastGroup = '';
  filtered.forEach((c, i) => {
    // Gruppen-Header anzeigen wenn neue Gruppe beginnt
    if (c.group && c.group !== lastGroup) {
      lastGroup = c.group;
      var hdr = document.createElement('div');
      hdr.className = 'slash-ac-group';
      hdr.textContent = c.group;
      dd.appendChild(hdr);
    }
    const item = document.createElement('div');
    item.className = 'slash-ac-item';
    item.dataset.cmd = c.cmd;
    item.innerHTML = '<strong>' + c.label + '</strong><span style="margin-left:8px;color:#888;font-size:12px">' + c.desc + '</span>';
    item.onmousedown = (e) => { e.preventDefault(); selectSlashCmd(inputEl, c.cmd); };
    dd.appendChild(item);
  });
  dd.style.display = 'block';
}

function hideSlashAutocomplete() {
  const dd = document.getElementById('slash-ac-dropdown');
  if (dd) dd.style.display = 'none';
  _slashAcIdx = -1;
}

// ─── EMAIL SEARCH MODAL + CHAT FLOW ───────────────────────────────────
let _emailContext = null;
let _esmDebounce = null;

function showEmailSearchModal() {
  document.getElementById('esm-from').value = '';
  document.getElementById('esm-subject').value = '';
  document.getElementById('esm-to').value = '';
  document.getElementById('esm-body').value = '';
  document.getElementById('esm-results').innerHTML = '<div class="esm-hint">Mindestens 2 Zeichen in ein Feld eingeben...</div>';
  var m = document.getElementById('email-search-modal');
  m.classList.add('show'); m.style.display = 'flex';
  setTimeout(function() { document.getElementById('esm-from').focus(); }, 100);
}

function closeEmailSearchModal() {
  var m = document.getElementById('email-search-modal');
  m.classList.remove('show'); m.style.display = 'none';
}

function _esmDoSearch() {
  var from = document.getElementById('esm-from').value.trim();
  var subj = document.getElementById('esm-subject').value.trim();
  var to = document.getElementById('esm-to').value.trim();
  var body = document.getElementById('esm-body').value.trim();
  // Need at least 2 chars in any field
  var q = from || subj || to || body;
  if (!q || q.length < 2) { document.getElementById('esm-results').innerHTML = '<div class="esm-hint">Mindestens 2 Zeichen in ein Feld eingeben...</div>'; return; }
  document.getElementById('esm-results').innerHTML = '<div class="esm-loading">Suche...</div>';
  var agent = getAgentName();
  var params = 'agent=' + encodeURIComponent(agent);
  if (from) params += '&from=' + encodeURIComponent(from);
  if (subj) params += '&subject=' + encodeURIComponent(subj);
  if (to) params += '&to=' + encodeURIComponent(to);
  if (body) params += '&body=' + encodeURIComponent(body);
  // Also send combined q for backwards compat
  params += '&q=' + encodeURIComponent([from, subj, to, body].filter(Boolean).join(' '));
  fetch('/api/email-search?' + params)
    .then(function(r) { return r.json(); })
    .then(function(results) {
      var container = document.getElementById('esm-results');
      if (!results || results.length === 0) {
        container.innerHTML = '<div class="esm-hint">Keine E-Mails gefunden</div>';
        return;
      }
      container.innerHTML = '';
      results.forEach(function(item, idx) {
        var div = document.createElement('div');
        div.className = 'esm-result';
        div.innerHTML = '<div><span class="esm-result-from">' + escHtml(item.from_name || item.from_email) + '</span>'
          + (item.from_name ? '<span class="esm-result-email">&lt;' + escHtml(item.from_email) + '&gt;</span>' : '')
          + '<span class="esm-result-date">' + escHtml(item.date) + '</span></div>'
          + '<div class="esm-result-subject">' + escHtml(item.subject || '(kein Betreff)') + '</div>'
          + (item.to ? '<div class="esm-result-meta">An: ' + escHtml(item.to).substring(0,60) + '</div>' : '');
        div.onclick = function() { _esmSelectEmail(item); };
        container.appendChild(div);
      });
    })
    .catch(function(e) { document.getElementById('esm-results').innerHTML = '<div class="esm-hint">Fehler: ' + escHtml(e.message) + '</div>'; });
}

function _esmSelectEmail(item) {
  closeEmailSearchModal();
  _openEmailInChat(item);
}

// Attach search listeners on DOMContentLoaded
document.addEventListener('DOMContentLoaded', function() {
  ['esm-from','esm-subject','esm-to','esm-body'].forEach(function(id) {
    var el = document.getElementById(id);
    if (el) {
      el.addEventListener('input', function() {
        clearTimeout(_esmDebounce);
        _esmDebounce = setTimeout(_esmDoSearch, 250);
      });
      el.addEventListener('keydown', function(e) {
        if (e.key === 'Escape') closeEmailSearchModal();
      });
    }
  });
  var modal = document.getElementById('email-search-modal');
  if (modal) modal.addEventListener('click', function(e) { if (e.target === modal) closeEmailSearchModal(); });
});

function _openEmailInChat(emailMeta) {
  addStatusMsg('Lade E-Mail...');
  var agent = getAgentName();
  fetch('/api/email-content?agent=' + encodeURIComponent(agent) + '&message_id=' + encodeURIComponent(emailMeta.message_id || '') + '&from_email=' + encodeURIComponent(emailMeta.from_email || '') + '&subject=' + encodeURIComponent(emailMeta.subject || ''))
    .then(function(r) { return r.json(); })
    .then(function(data) {
      if (!data.ok) { addStatusMsg('E-Mail konnte nicht geladen werden: ' + (data.error || 'unbekannt')); return; }
      _showEmailCard(data);
    })
    .catch(function(e) { addStatusMsg('Fehler: ' + e.message); });
}

function _showEmailCard(email) {
  var msgs = document.getElementById('messages');
  var div = document.createElement('div');
  div.className = 'msg assistant';
  var time = new Date().toLocaleTimeString('de-DE', {hour:'2-digit', minute:'2-digit'});
  var fromDisplay = email.from_name ? escHtml(email.from_name) + ' &lt;' + escHtml(email.from_email) + '&gt;' : escHtml(email.from_email);
  var bodyText = email.body || '(kein Inhalt)';
  var ccDisplay = email.cc ? escHtml(email.cc) : '';
  var cardId = 'ec-' + Date.now();
  var html = '<div class="bubble"><div class="email-card">'
    + '<div class="email-card-header">'
    + '<div class="email-card-label">\\u2709 E-Mail</div>'
    + '<div class="email-card-subject">' + escHtml(email.subject || '(kein Betreff)') + '</div>'
    + '<div class="email-card-row"><strong>Von:</strong> ' + fromDisplay + '</div>'
    + '<div class="email-card-row"><strong>An:</strong> ' + escHtml(email.to || '') + '</div>'
    + (ccDisplay ? '<div class="email-card-row"><strong>CC:</strong> ' + ccDisplay + '</div>' : '')
    + '<div class="email-card-date">' + escHtml(email.date || '') + '</div>'
    + '</div>'
    + '<div class="email-card-body">' + escHtml(bodyText) + '</div>'
    + (email.message_id ? '<div class="email-card-msgid">Message-ID: ' + escHtml(email.message_id) + '</div>' : '')
    + '<div class="email-card-actions">'
    + '<button class="email-card-btn email-card-btn-reply" data-card="' + cardId + '">Antworten</button>'
    + '<button class="email-card-btn email-card-btn-close" data-card="' + cardId + '">Schliessen</button>'
    + '</div>'
    + '</div></div><div class="meta">' + time + '</div>';
  div.innerHTML = html;
  div.id = cardId;
  msgs.appendChild(div);
  scrollDown();

  div.querySelector('.email-card-btn-reply').onclick = function() {
    var ownAddrs = ['moritz.cremer@me.com', 'londoncityfox@gmail.com', 'moritz.cremer@signicat.com'];
    var ccParts = [];
    [email.to || '', email.cc || ''].forEach(function(field) {
      if (!field) return;
      field.split(',').forEach(function(addr) {
        addr = addr.trim();
        var em = addr.includes('<') ? addr.substring(addr.indexOf('<')+1, addr.indexOf('>')) : addr;
        if (em && !ownAddrs.includes(em.toLowerCase().trim())) ccParts.push(addr);
      });
    });
    _emailContext = {
      from_email: email.from_email || '',
      from_name: email.from_name || '',
      subject: email.subject || '',
      message_id: email.message_id || '',
      cc: ccParts.join(', ')
    };
    addStatusMsg('E-Mail-Kontext gesetzt. Schreibe jetzt deine Antwort-Anweisung.');
    document.getElementById('msg-input').focus();
    this.textContent = '\\u2713 Kontext gesetzt';
    this.style.background = '#2a5a2a';
    this.onclick = null;
  };
  div.querySelector('.email-card-btn-close').onclick = function() { div.remove(); };
}

function _getAndClearEmailContext() {
  if (!_emailContext) return null;
  var ctx = _emailContext;
  _emailContext = null;
  return ctx;
}

function selectSlashCmd(inputEl, cmd) {
  var entry = _SLASH_COMMANDS.find(c => c.cmd === cmd);
  if (cmd === '/create-email-reply' || cmd === '/reply') {
    inputEl.value = '';
    hideSlashAutocomplete();
    showEmailSearchModal();
    return;
  }
  if (entry && entry.template) {
    inputEl.value = entry.template;
    hideSlashAutocomplete();
    inputEl.focus();
    onInputHandler(inputEl);
  } else {
    inputEl.value = cmd + ' ';
    hideSlashAutocomplete();
    inputEl.focus();
    onInputHandler(inputEl);
  }
}

function onSlashAcKey(e, inputEl) {
  const dd = document.getElementById('slash-ac-dropdown');
  if (!dd || dd.style.display === 'none') return false;
  const items = dd.querySelectorAll('.slash-ac-item');
  if (e.key === 'ArrowDown') { e.preventDefault(); _slashAcIdx = Math.min(_slashAcIdx+1, items.length-1); items.forEach((it,i) => it.classList.toggle('active', i===_slashAcIdx)); return true; }
  if (e.key === 'ArrowUp') { e.preventDefault(); _slashAcIdx = Math.max(_slashAcIdx-1, 0); items.forEach((it,i) => it.classList.toggle('active', i===_slashAcIdx)); return true; }
  if (e.key === 'Tab' && items.length > 0) { e.preventDefault(); var idx = _slashAcIdx >= 0 ? _slashAcIdx : 0; var selItem = items[idx]; if (selItem) selectSlashCmd(inputEl, selItem.dataset.cmd); return true; }
  if (e.key === 'Enter' && _slashAcIdx >= 0) { e.preventDefault(); var selItem = items[_slashAcIdx]; if (selItem) selectSlashCmd(inputEl, selItem.dataset.cmd); return true; }
  if (e.key === 'Escape') { hideSlashAutocomplete(); return true; }
  return false;
}

// ─── TYPE FILTER DROPDOWN ────────────────────────────────────────────────────
const _SEARCH_TYPES = [
  {key: 'all', label: '\U0001f50d Alles', shortcut: 'A'},
  {key: 'email', label: '\u2709 E-Mail', shortcut: 'E'},
  {key: 'webclip', label: '\U0001f310 Web Clip', shortcut: 'W'},
  {key: 'screenshot', label: '\U0001f4f8 Screenshot', shortcut: 'S'},
  {key: 'contact', label: '\U0001f464 Kontakt', shortcut: 'K'},
  {key: 'document', label: '\U0001f4c4 Dokument', shortcut: 'D'},
  {key: 'conversation', label: '\U0001f4ac Konversation', shortcut: 'G'},
];
let _typeAcIdx = -1;
let _typeAcVisible = false;

function showTypeAutocomplete(inputEl) {
  let dd = document.getElementById('type-ac-dropdown');
  if (!dd) {
    dd = document.createElement('div');
    dd.id = 'type-ac-dropdown';
    inputEl.parentElement.style.position = 'relative';
    inputEl.parentElement.appendChild(dd);
  }
  dd.innerHTML = '';
  _typeAcIdx = -1;
  _typeAcVisible = true;
  _SEARCH_TYPES.forEach((t, i) => {
    const item = document.createElement('div');
    item.className = 'type-ac-item';
    item.innerHTML = '<span>' + t.label + '</span><span class="type-shortcut">' + t.shortcut + '</span>';
    item.onmousedown = (e) => { e.preventDefault(); selectTypeCmd(inputEl, t.key); };
    dd.appendChild(item);
  });
  dd.style.display = 'block';
}

function hideTypeAutocomplete() {
  const dd = document.getElementById('type-ac-dropdown');
  if (dd) dd.style.display = 'none';
  _typeAcIdx = -1;
  _typeAcVisible = false;
}

function selectTypeCmd(inputEl, typeKey) {
  const val = inputEl.value;
  // Find the /find or /find_global prefix
  const prefix = val.match(/^\/find_global\s*|\/find\s*/);
  if (prefix) {
    if (typeKey === 'all') {
      inputEl.value = prefix[0];
    } else {
      inputEl.value = prefix[0] + typeKey + ' ';
    }
  }
  hideTypeAutocomplete();
  inputEl.focus();
}

function onTypeAcKey(e, inputEl) {
  if (!_typeAcVisible) return false;
  const dd = document.getElementById('type-ac-dropdown');
  if (!dd || dd.style.display === 'none') return false;
  const items = dd.querySelectorAll('.type-ac-item');

  if (e.key === 'ArrowDown') { e.preventDefault(); _typeAcIdx = Math.min(_typeAcIdx+1, items.length-1); items.forEach((it,i) => it.classList.toggle('active', i===_typeAcIdx)); return true; }
  if (e.key === 'ArrowUp') { e.preventDefault(); _typeAcIdx = Math.max(_typeAcIdx-1, 0); items.forEach((it,i) => it.classList.toggle('active', i===_typeAcIdx)); return true; }
  if (e.key === 'Tab' && items.length > 0) { e.preventDefault(); var tidx = _typeAcIdx >= 0 ? _typeAcIdx : 0; selectTypeCmd(inputEl, _SEARCH_TYPES[tidx].key); return true; }
  if (e.key === 'Enter' && _typeAcIdx >= 0) { e.preventDefault(); selectTypeCmd(inputEl, _SEARCH_TYPES[_typeAcIdx].key); return true; }
  if (e.key === 'Escape') { hideTypeAutocomplete(); return true; }

  // Shortcut keys
  const key = e.key.toUpperCase();
  const match = _SEARCH_TYPES.find(t => t.shortcut === key);
  if (match && key.length === 1) { e.preventDefault(); selectTypeCmd(inputEl, match.key); return true; }

  return false;
}

// ─── SEARCH DIALOG ──────────────────────────────────────────────────────────
let _pendingSearchMsg = null;
let _searchResults = [];
let _currentSearchFilter = 'all';
let _isGlobalSearch = false;

const _SOURCE_SUBTYPES = {
  'webclip': ['webclip_salesforce', 'webclip_slack', 'webclip_general'],
  'document': ['document_word', 'document_excel', 'document_pdf', 'document_pptx'],
  'whatsapp': ['whatsapp_direct', 'whatsapp_group'],
};
const _SOURCE_SUBLABELS = {
  'webclip_salesforce': 'Salesforce', 'webclip_slack': 'Slack', 'webclip_general': 'Web',
  'document_word': 'Word', 'document_excel': 'Excel', 'document_pdf': 'PDF', 'document_pptx': 'PowerPoint',
  'whatsapp_direct': 'Direktnachricht', 'whatsapp_group': 'Gruppenchat',
};
const _SOURCE_PARENTS = {
  'notification': 'email', 'webclip_salesforce': 'webclip', 'webclip_slack': 'webclip',
  'webclip_general': 'webclip', 'document_word': 'document', 'document_excel': 'document',
  'document_pdf': 'document', 'document_pptx': 'document',
  'whatsapp_direct': 'whatsapp', 'whatsapp_group': 'whatsapp',
};

function applySearchFilter(filter) {
  _currentSearchFilter = filter;
  // Update button states
  document.querySelectorAll('.search-filter-btn').forEach(b => b.classList.toggle('active', b.dataset.filter === filter));
  // Show subfilter bar if applicable
  const subbar = document.getElementById('search-subfilter-bar');
  if (_SOURCE_SUBTYPES[filter]) {
    subbar.innerHTML = '';
    const allBtn = document.createElement('button');
    allBtn.className = 'search-subfilter-btn active';
    allBtn.textContent = 'Alle';
    allBtn.onclick = () => { _currentSearchFilter = filter; rerenderSearchResults(); document.querySelectorAll('.search-subfilter-btn').forEach(b => b.classList.remove('active')); allBtn.classList.add('active'); };
    subbar.appendChild(allBtn);
    _SOURCE_SUBTYPES[filter].forEach(sub => {
      const btn = document.createElement('button');
      btn.className = 'search-subfilter-btn';
      btn.textContent = _SOURCE_SUBLABELS[sub] || sub;
      btn.onclick = () => { _currentSearchFilter = sub; rerenderSearchResults(); document.querySelectorAll('.search-subfilter-btn').forEach(b => b.classList.remove('active')); btn.classList.add('active'); };
      subbar.appendChild(btn);
    });
    subbar.classList.add('show');
  } else {
    subbar.classList.remove('show');
  }
  rerenderSearchResults();
}

function getFilteredResults() {
  if (_currentSearchFilter === 'all') {
    // Hide notifications in 'all' view
    return _searchResults.filter(r => !r.is_notification && r.source_type !== 'notification');
  }
  const filter = _currentSearchFilter;
  const subs = _SOURCE_SUBTYPES[filter];
  const validTypes = subs ? new Set([filter, ...subs]) : new Set([filter]);
  // For email filter, include notifications
  if (filter === 'email') validTypes.add('notification');
  return _searchResults.filter(r => {
    const st = r.source_type || r.type || 'file';
    const parent = _SOURCE_PARENTS[st] || st;
    return validTypes.has(st) || validTypes.has(parent);
  });
}

function rerenderSearchResults() {
  const filtered = getFilteredResults();
  const list = document.getElementById('search-results-list');
  list.innerHTML = '';
  let checkedCount = 0;

  function addItem(r) {
    const origIdx = _searchResults.indexOf(r);
    const div = document.createElement('div');
    const isNotif = r.is_notification || r.source_type === 'notification';
    div.className = 'search-result-item' + (isNotif ? ' notification' : '');
    const cb = document.createElement('input');
    cb.type = 'checkbox'; cb.dataset.idx = origIdx;
    // All unchecked by default — user selects manually
    cb.onchange = function() { onSearchCheckboxChange(this); };
    if (r.is_notification) cb.dataset.notif = '1';
    div.onclick = (e) => { if (e.target.tagName !== 'INPUT') { cb.checked = !cb.checked; onSearchCheckboxChange(cb); }};
    const info = document.createElement('div');
    info.className = 'search-result-info';
    const nameEl = document.createElement('div');
    nameEl.className = 'search-result-name';
    let prefix = '';
    if (isNotif) prefix = '[Notif] ';
    else if ((r.source_type||r.type) === 'email') prefix = '\u2709 ';
    else if ((r.source_type||r.type||'').startsWith('webclip')) prefix = '\U0001f310 ';
    else if ((r.source_type||r.type||'').startsWith('document')) prefix = '\U0001f4c4 ';
    else if ((r.source_type||r.type) === 'conversation') prefix = '\U0001f4ac ';
    else if ((r.source_type||r.type) === 'screenshot') prefix = '\U0001f4f8 ';
    else prefix = '\U0001f4c4 ';
    if (_isGlobalSearch && r.agent) prefix = '\U0001f310[' + r.agent + '] ' + prefix;
    nameEl.textContent = prefix + r.name;
    const meta = document.createElement('div');
    meta.className = 'search-result-meta';
    let metaParts = [];
    if (isNotif) metaParts.push('<span style="color:#999;font-style:italic">Notifikation</span>');
    if (r.from) { const cls = r.from_person ? 'from-person' : 'from-auto'; metaParts.push('<span class="'+cls+'">Von: '+r.from.substring(0,40)+'</span>'); }
    if (r.subject) metaParts.push(r.subject.substring(0,50));
    if (r.date) metaParts.push(r.date);
    meta.innerHTML = metaParts.join(' \u00b7 ');
    const prev = document.createElement('div');
    prev.className = 'search-result-preview';
    prev.textContent = r.preview || '';
    info.appendChild(nameEl);
    info.appendChild(meta);
    info.appendChild(prev);
    div.appendChild(cb);
    div.appendChild(info);
    list.appendChild(div);
  }

  if (_isGlobalSearch && _currentSearchFilter === 'all') {
    // Group by agent
    const groups = {};
    filtered.forEach(r => {
      const key = r.agent || 'global';
      if (!groups[key]) groups[key] = [];
      groups[key].push(r);
    });
    Object.keys(groups).sort().forEach(agent => {
      const items = groups[agent];
      const d = document.createElement('div');
      d.className = 'search-section-divider';
      d.textContent = '\u2500\u2500 ' + agent + ' (' + items.length + ' Treffer) \u2500\u2500';
      list.appendChild(d);
      items.forEach(r => addItem(r));
    });
  } else if (_currentSearchFilter === 'email') {
    // Email filter: show emails first, notifications last
    const emails = filtered.filter(r => !r.is_notification && r.source_type !== 'notification');
    const notifs = filtered.filter(r => r.is_notification || r.source_type === 'notification');
    emails.forEach(r => addItem(r));
    if (notifs.length > 0) {
      const d = document.createElement('div');
      d.className = 'search-section-divider';
      d.textContent = '\u2500\u2500 Notifikationen (' + notifs.length + ') \u2500\u2500';
      list.appendChild(d);
      notifs.forEach(r => addItem(r));
    }
  } else {
    // Group by source_type subcategory
    const groups = {};
    filtered.forEach(r => {
      const st = r.source_type || r.type || 'file';
      if (!groups[st]) groups[st] = [];
      groups[st].push(r);
    });
    const keys = Object.keys(groups).sort();
    if (keys.length > 1) {
      keys.forEach(st => {
        const label = _SOURCE_SUBLABELS[st] || st;
        const d = document.createElement('div');
        d.className = 'search-section-divider';
        d.textContent = '\u2500\u2500 ' + label + ' (' + groups[st].length + ') \u2500\u2500';
        list.appendChild(d);
        groups[st].forEach(r => addItem(r));
      });
    } else {
      filtered.forEach(r => addItem(r));
    }
  }
  updateSelectionCounter();
}

function updateSelectionCounter() {
  const checked = document.querySelectorAll('#search-results-list input[type=checkbox]:checked').length;
  const counter = document.getElementById('search-selection-counter');
  if (counter) counter.textContent = checked + ' ausgewaehlt (max 50)';
  updateToggleAllBtn();
}

function onSearchCheckboxChange(cb) {
  const checked = document.querySelectorAll('#search-results-list input[type=checkbox]:checked').length;
  if (cb.checked && checked > 50) { cb.checked = false; return; }
  updateSelectionCounter();
}

function showSearchDialog(results, query, isGlobal) {
  _searchResults = results;
  _isGlobalSearch = isGlobal || false;
  _currentSearchFilter = 'all';
  const overlay = document.getElementById('search-overlay');
  if (isGlobal) {
    document.getElementById('search-panel-title').textContent = '🌐 Globale Suche \u2014 ' + results.length + ' Datei(en)';
  } else {
    document.getElementById('search-panel-title').textContent = '🔍 ' + results.length + ' Datei(en) gefunden';
  }
  document.getElementById('search-panel-count').textContent = 'Suche: ' + query;
  var hitCount = document.getElementById('search-hit-count');
  if (hitCount) hitCount.textContent = results.length + ' Dateien gefunden';
  document.querySelectorAll('.search-filter-btn').forEach(b => b.classList.toggle('active', b.dataset.filter === 'all'));
  document.getElementById('search-subfilter-bar').classList.remove('show');
  rerenderSearchResults();
  overlay.classList.add('show');
  document.addEventListener('keydown', _searchEscHandler);
}

function closeSearchDialog(sendAnyway) {
  document.getElementById('search-overlay').classList.remove('show');
  document.removeEventListener('keydown', _searchEscHandler);
  if (sendAnyway && _pendingSearchMsg) {
    doSendChat(_pendingSearchMsg);
  }
  _pendingSearchMsg = null;
  _searchResults = [];
}

function _searchEscHandler(e) {
  if (e.key === 'Escape') { e.preventDefault(); closeSearchDialog(true); }
}

function toggleAllSearchCheckboxes() {
  var cbs = document.querySelectorAll('#search-results-list input[type=checkbox]');
  var nonNotif = Array.from(cbs).filter(function(cb) { return !cb.dataset.notif; });
  var allChecked = nonNotif.length > 0 && nonNotif.every(function(cb) { return cb.checked; });
  var count = 0;
  nonNotif.forEach(function(cb) {
    if (allChecked) { cb.checked = false; }
    else if (count < 50) { cb.checked = true; count++; }
  });
  updateSelectionCounter();
}

function updateToggleAllBtn() {
  var btn = document.getElementById('search-toggle-all-btn');
  if (!btn) return;
  var cbs = document.querySelectorAll('#search-results-list input[type=checkbox]');
  var nonNotif = Array.from(cbs).filter(function(cb) { return !cb.dataset.notif; });
  var allChecked = nonNotif.length > 0 && nonNotif.every(function(cb) { return cb.checked; });
  btn.textContent = allChecked ? 'Alle abwaehlen' : 'Alle markieren';
}

async function loadAllResults() {
  const paths = _searchResults.filter(r => !r.is_notification).slice(0, 50).map(r => r.path);
  await doLoadFiles(paths);
}

async function loadSelectedResults() {
  const checkboxes = document.querySelectorAll('#search-results-list input[type=checkbox]:checked');
  const paths = [];
  checkboxes.forEach(cb => {
    const idx = parseInt(cb.dataset.idx);
    if (_searchResults[idx] && paths.length < 50) paths.push(_searchResults[idx].path);
  });
  if (paths.length === 0) { closeSearchDialog(true); return; }
  await doLoadFiles(paths);
}

async function doLoadFiles(paths) {
  try {
    const r = await fetch('/load_selected_files', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({paths, session_id:SESSION_ID})});
    const data = await r.json();
    if (data.loaded) {
      const agent = getAgentName();
      data.loaded.forEach(n => { addCtxItem(n,'file',true); addFinderLink(agent,n); });
      addStatusMsg('🔍 ' + data.loaded.length + ' Datei(en) in Kontext geladen');
      if (data.agents) { data.agents.forEach(a => { if (a !== agent) addStatusMsg('🌐 Dateien aus Agent: ' + a); }); }
    }
  } catch(e) { addStatusMsg('Fehler beim Laden: '+e.message); }
  document.getElementById('search-overlay').classList.remove('show');
  if (_pendingSearchMsg) {
    doSendChat(_pendingSearchMsg);
  }
  _pendingSearchMsg = null;
  _searchResults = [];
}

// ─── SEND MESSAGE ───────────────────────────────────────────────────────────
async function sendMessage() {
  const input = document.getElementById('msg-input');
  const text = input.value.trim();
  if (!text) return;
  input.value = ''; input.style.height = 'auto';
  var _sendAgent = document.getElementById('agent-label').dataset.agentName;
  if (_sendAgent) { localStorage.removeItem('draft_' + _sendAgent); }
  hideSlashAutocomplete();

  // /find and /find_global commands — intercept, do NOT send to AI
  hideTypeAutocomplete();
  hideFindChips();
  // Match: /find-TYPE query, /find_global-TYPE query, /find query, /find_global query
  const typedFindMatch = text.match(/^\/find(_global)?(?:-(email|whatsapp|webclip|slack|salesforce|screenshot|contact|document|conversation))?(?:\s+(.*))?$/i);
  if (typedFindMatch && getAgentName() !== 'Kein Agent') {
    const isGlobal = !!typedFindMatch[1];
    let searchType = typedFindMatch[2] ? typedFindMatch[2].toLowerCase() : null;
    let rawQuery = (typedFindMatch[3] || '').trim();
    // Also check for old-style /find email query (type as first word)
    if (!searchType) {
      const knownTypes = ['email','whatsapp','webclip','slack','salesforce','screenshot','contact','document','conversation'];
      const firstWord = rawQuery.split(/\s+/)[0].toLowerCase();
      if (knownTypes.includes(firstWord)) {
        searchType = firstWord;
        rawQuery = rawQuery.substring(firstWord.length).trim();
      }
    }
    const query = rawQuery;
    // Empty query — show recent files instead
    if (!query) {
      addMessage('user', text);
      scrollDown();
      const typeLabels2 = {email:'E-Mail',whatsapp:'WhatsApp',webclip:'Web Clip',slack:'Slack',salesforce:'Salesforce',screenshot:'Screenshot',document:'Dokument',conversation:'Konversation'};
      const typeInfo2 = searchType ? ' | Typ: ' + (typeLabels2[searchType]||searchType) : '';
      startTyping('Lade neueste Dateien' + typeInfo2 + '...');
      try {
        const payload2 = {query:'', session_id:SESSION_ID, recent:true};
        if (searchType) payload2.type = searchType;
        payload2.agent = getAgentName();
        const endpoint2 = '/search_preview';
        const r2 = await fetch(endpoint2, {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(payload2)});
        const data2 = await r2.json();
        stopTyping();
        if (data2.ok && data2.results && data2.results.length > 0) {
          _pendingSearchMsg = null;
          showSearchDialog(data2.results, 'Neueste' + typeInfo2, isGlobal);
        } else {
          addStatusMsg('Keine Dateien gefunden' + typeInfo2 + ' (Index leer oder nicht aufgebaut)');
        }
      } catch(e2) { stopTyping(); addStatusMsg('Fehler: '+e2.message); }
      return;
    }
    addMessage('user', text);
    scrollDown();
    const typeLabels = {email:'E-Mail',webclip:'Web Clip',screenshot:'Screenshot',contact:'Kontakt',document:'Dokument',conversation:'Konversation'};
    const typeInfo = searchType ? ' | Typ: ' + (typeLabels[searchType]||searchType) : '';
    startTyping((isGlobal ? 'Globale Suche' : 'Suche') + typeInfo + '...');
    try {
      let r;
      const payload = {query:query, session_id:SESSION_ID};
      if (searchType) payload.type = searchType;
      if (isGlobal) {
        payload.requesting_agent = getAgentName();
        r = await fetch('/global_search_preview', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(payload)});
      } else {
        payload.agent = getAgentName();
        r = await fetch('/search_preview', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(payload)});
      }
      const data = await r.json();
      stopTyping();
      if (data.ok && data.results && data.results.length > 0) {
        _pendingSearchMsg = null;  // /find results stay in dialog, no auto-send to AI
        showSearchDialog(data.results, data.query || query, isGlobal || data.global);
        return;
      } else {
        addStatusMsg('Keine Ergebnisse fuer: ' + query + typeInfo);
      }
    } catch(e) { stopTyping(); addStatusMsg('Suchfehler: ' + e.message); }
    return;
  }

  // CALENDAR_SLASH_V1: Kalender-Befehle intercepten
  if (text.startsWith('/calendar-')) {
    addMessage('user', text);
    scrollDown();
    handleCalendarCommand(text);
    return;
  }

  // CANVA_SLASH_V1: Canva-Befehle intercepten
  if (text.startsWith('/canva-')) {
    addMessage('user', text);
    scrollDown();
    handleCanvaCommand(text);
    return;
  }

  addMessage('user', text);
  scrollDown();
  doSendChat(text);
}

async function handleCalendarCommand(text) {
  startTyping('Kalender...');
  try {
    const parts = text.split(/\s+/);
    const cmd = parts[0];
    const arg = parts.slice(1).join(' ').trim();
    let daysBack = 0, daysAhead = 1, label = 'Heute';

    if (cmd === '/calendar-today') { daysBack = 0; daysAhead = 1; label = 'Heute'; }
    else if (cmd === '/calendar-tomorrow') { daysBack = 0; daysAhead = 2; label = 'Morgen'; }
    else if (cmd === '/calendar-week') { daysBack = 0; daysAhead = 7; label = 'Diese Woche'; }
    else if (cmd === '/calendar-search') {
      if (!arg) { stopTyping(); addStatusMsg('Bitte Suchbegriff angeben: /calendar-search Meeting'); return; }
      daysBack = 30; daysAhead = 30; label = 'Suche: ' + arg;
    }

    const payload = {days_back: daysBack, days_ahead: daysAhead};
    if (cmd === '/calendar-search' && arg) payload.search = arg;

    const r = await fetch('/api/calendar', {method:'POST', headers:{'Content-Type':'application/json'}, body: JSON.stringify(payload)});
    const d = await r.json();
    stopTyping();

    if (d.error) { addStatusMsg('Kalender-Fehler: ' + d.error); return; }

    let msg = '**\U0001f4c5 Kalender — ' + label + '** (' + d.count + ' Termine)\\n';
    if (!d.events || !d.events.length) {
      msg += '\\nKeine Termine gefunden.';
    } else {
      let lastDate = '';
      d.events.forEach(function(e) {
        const startStr = e.start || '';
        const dateOnly = startStr.substring(0, 10);
        if (dateOnly !== lastDate) {
          msg += '\\n**' + dateOnly + '**';
          lastDate = dateOnly;
        }
        const time = e.all_day ? 'Ganztaegig' : startStr.substring(11, 16);
        msg += '\\n- ' + time + ' — **' + e.title + '**';
        if (e.calendar_name) msg += ' _(' + e.calendar_name + ')_';
        if (e.location) msg += ' \U0001f4cd ' + e.location;
      });
    }
    addMessage('assistant', msg);
  } catch(e) {
    stopTyping();
    addStatusMsg('Kalender-Fehler: ' + e.message);
  }
}

async function handleCanvaCommand(text) {
  startTyping('Canva...');
  try {
    const parts = text.split(/\s+/);
    const cmd = parts[0];
    const arg = parts.slice(1).join(' ').trim();

    if (cmd === '/canva-search') {
      const r = await fetch('/api/canva', {method:'POST', headers:{'Content-Type':'application/json'},
        body: JSON.stringify({action:'search', query: arg || '', count: 10})});
      const d = await r.json();
      stopTyping();
      if (d.ok && d.data && d.data.items) {
        let msg = '**Canva Designs' + (arg ? ' fuer "'+arg+'"' : '') + ':**\\n';
        d.data.items.forEach(function(it) {
          msg += '\\n- **' + (it.title || '(Ohne Titel)') + '** (ID: `' + it.id + '`)';
          if (it.urls && it.urls.edit_url) msg += ' — [Bearbeiten](' + it.urls.edit_url + ')';
        });
        if (!d.data.items.length) msg += '\\nKeine Designs gefunden.';
        addMessage('assistant', msg);
      } else { addStatusMsg('Canva-Fehler: ' + (d.error || d.data?.error || 'Unbekannt')); }

    } else if (cmd === '/canva-create') {
      if (!arg) { stopTyping(); addStatusMsg('Bitte Titel angeben: /canva-create Mein Design'); return; }
      const r = await fetch('/api/canva', {method:'POST', headers:{'Content-Type':'application/json'},
        body: JSON.stringify({action:'create', title: arg, design_type: 'doc'})});
      const d = await r.json();
      stopTyping();
      if (d.ok && d.data) {
        let msg = '**Canva Design erstellt:** ' + arg;
        if (d.data.design && d.data.design.urls) msg += '\\n[Design oeffnen](' + d.data.design.urls.edit_url + ')';
        addMessage('assistant', msg);
      } else { addStatusMsg('Canva-Fehler: ' + (d.error || d.data?.error || 'Unbekannt')); }

    } else if (cmd === '/canva-templates') {
      const r = await fetch('/api/canva', {method:'POST', headers:{'Content-Type':'application/json'},
        body: JSON.stringify({action:'brand_templates', query: arg || '', count: 20})});
      const d = await r.json();
      stopTyping();
      if (d.ok && d.data && d.data.items) {
        let msg = '**Canva Brand Templates:**\\n';
        d.data.items.forEach(function(it) {
          msg += '\\n- **' + (it.title || it.id) + '** (ID: `' + it.id + '`)';
        });
        if (!d.data.items.length) msg += '\\nKeine Brand Templates gefunden. Erstelle Templates in Canva und markiere sie als Brand Template.';
        addMessage('assistant', msg);
      } else { addStatusMsg('Canva-Fehler: ' + (d.error || d.data?.error || 'Unbekannt')); }

    } else if (cmd === '/canva-export') {
      const exportParts = arg.split(/\s+/);
      const designId = exportParts[0] || '';
      const fmt = exportParts[1] || 'pdf';
      if (!designId) { stopTyping(); addStatusMsg('Bitte Design-ID angeben: /canva-export DAGxxxx pdf'); return; }
      const r = await fetch('/api/canva', {method:'POST', headers:{'Content-Type':'application/json'},
        body: JSON.stringify({action:'export', design_id: designId, format: fmt})});
      const d = await r.json();
      stopTyping();
      if (d.ok && d.data) {
        addMessage('assistant', '**Export gestartet** (Design: `' + designId + '`, Format: ' + fmt + ')\\nStatus: ' + JSON.stringify(d.data).substring(0, 200));
      } else { addStatusMsg('Export-Fehler: ' + (d.error || d.data?.error || 'Unbekannt')); }

    } else if (cmd === '/canva-campaign') {
      // Kampagne wird ans LLM delegiert — der Agent erstellt die Varianten
      stopTyping();
      doSendChat('Erstelle eine Canva-Kampagne mit mehreren Varianten. ' + (arg || 'Ich moechte verschiedene Ad-Varianten aus einem Brand Template generieren. Zeige mir erst die verfuegbaren Templates.'));

    } else {
      stopTyping();
      addStatusMsg('Unbekannter Canva-Befehl: ' + cmd);
    }
  } catch(e) {
    stopTyping();
    addStatusMsg('Canva-Fehler: ' + e.message);
  }
}

async function doSendChat(text) {
  // Session-ID zum Sendezeitpunkt CAPTUREN. Wenn der Nutzer waehrend des
  // await Tab wechselt, darf die Response nicht ins DOM des neuen Tabs
  // fallen — sie gehoert zu mySid. Deshalb alles ueber mySid routen und
  // nur dann DOM anfassen, wenn mySid aktuell aktiv ist. Sonst puffern.
  var mySid = SESSION_ID;
  var mySt = _tabState(mySid);
  // Inject email context if set (one-time use)
  var ectx = _getAndClearEmailContext();
  if (ectx) {
    var ctxBlock = '[E-MAIL KONTEXT: Von: ' + ectx.from_email + (ectx.from_name ? ' (' + ectx.from_name + ')' : '') + ', Betreff: ' + ectx.subject + ', Message-ID: ' + ectx.message_id + (ectx.cc ? ', CC: ' + ectx.cc : '') + ']\\n\\nUser-Anweisung: ';
    text = ctxBlock + text;
  }
  startTyping(text.substring(0,50), mySid);
  startTaskDiscovery();
  let data;
  try {
    const r = await fetch('/chat', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({message:text, session_id: mySid})});
    data = await r.json();
  } catch(err) {
    stopTaskDiscovery();
    stopTyping(mySid);
    showStopBtn(false, mySid);
    if (_isActiveSession(mySid)) addStatusMsg('Verbindungsfehler: '+err.message);
    return;
  }

  if (data.queued) {
    if (_isActiveSession(mySid)) {
      const ph = addQueuedMessage(text, data.position, data.queue_id);
      mySt.queuedPlaceholders[data.queue_id] = ph;
    }
    // Wenn Tab inaktiv: Placeholder wird spaeter beim Flushen/Tab-Switch
    // anhand der queue_id im DOM gefunden — wir muessen nichts buffern.
    updateQueueDisplay(data.position, mySid);
    startPolling(mySid);
    startTyping(null, mySid);
    return;
  }

  // Direct response
  stopTaskDiscovery();
  stopTyping(mySid);
  if (_isActiveSession(mySid)) {
    handleResponse(data);
  } else {
    // Response gehoert zu inaktivem Tab — puffern und auf Tab-Switch
    // flushen (renderActiveTabState uebernimmt das).
    mySt.pendingResponses.push(data);
  }

  if (data.queue_active) {
    startPolling(mySid);
    startTyping(null, mySid);
  } else {
    showStopBtn(false, mySid);
  }
}

// ── Task Discovery Poller ────────────────────────────────────────────────────
// Discovers server-side video/image generation tasks via /queue_status and
// spawns a progress bar for each new task_id. Runs while a /chat request is
// in flight so the user sees progress during long-running generation.
var taskDiscoveryInterval = null;

async function discoverTasksOnce() {
  try {
    const r = await fetch('/queue_status?session_id=' + SESSION_ID);
    const st = await r.json();
    if (st && Array.isArray(st.active_tasks)) {
      st.active_tasks.forEach(function(t){
        if (!progressBars[t.task_id]) {
          createProgressBar(t.task_id, t.kind || 'video');
        }
        updateProgressBar(t.task_id, t);
      });
    }
  } catch(e) {}
}

function startTaskDiscovery() {
  if (taskDiscoveryInterval) return;
  discoverTasksOnce();
  taskDiscoveryInterval = setInterval(discoverTasksOnce, 2000);
}

function stopTaskDiscovery() {
  if (taskDiscoveryInterval) { clearInterval(taskDiscoveryInterval); taskDiscoveryInterval = null; }
}

// ─── URL / FILE / SEARCH ────────────────────────────────────────────────────────
async function addUrl() {
  let url = document.getElementById('url-input').value.trim();
  if (!url) return;
  if (!url.startsWith('http')) url = 'https://' + url;
  document.getElementById('url-input').value='';
  addStatusMsg('Lade: '+url);
  const r = await fetch('/add_url', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({url, session_id:SESSION_ID})});
  const data = await r.json();
  if (data.ok) { addCtxItem(data.title||url,'url'); addStatusMsg('Bereit: '+url); }
  else addStatusMsg('Fehler: '+data.error);
}

async function addFileFromInput() {
  const files = document.getElementById('file-input').files;
  for (const file of files) {
    const fd = new FormData(); fd.append('file', file);
    addStatusMsg('Lade: '+file.name+'...');
    fd.append('session_id', SESSION_ID);
    const r = await fetch('/add_file', {method:'POST', body:fd});
    const data = await r.json();
    if (data.ok) { addCtxItem(file.name,'file'); addStatusMsg('✓ '+file.name+' — Kontext + Memory gespeichert'); }
    else addStatusMsg('Fehler: '+data.error);
  }
  document.getElementById('file-input').value='';
}

function toggleSearch() {
  searchVisible = !searchVisible;
  document.getElementById('memory-search').style.display = searchVisible?'block':'none';
  document.getElementById('search-toggle').classList.toggle('active', searchVisible);
  if (searchVisible) document.getElementById('memory-search').focus();
}

async function searchMemory() {
  const q = document.getElementById('memory-search').value.trim();
  if (!q) return;
  addStatusMsg('Suche im Memory: '+q+'...');
  const r = await fetch('/search_memory', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({query:q, session_id:SESSION_ID})});
  const data = await r.json();
  if (data.results && data.results.length) {
    const agent = getAgentName();
    data.results.forEach(f => { addCtxItem(f.name,'file'); addFinderLink(agent,f.name); });
    addStatusMsg(data.results.length+' Datei(en) gefunden und geladen.');
  } else addStatusMsg('Keine Ergebnisse fuer: '+q);
}

// ─── FILE AUTOCOMPLETE ──────────────────────────────────────────────────────────
var _fileAcTimer = null;
var _fileAcIdx = -1;
var _fileAcResults = [];

async function onFileAcInput(val) {
  clearTimeout(_fileAcTimer);
  var dd = document.getElementById('file-ac-dropdown');
  if (!val || val.length < 2) { dd.style.display = 'none'; _fileAcResults = []; return; }
  _fileAcTimer = setTimeout(async function() {
    var agent = getAgentName();
    if (!agent || agent === 'Kein Agent') return;
    try {
      var r = await fetch('/api/memory-files-search?q=' + encodeURIComponent(val) + '&agent=' + encodeURIComponent(agent));
      var data = await r.json();
      _fileAcResults = data.results || data || [];
      _fileAcIdx = -1;
      if (_fileAcResults.length === 0) { dd.style.display = 'none'; return; }
      var html = '';
      for (var i = 0; i < Math.min(_fileAcResults.length, 8); i++) {
        var item = _fileAcResults[i];
        html += '<div class="file-ac-item" data-idx="' + i + '" onmousedown="selectFileAc(' + i + ')">';
        html += '<span class="ac-filename">' + escHtml(item.filename || item.name || '') + '</span>';
        if (item.snippet) html += '<span class="ac-snippet">' + escHtml(item.snippet) + '</span>';
        html += '</div>';
      }
      dd.innerHTML = html;
      dd.style.display = 'block';
    } catch(e) { dd.style.display = 'none'; }
  }, 250);
}

function onFileAcKey(e) {
  var dd = document.getElementById('file-ac-dropdown');
  if (dd.style.display === 'none') return;
  var items = dd.querySelectorAll('.file-ac-item');
  if (e.key === 'ArrowDown') { e.preventDefault(); _fileAcIdx = Math.min(_fileAcIdx + 1, items.length - 1); updateFileAcHighlight(items); }
  else if (e.key === 'ArrowUp') { e.preventDefault(); _fileAcIdx = Math.max(_fileAcIdx - 1, 0); updateFileAcHighlight(items); }
  else if (e.key === 'Enter' && _fileAcIdx >= 0) { e.preventDefault(); selectFileAc(_fileAcIdx); }
  else if (e.key === 'Escape') { dd.style.display = 'none'; }
}

function updateFileAcHighlight(items) {
  for (var i = 0; i < items.length; i++) items[i].classList.toggle('selected', i === _fileAcIdx);
}

async function selectFileAc(idx) {
  var item = _fileAcResults[idx];
  if (!item) return;
  var dd = document.getElementById('file-ac-dropdown');
  dd.style.display = 'none';
  document.getElementById('file-ac-input').value = '';
  var fname = item.filename || item.name;
  addStatusMsg('Lade: ' + fname + '...');
  try {
    var r = await fetch('/load_selected_files', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({files:[fname], session_id:SESSION_ID})});
    var data = await r.json();
    if (data.ok) {
      addCtxItem(fname, 'file');
      addStatusMsg('Geladen: ' + fname);
    } else {
      addStatusMsg('Fehler: ' + (data.error || 'unbekannt'));
    }
  } catch(e) { addStatusMsg('Fehler beim Laden: ' + e.message); }
}

// Close dropdown when clicking elsewhere
document.addEventListener('click', function(e) {
  if (!e.target.closest('#file-ac-wrap')) {
    var dd = document.getElementById('file-ac-dropdown');
    if (dd) dd.style.display = 'none';
  }
});

// ─── DRAG & DROP ────────────────────────────────────────────────────────────────
let dragCounter = 0;
document.addEventListener('dragenter', e => { e.preventDefault(); if(++dragCounter===1) document.getElementById('drop-overlay').classList.add('active'); });
document.addEventListener('dragleave', e => { if(--dragCounter===0) document.getElementById('drop-overlay').classList.remove('active'); });
document.addEventListener('dragover', e => e.preventDefault());
document.addEventListener('drop', async e => {
  e.preventDefault(); dragCounter=0; document.getElementById('drop-overlay').classList.remove('active');
  const files = Array.from(e.dataTransfer.files);
  if (!files.length) return;
  if (document.getElementById('agent-label').textContent==='Kein Agent') { addStatusMsg('Bitte zuerst einen Agenten auswaehlen.'); return; }
  for (const file of files) {
    const fd = new FormData(); fd.append('file', file);
    addStatusMsg('Lade: '+file.name+'...');
    try {
      fd.append('session_id', SESSION_ID);
    const r = await fetch('/add_file', {method:'POST', body:fd});
      const data = await r.json();
      if (data.ok) { addCtxItem(file.name,'file'); addStatusMsg('✓ '+file.name+' — Kontext + Memory gespeichert'); }
      else addStatusMsg('Fehler: '+data.error);
    } catch(err) { addStatusMsg('Fehler: '+err.message); }
  }
});

function mdEscHtml(s){
  if (s === null || s === undefined) return '';
  return String(s).replace(/[&<>"']/g, function(c){
    return ({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'})[c];
  });
}
async function handlePreloadMessage(msgId){
  try {
    var r = await fetch('/api/messages/' + encodeURIComponent(msgId));
    var d = await r.json();
    if (!d.ok) return;
    var m = d.message;
    var body = m.full_content || m.preview || '';
    if (body.length > 4000) body = body.substring(0, 4000) + '\\n[... gekuerzt]';
    var previewHtml = mdEscHtml(body.substring(0, 700));
    var banner = document.createElement('div');
    banner.className = 'preload-banner';
    banner.style.cssText = 'background:#1a1a14;border:1px solid #3a3020;border-radius:8px;padding:12px 14px;margin:10px;font-size:12px;color:#d0d0d0;';
    banner.innerHTML =
      '<div style="font-weight:700;color:#f0c060;margin-bottom:6px">\u2709\uFE0E Antwort auf eingehende Nachricht</div>' +
      '<div style="color:#c0c0c0;font-size:11px"><b>Von:</b> ' + mdEscHtml(m.sender_name) +
        (m.sender_address ? ' &lt;' + mdEscHtml(m.sender_address) + '&gt;' : '') + '</div>' +
      '<div style="color:#c0c0c0;font-size:11px"><b>Betreff:</b> ' + mdEscHtml(m.subject || '(kein Betreff)') + '</div>' +
      '<div style="color:#888;font-size:11px;margin-top:6px;max-height:140px;overflow-y:auto;white-space:pre-wrap;background:#111;border:1px solid #262626;border-radius:6px;padding:8px">' +
        previewHtml + '</div>';
    var msgsEl = document.getElementById('messages');
    if (msgsEl) msgsEl.insertBefore(banner, msgsEl.firstChild);
    var input = document.getElementById('msg-input');
    if (input) {
      var sep = '---';
      var NL = String.fromCharCode(10);
      var parts = [
        sep + ' EINGEHENDE NACHRICHT ' + sep,
        'Von: ' + (m.sender_name || '') + (m.sender_address ? ' <' + m.sender_address + '>' : ''),
        'Datum: ' + (m.timestamp || ''),
        'Betreff: ' + (m.subject || ''),
        '',
        body,
        '',
        sep + ' ENDE NACHRICHT ' + sep,
        '',
        'Bitte hilf mir, auf diese Nachricht zu antworten.'
      ];
      var quote = parts.join(NL);
      input.value = quote;
      if (typeof autoResize === 'function') autoResize(input);
      input.focus();
      try { input.setSelectionRange(quote.length, quote.length); } catch(e) {}
    }
  } catch(e){ console.log('Preload-Fehler:', e); }
}
window.onload = async () => {
  try { await loadProviders(); } catch(e) {}
  // URL-Parameter haben Vorrang vor localStorage (Message-Dashboard Deep-Link)
  var urlParams = new URLSearchParams(window.location.search);
  var agentParam = urlParams.get('agent');
  var preloadId = urlParams.get('preload_message');
  if (agentParam) {
    try {
      await selectAgent(agentParam);
      if (preloadId) { await handlePreloadMessage(preloadId); }
      // URL bereinigen damit Refresh nicht nochmal preloadet
      try { history.replaceState({}, '', '/'); } catch(e) {}
      return;
    } catch(e) { console.log('Deep-link agent load failed:', e); }
  }
  // Auto-restore last active agent (persists across page reloads and server restarts)
  const savedAgent = localStorage.getItem('last_active_agent');
  if (savedAgent) {
    try { await selectAgent(savedAgent); return; } catch(e) { console.log('Auto-restore failed:', e); }
  }
  showAgentModal();
};
</script>
</body>
</html>
"""

@app.route("/")
def index():
    # No-store verhindert dass Browser veraltetes HTML/JS nach einem Code-Update
    # weiterverwenden (sonst laeuft nach Server-Restart die alte UI-Version weiter).
    resp = make_response(render_template_string(HTML))
    resp.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate, max-age=0'
    resp.headers['Pragma'] = 'no-cache'
    resp.headers['Expires'] = '0'
    return resp

@app.route('/models')
def get_models():
    config = load_models()
    provider_names = {'anthropic': 'Anthropic', 'openai': 'OpenAI', 'mistral': 'Mistral', 'gemini': 'Google Gemini'}
    result = []
    for pkey, pdata in config.get('providers', {}).items():
        result.append({
            'provider': pkey,
            'name': provider_names.get(pkey, pkey.title()),
            'models': [dict(m, capabilities=[CAPABILITY_EMOJI.get(c,'') for c in MODEL_CAPABILITIES.get(m['id'], [])]) for m in pdata.get('models', [])]
        })
    if not result:
        result = [{'provider': 'anthropic', 'name': 'Anthropic', 'models': [
            {'id': 'claude-sonnet-4-6', 'name': 'Claude Sonnet 4.6'},
            {'id': 'claude-opus-4-6', 'name': 'Claude Opus 4.6'},
            {'id': 'claude-haiku-4-5-20251001', 'name': 'Claude Haiku 4.5'}
        ]}]
    return jsonify(result)

def get_parent_agent(name):
    """If name is a sub-agent (contains '_'), return parent name. Otherwise None."""
    if '_' in name:
        return name.split('_', 1)[0]
    return None

def get_agent_speicher(name):
    """Return the storage directory for an agent. Sub-agents use parent's directory."""
    parent = get_parent_agent(name)
    return os.path.join(BASE, parent if parent else name)

def get_agent_display_name(name):
    """Return display name: 'signicat > outbound' for sub-agents."""
    parent = get_parent_agent(name)
    if parent:
        sub = name.split('_', 1)[1]
        return parent + ' \u203a ' + sub
    return name

def _agent_description(name):
    """TOOLTIPS_V1: Liest die ersten zwei Saetze aus dem Agent-System-Prompt
    als knappe Beschreibung fuer Frontend-Tooltips. Max 180 Zeichen."""
    try:
        fpath = os.path.join(AGENTS_DIR, name + '.txt')
        if not os.path.exists(fpath):
            return ''
        with open(fpath, 'r', encoding='utf-8', errors='ignore') as f:
            raw = f.read(800)
        # Memory-/System-Marker entfernen
        for marker in ['--- GEDAECHTNIS:', '--- DATEI-ERSTELLUNG ---', '--- WEITERE FAEHIGKEITEN ---', 'VERGANGENE KONVERSATIONEN:']:
            idx = raw.find(marker)
            if idx != -1:
                raw = raw[:idx]
        raw = raw.strip()
        # Erste ~2 Saetze / 180 Zeichen
        snippet = raw.split('\n\n')[0] if '\n\n' in raw else raw
        snippet = ' '.join(snippet.split())
        if len(snippet) > 180:
            snippet = snippet[:177].rstrip() + '...'
        return snippet
    except Exception:
        return ''


@app.route('/agents')
def get_agents():
    files = sorted([f.replace('.txt','') for f in os.listdir(AGENTS_DIR) if f.endswith('.txt')])
    # Group: parents first, then sub-agents nested
    parents = []
    children = {}  # parent_name -> [sub_agent_names]
    for name in files:
        parent = get_parent_agent(name)
        if parent:
            children.setdefault(parent, []).append(name)
        else:
            parents.append(name)
    # Build hierarchical structure
    result = []
    for p in parents:
        subs = []
        for c in children.get(p, []):
            sub_label = c.split('_', 1)[1] if '_' in c else c
            subs.append({'name': c, 'label': sub_label, 'description': _agent_description(c)})
        result.append({
            'name': p,
            'label': p,
            'description': _agent_description(p),
            'has_subagents': len(subs) > 0,
            'subagents': subs
        })
    # Orphan sub-agents (parent has no .txt)
    orphan_parents = set(children.keys()) - set(parents)
    for op in sorted(orphan_parents):
        subs = []
        for c in children[op]:
            sub_label = c.split('_', 1)[1] if '_' in c else c
            subs.append({'name': c, 'label': sub_label, 'description': _agent_description(c)})
        result.append({
            'name': op,
            'label': op,
            'description': '',
            'has_subagents': True,
            'subagents': subs
        })
    return jsonify(result)

@app.route('/close_session', methods=['POST'])
def close_session():
    session_id = request.json.get('session_id', 'default') if request.is_json else 'default'
    state = get_session(session_id)
    close_current_session(state)
    return jsonify({'ok': True})

def parse_konversation_file(pfad):
    """Parse a konversation_*.txt file into a verlauf list (list of dicts with role/content)."""
    try:
        with open(pfad, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()
        messages = []
        lines = raw.split('\n')
        i = 0
        while i < len(lines) and not lines[i].startswith('[') and not lines[i].startswith('Du: '):
            i += 1
        current_role = None
        current_content = []
        for line in lines[i:]:
            if line.startswith('Du: '):
                if current_role and current_content:
                    messages.append({'role': current_role, 'content': '\n'.join(current_content).strip()})
                current_role = 'user'
                current_content = [line[4:]]
            elif line.startswith('Assistant: '):
                if current_role and current_content:
                    messages.append({'role': current_role, 'content': '\n'.join(current_content).strip()})
                current_role = 'assistant'
                current_content = [line[11:]]
            elif current_role:
                current_content.append(line)
        if current_role and current_content:
            messages.append({'role': current_role, 'content': '\n'.join(current_content).strip()})
        return [m for m in messages if m['content'].strip()]
    except Exception as e:
        print(f'[PARSE] Error parsing {pfad}: {e}')
        return []


def find_latest_konversation(speicher, name):
    """Find the most recent konversation file for today. Returns (path, verlauf) or (None, [])."""
    import glob as _glob_mod
    today = datetime.datetime.now().strftime("%Y-%m-%d")
    parent = get_parent_agent(name)
    if parent:
        sub_label = name.split('_', 1)[1]
        pattern = os.path.join(speicher, 'konversation_' + today + '_*_' + sub_label + '.txt')
    else:
        pattern = os.path.join(speicher, 'konversation_' + today + '_*.txt')
    files = sorted(_glob_mod.glob(pattern))
    if parent:
        pass
    else:
        files = [f for f in files if '_' not in os.path.basename(f).replace('konversation_' + today + '_', '', 1).replace('.txt', '') or os.path.basename(f).count('_') == 2]
    if not files:
        return None, []
    latest = files[-1]
    verlauf = parse_konversation_file(latest)
    return latest, verlauf


@app.route('/select_agent', methods=['POST'])
def select_agent():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    name = request.json['agent']

    # Guard: block agent switch while processing
    if state.get('processing'):
        return jsonify({'ok': False, 'error': 'Agent-Wechsel nicht moeglich waehrend eine Antwort generiert wird. Bitte warte bis die Antwort fertig ist.'})

    prompt_file = os.path.join(AGENTS_DIR, name + '.txt')
    if not os.path.exists(prompt_file):
        return jsonify({'ok': False, 'error': f'Agent "{name}" nicht gefunden (keine .txt Datei)'})

    # Sub-agent: use parent's storage directory for memory sharing
    speicher = get_agent_speicher(name)
    os.makedirs(speicher, exist_ok=True)

    with open(prompt_file) as f:
        raw_prompt = f.read()

    # Strip ANY memory/capability blocks from the file - keep only the user's base prompt
    for marker in ['\n\n--- GEDAECHTNIS:', '\n--- GEDAECHTNIS:',
                   '\n\n--- DATEI-ERSTELLUNG ---', '\n--- DATEI-ERSTELLUNG ---',
                   '\n\n--- WEITERE FAEHIGKEITEN ---']:
        raw_prompt = raw_prompt.split(marker)[0]
    base_prompt = raw_prompt.strip()

    # Write clean base prompt back to txt file
    with open(prompt_file, 'w') as f:
        f.write(base_prompt)

    # Migrate old conversations not yet indexed
    migrated = migrate_old_conversations(speicher)

    # Build working memory (persistent agent knowledge base)
    wm_ctx = load_working_memory(name)

    # Build memory context (always from parent's storage)
    display_name = get_agent_display_name(name)
    memory_ctx = build_memory_context(speicher, display_name)
    system_prompt = base_prompt + wm_ctx + memory_ctx

    # Memory info for UI
    index = load_index(speicher)
    if index:
        n = len(index)
        last = index[-1]['date']
        migration_note = " (+" + str(migrated) + " importiert)" if migrated > 0 else ""
        memory_info = str(n) + " vergangene Session(s) geladen. Letzte: " + last + migration_note
    else:
        memory_info = None

    # Load saved provider/model preference for this agent
    agent_prefs = _load_agent_prefs()
    agent_pref = agent_prefs.get(name, {})
    pref_provider = agent_pref.get('provider', '')
    pref_model = agent_pref.get('model', '')

    # Re-click same agent: keep current session, don't reset
    if name == state.get('agent'):
        state['system_prompt'] = system_prompt
        return jsonify({'ok': True, 'memory_info': memory_info, 'base_prompt': system_prompt,
                        'pref_provider': pref_provider, 'pref_model': pref_model})

    # Save current session before switching (defensive)
    if state.get('agent') and state.get('verlauf'):
        auto_save_session(session_id)

    # Neuer Tab / neue Session / Agent-Wechsel: IMMER frische Konversation.
    # Die Datei wird bewusst noch NICHT angelegt — sie entsteht erst, wenn
    # der Nutzer den ersten Prompt abschickt (siehe auto_save_session). So
    # werden leere Konversationsdateien vermieden, die nur Rauschen in der
    # History-Sidebar erzeugen wuerden.
    # Fuer das Wiederaufnehmen einer alten Konversation benutzt der Nutzer
    # die History-Sidebar (onHistoryClick → /load_conversation).
    dateiname = None
    new_verlauf = []

    # Add file creation capability to system prompt
    file_capability = """

--- DATEI-ERSTELLUNG ---
Du kannst Word-, Excel- und PDF-Dateien erstellen. Wenn der Nutzer eine Datei moechte, antworte mit einem JSON-Block am Ende deiner Antwort:

Fuer Word (.docx):
[CREATE_FILE:docx:{"title":"Dateiname","content":[{"type":"heading","level":1,"text":"Titel"},{"type":"paragraph","text":"Text"},{"type":"bullet","text":"Punkt"},{"type":"table","rows":[["Spalte1","Spalte2"],["Wert1","Wert2"]]}]}]

Fuer Excel (.xlsx):
[CREATE_FILE:xlsx:{"title":"Dateiname","sheets":[{"name":"Tabelle1","rows":[["Kopf1","Kopf2"],["Wert1","Wert2"]]}]}]

Fuer PDF (.pdf):
[CREATE_FILE:pdf:{"title":"Dateiname","content":[{"type":"heading","level":1,"text":"Titel"},{"type":"paragraph","text":"Text"}]}]

Fuer PowerPoint (.pptx):
[CREATE_FILE:pptx:{"title":"Praesentation","slides":[{"type":"title","heading":"Untertitel"},{"type":"content","heading":"Folientitel","body":"Einleitungstext","bullets":["Punkt 1","Punkt 2"],"footer":"Fusszeile"}]}]

Fuer E-Mail-Draft (oeffnet Apple Mail):
[CREATE_EMAIL:{"to":"empfaenger@example.com","cc":"","subject":"Betreff","body":"E-Mail Text hier","from":"optionale-absender@example.com"}]

Fuer E-Mail-Antwort (oeffnet Apple Mail Reply mit korrektem Threading):
[CREATE_EMAIL_REPLY:{"message_id":"<original-message-id@domain.com>","to":"absender@example.com","cc":"andere@example.com","subject":"Re: Betreff","body":"Antworttext hier","quote_original":true,"from":"optionale-absender@example.com"}]

WICHTIG: Du schickst NIEMALS eine E-Mail direkt ab. Du erstellst IMMER nur einen Draft der in Apple Mail geoeffnet wird. Das gilt fuer alle Formulierungen: 'schreibe', 'sende', 'schick', 'antworte', 'reply', 'forward' — immer CREATE_EMAIL oder CREATE_EMAIL_REPLY, niemals direkt senden.

Wenn du auf eine E-Mail antwortest (z.B. "antworte auf diese E-Mail", "reply", "Antwort verfassen", "schreib eine Antwort"):
- Verwende CREATE_EMAIL_REPLY statt CREATE_EMAIL fuer korrekte Threading
- message_id: Lies den Message-ID Header aus der Original-E-Mail (z.B. aus der .eml Datei im Kontext)
- to: Original-Absender (aus "Von:" Feld)
- cc: Alle Original-Empfaenger aus "An:" AUSSER moritz.cremer@me.com und londoncityfox@gmail.com
- subject: "Re: " + Original-Betreff (nur wenn nicht bereits "Re:" vorhanden)
- body: Dein Antworttext
- quote_original: true wenn das Original zitiert werden soll (Standard), false wenn nicht
- Falls keine message_id verfuegbar: verwende CREATE_EMAIL als Fallback
- from: (Optional) Absender-E-Mail-Adresse. Wenn angegeben, wird dieser Account in Apple Mail als Sender verwendet. Nur setzen wenn ein bestimmter Absender-Account gewuenscht ist.

Verwende dies immer wenn der Nutzer ein Dokument, eine Tabelle, PDF, Praesentation oder E-Mail benoetigt.

Fuer WhatsApp-Nachricht (oeffnet WhatsApp mit vorausgefuellter Nachricht):
[CREATE_WHATSAPP:{"to":"Vorname oder Name","message":"Nachrichtentext hier"}]

WICHTIG: WhatsApp-Nachrichten werden NIEMALS automatisch gesendet. Die App wird nur geoeffnet mit vorausgefuelltem Text. Der Nutzer muss manuell auf Senden klicken. Verwende dies wenn der Nutzer sagt: 'schreib auf WhatsApp', 'WhatsApp an', 'schick per WhatsApp'.

KEINE WIEDERHOLUNG: Wenn im Verlauf bereits eine Aktion ausgefuehrt wurde (erkennbar an '[... — Aktion ausgefuehrt]'), erzeuge diese Aktion NICHT erneut. Jede CREATE_EMAIL, CREATE_EMAIL_REPLY, CREATE_WHATSAPP, CREATE_SLACK und CREATE_VIDEO Aktion darf nur EINMAL pro expliziter Nutzer-Anfrage erzeugt werden. Wenn der Nutzer eine NEUE Anfrage stellt (z.B. Video generieren statt WhatsApp), fuehre NUR die neue Aktion aus.

Fuer Slack-Nachricht (oeffnet Slack Desktop mit vorausgefuelltem Text):
[CREATE_SLACK:{"channel":"#kanalname","message":"Nachrichtentext hier"}]

Fuer Slack-DM:
[CREATE_SLACK:{"to":"Vorname Nachname","message":"Nachrichtentext hier"}]

WICHTIG: Slack-Nachrichten werden NIEMALS automatisch gesendet. Die App wird geoeffnet, Text wird eingefuegt. Der Nutzer muss manuell auf Senden klicken. Verwende dies wenn der Nutzer sagt: 'schreib auf Slack', 'Slack an', 'schick per Slack', 'poste in #channel'.
--- WEITERE FAEHIGKEITEN ---
- Web-Suche: Du kannst aktuelle Informationen aus dem Internet abrufen. Nutze dies automatisch wenn der Nutzer nach aktuellen Infos, Preisen, Nachrichten oder Website-Inhalten fragt.
- Bilder lesen: Hochgeladene Screenshots und Fotos kannst du analysieren und beschreiben.
- Dateien lesen: PDF, Word, Excel werden automatisch extrahiert und stehen als Kontext zur Verfuegung.
- Bilder erstellen: Wenn der Nutzer ein Bild moechte, verwende CREATE_IMAGE. Das System wechselt automatisch auf das passende Bildmodell des aktiven Anbieters (Imagen 4 bei Google Gemini, gpt-image-1 bei OpenAI). Sage NIEMALS dass du keine Bilder erstellen kannst.
- Videos erstellen: Wenn der Nutzer ein Video moechte, verwende CREATE_VIDEO. Das System nutzt automatisch Google Veo. Sage NIEMALS dass du keine Videos erstellen kannst.
- Kalender: Du hast Zugriff auf den Kalender (Fantastical/Apple Calendar). Befehle: /calendar-today, /calendar-tomorrow, /calendar-week, /calendar-search [query]. Termine werden auch automatisch eingeblendet wenn der Nutzer danach fragt.
- Canva Designs: Du kannst Canva-Designs suchen, erstellen und exportieren. Verwende die /canva-Befehle oder sage einfach "erstelle ein Canva Design", "suche in Canva", "exportiere als PDF". Fuer Kampagnen mit Brand Templates: /canva-campaign.
--- ENDE DATEI-ERSTELLUNG ---"""

    system_prompt = system_prompt + file_capability

    update_dict = {
        'agent': name, 'system_prompt': system_prompt, 'speicher': speicher,
        'verlauf': new_verlauf, 'dateiname': dateiname, 'kontext_items': [], 'session_files': []
    }
    # Set provider/model from saved preference (if available)
    if pref_provider and pref_provider in VALID_PROVIDERS:
        update_dict['provider'] = pref_provider
    if pref_model:
        update_dict['model_id'] = pref_model
    state.update(update_dict)

    # txt file stays clean (base prompt only) - sidebar shows full prompt with memory
    # Build/update search index in background
    if build_index_async and state.get('speicher'):
        build_index_async(state['speicher'])
    return jsonify({'ok': True, 'memory_info': memory_info, 'base_prompt': system_prompt,
                    'pref_provider': pref_provider, 'pref_model': pref_model,
                    'recovered_messages': new_verlauf if new_verlauf else []})


@app.route('/new_conversation', methods=['POST'])
def new_conversation():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    name = request.json['agent']
    if not state['agent']:
        return jsonify({'ok': False})
    datum = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M")
    dateiname = os.path.join(state['speicher'], 'konversation_' + datum + '.txt')
    tmp_path = dateiname + '.tmp'
    with open(tmp_path, 'w') as f:
        f.write('Agent: ' + name + '\nDatum: ' + datum + '\n\n')
    os.replace(tmp_path, dateiname)
    state['verlauf'] = []
    state['dateiname'] = dateiname
    state['kontext_items'] = []
    state['session_files'] = []
    return jsonify({'ok': True})

# ─── ACCESS CONTROL ───────────────────────────────────────────────────────
ACCESS_CONTROL_FILE = os.path.join(BASE, "config", "access_control.json")

# Built-in shared sources (nicht entfernbar). Pfade relativ zu BASE, soweit
# sinnvoll — whatsapp haengt als ordnerbasierte Quelle im privat-Memory.
BUILTIN_SHARED_SOURCES = [
    {"key": "email_inbox",    "label": "E-Mail Inbox",    "icon": "\u2709",      "path": os.path.join(BASE, "email_inbox")},
    {"key": "webclips",       "label": "Webclips",        "icon": "\U0001F310",  "path": os.path.join(BASE, "webclips")},
    {"key": "calendar",       "label": "Kalender",        "icon": "\U0001F4C5",  "path": os.path.join(BASE, "calendar")},
    {"key": "working_memory", "label": "Working Memory",  "icon": "\U0001F9E0",  "path": ""},   # pro Agent, kein globaler Pfad
    {"key": "whatsapp",       "label": "WhatsApp Chats",  "icon": "\U0001F4AC",  "path": os.path.join(BASE, "privat", "memory", "whatsapp")},
]


def _load_access_control():
    try:
        with open(ACCESS_CONTROL_FILE, 'r') as f:
            data = json.load(f)
    except Exception:
        data = {"agents": {}, "last_modified": "", "version": "1.0"}
    data.setdefault("agents", {})
    data.setdefault("custom_sources", [])
    return data


def _save_access_control(data):
    import datetime as _dt
    data['last_modified'] = _dt.datetime.now().isoformat()
    data.setdefault('version', '1.0')
    os.makedirs(os.path.dirname(ACCESS_CONTROL_FILE), exist_ok=True)
    with open(ACCESS_CONTROL_FILE, 'w') as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


def _slugify_source_key(label: str, existing: set) -> str:
    """Erzeugt einen stabilen Key fuer eine Custom-Source aus dem Label.
    Nur Kleinbuchstaben, Zahlen, Unterstriche; erzwingt Eindeutigkeit."""
    import re as _re
    base = _re.sub(r'[^a-z0-9]+', '_', (label or '').lower()).strip('_')
    if not base:
        base = 'custom'
    base = 'custom_' + base
    key = base
    i = 2
    while key in existing:
        key = f"{base}_{i}"
        i += 1
    return key


def _source_status(path: str, key: str = "") -> dict:
    """Gibt {exists, count} fuer einen Ordner-Pfad zurueck.

    Spezial-Behandlung fuer Quellen, deren Dateien nicht in einem zentralen
    Ordner liegen, sondern ueber alle Agent-Memorys verteilt sind:
    - email_inbox: summiert IN_/OUT_-Dateien in <agent>/memory/ ueber alle
      Agent-Ordner (der Staging-Ordner `email_inbox/` ist fast immer leer,
      weil der Email-Watcher verarbeitet und wegraeumt).
    - whatsapp: wenn der per-Agent-Pfad leer ist, Fallback auf globalen Count.
    """
    if key == "email_inbox":
        total = 0
        try:
            for entry in os.listdir(BASE):
                mem = os.path.join(BASE, entry, "memory")
                if not os.path.isdir(mem):
                    continue
                try:
                    for f in os.listdir(mem):
                        if f.endswith('.txt') and ('_IN_' in f or '_OUT_' in f):
                            total += 1
                except OSError:
                    continue
        except OSError:
            return {"exists": False, "count": 0}
        return {"exists": total > 0, "count": total}
    if not path:
        return {"exists": False, "count": 0}
    try:
        if not os.path.isdir(path):
            return {"exists": False, "count": 0}
        cnt = sum(1 for _ in os.listdir(path))
        return {"exists": True, "count": cnt}
    except Exception:
        return {"exists": False, "count": 0}


@app.route('/api/access-control', methods=['GET'])
def api_access_control_get():
    data = _load_access_control()
    # Enrich: fuege statische Shared-Source-Metadaten (Pfad, Status) bei,
    # damit das Frontend keine Pfade kennen muss.
    shared = []
    for src in BUILTIN_SHARED_SOURCES:
        entry = dict(src)
        entry['builtin'] = True
        if src['key'] == 'email_inbox':
            entry['status'] = _source_status(src['path'], key='email_inbox')
        elif src['path']:
            entry['status'] = _source_status(src['path'])
        else:
            entry['status'] = {"exists": True, "count": 0}
        shared.append(entry)
    for src in data.get('custom_sources', []):
        entry = dict(src)
        entry['builtin'] = False
        entry['status'] = _source_status(src.get('path', ''))
        shared.append(entry)
    data['shared_sources'] = shared
    return jsonify(data)


@app.route('/api/access-control/custom-sources', methods=['POST'])
def api_access_control_add_custom_source():
    """Body: {label: str, path: str, icon?: str}. Pfad muss existieren."""
    payload = request.get_json(silent=True) or {}
    label = (payload.get('label') or '').strip()
    path = (payload.get('path') or '').strip()
    icon = (payload.get('icon') or '').strip() or "\U0001F4C1"
    if not label:
        return jsonify({'success': False, 'error': 'Label fehlt'}), 400
    if not path:
        return jsonify({'success': False, 'error': 'Pfad fehlt'}), 400
    # Expand ~ und pruefen
    path = os.path.expanduser(path)
    if not os.path.isdir(path):
        return jsonify({'success': False, 'error': f'Pfad ist kein Ordner oder existiert nicht: {path}'}), 400

    data = _load_access_control()
    existing_keys = {s['key'] for s in data.get('custom_sources', [])}
    existing_keys |= {s['key'] for s in BUILTIN_SHARED_SOURCES}
    key = _slugify_source_key(label, existing_keys)
    data.setdefault('custom_sources', []).append({
        'key': key,
        'label': label,
        'path': path,
        'icon': icon,
    })
    try:
        _save_access_control(data)
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500
    return jsonify({'success': True, 'key': key, 'saved_at': data['last_modified']})


@app.route('/api/access-control/custom-sources/<key>', methods=['DELETE'])
def api_access_control_remove_custom_source(key):
    data = _load_access_control()
    sources = data.get('custom_sources', [])
    new_sources = [s for s in sources if s.get('key') != key]
    if len(new_sources) == len(sources):
        return jsonify({'success': False, 'error': f'Custom-Source nicht gefunden: {key}'}), 404
    data['custom_sources'] = new_sources
    # Referenzen aus allen Agents entfernen
    for agent_cfg in data.get('agents', {}).values():
        if isinstance(agent_cfg, dict) and key in (agent_cfg.get('shared_memory') or []):
            agent_cfg['shared_memory'] = [k for k in agent_cfg['shared_memory'] if k != key]
    try:
        _save_access_control(data)
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500
    return jsonify({'success': True, 'saved_at': data['last_modified']})


@app.route('/api/access-control', methods=['POST'])
def api_access_control_post():
    data = request.get_json(silent=True)
    if not data or 'agents' not in data:
        return jsonify({'success': False, 'error': 'Ungueltige Eingabe: agents fehlt'}), 400
    # Validate: each agent must exist as .txt in AGENTS_DIR
    valid_agents = set()
    if os.path.exists(AGENTS_DIR):
        for fname in os.listdir(AGENTS_DIR):
            if fname.endswith('.txt') and '.backup_' not in fname:
                valid_agents.add(fname[:-4])
    for agent in data['agents'].keys():
        if agent not in valid_agents:
            return jsonify({'success': False, 'error': f'Unbekannter Agent: {agent}'}), 400
    # Custom-Sources erhalten: wenn der Client sie nicht mitschickt, aus
    # existierender Datei uebernehmen (Frontend re-schickt sie aber auch).
    if 'custom_sources' not in data:
        existing = _load_access_control()
        data['custom_sources'] = existing.get('custom_sources', [])
    # Enriched 'shared_sources' aus GET nicht persistieren
    data.pop('shared_sources', None)
    try:
        _save_access_control(data)
        return jsonify({'success': True, 'saved_at': data['last_modified']})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/admin/access-control', methods=['GET'])
def admin_access_control_page():
    html = """<!DOCTYPE html>
<html><head><meta charset="UTF-8"><title>Access Control — AssistantDev</title>
<style>
* { box-sizing:border-box; margin:0; padding:0; }
body { background:#1a1a2e; color:#e0e0e0; font-family:-apple-system,Inter,sans-serif; padding:24px; }
.container { max-width:1400px; margin:0 auto; }
.admin-topbar { display:flex; align-items:center; gap:12px; padding:10px 0 16px; border-bottom:1px solid #334; margin-bottom:20px; flex-wrap:wrap; }
.admin-topbar a { color:#aaa; text-decoration:none; font-size:12px; padding:6px 14px; border:1px solid #444; border-radius:6px; transition:all .15s; }
.admin-topbar a:hover { border-color:#f0c060; color:#f0c060; text-decoration:none; }
.admin-topbar .back-btn { background:#222; border-color:#555; color:#ccc; font-weight:600; }
.admin-topbar .back-btn:hover { background:#333; border-color:#f0c060; color:#f0c060; }
h1 { color:#f0c060; font-size:24px; margin:8px 0 4px; }
.subtitle { color:#888; font-size:13px; margin-bottom:20px; }
.msg { padding:10px 14px; border-radius:6px; margin:10px 0; font-size:13px; display:none; }
.msg.success { background:#1f4a1f; border:1px solid #4a8a4a; color:#a0d090; display:block; }
.msg.error { background:#4a1f1f; border:1px solid #8a4a4a; color:#d09090; display:block; }

.matrix-wrap { overflow-x:auto; margin:16px 0; border:1px solid #334; border-radius:10px; background:#16162e; max-height:75vh; overflow-y:auto; }
.matrix-table { border-collapse:separate; border-spacing:0; width:100%; min-width:600px; }
.matrix-table thead th { position:sticky; top:0; z-index:20; background:#1e1e40; padding:0; border-bottom:2px solid #445; }
.matrix-table thead th.corner { z-index:30; left:0; min-width:260px; width:260px; background:#1e1e40; }
.agent-header { writing-mode:vertical-rl; transform:rotate(180deg); padding:12px 6px 8px; font-size:12px; font-weight:600; color:#c0c0d0; white-space:nowrap; text-align:left; min-height:80px; letter-spacing:0.3px; }
.matrix-table td.source-cell, .matrix-table th.corner { position:sticky; left:0; z-index:10; background:#1a1a34; }
.matrix-table td.source-cell { min-width:300px; width:300px; padding:8px 14px; border-right:1px solid #334; font-size:13px; }
.source-top { display:flex; align-items:center; white-space:nowrap; }
.source-icon { margin-right:8px; font-size:15px; display:inline-block; }
.source-name { font-weight:500; color:#d0d0e0; }
.source-badge { display:inline-block; font-size:10px; padding:1px 7px; border-radius:8px; margin-left:8px; font-weight:600; vertical-align:middle; }
.source-path { margin-top:3px; margin-left:23px; font-size:10px; color:#6a6a80; font-family:ui-monospace,Menlo,monospace; overflow:hidden; text-overflow:ellipsis; white-space:nowrap; max-width:270px; }
.source-path.missing { color:#c07070; }
.source-count { color:#5a7a9a; }
.source-remove { background:transparent; border:1px solid #6a3a3a; color:#c08080; border-radius:4px; font-size:11px; line-height:1; padding:2px 6px; margin-left:8px; cursor:pointer; }
.source-remove:hover { background:#4a1f1f; color:#fff; border-color:#a04040; }
.badge-shared { background:#2a4a6a; color:#7ab8f5; }
.badge-exclusive { background:#3a2a5a; color:#b08ae0; }
.badge-custom { background:#3a4a2a; color:#a0d070; }
.section-row td { background:#12122a !important; color:#888; font-size:11px; font-weight:700; text-transform:uppercase; letter-spacing:1.5px; padding:10px 14px !important; border-bottom:1px solid #334; }
.section-row td.source-cell { background:#12122a !important; }
.matrix-table td.data-cell { text-align:center; padding:6px 4px; border-bottom:1px solid #222240; border-right:1px solid #222240; min-width:44px; }
.matrix-table tr:hover td.data-cell { background:#22224a; }
.matrix-table tr:hover td.source-cell { background:#1e1e44; }
.ac-cb { width:18px; height:18px; accent-color:#1B6FD8; cursor:pointer; border-radius:3px; }
.ac-cb:checked { box-shadow:0 0 4px rgba(27,111,216,0.4); }
.btn-row { margin-top:16px; padding:16px 0; border-top:1px solid #334; display:flex; align-items:center; gap:12px; flex-wrap:wrap; }
.btn { padding:10px 28px; border:none; border-radius:6px; cursor:pointer; font-size:14px; font-family:Inter,sans-serif; font-weight:600; }
.btn-primary { background:#1B6FD8; color:#fff; }
.btn-primary:hover { background:#2580e8; }
.btn-secondary { background:#333; color:#aaa; }
.btn-secondary:hover { background:#444; }
.last-mod { color:#666; font-size:11px; margin-left:auto; }
.loading { text-align:center; padding:60px; color:#888; font-size:14px; }
.modal-back { display:none; position:fixed; inset:0; background:rgba(0,0,0,0.6); z-index:100; align-items:center; justify-content:center; }
.modal-back.open { display:flex; }
.modal { background:#1e1e3a; border:1px solid #445; border-radius:10px; padding:24px; min-width:460px; max-width:620px; }
.modal h2 { color:#f0c060; font-size:17px; margin:0 0 14px; }
.modal label { display:block; color:#aaa; font-size:12px; margin:10px 0 4px; }
.modal input[type=text] { width:100%; padding:8px 10px; background:#11112a; border:1px solid #445; border-radius:5px; color:#e0e0e0; font-size:13px; font-family:inherit; }
.modal input[type=text]:focus { outline:none; border-color:#4a8aca; }
.modal .hint { color:#666; font-size:11px; margin-top:4px; }
.modal .btn-row { border-top:none; padding-top:16px; margin-top:16px; }
.btn-ghost { background:transparent; color:#aaa; border:1px solid #445; padding:10px 28px; border-radius:6px; cursor:pointer; font-size:14px; font-weight:600; font-family:inherit; }
.btn-ghost:hover { background:#22224a; color:#fff; }
.modal-err { color:#d09090; font-size:12px; margin-top:8px; min-height:16px; }
</style></head><body>
<div class="container">
<div class="admin-topbar">
  <a class="back-btn" href="/">&#8592; Chat</a>
  <a href="/admin">Admin</a>
  <a href="/admin/access-control" style="border-color:#4a8aca;color:#4a8aca;background:#1a2a4a;">Access Control</a>
  <a href="/admin/permissions">Berechtigungen</a>
  <a href="/admin/docs">Docs</a>
  <a href="/admin/changelog">Changelog</a>
</div>
<h1>&#9881; Access Control</h1>
<div class="subtitle">Zugriffsrechte als Matrix &mdash; Datenquellen (Zeilen) &times; Agenten (Spalten)</div>
<div id="msg" class="msg"></div>
<div id="matrix-container" class="loading">Lade Konfiguration...</div>
<div class="btn-row">
  <button class="btn btn-primary" onclick="saveMatrix()">&#128190; Speichern</button>
  <button class="btn btn-secondary" onclick="loadMatrix()">Verwerfen</button>
  <button class="btn btn-ghost" onclick="openAddSourceModal()">&#128193; Ordner hinzufuegen</button>
  <span class="last-mod" id="last-mod"></span>
</div>
</div>

<div id="add-source-modal" class="modal-back" onclick="if(event.target===this)closeAddSourceModal()">
  <div class="modal">
    <h2>&#128193; Neue Datenquelle hinzufuegen</h2>
    <label>Anzeigename</label>
    <input type="text" id="new-src-label" placeholder="z.B. Slack Export 2026">
    <label>Ordner-Pfad (absolut oder mit ~)</label>
    <input type="text" id="new-src-path" placeholder="/Users/.../Downloads/slack-export oder ~/Documents/...">
    <div class="hint">Der Ordner muss existieren. Nach dem Hinzufuegen kannst du ihn pro Agent per Checkbox freigeben.</div>
    <div class="modal-err" id="add-src-err"></div>
    <div class="btn-row">
      <button class="btn btn-primary" onclick="submitAddSource()">Hinzufuegen</button>
      <button class="btn btn-ghost" onclick="closeAddSourceModal()">Abbrechen</button>
    </div>
  </div>
</div>

<script>
var _acData = null;
var _agents = [];
var _sharedSources = [];

function escH(s){ return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;'); }

async function loadMatrix(){
  var acRes = await fetch('/api/access-control');
  var agRes = await fetch('/agents');
  _acData = await acRes.json();
  var agList = await agRes.json();
  _agents = [];
  agList.forEach(function(a){
    _agents.push(a.name);
    if(a.subagents) a.subagents.forEach(function(s){ _agents.push(s.name); });
  });
  _agents.sort();
  if(!_acData.agents) _acData.agents = {};
  _agents.forEach(function(n){
    if(!_acData.agents[n]) _acData.agents[n] = {own_memory:true, shared_memory:[], cross_agent_read:[], description:''};
  });
  _sharedSources = _acData.shared_sources || [];
  renderMatrix();
  var lm = document.getElementById('last-mod');
  lm.textContent = _acData.last_modified ? 'Zuletzt gespeichert: '+_acData.last_modified : '';
}

function countAccess(sourceType, sourceKey){
  var c=0;
  _agents.forEach(function(ag){
    var a = _acData.agents[ag];
    if(!a) return;
    if(sourceType==='own' && sourceKey===ag && a.own_memory) c++;
    if(sourceType==='shared' && (a.shared_memory||[]).indexOf(sourceKey)>=0) c++;
    if(sourceType==='cross' && (a.cross_agent_read||[]).indexOf(sourceKey)>=0) c++;
  });
  return c;
}

function badgeHtml(count){
  if(count<=0) return '<span class="source-badge" style="background:#333;color:#666;">Kein Zugriff</span>';
  if(count===1) return '<span class="source-badge badge-exclusive">Exklusiv</span>';
  return '<span class="source-badge badge-shared">Geteilt von '+count+'</span>';
}

function renderMatrix(){
  var container = document.getElementById('matrix-container');
  container.className = '';
  container.innerHTML = '';
  var wrap = document.createElement('div');
  wrap.className = 'matrix-wrap';
  var h = '<table class="matrix-table"><thead><tr><th class="corner"></th>';
  _agents.forEach(function(ag){
    h += '<th><div class="agent-header">'+escH(ag)+'</div></th>';
  });
  h += '</tr></thead><tbody>';
  h += '<tr class="section-row"><td class="source-cell" colspan="'+(_agents.length+1)+'">&#128274; Eigenes Memory</td></tr>';
  _agents.forEach(function(ag){
    var cnt = _acData.agents[ag] && _acData.agents[ag].own_memory ? 1 : 0;
    h += '<tr><td class="source-cell"><span class="source-icon">&#129504;</span><span class="source-name">'+escH(ag)+'</span>'+badgeHtml(cnt)+'</td>';
    _agents.forEach(function(col){
      if(col===ag){
        var chk = (_acData.agents[ag] && _acData.agents[ag].own_memory) ? ' checked' : '';
        h += '<td class="data-cell"><input type="checkbox" class="ac-cb" data-type="own" data-source="'+escH(ag)+'" data-agent="'+escH(col)+'"'+chk+' onchange="updateBadges()"></td>';
      } else {
        h += '<td class="data-cell" style="background:#14142a;"></td>';
      }
    });
    h += '</tr>';
  });
  h += '<tr class="section-row"><td class="source-cell" colspan="'+(_agents.length+1)+'">&#128279; Shared Memory</td></tr>';
  _sharedSources.forEach(function(src){
    var cnt = countAccess('shared', src.key);
    var pathHtml = '';
    if(src.path){
      var st = src.status || {exists:false, count:0};
      var cls = st.exists ? '' : ' missing';
      var cnt2 = st.exists ? ('<span class="source-count"> &middot; '+st.count+' Eintr&auml;ge</span>') : ' &middot; nicht gefunden';
      pathHtml = '<div class="source-path'+cls+'"><code>'+escH(src.path)+'</code>'+cnt2+'</div>';
    } else if(src.key === 'working_memory'){
      pathHtml = '<div class="source-path">&lt;agent&gt;/working_memory/</div>';
    }
    var removeBtn = src.builtin ? '' :
      '<button class="source-remove" title="Quelle entfernen" onclick="removeCustomSource('+"'"+escH(src.key)+"'"+',event)">&times;</button>';
    var customBadge = src.builtin ? '' : '<span class="source-badge badge-custom">Custom</span>';
    h += '<tr><td class="source-cell">' +
         '<div class="source-top"><span class="source-icon">'+(src.icon||'&#128193;')+'</span>' +
         '<span class="source-name">'+escH(src.label)+'</span>'+customBadge+badgeHtml(cnt)+removeBtn+'</div>' +
         pathHtml +
         '</td>';
    _agents.forEach(function(ag){
      var a = _acData.agents[ag];
      var chk = (a && (a.shared_memory||[]).indexOf(src.key)>=0) ? ' checked' : '';
      h += '<td class="data-cell"><input type="checkbox" class="ac-cb" data-type="shared" data-source="'+escH(src.key)+'" data-agent="'+escH(ag)+'"'+chk+' onchange="updateBadges()"></td>';
    });
    h += '</tr>';
  });
  h += '<tr class="section-row"><td class="source-cell" colspan="'+(_agents.length+1)+'">&#128101; Cross-Agent Read</td></tr>';
  _agents.forEach(function(srcAg){
    var cnt = countAccess('cross', srcAg);
    h += '<tr><td class="source-cell"><span class="source-icon">&#128100;</span><span class="source-name">'+escH(srcAg)+'</span>'+badgeHtml(cnt)+'</td>';
    _agents.forEach(function(col){
      if(col===srcAg){
        h += '<td class="data-cell" style="background:#14142a;"></td>';
      } else {
        var a = _acData.agents[col];
        var chk = (a && (a.cross_agent_read||[]).indexOf(srcAg)>=0) ? ' checked' : '';
        h += '<td class="data-cell"><input type="checkbox" class="ac-cb" data-type="cross" data-source="'+escH(srcAg)+'" data-agent="'+escH(col)+'"'+chk+' onchange="updateBadges()"></td>';
      }
    });
    h += '</tr>';
  });
  h += '</tbody></table>';
  wrap.innerHTML = h;
  container.appendChild(wrap);
}

function updateBadges(){
  collectFromDOM();
  var rows = document.querySelectorAll('tr:not(.section-row)');
  rows.forEach(function(row){
    var cb = row.querySelector('.ac-cb');
    if(!cb) return;
    var type = cb.dataset.type;
    var source = cb.dataset.source;
    var cnt;
    if(type==='own'){
      cnt = (_acData.agents[source] && _acData.agents[source].own_memory) ? 1 : 0;
    } else if(type==='shared'){
      cnt = countAccess('shared', source);
    } else {
      cnt = countAccess('cross', source);
    }
    var badgeEl = row.querySelector('.source-badge');
    if(badgeEl) badgeEl.outerHTML = badgeHtml(cnt);
  });
}

function collectFromDOM(){
  if(!_acData) return;
  _agents.forEach(function(ag){
    if(!_acData.agents[ag]) _acData.agents[ag] = {own_memory:false, shared_memory:[], cross_agent_read:[], description:''};
    _acData.agents[ag].shared_memory = [];
    _acData.agents[ag].cross_agent_read = [];
  });
  document.querySelectorAll('.ac-cb').forEach(function(cb){
    var type = cb.dataset.type;
    var src  = cb.dataset.source;
    var ag   = cb.dataset.agent;
    if(!_acData.agents[ag]) return;
    if(type==='own'){
      _acData.agents[ag].own_memory = cb.checked;
    } else if(type==='shared'){
      if(cb.checked) _acData.agents[ag].shared_memory.push(src);
    } else if(type==='cross'){
      if(cb.checked) _acData.agents[ag].cross_agent_read.push(src);
    }
  });
}

async function saveMatrix(){
  if(!_acData) return;
  collectFromDOM();
  var msg = document.getElementById('msg');
  try {
    var r = await fetch('/api/access-control', {
      method:'POST', headers:{'Content-Type':'application/json'},
      body:JSON.stringify(_acData)
    });
    var d = await r.json();
    if(d.success){
      msg.className='msg success';
      msg.textContent='Gespeichert um '+d.saved_at;
      msg.style.display='block';
      document.getElementById('last-mod').textContent='Zuletzt gespeichert: '+d.saved_at;
    } else {
      msg.className='msg error';
      msg.textContent='Fehler: '+(d.error||'unbekannt');
      msg.style.display='block';
    }
  } catch(e){
    msg.className='msg error';
    msg.textContent='Netzwerk-Fehler: '+e.message;
    msg.style.display='block';
  }
  setTimeout(function(){ msg.style.display='none'; }, 5000);
}

function openAddSourceModal(){
  document.getElementById('new-src-label').value = '';
  document.getElementById('new-src-path').value = '';
  document.getElementById('add-src-err').textContent = '';
  document.getElementById('add-source-modal').classList.add('open');
  document.getElementById('new-src-label').focus();
}
function closeAddSourceModal(){
  document.getElementById('add-source-modal').classList.remove('open');
}
async function submitAddSource(){
  var label = document.getElementById('new-src-label').value.trim();
  var path = document.getElementById('new-src-path').value.trim();
  var err = document.getElementById('add-src-err');
  err.textContent = '';
  if(!label){ err.textContent = 'Bitte Anzeigename eingeben.'; return; }
  if(!path){ err.textContent = 'Bitte Pfad eingeben.'; return; }
  // Zuerst aktuelle Matrix-Checkboxen ins Modell uebernehmen, damit unsave
  // Aenderungen beim Reload nicht verloren gehen.
  collectFromDOM();
  try {
    await fetch('/api/access-control', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(_acData)});
  } catch(e){ /* best effort */ }
  try {
    var r = await fetch('/api/access-control/custom-sources', {
      method:'POST', headers:{'Content-Type':'application/json'},
      body:JSON.stringify({label:label, path:path})
    });
    var d = await r.json();
    if(!d.success){ err.textContent = d.error || 'Fehler'; return; }
    closeAddSourceModal();
    await loadMatrix();
    var msg = document.getElementById('msg');
    msg.className='msg success';
    msg.textContent='Quelle hinzugefuegt: '+label;
    msg.style.display='block';
    setTimeout(function(){ msg.style.display='none'; }, 4000);
  } catch(e){
    err.textContent = 'Netzwerk-Fehler: '+e.message;
  }
}
async function removeCustomSource(key, ev){
  if(ev){ ev.stopPropagation(); ev.preventDefault(); }
  if(!confirm('Quelle "'+key+'" wirklich entfernen? Zugriffsrechte werden aus allen Agents entfernt.')) return;
  collectFromDOM();
  try {
    await fetch('/api/access-control', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(_acData)});
  } catch(e){ /* best effort */ }
  try {
    var r = await fetch('/api/access-control/custom-sources/'+encodeURIComponent(key), {method:'DELETE'});
    var d = await r.json();
    var msg = document.getElementById('msg');
    if(d.success){
      await loadMatrix();
      msg.className='msg success';
      msg.textContent='Quelle entfernt.';
    } else {
      msg.className='msg error';
      msg.textContent='Fehler: '+(d.error||'unbekannt');
    }
    msg.style.display='block';
    setTimeout(function(){ msg.style.display='none'; }, 4000);
  } catch(e){
    alert('Netzwerk-Fehler: '+e.message);
  }
}

loadMatrix();
</script>
</body></html>"""
    return html

@app.route('/get_prompt', methods=['GET'])
def get_prompt():
    session_id = request.args.get('session_id', 'default')
    state = get_session(session_id)
    agent = request.args.get('agent', '')
    prompt_file = os.path.join(AGENTS_DIR, agent + '.txt')
    if not agent or not os.path.exists(prompt_file):
        return jsonify({'ok': False, 'prompt': ''})
    # If this is the active agent, return the full in-memory system prompt
    if state['agent'] == agent and state.get('system_prompt'):
        return jsonify({'ok': True, 'prompt': state['system_prompt'], 'has_memory': True})
    # For inactive agents, return only the base prompt from .txt file
    with open(prompt_file) as f:
        raw = f.read()
    for marker in ['\n\n--- GEDAECHTNIS:', '\n--- GEDAECHTNIS:',
                   '\n\n--- DATEI-ERSTELLUNG ---', '\n--- DATEI-ERSTELLUNG ---',
                   '\n\n--- WEITERE FAEHIGKEITEN ---']:
        raw = raw.split(marker)[0]
    return jsonify({'ok': True, 'prompt': raw.strip(), 'has_memory': False})

@app.route('/save_prompt', methods=['POST'])
def save_prompt():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    agent = request.json['agent']
    new_prompt = request.json['prompt']
    prompt_file = os.path.join(AGENTS_DIR, agent + '.txt')
    try:
        # Strip all auto-generated blocks, keep only user's base text
        base_part = new_prompt
        for marker in ['\n\n--- GEDAECHTNIS:', '\n--- GEDAECHTNIS:',
                       '\n\n--- DATEI-ERSTELLUNG ---', '\n--- DATEI-ERSTELLUNG ---',
                       '\n\n--- WEITERE FAEHIGKEITEN ---']:
            base_part = base_part.split(marker)[0]
        base_part = base_part.strip()
        # Always write ONLY the base prompt to the .txt file
        with open(prompt_file, 'w') as f:
            f.write(base_part)
        # If this agent is active, rebuild the in-memory system prompt
        if state['agent'] == agent:
            memory_ctx = build_memory_context(state['speicher'], agent)
            state['system_prompt'] = base_part + memory_ctx
        return jsonify({'ok': True, 'prompt': base_part})
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})

@app.route('/create_agent', methods=['POST'])
def create_agent():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    name = request.json.get('name', '').strip()
    prompt = request.json.get('prompt', 'Du bist ein hilfreicher Assistent.')
    if not name:
        return jsonify({'ok': False, 'error': 'Name erforderlich'})
    prompt_file = os.path.join(AGENTS_DIR, name + '.txt')
    if os.path.exists(prompt_file):
        return jsonify({'ok': False, 'error': 'Agent "' + name + '" existiert bereits'})
    try:
        with open(prompt_file, 'w') as f:
            f.write(prompt)
        # Sub-agents use parent's directory, no new folder needed
        speicher = get_agent_speicher(name)
        os.makedirs(speicher, exist_ok=True)
        return jsonify({'ok': True})
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})

@app.route('/open_in_finder', methods=['POST'])
def open_in_finder():
    import subprocess
    body = request.json or {}
    agent = body.get('agent', '')
    filename = body.get('filename', '')
    direct_path = body.get('path', '')
    memory_dir = os.path.join(BASE, agent, 'memory') if agent else ''
    try:
        # Direkter Pfad (z.B. aus Message-Dashboard) — nur zulassen wenn er
        # innerhalb BASE liegt, damit die Route nicht zum allgemeinen
        # Filesystem-Opener wird.
        if direct_path:
            real = os.path.realpath(direct_path)
            real_base = os.path.realpath(BASE)
            if not real.startswith(real_base + os.sep) and real != real_base:
                return jsonify({'ok': False, 'error': 'Pfad ausserhalb des Datalake'})
            if os.path.isfile(real):
                subprocess.run(['open', '-R', real], check=True)
                return jsonify({'ok': True})
            if os.path.isdir(real):
                subprocess.run(['open', real], check=True)
                return jsonify({'ok': True})
            return jsonify({'ok': False, 'error': 'Pfad nicht gefunden'})
        if filename:
            fpath = os.path.join(memory_dir, filename)
            if os.path.exists(fpath):
                # -R reveals file in Finder (select it)
                subprocess.run(['open', '-R', fpath], check=True)
                return jsonify({'ok': True})
            else:
                # File not found - open folder instead
                subprocess.run(['open', memory_dir], check=True)
                return jsonify({'ok': True, 'note': 'Datei nicht gefunden, Ordner geoeffnet'})
        elif memory_dir and os.path.exists(memory_dir):
            subprocess.run(['open', memory_dir], check=True)
            return jsonify({'ok': True})
        return jsonify({'ok': False, 'error': 'Ordner nicht gefunden: ' + (memory_dir or '(leer)')})
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})

@app.route('/get_history', methods=['GET'])
def get_history():
    session_id = request.args.get('session_id', 'default')
    state = get_session(session_id)
    agent = request.args.get('agent', '')
    if not agent:
        return jsonify({'sessions': []})
    # SUBAGENT_HISTORY_V1: Sub-Agents nutzen Parent-Ordner
    speicher = get_agent_speicher(agent)
    if not os.path.exists(speicher):
        return jsonify({'sessions': []})
    
    # Scan conversation files directly from disk (ignore _index.json for listing)
    # Sub-Agent-Suffix (z.B. '_outbound' fuer 'signicat_outbound')
    parent = get_parent_agent(agent)
    sub_suffix = '_' + agent.split('_', 1)[1] if parent else None
    all_conv = [f for f in os.listdir(speicher)
                if f.startswith('konversation_') and f.endswith('.txt')]
    if sub_suffix:
        # Sub-Agent: nur eigene Konversationen (mit Suffix im Dateinamen)
        conv_files = [f for f in all_conv if f.endswith(sub_suffix + '.txt')]
    else:
        # Parent: Konversationen ohne bekannte Sub-Agent-Suffixe
        known_subs = set()
        for afile in os.listdir(AGENTS_DIR):
            if afile.endswith('.txt') and '_' in afile:
                aname = afile.replace('.txt', '')
                if get_parent_agent(aname) == agent:
                    known_subs.add('_' + aname.split('_', 1)[1] + '.txt')
        conv_files = [f for f in all_conv
                      if not any(f.endswith(s) for s in known_subs)] if known_subs else all_conv
    # Build result with file metadata
    result = []
    for fname in conv_files:
        fpath = os.path.join(speicher, fname)
        # Skip tiny files (empty sessions with only header, <=50 bytes)
        try:
            fsize = os.path.getsize(fpath)
        except Exception:
            continue
        if fsize <= 50:
            continue
        # Auto-title + Test-Artefakt-Check: extrahiere ALLE User-Messages.
        # Wenn jede User-Message ein bekanntes Test-Muster ist UND die Datei
        # klein ist, markieren wir sie als Test-Artefakt und blenden sie aus
        # (der Server-Auto-Save erzeugt sonst UI-Rauschen aus der Test-Suite).
        _HISTORY_TEST_PATTERNS = {
            "Sag nur das Wort: TESTOK", "/find test",
            "Antworte NUR mit: TEST_OK", "Say hello", "TEST", "test",
        }
        all_du = []
        try:
            with open(fpath, encoding='utf-8') as f:
                for line in f:
                    if line.startswith('Du: '):
                        all_du.append(line[4:].strip())
        except Exception:
            pass
        # Skip if no user message (test artifacts / empty files)
        if not all_du:
            continue
        # Skip wenn alle User-Messages Test-Muster sind UND Datei klein
        if fsize <= 3500 and all(m in _HISTORY_TEST_PATTERNS for m in all_du):
            continue
        # Titel: erste ECHTE (nicht-Test) User-Message bevorzugen — sonst
        # erste Du-Zeile. So erscheinen gemischte Dateien mit dem sinnvollen
        # Prompt als Label statt mit "/find test".
        real_msgs = [m for m in all_du if m not in _HISTORY_TEST_PATTERNS]
        title = (real_msgs[0] if real_msgs else all_du[0])[:60]
        # Sort key: file modification time (last save = most recent activity)
        try:
            mtime = os.path.getmtime(fpath)
        except Exception:
            mtime = 0
        # Readable date from modification time
        import datetime as _dt
        date_display = _dt.datetime.fromtimestamp(mtime).strftime('%d.%m.%Y %H:%M')
        result.append({
            'date': date_display,
            'file': fname,
            'title': title,
            'mtime': mtime,
        })
    # Sort by modification time, newest first
    result.sort(key=lambda x: x['mtime'], reverse=True)
    # Remove mtime from response
    for r in result:
        del r['mtime']
    return jsonify({'sessions': result})

@app.route('/load_conversation', methods=['POST'])
def load_conversation():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    agent = request.json.get('agent', '')
    filename = request.json.get('file', '')
    resume = request.json.get('resume', False)
    speicher = get_agent_speicher(agent)  # SUBAGENT_HISTORY_V1
    fpath = os.path.join(speicher, filename)

    if not filename or not os.path.exists(fpath):
        return jsonify({'ok': False, 'error': 'Datei nicht gefunden: ' + filename})

    try:
        with open(fpath) as f:
            raw = f.read()

        # Parse conversation file into messages
        messages = []
        lines = raw.split('\n')
        i = 0
        # Skip header lines (Agent:, Datum:)
        while i < len(lines) and not lines[i].startswith('['):
            i += 1

        current_role = None
        current_content = []

        for line in lines[i:]:
            if line.startswith('Du: '):
                if current_role and current_content:
                    messages.append({'role': current_role, 'content': '\n'.join(current_content).strip()})
                current_role = 'user'
                current_content = [line[4:]]
            elif line.startswith('Assistant: '):
                if current_role and current_content:
                    messages.append({'role': current_role, 'content': '\n'.join(current_content).strip()})
                current_role = 'assistant'
                current_content = [line[11:]]
            elif current_role:
                current_content.append(line)

        if current_role and current_content:
            messages.append({'role': current_role, 'content': '\n'.join(current_content).strip()})

        # Filter empty
        messages = [m for m in messages if m['content'].strip()]

        # Parse KONTEXT_DATEIEN block if present
        restored_ctx = []
        missing_ctx = []
        import re as _ctx_re
        ctx_match = _ctx_re.search(r'\[KONTEXT_DATEIEN:(\[.*?\])\]', raw)
        if ctx_match:
            try:
                ctx_entries = json.loads(ctx_match.group(1))
                for entry in ctx_entries:
                    fname = entry.get('name', '')
                    if not fname:
                        continue
                    # Try to find the file in agent memory
                    ctx_fpath = os.path.join(speicher, 'memory', fname)
                    if os.path.exists(ctx_fpath):
                        try:
                            with open(ctx_fpath, 'r', encoding='utf-8', errors='replace') as cf:
                                ctx_content = cf.read(50000)
                            restored_ctx.append({'name': fname, 'content': ctx_content})
                        except Exception:
                            missing_ctx.append(fname)
                    else:
                        missing_ctx.append(fname)
            except Exception as ctx_err:
                print(f'[LOAD] KONTEXT_DATEIEN parse error: {ctx_err}')

        # Resume mode: set session state so new messages continue in this file
        if resume and state.get('agent'):
            state['dateiname'] = fpath
            state['verlauf'] = messages[:]
            state['kontext_items'] = restored_ctx[:]
            state['session_files'] = [c['name'] for c in restored_ctx]

        # Extract provider/model from conversation (format: [provider/model_id])
        conv_provider = None
        conv_model_id = None
        import re as _re
        for line in raw.split('\n'):
            pm = _re.match(r'\[([^/]+)/([^\]]+)\]$', line.strip())
            if pm and pm.group(1) in VALID_PROVIDERS:
                conv_provider = pm.group(1)
                conv_model_id = pm.group(2)
        # Set on session state if resuming
        if resume and conv_provider and conv_model_id:
            state['provider'] = conv_provider
            state['model_id'] = conv_model_id

        return jsonify({'ok': True, 'messages': messages, 'resumed': bool(resume),
                        'provider': conv_provider, 'model_id': conv_model_id,
                        'restored_ctx': [c['name'] for c in restored_ctx],
                        'missing_ctx': missing_ctx})
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})

@app.route('/create_file', methods=['POST'])
def create_file():
    session_id = request.json.get('session_id', 'default') if request.is_json else 'default'
    state = get_session(session_id)
    try:
        filetype = request.json.get('type', 'docx')
        spec = request.json.get('spec', {})
        if filetype == 'docx':
            fname, fpath = create_docx_from_spec(spec)
        elif filetype == 'xlsx':
            fname, fpath = create_xlsx_from_spec(spec)
        elif filetype == 'pdf':
            fname, fpath = create_pdf_from_spec(spec)
        elif filetype == 'pptx':
            fname, fpath = create_pptx_from_spec(spec)
        else:
            return jsonify({'ok': False, 'error': 'Unbekannter Dateityp: ' + filetype})
        # Also save to agent memory
        if state.get('speicher'):
            memory_dir = os.path.join(state['speicher'], 'memory')
            os.makedirs(memory_dir, exist_ok=True)
            import shutil
            shutil.copy2(fpath, os.path.join(memory_dir, fname))
        return jsonify({'ok': True, 'filename': fname, 'path': fpath})
    except Exception as e:
        import traceback
        return jsonify({'ok': False, 'error': str(e), 'trace': traceback.format_exc()})

@app.route('/send_email_draft', methods=['POST'])
def send_email_draft_route():
    session_id = request.json.get('session_id', 'default') if request.is_json else 'default'
    state = get_session(session_id)
    try:
        spec = request.json
        send_email_draft(spec)
        return jsonify({'ok': True, 'subject': spec.get('subject',''), 'to': spec.get('to','')})
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})


@app.route('/send_email_reply', methods=['POST'])
def send_email_reply_route():
    try:
        spec = request.json
        send_email_reply(spec)
        return jsonify({'ok': True, 'subject': spec.get('subject',''), 'to': spec.get('to','')})
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})


# ─── EMAIL HEADER CACHE (in-memory) ────────────────────────────────────────
_email_header_cache = {}  # key: dir_path -> list of parsed headers
_email_cache_mtime = {}   # key: dir_path -> last build time
_EMAIL_CACHE_TTL = 300    # rebuild cache every 5 minutes

def _parse_filename_timestamp(fname):
    """Extract timestamp from filename like 2025-10-30_16-46-32_... Returns (display, ts)."""
    import re as _re
    import datetime as _dt
    m = _re.match(r'^(\d{4})-(\d{2})-(\d{2})_(\d{2})-(\d{2})-(\d{2})', fname)
    if m:
        try:
            dt = _dt.datetime(int(m.group(1)), int(m.group(2)), int(m.group(3)),
                              int(m.group(4)), int(m.group(5)), int(m.group(6)))
            return dt.strftime('%d.%m.%Y %H:%M'), dt.timestamp()
        except Exception:
            pass
    return '', 0


def _parse_txt_email(fpath, fname):
    """Parse .txt email with German header format (Von:, An:, Betreff:, Datum:)."""
    from email.utils import parsedate_to_datetime as _pdt
    import re as _re
    try:
        with open(fpath, 'r', errors='replace') as f:
            lines = []
            for i, line in enumerate(f):
                if i >= 30:
                    break
                lines.append(line)
        from_raw = subject = date_str = to_raw = cc_raw = message_id = ''
        for line in lines:
            low = line.lower()
            if low.startswith('von:') or low.startswith('from:'):
                from_raw = line.split(':', 1)[1].strip()
            elif low.startswith('betreff:') or low.startswith('subject:'):
                subject = line.split(':', 1)[1].strip()
            elif low.startswith('datum:') or low.startswith('date:'):
                date_str = line.split(':', 1)[1].strip()
            elif low.startswith('an:') or low.startswith('to:'):
                to_raw = line.split(':', 1)[1].strip()
            elif low.startswith('cc:') or low.startswith('kopie:'):
                cc_raw = line.split(':', 1)[1].strip()
            elif low.startswith('message-id:'):
                message_id = line.split(':', 1)[1].strip()
        from_name = ''
        from_email = from_raw
        if '<' in from_raw and '>' in from_raw:
            from_name = from_raw[:from_raw.index('<')].strip().strip('"').replace(',', ' ')
            from_email = from_raw[from_raw.index('<')+1:from_raw.index('>')]
        date_display = ''
        date_ts = 0
        if date_str:
            try:
                dt = _pdt(date_str)
                date_display = dt.strftime('%d.%m.%Y %H:%M')
                date_ts = dt.timestamp()
            except Exception:
                pass
        if not date_ts:
            fn_display, fn_ts = _parse_filename_timestamp(fname)
            if fn_ts:
                date_display = date_display or fn_display
                date_ts = fn_ts
        return {
            'message_id': message_id,
            'from_name': from_name, 'from_email': from_email,
            'subject': subject, 'date': date_display, 'date_ts': date_ts,
            'to': to_raw, 'cc': cc_raw,
        }
    except Exception:
        return None


def _build_email_cache(sdir):
    import email as _eml
    from email.header import decode_header as _dh
    from email.utils import parsedate_to_datetime as _pdt

    def _dec(val):
        if not val:
            return ''
        try:
            parts = _dh(val)
            decoded = []
            for part, charset in parts:
                if isinstance(part, bytes):
                    decoded.append(part.decode(charset or 'utf-8', errors='replace'))
                else:
                    decoded.append(part)
            return ' '.join(decoded)
        except Exception:
            return str(val)

    entries = []
    try:
        for fname in os.scandir(sdir):
            if not (fname.name.endswith('.eml') or fname.name.endswith('.txt')):
                continue
            try:
                parsed = None
                if fname.name.endswith('.eml'):
                    with open(fname.path, 'r', errors='replace') as f:
                        header_lines = []
                        for i, line in enumerate(f):
                            if i >= 40:
                                break
                            header_lines.append(line)
                    msg = _eml.message_from_string(''.join(header_lines))
                    from_raw = _dec(msg.get('From', ''))
                    subject = _dec(msg.get('Subject', ''))
                    date_str = msg.get('Date', '')
                    message_id = msg.get('Message-ID', '').strip()
                    to_raw = _dec(msg.get('To', ''))
                    cc_raw = _dec(msg.get('Cc', ''))
                    from_name = ''
                    from_email = from_raw
                    if '<' in from_raw and '>' in from_raw:
                        from_name = from_raw[:from_raw.index('<')].strip().strip('"').replace(',', ' ')
                        from_email = from_raw[from_raw.index('<')+1:from_raw.index('>')]
                    date_display = ''
                    date_ts = 0
                    try:
                        dt = _pdt(date_str)
                        date_display = dt.strftime('%d.%m.%Y %H:%M')
                        date_ts = dt.timestamp()
                    except Exception:
                        pass
                    if not date_ts:
                        fn_display, fn_ts = _parse_filename_timestamp(fname.name)
                        if fn_ts:
                            date_display = date_display or fn_display
                            date_ts = fn_ts
                    if not date_ts:
                        try:
                            date_ts = fname.stat().st_mtime
                        except Exception:
                            pass
                    parsed = {
                        'message_id': message_id,
                        'from_name': from_name, 'from_email': from_email,
                        'subject': subject, 'date': date_display, 'date_ts': date_ts,
                        'to': to_raw, 'cc': cc_raw,
                    }
                else:
                    # .txt with German header format
                    parsed = _parse_txt_email(fname.path, fname.name)
                    if parsed and not parsed.get('date_ts'):
                        try:
                            parsed['date_ts'] = fname.stat().st_mtime
                        except Exception:
                            pass

                if not parsed:
                    continue
                # Only include entries that actually look like emails
                if not (parsed['from_email'] or parsed['subject']):
                    continue

                entries.append({
                    **parsed,
                    '_filename': fname.name,
                    '_fpath': fname.path,
                    '_s_from': (parsed['from_name'].replace(',', ' ') + ' ' + parsed['from_email'] + ' ' + fname.name).lower(),
                    '_s_subj': parsed['subject'].lower(),
                    '_s_to': (parsed.get('to','') + ' ' + parsed.get('cc','')).lower(),
                })
            except Exception:
                continue
    except Exception:
        pass
    entries.sort(key=lambda e: e.get('date_ts', 0), reverse=True)
    return entries


def _get_email_cache(sdir):
    import time
    now = time.time()
    if sdir in _email_header_cache and (now - _email_cache_mtime.get(sdir, 0)) < _EMAIL_CACHE_TTL:
        return _email_header_cache[sdir]
    entries = _build_email_cache(sdir)
    _email_header_cache[sdir] = entries
    _email_cache_mtime[sdir] = now
    print(f"[EMAIL_CACHE] Built cache for {sdir}: {len(entries)} entries", flush=True)
    return entries


@app.route('/api/email-search')
def email_search_route():
    agent = request.args.get('agent', 'standard')
    q = request.args.get('q', '').strip().lower()
    from_filter = request.args.get('from', '').strip().lower()
    subj_filter = request.args.get('subject', '').strip().lower()
    to_filter = request.args.get('to', '').strip().lower()
    body_filter = request.args.get('body', '').strip().lower()
    has_field_filter = bool(from_filter or subj_filter or to_filter or body_filter)
    if len(q) < 2 and not has_field_filter:
        return jsonify([])

    search_dirs = []
    speicher = get_agent_speicher(agent)
    memory_dir = os.path.join(speicher, 'memory')
    if os.path.exists(memory_dir):
        search_dirs.append(memory_dir)
    inbox_dir = os.path.join(BASE, 'email_inbox')
    if os.path.exists(inbox_dir):
        search_dirs.append(inbox_dir)

    all_entries = []
    for sdir in search_dirs:
        all_entries.extend(_get_email_cache(sdir))
    # Re-sort merged list
    all_entries.sort(key=lambda e: e.get('date_ts', 0), reverse=True)

    def _clean_email(s):
        # Strip mailto: wrappers and angle brackets
        import re as _re
        if not s: return ''
        # Common pattern in .txt files: "user@domain.de<mailto:user@domain.de>"
        m = _re.search(r'([\w\.\-]+@[\w\.\-]+)', s)
        return m.group(1) if m else s

    results = []
    seen_ids = set()
    seen_dedup = set()
    for entry in all_entries:
        if len(results) >= 8:
            break
        mid = entry['message_id']
        # Primary dedup: message_id
        if mid and mid in seen_ids:
            continue
        # Secondary dedup: (clean_from_email, subject, date_ts) — handles iCloud duplicates
        clean_from = _clean_email(entry['from_email']).lower()
        dedup_key = (clean_from, entry['subject'].strip(), entry.get('date_ts', 0))
        if dedup_key in seen_dedup:
            continue
        match = True
        if from_filter:
            if from_filter not in entry['_s_from']:
                match = False
        if subj_filter and match:
            if subj_filter not in entry['_s_subj']:
                match = False
        if to_filter and match:
            if to_filter not in entry['_s_to']:
                match = False
        if not from_filter and not subj_filter and not to_filter and q:
            if q not in (entry['_s_from'] + ' ' + entry['_s_subj']):
                match = False
        if match:
            if mid:
                seen_ids.add(mid)
            seen_dedup.add(dedup_key)
            results.append({
                'message_id': entry['message_id'],
                'from_name': entry['from_name'],
                'from_email': _clean_email(entry['from_email']),
                'subject': entry['subject'], 'date': entry['date'],
                'date_ts': entry['date_ts'],
                'to': entry['to'], 'cc': entry['cc'],
                'file': entry.get('_filename', ''),
                'fpath': entry.get('_fpath', ''),
            })
    return jsonify(results[:8])


@app.route('/api/email-content')
def email_content_route():
    """Load full email content for display in chat."""
    import email as _email_mod
    from email.header import decode_header as _dec_hdr
    agent = request.args.get('agent', 'standard')
    target_mid = request.args.get('message_id', '').strip()
    target_from = request.args.get('from_email', '').strip().lower()
    target_subj = request.args.get('subject', '').strip().lower()

    if not target_mid and not target_from:
        return jsonify({'ok': False, 'error': 'message_id or from_email required'})

    def _dec(val):
        if not val:
            return ''
        parts = _dec_hdr(val)
        decoded = []
        for part, charset in parts:
            if isinstance(part, bytes):
                decoded.append(part.decode(charset or 'utf-8', errors='replace'))
            else:
                decoded.append(part)
        return ' '.join(decoded)

    def _get_body(msg):
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == 'text/plain':
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or 'utf-8'
                        return payload.decode(charset, errors='replace')
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == 'text/html':
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or 'utf-8'
                        import re as _re
                        html = payload.decode(charset, errors='replace')
                        return _re.sub(r'<[^>]+>', '', html)[:3000]
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                charset = msg.get_content_charset() or 'utf-8'
                return payload.decode(charset, errors='replace')
        return ''

    search_dirs = []
    speicher = get_agent_speicher(agent)
    memory_dir = os.path.join(speicher, 'memory')
    if os.path.exists(memory_dir):
        search_dirs.append(memory_dir)
    inbox_dir = os.path.join(BASE, 'email_inbox')
    if os.path.exists(inbox_dir):
        search_dirs.append(inbox_dir)

    def _read_txt_email(fpath, fname):
        try:
            with open(fpath, 'r', errors='replace') as f:
                full = f.read()
            # Split headers from body (first blank line)
            header_end = full.find('\n\n')
            if header_end == -1:
                header_end = len(full)
            header_block = full[:header_end]
            body = full[header_end+2:] if header_end < len(full) else ''
            headers = {}
            for line in header_block.split('\n'):
                for key_de, key_en in [('von:','from'),('an:','to'),('betreff:','subject'),('datum:','date'),('cc:','cc'),('kopie:','cc'),('message-id:','message_id')]:
                    if line.lower().startswith(key_de) or line.lower().startswith(key_en + ':'):
                        headers[key_en] = line.split(':', 1)[1].strip()
                        break
            return headers, body
        except Exception:
            return None, None

    for sdir in search_dirs:
        try:
            for fname in os.listdir(sdir):
                if not (fname.endswith('.eml') or fname.endswith('.txt')):
                    continue
                fpath = os.path.join(sdir, fname)
                try:
                    if fname.endswith('.txt'):
                        headers, body = _read_txt_email(fpath, fname)
                        if not headers:
                            continue
                        mid = headers.get('message_id', '').strip()
                        from_raw_txt = headers.get('from', '')
                        subj_txt = headers.get('subject', '')
                        # Match
                        if target_mid and mid == target_mid:
                            pass
                        elif target_from:
                            if target_from not in from_raw_txt.lower():
                                continue
                            if target_subj and target_subj not in subj_txt.lower():
                                continue
                        else:
                            continue
                        # Build result
                        from_name = ''
                        from_email = from_raw_txt
                        if '<' in from_raw_txt and '>' in from_raw_txt:
                            from_name = from_raw_txt[:from_raw_txt.index('<')].strip().strip('"')
                            from_email = from_raw_txt[from_raw_txt.index('<')+1:from_raw_txt.index('>')]
                        date_display = headers.get('date', '')
                        try:
                            from email.utils import parsedate_to_datetime
                            dt = parsedate_to_datetime(date_display)
                            date_display = dt.strftime('%d.%m.%Y %H:%M')
                        except Exception:
                            pass
                        if len(body) > 5000:
                            body = body[:5000] + '\n\n[... gekuerzt, ' + str(len(body)) + ' Zeichen gesamt]'
                        return jsonify({
                            'ok': True,
                            'from_name': from_name, 'from_email': from_email,
                            'to': headers.get('to', ''), 'cc': headers.get('cc', ''),
                            'subject': subj_txt, 'date': date_display,
                            'message_id': mid, 'body': body, 'file': fname,
                        })
                    with open(fpath, 'r', errors='replace') as f:
                        msg = _email_mod.message_from_file(f)
                    mid = msg.get('Message-ID', '').strip()
                    if target_mid and mid == target_mid:
                        pass
                    elif target_from:
                        from_raw = _dec(msg.get('From', '')).lower()
                        subj_raw = _dec(msg.get('Subject', '')).lower()
                        if target_from not in from_raw:
                            continue
                        if target_subj and target_subj not in subj_raw:
                            continue
                    else:
                        continue

                    from_raw = _dec(msg.get('From', ''))
                    from_name = ''
                    from_email = from_raw
                    if '<' in from_raw and '>' in from_raw:
                        from_name = from_raw[:from_raw.index('<')].strip().strip('"')
                        from_email = from_raw[from_raw.index('<')+1:from_raw.index('>')]

                    date_display = ''
                    try:
                        from email.utils import parsedate_to_datetime
                        dt = parsedate_to_datetime(msg.get('Date', ''))
                        date_display = dt.strftime('%d.%m.%Y %H:%M')
                    except Exception:
                        date_display = (msg.get('Date', '') or '')[:30]

                    body = _get_body(msg)
                    if len(body) > 5000:
                        body = body[:5000] + '\n\n[... gekuerzt, ' + str(len(body)) + ' Zeichen gesamt]'

                    return jsonify({
                        'ok': True,
                        'from_name': from_name,
                        'from_email': from_email,
                        'to': _dec(msg.get('To', '')),
                        'cc': _dec(msg.get('Cc', '')),
                        'subject': _dec(msg.get('Subject', '')),
                        'date': date_display,
                        'message_id': mid,
                        'body': body,
                        'file': fname,
                    })
                except Exception:
                    continue
        except Exception:
            continue

    return jsonify({'ok': False, 'error': 'E-Mail nicht gefunden'})


@app.route('/send_whatsapp_draft', methods=['POST'])
def send_whatsapp_draft_route():
    session_id = request.json.get('session_id', 'default') if request.is_json else 'default'
    state = get_session(session_id)
    try:
        spec = request.json
        agent_name = state.get('agent', 'standard')
        to_name, phone = send_whatsapp_draft(spec, agent_name)
        return jsonify({'ok': True, 'to': to_name, 'phone': phone, 'clipboard_fallback': phone is None})
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})


@app.route('/open_slack_draft', methods=['POST'])
def open_slack_draft_route():
    try:
        spec = request.json
        target, clipboard_only = send_slack_draft(spec)
        return jsonify({'ok': True, 'target': target, 'clipboard_only': clipboard_only})
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})


@app.route('/download_file', methods=['GET'])
def download_file():
    session_id = request.args.get('session_id', 'default')
    state = get_session(session_id)
    from flask import send_file
    fpath = request.args.get('path', '')
    if not fpath or not os.path.exists(fpath):
        return "Datei nicht gefunden", 404
    return send_file(fpath, as_attachment=True)

@app.route('/select_model', methods=['POST'])
def select_model():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    new_provider = request.json['provider']
    if new_provider not in VALID_PROVIDERS:
        return jsonify({'ok': False, 'error': 'Ungueltiger Provider: ' + str(new_provider)})
    state['provider'] = new_provider
    state['model_id'] = request.json['model_id']
    return jsonify({'ok': True})


AGENT_PREFS_FILE = os.path.join(BASE, "config", "agent_model_preferences.json")

def _load_agent_prefs():
    if os.path.exists(AGENT_PREFS_FILE):
        try:
            with open(AGENT_PREFS_FILE) as f:
                return json.load(f)
        except Exception:
            pass
    return {}

def _save_agent_prefs(prefs):
    os.makedirs(os.path.dirname(AGENT_PREFS_FILE), exist_ok=True)
    with open(AGENT_PREFS_FILE, 'w') as f:
        json.dump(prefs, f, indent=2)

@app.route('/api/agent-model-preference', methods=['GET', 'POST'])
def agent_model_preference():
    if request.method == 'POST':
        data = request.json
        agent = data.get('agent', '')
        if not agent:
            return jsonify({'ok': False, 'error': 'No agent specified'})
        prefs = _load_agent_prefs()
        prefs[agent] = {'provider': data.get('provider', ''), 'model': data.get('model', '')}
        _save_agent_prefs(prefs)
        return jsonify({'ok': True})
    else:
        agent = request.args.get('agent', '')
        prefs = _load_agent_prefs()
        pref = prefs.get(agent, {})
        return jsonify({'ok': True, 'provider': pref.get('provider', ''), 'model': pref.get('model', '')})


SEARCH_STOPWORDS = {
    # Deutsch — Artikel, Präpositionen, Pronomen
    'die', 'der', 'das', 'ein', 'eine', 'und', 'oder', 'von', 'für', 'fuer',
    'mit', 'den', 'dem', 'des', 'auf', 'aus', 'bei', 'bis',
    'ist', 'mir', 'ich', 'er', 'sie', 'wir', 'ihr',
    'über', 'ueber', 'wie', 'was', 'wer', 'zum', 'zur',
    'meine', 'meinen', 'meinem', 'dein', 'sein', 'ihre',
    'bitte', 'nochmal', 'noch', 'mal', 'auch', 'nur', 'sehr',
    'dann', 'jetzt', 'hier', 'dort', 'aber', 'wenn', 'weil',
    'kann', 'dass', 'nach', 'suche', 'suchen',
    # Englisch
    'the', 'and', 'from', 'for',
    'that', 'this', 'also', 'please', 'just', 'load',
    # Portugiesisch
    'uma', 'uns', 'umas', 'para', 'com', 'por', 'que',
    'mas', 'não', 'nao', 'ao', 'na', 'no', 'nas', 'nos',
    'pelo', 'pela', 'este', 'essa', 'esse', 'isso',
    'algum', 'alguma', 'alguns', 'algumas',
    'dos', 'tem',
    # Kontext-spezifisch (Memory-Suche Trigger)
    'memory', 'folder', 'ordner', 'prompt', 'kontext',
    'schau', 'guck', 'lies', 'read', 'find', 'such', 'hol', 'lad',
    'e-mail', 'email', 'mail',
}

# Note: _SEARCH_ACTIONS, _SEARCH_PHRASES, _SEARCH_OBJECTS removed.
# Search is now triggered only by explicit /find and /find_global commands.
BINARY_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.mp4', '.mov', '.mp3', '.wav'}


def parse_search_keywords(query):
    """Split query into normalized keywords, removing stopwords and short tokens."""
    import re
    query_lower = query.lower().strip()
    raw_tokens = query_lower.split()
    keywords = []
    for token in raw_tokens:
        token = token.strip('.,;:!?()[]{}\"\'/') 
        if len(token) < 3:
            continue
        if token in SEARCH_STOPWORDS:
            continue
        keywords.append(token)
        # Also add version without hyphens for fuzzy matching (wave-x -> wavex)
        no_hyphen = token.replace('-', '')
        if no_hyphen != token and len(no_hyphen) >= 3:
            keywords.append(no_hyphen)
    return list(dict.fromkeys(keywords))  # dedupe, preserve order


def _is_binary(fname):
    return any(fname.lower().endswith(ext) for ext in BINARY_EXTS)


def _read_file_text(fpath, fname, max_chars=20000):
    """Read file content as text, skipping binaries. Returns empty string on failure."""
    if _is_binary(fname):
        return ''
    try:
        with open(fpath, 'rb') as f:
            raw = f.read()
        return extract_file_content(raw, fname)[:max_chars]
    except Exception:
        return ''


def scored_memory_search(memory_dir, keywords, max_results=3):
    """Search memory files by keyword scoring. Returns list of {name, content, score}."""
    if not keywords or not os.path.exists(memory_dir):
        return []

    all_files = os.listdir(memory_dir)
    scored = []

    # --- Pass 1: score by filename ---
    for fname in all_files:
        if _is_binary(fname):
            continue
        fname_lower = fname.lower()
        fname_no_hyphen = fname_lower.replace('-', '')
        score = 0
        matched_kws = 0
        for kw in keywords:
            if kw in fname_lower or kw in fname_no_hyphen:
                score += 3
                matched_kws += 1
        if matched_kws == len(keywords) and len(keywords) > 1:
            score += 5
        if score > 0:
            scored.append((fname, score))

    # --- Pass 2: if fewer than 3 filename hits, search file contents ---
    if len([s for s in scored if s[1] > 0]) < max_results:
        files_by_mtime = sorted(
            [f for f in all_files if not _is_binary(f)],
            key=lambda f: os.path.getmtime(os.path.join(memory_dir, f)),
            reverse=True,
        )[:1000]

        already_scored = {s[0] for s in scored}
        for fname in files_by_mtime:
            fpath = os.path.join(memory_dir, fname)
            text = _read_file_text(fpath, fname)
            if not text:
                continue
            text_lower = text.lower()

            content_score = 0
            matched_content_kws = 0
            for kw in keywords:
                count = text_lower.count(kw)
                # Also count without hyphens
                kw_no_h = kw.replace('-', '')
                if kw_no_h != kw:
                    count += text_lower.replace('-', '').count(kw_no_h)
                if count >= 3:
                    content_score += 2
                    matched_content_kws += 1
                elif count >= 1:
                    content_score += 1
                    matched_content_kws += 1
            if matched_content_kws == len(keywords) and len(keywords) > 1:
                content_score += 3

            if fname in already_scored:
                scored = [(f, s + content_score) if f == fname else (f, s) for f, s in scored]
            elif content_score > 0:
                scored.append((fname, content_score))

            if len([s for s in scored if s[1] > 0]) >= max_results:
                break

    # Sort by score descending, take top results
    scored.sort(key=lambda x: x[1], reverse=True)
    results = []
    for fname, score in scored[:max_results]:
        if score <= 0:
            break
        fpath = os.path.join(memory_dir, fname)
        content = _read_file_text(fpath, fname)
        if content:
            results.append({'name': fname, 'content': content, 'score': score})
    return results


def auto_search_memory(msg, speicher):
    """Auto-search memory. Only triggers on legacy 'memory folder/ordner' keyword.
    All other searches now use explicit /find and /find_global commands."""
    try:
        import re
        msg_lower = msg.lower()

        # --- Legacy trigger: "memory folder" / "memory ordner" ---
        if 'memory folder' in msg_lower or 'memory ordner' in msg_lower:
            memory_dir = os.path.join(speicher, 'memory')
            if not os.path.exists(memory_dir):
                return []
            after = re.split(r'memory folder|memory ordner', msg_lower, maxsplit=1)[-1]
            after = re.split(r'\bnach\b|\bfor\b|\bafter\b', after)[-1]
            keywords = parse_search_keywords(after)
            if not keywords:
                return []
            return scored_memory_search(memory_dir, keywords, max_results=3)

        return []

    except Exception as e:
        print("auto_search_memory error: " + str(e))
        return []

def process_single_message(msg, kontext_override=None, state=None, **kwargs):
    """Process a single chat message through the LLM. Returns result dict.
    Does not use Flask request/jsonify — can be called from queue worker thread."""
    if state is None:
        state = get_session()
    kontext_items = kontext_override if kontext_override is not None else state['kontext_items']

    # Auto-search memory (via search_engine.py)
    auto_loaded_names = kwargs.get('auto_loaded_override', [])
    auto_search_info = kwargs.get('auto_search_info_override', '')
    if auto_loaded_names or auto_search_info:
        pass  # Overrides provided — skip auto-search
    elif auto_search and state.get('speicher'):
        try:
            # Dedup against files already pinned in working memory so we don't
            # double-inject the same file into the prompt.
            wm_names = set()
            try:
                agent_name_for_wm = state.get('agent', '')
                wm_manifest = working_memory_list(agent_name_for_wm)
                for f in wm_manifest.get('files', []):
                    if f.get('filename'):
                        wm_names.add(f['filename'])
            except Exception:
                pass

            search_results, search_feedback = auto_search(msg, state['speicher'])
            for item in search_results:
                fname = item['name']
                if fname in wm_names:
                    continue  # already pinned in working memory
                if not any(k['name'] == fname for k in kontext_items):
                    kontext_items.append(item)
                    if fname not in state['session_files']:
                        state['session_files'].append(fname)
                    auto_loaded_names.append(fname)
            if search_feedback:
                auto_search_info = format_search_feedback(search_feedback, len(auto_loaded_names))
        except Exception as e:
            print(f"auto_search error: {e}")
            import traceback; traceback.print_exc()

    # Fallback: deep_memory_search only for explicit /find command (passed through from frontend)
    if not auto_loaded_names and state.get('speicher') and not kontext_override:
        msg_lower = msg.lower()
        if msg_lower.startswith('/find '):
            try:
                memory_dir = os.path.join(state['speicher'], 'memory')
                ds_results = deep_memory_search(memory_dir, msg, max_results=15)
                hits = ds_results.get('results', [])
                if hits:
                    if len(hits) <= 5:
                        # Auto-load up to 5
                        for r in hits:
                            fname = r['filename']
                            if not any(k['name'] == fname for k in kontext_items):
                                fpath = os.path.join(memory_dir, fname)
                                try:
                                    with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                                        content = f.read(50000)
                                    kontext_items.append({'name': fname, 'content': content})
                                    if fname not in state['session_files']:
                                        state['session_files'].append(fname)
                                    auto_loaded_names.append(fname)
                                except Exception:
                                    pass
                        auto_search_info = f'{len(auto_loaded_names)} Datei(en) aus Deep-Search geladen'
                    elif len(hits) <= 15:
                        # Inject hint into message for the AI
                        filelist = ', '.join(r['filename'][:40] for r in hits[:10])
                        auto_search_info = f'{len(hits)} relevante Dateien gefunden. Dateien: {filelist}'
                    else:
                        auto_search_info = f'Zu viele Treffer ({ds_results["candidates_after_filter"]}). Bitte praezisiere deinen Suchbegriff.'
            except Exception as e:
                print(f"deep_search fallback error: {e}")

    # Check for sub-agent delegation (requires user confirmation)
    # Skip delegation check when email reply context is active
    _skip_deleg = kwargs.get('skip_delegation') or msg.startswith('[E-MAIL KONTEXT:')
    if state['agent'] and not _skip_deleg:
        deleg_info = detect_delegation(msg, state['agent'])
        if deleg_info:
            import uuid as _deleg_uuid
            import time as _deleg_time
            # Clean expired pending delegations (>5 min)
            expired = [k for k, v in _pending_delegations.items()
                       if _deleg_time.time() - v.get('timestamp', 0) > 300]
            for k in expired:
                _pending_delegations.pop(k, None)
            # Also expire any previous pending for this session
            old_pending = [k for k, v in _pending_delegations.items()
                          if v.get('session_id') == kwargs.get('_session_id', '')]
            for k in old_pending:
                _pending_delegations.pop(k, None)
            conf_id = str(_deleg_uuid.uuid4())[:12]
            _pending_delegations[conf_id] = {
                'session_id': kwargs.get('_session_id', ''),
                'sub_agent': deleg_info['full_name'],
                'msg': msg,
                'kontext': kontext_items[:],
                'auto_loaded': auto_loaded_names,
                'auto_search_info': auto_search_info,
                'timestamp': _deleg_time.time(),
            }
            return {
                'type': 'subagent_confirmation_required',
                'suggested_subagent': deleg_info['full_name'],
                'subagent_display': deleg_info['display_name'],
                'matched_keywords': deleg_info['matched_keywords'],
                'score': deleg_info['score'],
                'confirmation_id': conf_id,
                'original_message': msg[:100],
            }

    # Build context string
    text_ctx = ''
    image_items = []
    if kontext_items:
        text_ctx = '\n\n--- KONTEXT ---\n'
        for item in kontext_items:
            if item.get('image_b64'):
                image_items.append(item)
                text_ctx += '\n[Bild: ' + item['name'] + ']\n'
            else:
                text_ctx += '\n[' + item['name'] + ']:\n' + item['content'][:10000] + '\n'
        text_ctx += '\n--- ENDE KONTEXT ---\n'

    full_text = msg + text_ctx if text_ctx else msg

    # CALENDAR_INTEGRATION_V1: Automatisch Kalender-Daten injizieren wenn Intent erkannt
    if _has_calendar_intent(msg):
        try:
            _cal_events, _cal_cals, _cal_err = get_calendar_events(days_back=1, days_ahead=7)
            if _cal_events:
                _cal_ctx = format_calendar_context(_cal_events)
                full_text = full_text + '\n\n' + _cal_ctx
                print(f"[CALENDAR] {len(_cal_events)} Events injiziert fuer Intent in: {msg[:50]}", flush=True)
            elif _cal_err:
                print(f"[CALENDAR] Fehler: {_cal_err}", flush=True)
        except Exception as _cal_ex:
            print(f"[CALENDAR] Exception: {_cal_ex}", flush=True)

    provider_key = state.get('provider', 'anthropic')
    if image_items and provider_key == 'anthropic':
        user_content = []
        for img in image_items:
            safe_b64, safe_mime = downscale_image_b64_if_needed(img.get('image_b64'), img.get('image_type'))
            if not safe_b64:
                print(f"[IMG] Oversize-Bild '{img.get('name','?')}' wird nicht gesendet", flush=True)
                continue
            if safe_b64 is not img.get('image_b64'):
                img['image_b64'] = safe_b64
                img['image_type'] = safe_mime
            user_content.append({
                'type': 'image',
                'source': {'type': 'base64', 'media_type': safe_mime, 'data': safe_b64}
            })
        user_content.append({'type': 'text', 'text': full_text})
    else:
        user_content = full_text

    state['verlauf'].append({'role': 'user', 'content': user_content})
    # Write-Through: User-Nachricht + Pending-Marker sofort auf Disk schreiben
    state['verlauf'].append({'role': 'assistant', 'content': PENDING_MARKER})
    for _sid_imm, _st_imm in sessions.items():
        if _st_imm is state:
            auto_save_session(_sid_imm)
            break
    state['verlauf'].pop()  # Pending-Marker aus In-Memory-Verlauf entfernen
    try:
        config = load_models()
        provider_key = state.get('provider', 'anthropic')
        model_id = state.get('model_id', 'claude-sonnet-4-6')
        provider_cfg = config['providers'].get(provider_key, {})
        api_key = provider_cfg.get('api_key', '')
        model_name = next((m['name'] for m in provider_cfg.get('models', []) if m['id'] == model_id), model_id)
        adapter = ADAPTERS.get(provider_key)
        if not adapter:
            raise ValueError('Unbekannter Anbieter: ' + provider_key)
        text = adapter(api_key, model_id, state['system_prompt'], state['verlauf'])

        # Parse MEMORY_SEARCH in agent response and re-query if found
        ms_match = re.search(r'MEMORY_SEARCH:\s*(\{[^}]+\})', text)
        if ms_match and state.get('speicher'):
            try:
                ms_params = json.loads(ms_match.group(1))
                ms_results = deep_memory_search(
                    os.path.join(state['speicher'], 'memory'),
                    ms_params.get('query', ''),
                    date_from=ms_params.get('date_from'),
                    date_to=ms_params.get('date_to'),
                    direction=ms_params.get('direction'),
                    contact=ms_params.get('contact'),
                )
                if ms_results.get('results'):
                    # Inject results and re-query
                    search_ctx = '\n\n--- MEMORY_SEARCH ERGEBNIS ---\n'
                    for r in ms_results['results'][:5]:
                        search_ctx += f"\n[{r['filename']}]\nVon: {r.get('von','')}\nAn: {r.get('an','')}\nBetreff: {r.get('betreff','')}\nDatum: {r.get('datum','')}\n{r.get('preview','')[:500]}\n"
                    search_ctx += '\n--- ENDE MEMORY_SEARCH ---\n'
                    # Replace MEMORY_SEARCH block with results and re-send
                    state['verlauf'].append({'role': 'assistant', 'content': text})
                    state['verlauf'].append({'role': 'user', 'content': search_ctx + '\nBitte beantworte die urspruengliche Frage mit diesen Suchergebnissen.'})
                    text = adapter(api_key, model_id, state['system_prompt'], state['verlauf'])
                    # Remove the injected messages from verlauf (we'll add the final response below)
                    state['verlauf'].pop()  # remove injected user msg
                    state['verlauf'].pop()  # remove MEMORY_SEARCH response
            except Exception as mse:
                print(f"MEMORY_SEARCH error: {mse}")

        # Parse WORKING_MEMORY commands in agent response
        agent_name = state.get('agent', '')
        wm_add_pattern = re.compile(r'WORKING_MEMORY_ADD:\s*(\{.*?\})', re.DOTALL)
        for wm_m in wm_add_pattern.finditer(text):
            try:
                wm_spec = json.loads(wm_m.group(1))
                wm_manifest = working_memory_add(
                    agent_name, wm_spec['filename'], wm_spec['content'],
                    priority=wm_spec.get('priority', 5),
                    description=wm_spec.get('description', ''),
                )
                marker = f'\n[Working Memory: {wm_spec["filename"]} gespeichert]\n'
                text = text.replace(wm_m.group(0), marker)
            except Exception as wme:
                print(f"WORKING_MEMORY_ADD error: {wme}")

        wm_rm_pattern = re.compile(r'WORKING_MEMORY_REMOVE:\s*(\{.*?\})', re.DOTALL)
        for wm_m in wm_rm_pattern.finditer(text):
            try:
                wm_spec = json.loads(wm_m.group(1))
                working_memory_remove(agent_name, wm_spec['filename'])
                marker = f'\n[Working Memory: {wm_spec["filename"]} entfernt]\n'
                text = text.replace(wm_m.group(0), marker)
            except Exception as wme:
                print(f"WORKING_MEMORY_REMOVE error: {wme}")

        wm_list_pattern = re.compile(r'WORKING_MEMORY_LIST:\s*\{[^}]*\}')
        wm_list_match = wm_list_pattern.search(text)
        if wm_list_match:
            try:
                wm_info = working_memory_list(agent_name)
                wm_files = wm_info.get('files', [])
                if wm_files:
                    listing = '\n[Working Memory Inhalt:'
                    for wf in wm_files:
                        listing += f'\n  - {wf["filename"]} (Prio: {wf.get("priority", 5)}, {wf.get("added", "?")}) — {wf.get("description", "")}'
                    listing += f'\n  Token-Limit: {wm_info.get("max_tokens", 8000)}]\n'
                else:
                    listing = '\n[Working Memory ist leer]\n'
                text = text.replace(wm_list_match.group(0), listing)
            except Exception as wme:
                print(f"WORKING_MEMORY_LIST error: {wme}")

        created_files = []
        created_emails = []

        def extract_blocks(src, tag):
            results = []
            prefix = '[' + tag + ':'
            i = 0
            while i < len(src):
                idx = src.find(prefix, i)
                if idx == -1:
                    break
                after_tag = src[idx + len(prefix):]
                colon = after_tag.find(':')
                if colon == -1:
                    i = idx + 1
                    continue
                btype = after_tag[:colon]
                jstart = idx + len(prefix) + colon + 1
                depth = 0
                j = jstart
                jend = -1
                in_string = False
                escape_next = False
                while j < len(src):
                    c = src[j]
                    if escape_next:
                        escape_next = False
                    elif c == '\\':
                        escape_next = True
                    elif c == '"' and not escape_next:
                        in_string = not in_string
                    elif not in_string:
                        if c == '{':
                            depth += 1
                        elif c == '}':
                            depth -= 1
                            if depth == 0:
                                jend = j + 1
                                break
                    j += 1
                if jend != -1 and jend < len(src) and src[jend] == ']':
                    json_str = src[jstart:jend]
                    results.append((src[idx:jend+1], btype, json_str))
                    i = jend + 1
                else:
                    i = idx + 1
            return results

        # Parse CREATE_EMAIL (but not CREATE_EMAIL_REPLY)
        email_prefix = '[CREATE_EMAIL:'
        ei = 0
        while ei < len(text):
            eidx = text.find(email_prefix, ei)
            if eidx == -1:
                break
            # Skip if this is actually CREATE_EMAIL_REPLY
            if text[eidx:eidx+len('[CREATE_EMAIL_REPLY:')] == '[CREATE_EMAIL_REPLY:':
                ei = eidx + 1
                continue
            jstart = eidx + len(email_prefix)
            depth = 0
            j = jstart
            jend = -1
            in_str = False
            esc = False
            while j < len(text):
                c = text[j]
                if esc:
                    esc = False
                elif c == '\\':
                    esc = True
                elif c == '"' and not esc:
                    in_str = not in_str
                elif not in_str:
                    if c == '{':
                        depth += 1
                    elif c == '}':
                        depth -= 1
                        if depth == 0:
                            jend = j + 1
                            break
                j += 1
            if jend != -1 and jend < len(text) and text[jend] == ']':
                full_block = text[eidx:jend+1]
                json_str = text[jstart:jend]
                spec = {}
                try:
                    spec = json.loads(json_str)
                    send_email_draft(spec)
                    created_emails.append({'ok': True, 'subject': spec.get('subject',''), 'to': spec.get('to',''), 'body': spec.get('body','')})
                    marker = f'\n[E-Mail-Draft an {spec.get("to","")} erstellt — Aktion ausgefuehrt]\n'
                    text = text[:eidx] + marker + text[jend+1:]
                    ei = eidx + len(marker)
                except Exception as ee:
                    created_emails.append({'ok': False, 'subject': spec.get('subject',''), 'to': '', 'error': str(ee), 'body': spec.get('body','')})
                    marker = '\n[E-Mail-Erstellung fehlgeschlagen]\n'
                    text = text[:eidx] + marker + text[jend+1:]
                    ei = eidx + len(marker)
            else:
                ei = eidx + 1

        # Parse CREATE_EMAIL_REPLY
        er_prefix = '[CREATE_EMAIL_REPLY:'
        eri = 0
        while eri < len(text):
            eridx = text.find(er_prefix, eri)
            if eridx == -1:
                break
            erjstart = eridx + len(er_prefix)
            depth = 0
            erj = erjstart
            erjend = -1
            erin_str = False
            eresc = False
            while erj < len(text):
                erc = text[erj]
                if eresc:
                    eresc = False
                elif erc == '\\':
                    eresc = True
                elif erc == '"' and not eresc:
                    erin_str = not erin_str
                elif not erin_str:
                    if erc == '{':
                        depth += 1
                    elif erc == '}':
                        depth -= 1
                        if depth == 0:
                            erjend = erj + 1
                            break
                erj += 1
            if erjend != -1 and erjend < len(text) and text[erjend] == ']':
                full_block = text[eridx:erjend+1]
                json_str = text[erjstart:erjend]
                spec = {}
                try:
                    spec = json.loads(json_str)
                    send_email_reply(spec)
                    created_emails.append({'ok': True, 'subject': spec.get('subject',''), 'to': spec.get('to',''), 'body': spec.get('body',''), 'reply': True})
                    marker = f'\n[E-Mail-Reply an {spec.get("to","")} erstellt — Aktion ausgefuehrt]\n'
                    text = text[:eridx] + marker + text[erjend+1:]
                    eri = eridx + len(marker)
                except Exception as ere:
                    created_emails.append({'ok': False, 'subject': spec.get('subject',''), 'to': '', 'error': str(ere), 'body': spec.get('body',''), 'reply': True})
                    marker = '\n[E-Mail-Reply fehlgeschlagen]\n'
                    text = text[:eridx] + marker + text[erjend+1:]
                    eri = eridx + len(marker)
            else:
                eri = eridx + 1

        # Parse CREATE_WHATSAPP
        created_whatsapps = []
        wa_prefix = '[CREATE_WHATSAPP:'
        wi = 0
        while wi < len(text):
            widx = text.find(wa_prefix, wi)
            if widx == -1:
                break
            wjstart = widx + len(wa_prefix)
            depth = 0
            wj = wjstart
            wjend = -1
            win_str = False
            wesc = False
            while wj < len(text):
                wc = text[wj]
                if wesc:
                    wesc = False
                elif wc == '\\':
                    wesc = True
                elif wc == '"' and not wesc:
                    win_str = not win_str
                elif not win_str:
                    if wc == '{':
                        depth += 1
                    elif wc == '}':
                        depth -= 1
                        if depth == 0:
                            wjend = wj + 1
                            break
                wj += 1
            if wjend != -1 and wjend < len(text) and text[wjend] == ']':
                full_block = text[widx:wjend+1]
                json_str = text[wjstart:wjend]
                wspec = {}
                try:
                    wspec = json.loads(json_str)
                    agent_name = state.get('agent', 'standard')
                    wa_to, wa_phone = send_whatsapp_draft(wspec, agent_name)
                    created_whatsapps.append({'ok': True, 'to': wa_to, 'phone': wa_phone, 'clipboard_fallback': wa_phone is None})
                    marker = f'\n[WhatsApp an {wa_to} vorbereitet — Aktion ausgefuehrt]\n'
                    text = text[:widx] + marker + text[wjend+1:]
                    wi = widx + len(marker)
                except Exception as we:
                    created_whatsapps.append({'ok': False, 'to': wspec.get('to',''), 'error': str(we)})
                    marker = '\n[WhatsApp-Erstellung fehlgeschlagen]\n'
                    text = text[:widx] + marker + text[wjend+1:]
                    wi = widx + len(marker)
            else:
                wi = widx + 1


        # Parse CREATE_SLACK
        created_slacks = []
        sl_prefix = '[CREATE_SLACK:'
        si = 0
        while si < len(text):
            sidx = text.find(sl_prefix, si)
            if sidx == -1:
                break
            sjstart = sidx + len(sl_prefix)
            depth = 0
            sj = sjstart
            sjend = -1
            sin_str = False
            sesc = False
            while sj < len(text):
                sc = text[sj]
                if sesc:
                    sesc = False
                elif sc == '\\':
                    sesc = True
                elif sc == '"' and not sesc:
                    sin_str = not sin_str
                elif not sin_str:
                    if sc == '{':
                        depth += 1
                    elif sc == '}':
                        depth -= 1
                        if depth == 0:
                            sjend = sj + 1
                            break
                sj += 1
            if sjend != -1 and sjend < len(text) and text[sjend] == ']':
                json_str = text[sjstart:sjend]
                sspec = {}
                try:
                    sspec = json.loads(json_str)
                    sl_target, sl_clipboard = send_slack_draft(sspec)
                    created_slacks.append({'ok': True, 'target': sl_target, 'clipboard_only': sl_clipboard})
                    marker = f'\n[Slack-Nachricht an {sl_target} vorbereitet — Aktion ausgefuehrt]\n'
                    text = text[:sidx] + marker + text[sjend+1:]
                    si = sidx + len(marker)
                except Exception as se:
                    created_slacks.append({'ok': False, 'target': sspec.get('channel', sspec.get('to', '')), 'error': str(se)})
                    marker = '\n[Slack-Erstellung fehlgeschlagen]\n'
                    text = text[:sidx] + marker + text[sjend+1:]
                    si = sidx + len(marker)
            else:
                si = sidx + 1


        # Parse CREATE_FILE
        for full_block, ftype, json_str in extract_blocks(text, 'CREATE_FILE'):
            try:
                spec = sanitize_llm_json(json_str)
                if not isinstance(spec, dict):
                    raise ValueError('JSON-Spec muss ein Objekt sein, ist aber ' + type(spec).__name__)
                if ftype == 'docx':
                    fname, fpath = create_docx_from_spec(spec)
                elif ftype == 'xlsx':
                    fname, fpath = create_xlsx_from_spec(spec)
                elif ftype == 'pdf':
                    fname, fpath = create_pdf_from_spec(spec)
                elif ftype == 'pptx':
                    fname, fpath = create_pptx_from_spec(spec)
                else:
                    raise ValueError('Unbekannter Typ: ' + ftype)
                created_files.append({'filename': fname, 'path': fpath, 'type': ftype})
                # Copy to agent memory
                if state.get('speicher'):
                    _mem_dir = os.path.join(state['speicher'], 'memory')
                    os.makedirs(_mem_dir, exist_ok=True)
                    import shutil
                    try:
                        shutil.copy2(fpath, os.path.join(_mem_dir, fname))
                    except Exception:
                        pass
                text = text.replace(full_block, '')
            except Exception as fe:
                err_msg = str(fe)
                if 'double quotes' in err_msg or 'Expecting' in err_msg:
                    err_msg = 'JSON-Format ungueltig. Bitte versuche es erneut.'
                text = text.replace(full_block, f'\n*Datei-Erstellung fehlgeschlagen: {err_msg}*\n')

        # Resolve session_id for task-status attribution
        _task_session_id = kwargs.get('_session_id')
        if not _task_session_id:
            for _sid, _st in sessions.items():
                if _st is state:
                    _task_session_id = _sid
                    break
        _task_session_id = _task_session_id or 'default'

        # Parse CREATE_IMAGE
        created_images = []
        img_pattern = re.compile(r'\[?CREATE_IMAGE:\s*(.+?)\]?(?:\n|$)')
        img_matches = list(img_pattern.finditer(text))
        for m in reversed(img_matches):
            img_prompt = m.group(1).strip().rstrip(']')
            _img_task_id = task_create('image', _task_session_id, estimated_total=30)
            try:
                current_provider = state.get('provider', 'anthropic')
                fname, fpath, fallback_info = generate_image(
                    img_prompt, state['agent'] or 'standard', current_provider,
                    task_id=_img_task_id,
                )
                created_images.append({
                    'filename': fname, 'path': fpath,
                    'prompt': img_prompt[:100], 'task_id': _img_task_id,
                })
                replacement = f'\n\n*Bild erfolgreich generiert: {fname}{fallback_info}. Das Bild wird unten angezeigt.*'
                text = text[:m.start()] + replacement + text[m.end():]
            except Exception as ie:
                task_error(_img_task_id, str(ie))
                text = text[:m.start()] + f'\n\n*Bild-Generierung fehlgeschlagen: {str(ie)}. Du kannst es erneut versuchen.*' + text[m.end():]

        # Parse CREATE_VIDEO
        created_videos = []
        vid_pattern = re.compile(r'\[?CREATE_VIDEO:\s*(.+?)\]?(?:\n|$)')
        vid_matches = list(vid_pattern.finditer(text))
        for m in reversed(vid_matches):
            vid_prompt = m.group(1).strip().rstrip(']')
            _vid_task_id = task_create('video', _task_session_id, estimated_total=180)
            try:
                fname, fpath, _vid_aspect = generate_video(
                    vid_prompt, state['agent'] or 'standard',
                    state.get('provider', 'anthropic'),
                    task_id=_vid_task_id,
                )
                # Detect portrait: API aspect OR prompt keywords
                _portrait_kws = ['9:16', 'portrait', 'hochformat', 'vertical', 'senkrecht', 'vertikal', 'tiktok', 'reels', 'shorts']
                _is_portrait = (_vid_aspect == '9:16') or any(kw in vid_prompt.lower() for kw in _portrait_kws)
                created_videos.append({
                    'filename': fname, 'path': fpath,
                    'prompt': vid_prompt[:100], 'task_id': _vid_task_id,
                    'is_portrait': _is_portrait,
                })
                text = text[:m.start()] + f'\n\n*Video erfolgreich generiert: {fname}. Das Video wird unten angezeigt.*' + text[m.end():]
            except Exception as ve:
                task_error(_vid_task_id, str(ve))
                text = text[:m.start()] + f'\n\n*Video-Generierung fehlgeschlagen: {str(ve)}. Du kannst es erneut versuchen.*' + text[m.end():]

        text = text.strip()
        state['verlauf'].append({'role': 'assistant', 'content': text})
        # Auto-save full session after each message (overwrites entire file)
        for _sid, _st in sessions.items():
            if _st is state:
                auto_save_session(_sid)
                break
        return {
            'response': text,
            'model_name': model_name,
            'provider_display': PROVIDER_DISPLAY.get(provider_key, provider_key),
            'model_display': MODEL_DISPLAY.get(model_id, model_name),
            'auto_loaded': auto_loaded_names,
            'auto_search_info': auto_search_info,
            'agent': state['agent'],
            'created_files': created_files,
            'created_emails': created_emails,
            'created_whatsapps': created_whatsapps,
            'created_images': created_images,
            'created_videos': created_videos,
            'created_slacks': created_slacks,
        }
    except Exception as e:
        state['verlauf'].pop()
        raise


def process_queue_worker(state):
    """Background thread: processes queued messages one by one."""
    while True:
        with queue_lock:
            if state['stop_requested']:
                state['stop_requested'] = False
                state['queue'] = []
                state['processing'] = False
                state['current_prompt'] = ''
                return
            if not state['queue']:
                state['processing'] = False
                state['current_prompt'] = ''
                return
            item = state['queue'].pop(0)
            state['current_prompt'] = item['message'][:50]

        try:
            result = process_single_message(item['message'], kontext_override=item.get('kontext_snapshot'), state=state)
            result['queue_id'] = item['id']
            result['original_message'] = item['message'][:50]
            with queue_lock:
                state['completed_responses'].append(result)
            # Auto-save after queue item processed
            for _sid, _st in sessions.items():
                if _st is state:
                    auto_save_session(_sid)
                    break
        except Exception as e:
            with queue_lock:
                state['completed_responses'].append({
                    'queue_id': item['id'],
                    'original_message': item['message'][:50],
                    'error': str(e),
                })


# Pending sub-agent delegations awaiting user confirmation
_pending_delegations = {}  # confirmation_id -> {session_id, sub_agent, msg, kontext, timestamp, auto_loaded, auto_search_info}

@app.route('/api/subagent_confirm', methods=['POST'])
def subagent_confirm():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    confirmation_id = request.json.get('confirmation_id', '')
    confirmed = request.json.get('confirmed', False)

    pending = _pending_delegations.pop(confirmation_id, None)
    if not pending:
        return jsonify({'error': 'Confirmation abgelaufen oder nicht gefunden'})

    if not confirmed:
        # User declined — process with current agent (no delegation)
        try:
            # Re-process the original message without delegation
            result = process_single_message(
                pending['msg'], state=state, skip_delegation=True,
                kontext_override=pending.get('kontext'),
                auto_loaded_override=pending.get('auto_loaded', []),
                auto_search_info_override=pending.get('auto_search_info', ''),
            )
            return jsonify(result)
        except Exception as e:
            return jsonify({'error': str(e)})

    # User confirmed — execute delegation
    try:
        deleg_result = execute_delegation(
            pending['sub_agent'], pending['msg'],
            pending.get('kontext', []), state=state
        )
        if 'error' in deleg_result:
            return jsonify({'error': deleg_result['error']})
        state['verlauf'].append({'role': 'user', 'content': pending['msg']})
        state['verlauf'].append({'role': 'assistant', 'content': deleg_result['response']})
        for _sid, _st in sessions.items():
            if _st is state:
                auto_save_session(_sid)
                break
        model_name = deleg_result.get('model_name', '')
        return jsonify({
            'response': deleg_result['response'],
            'model_name': model_name,
            'provider_display': deleg_result.get('provider_display', ''),
            'model_display': deleg_result.get('model_display', ''),
            'auto_loaded': pending.get('auto_loaded', []),
            'auto_search_info': pending.get('auto_search_info', ''),
            'agent': state['agent'],
            'created_files': [], 'created_emails': [], 'created_whatsapps': [],
            'created_images': [], 'created_videos': [], 'created_slacks': [],
            'delegated_to': deleg_result.get('delegated_to', ''),
            'delegated_display': deleg_result.get('delegated_display', ''),
        })
    except Exception as e:
        return jsonify({'error': str(e)})


@app.route('/chat', methods=['POST'])
def chat():
    if _shutdown_event.is_set():
        return jsonify({'error': 'Server faehrt herunter — bitte warten'}), 503
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    if not state['agent']:
        return jsonify({'error': 'Kein Agent aktiv'})
    msg = request.json['message']

    with queue_lock:
        if state['processing']:
            item = {
                'id': str(uuid.uuid4()),
                'message': msg,
                'kontext_snapshot': copy.deepcopy(state['kontext_items']),
                'timestamp': datetime.datetime.now().isoformat(),
            }
            state['queue'].append(item)
            return jsonify({
                'queued': True,
                'position': len(state['queue']),
                'message': msg[:50],
                'queue_id': item['id'],
            })
        state['processing'] = True
        state['current_prompt'] = msg[:50]

    with _active_requests_lock:
        global _active_requests
        _active_requests += 1
    try:
        result = process_single_message(msg, state=state, _session_id=session_id)
    except Exception as e:
        with queue_lock:
            state['processing'] = False
            state['current_prompt'] = ''
        return jsonify({'error': str(e)})
    finally:
        with _active_requests_lock:
            _active_requests -= 1

    with queue_lock:
        has_queue = len(state['queue']) > 0
    if has_queue:
        threading.Thread(target=process_queue_worker, args=(state,), daemon=True).start()
    else:
        with queue_lock:
            state['processing'] = False
            state['current_prompt'] = ''

    result['queue_active'] = has_queue
    return jsonify(result)


@app.route('/stop_queue', methods=['POST'])
def stop_queue():
    session_id = request.json.get('session_id', 'default') if request.is_json else 'default'
    state = get_session(session_id)
    with queue_lock:
        cancelled = len(state['queue'])
        state['stop_requested'] = True
        state['queue'] = []
    return jsonify({'ok': True, 'cancelled': cancelled})


@app.route('/queue_status', methods=['GET'])
def queue_status():
    session_id = request.args.get('session_id', 'default')
    state = get_session(session_id)
    with queue_lock:
        payload = {
            'processing': state['processing'],
            'queue_length': len(state['queue']),
            'queue_preview': [{'id': q['id'], 'message': q['message'][:50]} for q in state['queue'][:3]],
            'current_prompt': state['current_prompt'],
        }
    # Piggyback active long-running tasks (video/image) so the frontend
    # can discover them without a separate endpoint.
    payload['active_tasks'] = tasks_for_session(session_id)
    return jsonify(payload)


@app.route('/task_status/<task_id>', methods=['GET'])
def task_status(task_id):
    """Return the current status of a long-running task (video/image generation)."""
    t = task_get(task_id)
    if not t:
        return jsonify({'error': 'Task nicht gefunden', 'task_id': task_id}), 404
    # Opportunistic cleanup of very old finished tasks
    try:
        tasks_cleanup()
    except Exception:
        pass
    return jsonify(t)


@app.route('/poll_responses', methods=['GET'])
def poll_responses():
    session_id = request.args.get('session_id', 'default')
    state = get_session(session_id)
    with queue_lock:
        responses = list(state['completed_responses'])
        state['completed_responses'] = []
    return jsonify({'responses': responses})


@app.route('/available_subagents', methods=['GET'])
def available_subagents():
    session_id = request.args.get('session_id', 'default')
    state = get_session(session_id)
    agent = request.args.get('agent', state.get('agent', ''))
    parent = get_parent_agent(agent) or agent
    if not parent:
        return jsonify({'ok': False, 'subagents': []})
    subs = _get_available_subagents(parent)
    kw_map = _load_subagent_keywords()
    result = []
    for s in subs:
        result.append({
            'name': s['full_name'],
            'label': s['sub_label'],
            'display': get_agent_display_name(s['full_name']),
            'keywords': kw_map.get(s['sub_label'], []),
        })
    return jsonify({'ok': True, 'parent': parent, 'subagents': result})


@app.route('/add_url', methods=['POST'])
def add_url():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    url = request.json['url'].strip()
    if url and not url.startswith('http://') and not url.startswith('https://'):
        url = 'https://' + url
    try:
        r = requests.get(url, timeout=15, headers={
            'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 Chrome/120.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,*/*',
            'Accept-Language': 'de-DE,de;q=0.9,en;q=0.8'
        })
        soup = BeautifulSoup(r.text, 'html.parser')
        for tag in soup(['script','style','nav','footer','header']): tag.decompose()
        text = soup.get_text(separator='\n', strip=True)
        text = '\n'.join(l for l in text.splitlines() if l.strip())[:20000]
        title = (soup.title.string if soup.title else url).strip()[:60]
        memory_dir = os.path.join(state['speicher'], 'memory')
        os.makedirs(memory_dir, exist_ok=True)
        safe = ''.join(c if c.isalnum() or c in '-_' else '_' for c in title)
        fname = safe + '.txt'
        with open(os.path.join(memory_dir, fname), 'w') as f:
            f.write('Quelle: ' + url + '\n\n' + text)
        state['kontext_items'].append({'name':title, 'content':text})
        if fname not in state['session_files']:
            state['session_files'].append(fname)
        return jsonify({'ok':True, 'title':title})
    except Exception as e:
        return jsonify({'ok':False, 'error':str(e)})


# ── Deep Memory Search ────────────────────────────────────────────────────────

def deep_memory_search(memory_dir, query, date_from=None, date_to=None, direction=None, contact=None, max_results=10):
    """3-stage deep search: filename filter -> content scan -> scored results."""
    if not os.path.isdir(memory_dir) or not query:
        return {'total_files_scanned': 0, 'candidates_after_filter': 0, 'results': []}

    all_files = os.listdir(memory_dir)
    total = len(all_files)

    # Stage 1: Filename filter (no file opening)
    candidates = []
    for fname in all_files:
        if not fname.endswith(('.txt', '.eml', '.json')):
            continue
        fname_lower = fname.lower()

        # Date filter (prefix YYYY-MM-DD)
        if date_from:
            file_date = fname[:10]
            if file_date < date_from:
                continue
        if date_to:
            file_date = fname[:10]
            if file_date > date_to:
                continue

        # Direction filter
        if direction:
            d = direction.upper()
            if d == 'IN' and '_IN_' not in fname:
                continue
            if d == 'OUT' and '_OUT_' not in fname:
                continue

        # Contact filter (substring in filename)
        if contact:
            if contact.lower().replace('@', '_at_').replace('.', '_') not in fname_lower and contact.lower() not in fname_lower:
                continue

        candidates.append(fname)

    # Stage 2: Content scan (only candidates)
    keywords = [kw.lower() for kw in query.split() if len(kw) >= 2]
    scored = []
    for fname in candidates:
        fpath = os.path.join(memory_dir, fname)
        try:
            with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                raw = f.read(30000)
            # For JSON webclips: extract searchable text fields
            if fname.endswith('.json'):
                try:
                    parsed = json.loads(raw)
                    content = (parsed.get('title', '') + '\n' +
                              parsed.get('url', '') + '\n' +
                              parsed.get('full_text', ''))[:5000]
                except (json.JSONDecodeError, TypeError):
                    content = raw[:5000]
            else:
                content = raw[:5000]
        except Exception:
            continue

        content_lower = content.lower()
        score = sum(content_lower.count(kw) for kw in keywords)
        # Bonus for filename match
        fname_lower = fname.lower()
        score += sum(2 for kw in keywords if kw in fname_lower)

        if score > 0:
            # Parse headers
            headers = {}
            for line in content.split('\n')[:10]:
                for key in ['Von', 'An', 'Betreff', 'Datum', 'Richtung', 'Kontakt']:
                    if line.startswith(key + ': '):
                        headers[key] = line[len(key) + 2:]

            # Preview: skip headers, get body start
            body_start = content.find('\u2500' * 10)
            if body_start > 0:
                preview = content[body_start:body_start + 300].strip().lstrip('\u2500').strip()
            else:
                preview = content[:300]

            scored.append({
                'filename': fname,
                'score': score,
                'preview': preview,
                'von': headers.get('Von', ''),
                'an': headers.get('An', ''),
                'betreff': headers.get('Betreff', ''),
                'datum': headers.get('Datum', ''),
                'richtung': headers.get('Richtung', ''),
                'kontakt': headers.get('Kontakt', ''),
            })

    scored.sort(key=lambda x: x['score'], reverse=True)

    return {
        'total_files_scanned': total,
        'candidates_after_filter': len(candidates),
        'results': scored[:max_results],
    }


@app.route('/api/working-memory/<agent>', methods=['POST'])
def api_working_memory(agent):
    """Manage agent working memory: add, remove, list."""
    data = request.get_json(force=True)
    action = data.get('action', 'list')
    try:
        if action == 'add':
            manifest = working_memory_add(
                agent, data['filename'], data['content'],
                priority=data.get('priority', 5),
                description=data.get('description', ''),
            )
            return jsonify({'ok': True, 'action': 'add', 'filename': data['filename'], 'manifest': manifest})
        elif action == 'remove':
            manifest = working_memory_remove(agent, data['filename'])
            if manifest is None:
                return jsonify({'ok': False, 'error': 'Manifest nicht gefunden'}), 404
            return jsonify({'ok': True, 'action': 'remove', 'filename': data['filename'], 'manifest': manifest})
        elif action == 'list':
            manifest = working_memory_list(agent)
            return jsonify({'ok': True, 'action': 'list', 'manifest': manifest})
        else:
            return jsonify({'ok': False, 'error': f'Unbekannte Aktion: {action}'}), 400
    except KeyError as ke:
        return jsonify({'ok': False, 'error': f'Fehlendes Feld: {ke}'}), 400
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)}), 500


@app.route('/api/memory-files-search', methods=['GET'])
def api_memory_files_search():
    """Autocomplete search for files in agent memory."""
    q = request.args.get('q', '').strip().lower()
    agent = request.args.get('agent', '').strip()
    if not q or len(q) < 2 or not agent:
        return jsonify([])

    speicher = get_agent_speicher(agent)
    if not speicher:
        return jsonify([])

    results = []
    # Search in search index first (fast)
    index_path = os.path.join(speicher, '.search_index.json')
    if os.path.exists(index_path):
        try:
            with open(index_path) as f:
                idx = json.load(f)
            entries = idx if isinstance(idx, dict) else {}
            for fname, entry in entries.items():
                fname_lower = fname.lower()
                subject = (entry.get('subject') or '').lower()
                preview = (entry.get('preview') or '').lower()
                keywords_str = ' '.join(entry.get('keywords', [])).lower()
                if q in fname_lower or q in subject or q in preview or q in keywords_str:
                    results.append({
                        'filename': fname,
                        'snippet': entry.get('subject') or entry.get('preview', '')[:80] or fname,
                    })
                    if len(results) >= 8:
                        break
        except Exception:
            pass

    # Fallback: scan filenames directly if no index results
    if not results:
        memory_dir = os.path.join(speicher, 'memory')
        if os.path.isdir(memory_dir):
            for fname in os.listdir(memory_dir):
                if q in fname.lower():
                    results.append({'filename': fname, 'snippet': fname})
                    if len(results) >= 8:
                        break

    return jsonify(results)


@app.route('/api/memory/search', methods=['POST'])
def api_memory_search():
    """Deep memory search with filename + content filters."""
    data = request.json or {}
    agent = data.get('agent', '')
    query = data.get('query', '')

    if not agent or not query:
        return jsonify({'error': 'agent and query required', 'results': []})

    speicher = get_agent_speicher(agent)
    if not speicher:
        return jsonify({'error': f'Agent "{agent}" not found', 'results': []})

    memory_dir = os.path.join(speicher, 'memory')
    result = deep_memory_search(
        memory_dir,
        query,
        date_from=data.get('date_from'),
        date_to=data.get('date_to'),
        direction=data.get('direction'),
        contact=data.get('contact'),
        max_results=int(data.get('max_results', 10)),
    )
    return jsonify(result)


@app.route('/api/context-info', methods=['GET'])
def api_context_info():
    """Returns token estimates for system prompt, conversation, and loaded memory files."""
    try:
        session_id = request.args.get('session_id', 'default')
        state = get_session(session_id)

        sp = state.get('system_prompt') or ''
        sp_tokens = len(sp) // 4

        conv_text = json.dumps(state.get('verlauf', []))
        conv_tokens = len(conv_text) // 4

        memory_files = []
        for item in state.get('kontext_items', []):
            content = item.get('content', '')
            tokens = len(content) // 4
            memory_files.append({
                'name': item.get('name', ''),
                'tokens': tokens,
                'removable': True,
            })

        total = sp_tokens + conv_tokens + sum(m['tokens'] for m in memory_files)

        return jsonify({
            'system_prompt_tokens': sp_tokens,
            'conversation_tokens': conv_tokens,
            'memory_files_loaded': memory_files,
            'total_tokens': total,
            'limit_warning': total > 25000,
        })
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e), 'total_tokens': 0, 'limit_warning': False,
                        'system_prompt_tokens': 0, 'conversation_tokens': 0, 'memory_files_loaded': []})


# CALENDAR_INTEGRATION_V1: Kalender-API Route
@app.route('/api/calendar', methods=['GET', 'POST'])
def api_calendar():
    """Gibt Kalender-Events im angegebenen Zeitraum zurueck."""
    if request.method == 'POST' and request.is_json:
        data = request.json or {}
    else:
        data = dict(request.args)
    days_back = int(data.get('days_back', 0))
    days_ahead = int(data.get('days_ahead', 7))
    cal_filter = data.get('calendar_filter', None)
    search = data.get('search', None)
    calendars = [cal_filter] if cal_filter else None

    events, cals_found, error = get_calendar_events(
        days_back=days_back, days_ahead=days_ahead,
        calendars=calendars, search=search,
    )
    if error:
        return jsonify({'error': error, 'events': [], 'count': 0})

    now = datetime.datetime.now()
    return jsonify({
        'events': events,
        'count': len(events),
        'range': {
            'from': (now - datetime.timedelta(days=days_back)).strftime('%Y-%m-%d'),
            'to': (now + datetime.timedelta(days=days_ahead)).strftime('%Y-%m-%d'),
        },
        'calendars_found': sorted(cals_found),
    })


# SLACK_API_V1: Slack-API Route
@app.route('/api/slack', methods=['POST'])
def api_slack():
    """Slack-API Proxy: send, channels, history, users."""
    data = request.json or {}
    action = data.get('action', 'send')

    if action == 'send':
        channel = data.get('channel', '')
        text = data.get('text', data.get('message', ''))
        if not channel or not text:
            return jsonify({'error': 'channel und text/message erforderlich'})
        # Channel-Name → ID aufloesen wenn noetig
        if channel.startswith('#'):
            ch_id = slack_find_channel_id(channel)
            if ch_id:
                channel = ch_id
        ok, resp = slack_send_message(channel, text, thread_ts=data.get('thread_ts'))
        return jsonify({'ok': ok, 'data': resp})

    elif action == 'channels':
        ok, resp = slack_list_channels(limit=data.get('limit', 100))
        channels = [{'id': c['id'], 'name': c['name'], 'topic': c.get('topic', {}).get('value', '')}
                    for c in resp.get('channels', [])] if ok else []
        return jsonify({'ok': ok, 'channels': channels})

    elif action == 'history':
        ch = data.get('channel_id', data.get('channel', ''))
        if not ch:
            return jsonify({'error': 'channel_id erforderlich'})
        ok, resp = slack_channel_history(ch, limit=data.get('limit', 20))
        return jsonify({'ok': ok, 'messages': resp.get('messages', []) if ok else [], 'error': resp.get('error')})

    elif action == 'users':
        ok, resp = slack_list_users(limit=data.get('limit', 200))
        users = [{'id': u['id'], 'name': u.get('real_name', u.get('name', '')),
                  'display': u.get('profile', {}).get('display_name', '')}
                 for u in resp.get('members', []) if not u.get('is_bot') and not u.get('deleted')] if ok else []
        return jsonify({'ok': ok, 'users': users})

    return jsonify({'error': f'Unbekannte action: {action}'})


# CANVA_API_V1: Canva-API Route
@app.route('/api/canva', methods=['POST'])
def api_canva():
    """Canva Connect API Proxy: designs, create, export, folders."""
    data = request.json or {}
    action = data.get('action', 'list')

    if action == 'list' or action == 'search':
        ok, resp = canva_list_designs(query=data.get('query'), count=data.get('count', 20))
        return jsonify({'ok': ok, 'data': resp})

    elif action == 'get':
        did = data.get('design_id', '')
        if not did:
            return jsonify({'error': 'design_id erforderlich'})
        ok, resp = canva_get_design(did)
        return jsonify({'ok': ok, 'data': resp})

    elif action == 'create':
        title = data.get('title', 'Neues Design')
        ok, resp = canva_create_design(
            title, design_type=data.get('design_type', 'doc'),
            width=data.get('width'), height=data.get('height'),
        )
        return jsonify({'ok': ok, 'data': resp})

    elif action == 'export':
        did = data.get('design_id', '')
        fmt = data.get('format', 'pdf')
        if not did:
            return jsonify({'error': 'design_id erforderlich'})
        ok, resp = canva_export_design(did, format_type=fmt)
        return jsonify({'ok': ok, 'data': resp})

    elif action == 'folders':
        ok, resp = canva_list_folders(count=data.get('count', 50))
        return jsonify({'ok': ok, 'data': resp})

    # CANVA_CAMPAIGNS_V1: Brand Templates + Autofill Actions
    elif action == 'brand_templates' or action == 'templates':
        ok, resp = canva_list_brand_templates(query=data.get('query'), count=data.get('count', 50))
        return jsonify({'ok': ok, 'data': resp})

    elif action == 'template_dataset':
        tid = data.get('template_id', '')
        if not tid:
            return jsonify({'error': 'template_id erforderlich'})
        ok, resp = canva_get_template_dataset(tid)
        return jsonify({'ok': ok, 'data': resp})

    elif action == 'autofill':
        tid = data.get('template_id', '')
        field_data = data.get('data', {})
        if not tid:
            return jsonify({'error': 'template_id erforderlich'})
        if not field_data:
            return jsonify({'error': 'data (Feld-Mappings) erforderlich'})
        ok, resp = canva_autofill(tid, field_data, title=data.get('title'))
        return jsonify({'ok': ok, 'data': resp})

    elif action == 'autofill_status':
        jid = data.get('job_id', '')
        if not jid:
            return jsonify({'error': 'job_id erforderlich'})
        ok, resp = canva_get_autofill_job(jid)
        return jsonify({'ok': ok, 'data': resp})

    elif action == 'batch_campaign':
        tid = data.get('template_id', '')
        rows = data.get('rows', [])
        if not tid or not rows:
            return jsonify({'error': 'template_id und rows erforderlich'})
        jobs = canva_batch_campaign(tid, rows, title_prefix=data.get('title_prefix', 'Campaign'))
        return jsonify({'ok': True, 'jobs': jobs, 'count': len(jobs)})

    elif action == 'upload_asset':
        url = data.get('url', '')
        name = data.get('name', 'uploaded_image')
        if not url:
            return jsonify({'error': 'url erforderlich'})
        ok, resp = canva_upload_asset(url, name=name)
        return jsonify({'ok': ok, 'data': resp})

    return jsonify({'error': f'Unbekannte action: {action}'})


@app.route('/search_memory', methods=['POST'])
def search_memory():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    query = request.json.get('query', '')
    if not state.get('speicher') or not query:
        return jsonify({'ok': False, 'results': []})
    if auto_search:
        results, feedback = auto_search(query, state['speicher'])
        return jsonify({'ok': True, 'results': [{'name': r['name'], 'preview': r['content'][:500]} for r in results], 'feedback': feedback})
    # Fallback to old search
    memory_dir = os.path.join(state['speicher'], 'memory')
    keywords = parse_search_keywords(query)
    results = scored_memory_search(memory_dir, keywords, max_results=5)
    return jsonify({'ok': True, 'results': [{'name': r['name'], 'preview': r['content'][:500]} for r in results]})


NOTIFICATION_PATTERNS = [
    'noreply', 'no-reply', 'no_reply', 'donotreply', 'do-not-reply',
    'mailer-daemon', 'postmaster', 'notifications@', 'notification@',
    'newsletter', 'digest', 'automated', 'auto-reply', 'autoreply',
    'bounce', 'info@', 'support@', 'service@', 'team@', 'hello@',
    'mailchimp', 'sendgrid', 'hubspot', 'salesforce', 'marketo',
    'intercom', 'zendesk', 'jira', 'confluence', 'atlassian',
    'github', 'gitlab', 'bitbucket', 'circleci', 'travis',
    'slack', 'notion', 'asana', 'trello', 'monday.com',
    'google-workspace', 'calendar-notification', 'accepted this invitation',
    'automatic reply', 'automatische antwort', 'abwesenheit',
    'out of office', 'ooo', 'delivery status', 'undeliverable',
    'unsubscribe', 'abmelden', 'cancelar inscricao',
]

def is_notification(from_field, subject):
    check = (from_field + ' ' + subject).lower()
    return any(p in check for p in NOTIFICATION_PATTERNS)


@app.route('/search_preview', methods=['POST'])
def search_preview():
    """Search memory and return rich preview results for interactive selection."""
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    query = request.json.get('query', '')
    search_type = request.json.get('type', None)
    # Map shorthand type names to SOURCE_TAXONOMY keys
    _type_aliases = {'slack': 'webclip_slack', 'salesforce': 'webclip_salesforce'}
    if search_type in _type_aliases:
        search_type = _type_aliases[search_type]
    is_recent = request.json.get('recent', False)
    if not state.get('speicher'):
        return jsonify({'ok': False, 'results': []})
    if not query and not is_recent:
        return jsonify({'ok': False, 'results': []})
    try:
        from search_engine import QueryParser, HybridSearch, SearchIndex, normalize_unicode, get_or_build_index, extract_search_keywords, search_contacts, get_recent_files

        # Empty query with recent flag — return latest files
        if is_recent and not query:
            if search_type:
                recent = get_recent_files(state['speicher'], category=search_type, limit=10)
            else:
                recent = get_recent_files(state['speicher'], per_category=True, limit=3)
            items = []
            for r in recent:
                items.append({
                    'name': r['name'], 'path': r['path'], 'type': r.get('source_type','file'),
                    'source_type': r.get('source_type','file'), 'date': r.get('date',''),
                    'from': r.get('from',''), 'subject': r.get('subject',''),
                    'preview': r.get('preview',''), 'score': 0,
                    'from_person': False, 'is_notification': False,
                })
            return jsonify({'ok': True, 'results': items, 'query': 'Neueste Dateien', 'feedback': None})

        # Contact search — special case
        if search_type == 'contact':
            contact_results = search_contacts(query, state['speicher'])
            items = []
            for c in contact_results:
                items.append({
                    'name': c['name'], 'path': '', 'type': 'contact', 'source_type': 'contact',
                    'date': '', 'from': '', 'subject': c['name'].replace('contact_','').replace('_',' '),
                    'preview': c['content'][:150], 'score': c['score'],
                    'from_person': False, 'is_notification': False,
                })
            return jsonify({'ok': True, 'results': items, 'query': query, 'feedback': {'query': query}})

        # NLP keyword extraction for long queries
        extracted_keywords = None
        effective_query = query
        words = query.split()
        if len(words) > 5:
            extracted_keywords = extract_search_keywords(query, search_type)
            if extracted_keywords:
                effective_query = ' '.join(extracted_keywords)

        intent = QueryParser.parse(effective_query, force_search=True)
        # Ensure index
        idx = get_or_build_index(state['speicher'])
        idx.update_index()
        results, feedback = HybridSearch.search(intent, state['speicher'], max_results=500, forced_type=search_type)
        if extracted_keywords and feedback:
            feedback['extracted_keywords'] = extracted_keywords
        # Build rich preview items
        items = []
        names_norm = [normalize_unicode(pn) for pn in intent.person_names]
        for r in results:
            fname = r['name']
            entry = idx.entries.get(fname, {})
            from_field = entry.get('from', '')
            from_norm = normalize_unicode(from_field) if from_field else ''
            from_person = any(pn in from_norm for pn in names_norm) if names_norm else False
            stype = entry.get('source_type') or entry.get('type', 'file')
            is_notif = is_notification(from_field, entry.get('subject', '')) or stype == 'notification'
            score = r.get('score', 0)
            if is_notif:
                score *= 0.1
            items.append({
                'name': fname,
                'path': entry.get('path', os.path.join(state['speicher'], 'memory', fname)),
                'type': entry.get('type', 'file'),
                'source_type': stype,
                'date': entry.get('date', ''),
                'from': from_field,
                'subject': entry.get('subject', ''),
                'preview': r['content'][:150].replace('\n', ' ').strip(),
                'score': score,
                'from_person': from_person,
                'is_notification': is_notif,
            })
        # Sort: from_person first, then highest score, then newest date
        if search_type == 'email':
            # For email searches: sort by date desc (newest first)
            items.sort(key=lambda x: x.get('date','') or '', reverse=True)
        else:
            items.sort(key=lambda x: (not x.get('from_person'), -x['score'], x.get('date','') or '0'), reverse=False)
        return jsonify({
            'ok': True,
            'results': items,
            'query': feedback.get('query', query) if feedback else query,
            'feedback': feedback,
        })
    except Exception as e:
        print(f"search_preview error: {e}")
        import traceback; traceback.print_exc()
        return jsonify({'ok': False, 'results': [], 'error': str(e)})


@app.route('/global_search_preview', methods=['POST'])
def global_search_preview():
    """Global search across all agents and Downloads shared."""
    query = request.json.get('query', '')
    search_type = request.json.get('type', None)
    _type_aliases2 = {'slack': 'webclip_slack', 'salesforce': 'webclip_salesforce'}
    if search_type in _type_aliases2:
        search_type = _type_aliases2[search_type]
    if not query:
        return jsonify({'ok': False, 'results': []})
    try:
        from search_engine import global_search, normalize_unicode, QueryParser, extract_search_keywords
        # NLP keyword extraction for long queries
        effective_query = query
        extracted_keywords = None
        if len(query.split()) > 5:
            extracted_keywords = extract_search_keywords(query, search_type)
            if extracted_keywords:
                effective_query = ' '.join(extracted_keywords)
        results, feedback = global_search(effective_query, max_results=500)
        if extracted_keywords and feedback:
            feedback['extracted_keywords'] = extracted_keywords
        return jsonify({
            'ok': True,
            'results': results,
            'query': feedback.get('query', query),
            'feedback': feedback,
            'global': True,
        })
    except Exception as e:
        print(f"global_search_preview error: {e}")
        import traceback; traceback.print_exc()
        return jsonify({'ok': False, 'results': [], 'error': str(e)})


@app.route('/load_selected_files', methods=['POST'])
def load_selected_files():
    """Load selected files into kontext_items."""
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    paths = request.json.get('paths', [])[:50]  # Max 50 files
    loaded = []
    agents_seen = set()
    for fpath in paths:
        fname = os.path.basename(fpath)
        if any(k['name'] == fname for k in state['kontext_items']):
            loaded.append(fname)
            continue
        try:
            ext = os.path.splitext(fname.lower())[1]
            # Screenshots: load as base64 image for vision
            if ext in ('.png', '.jpg', '.jpeg', '.gif', '.webp'):
                import base64
                with open(fpath, 'rb') as f:
                    img_data = base64.b64encode(f.read()).decode('utf-8')
                mime = {'png': 'image/png', 'jpg': 'image/jpeg', 'jpeg': 'image/jpeg',
                        'gif': 'image/gif', 'webp': 'image/webp'}.get(ext.lstrip('.'), 'image/png')
                img_data, mime = downscale_image_b64_if_needed(img_data, mime)
                if not img_data:
                    print(f"[IMG] Oversize-Screenshot '{fname}' uebersprungen", flush=True)
                    continue
                state['kontext_items'].append({
                    'name': fname,
                    'content': '[Screenshot]',
                    'image_b64': img_data,
                    'image_type': mime,
                    'path': fpath,
                    'auto_loaded': True,
                })
            else:
                with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                    text = f.read(20000)
                state['kontext_items'].append({
                    'name': fname,
                    'content': text,
                    'path': fpath,
                    'auto_loaded': True,
                })
            if fname not in state.get('session_files', []):
                state.setdefault('session_files', []).append(fname)
            loaded.append(fname)
            # Track agent from path
            fpath_lower = fpath.lower()
            for agent_name in ['signicat', 'privat', 'trustedcarrier']:
                if f'/{agent_name}/' in fpath_lower:
                    agents_seen.add(agent_name)
        except Exception as e:
            print(f"load_selected_files error for {fname}: {e}")
    return jsonify({'ok': True, 'loaded': loaded, 'agents': list(agents_seen)})


@app.route('/add_file', methods=['POST'])
def add_file():
    session_id = request.form.get('session_id', 'default')
    state = get_session(session_id)
    file = request.files.get('file')
    if not file:
        return jsonify({'ok': False, 'error': 'Keine Datei erhalten'})
    if not state['agent'] or not state['speicher']:
        return jsonify({'ok': False, 'error': 'Kein Agent aktiv - bitte zuerst Agenten auswaehlen'})
    try:
        raw = file.read()
        filename = file.filename
        fname_lower = filename.lower()

        # Detect image files
        image_exts = {'.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
                      '.gif': 'image/gif', '.webp': 'image/webp'}
        is_image = any(fname_lower.endswith(ext) for ext in image_exts)
        image_type = next((v for k,v in image_exts.items() if fname_lower.endswith(k)), None)

        # Save to memory folder
        memory_dir = os.path.join(state['speicher'], 'memory')
        os.makedirs(memory_dir, exist_ok=True)
        memory_path = os.path.join(memory_dir, filename)
        with open(memory_path, 'wb') as f:
            f.write(raw)

        if is_image:
            import base64
            b64 = base64.b64encode(raw).decode('utf-8')
            b64, image_type = downscale_image_b64_if_needed(b64, image_type)
            if not b64:
                return jsonify({'ok': False, 'error': f"Bild '{filename}' konnte nicht verarbeitet werden (zu gross oder beschaedigt)"})
            # Store as image item - will be sent as vision content
            state['kontext_items'].append({
                'name': filename,
                'content': f'[Bild: {filename}]',
                'image_b64': b64,
                'image_type': image_type
            })
        else:
            file_content = extract_file_content(raw, filename)
            state['kontext_items'].append({'name': filename, 'content': file_content})

        if filename not in state['session_files']:
            state['session_files'].append(filename)
        return jsonify({'ok': True, 'filename': filename, 'memory_path': memory_path})
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})

@app.route('/remove_ctx', methods=['POST'])
def remove_ctx():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    name = request.json['name']
    state['kontext_items'] = [i for i in state['kontext_items'] if i['name'] != name]
    return jsonify({'ok':True})


@app.route('/remove_all_ctx', methods=['POST'])
def remove_all_ctx():
    session_id = request.json.get('session_id', 'default')
    state = get_session(session_id)
    state['kontext_items'] = []
    return jsonify({'ok':True})


# ── ADMIN-BEREICH (added 2026-04-15) ───────────────────────────────────────
import html as _html_lib

_ADMIN_BASE = os.path.expanduser("~/AssistantDev")
_ADMIN_DOCS_DIR = os.path.join(_ADMIN_BASE, "docs")
_ADMIN_DATALAKE = os.path.expanduser(
    "~/Library/Mobile Documents/com~apple~CloudDocs/Downloads shared/claude_datalake"
)


def _admin_render_md(text):
    """Minimal markdown -> HTML: ## h2, ### h3, **bold**, lists, paragraphs.
    Escapes HTML first to prevent XSS from changelog content."""
    text = _html_lib.escape(text)
    out_lines = []
    in_list = False
    in_code = False
    code_buf = []
    for raw in text.split("\n"):
        line = raw.rstrip()

        # Fenced code blocks (```)
        if line.strip().startswith("```"):
            if in_code:
                out_lines.append('<pre><code>' + _html_lib.escape("\n".join(code_buf)) + '</code></pre>')
                code_buf = []
                in_code = False
            else:
                if in_list:
                    out_lines.append("</ul>")
                    in_list = False
                in_code = True
            continue
        if in_code:
            code_buf.append(line)
            continue

        # Headings
        if line.startswith("### "):
            if in_list:
                out_lines.append("</ul>")
                in_list = False
            out_lines.append(f"<h3>{line[4:].strip()}</h3>")
            continue
        if line.startswith("## "):
            if in_list:
                out_lines.append("</ul>")
                in_list = False
            out_lines.append(f"<h2>{line[3:].strip()}</h2>")
            continue
        if line.startswith("# "):
            if in_list:
                out_lines.append("</ul>")
                in_list = False
            out_lines.append(f"<h1>{line[2:].strip()}</h1>")
            continue

        # List items
        m_list = None
        stripped = line.lstrip()
        if stripped.startswith("- ") or stripped.startswith("* "):
            m_list = stripped[2:]
        if m_list is not None:
            if not in_list:
                out_lines.append("<ul>")
                in_list = True
            # bold + inline code
            item = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", m_list)
            item = re.sub(r"`([^`]+)`", r"<code>\1</code>", item)
            out_lines.append(f"<li>{item}</li>")
            continue

        if in_list:
            out_lines.append("</ul>")
            in_list = False

        # Horizontal rule
        if line.strip() == "---":
            out_lines.append("<hr>")
            continue

        # Empty -> paragraph break
        if not line.strip():
            out_lines.append("")
            continue

        # Bold + inline code in normal paragraph
        para = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", line)
        para = re.sub(r"`([^`]+)`", r"<code>\1</code>", para)
        out_lines.append(f"<p>{para}</p>")

    if in_list:
        out_lines.append("</ul>")
    if in_code:
        out_lines.append('<pre><code>' + _html_lib.escape("\n".join(code_buf)) + '</code></pre>')
    return "\n".join(out_lines)


_ADMIN_CSS = """
* { box-sizing: border-box; }
body { background:#1a1a2e; color:#e0e0e0; font-family:-apple-system,Inter,sans-serif; margin:0; padding:30px; line-height:1.55; }
.container { max-width:1000px; margin:0 auto; }
h1 { color:#f0c060; font-size:26px; margin:0 0 6px; }
h2 { color:#9ec5fe; font-size:20px; margin:24px 0 8px; border-bottom:1px solid #334; padding-bottom:4px; }
h3 { color:#cbb; font-size:16px; margin:16px 0 6px; }
p { margin:6px 0; color:#d0d0d0; }
ul { margin:6px 0 12px 22px; padding:0; }
li { margin:3px 0; color:#d0d0d0; }
code { background:#111; padding:1px 5px; border-radius:3px; font-family:ui-monospace,Menlo,monospace; font-size:12px; color:#9ec5fe; }
pre { background:#111; padding:12px; border-radius:6px; overflow-x:auto; border:1px solid #334; }
pre code { background:none; padding:0; color:#cbb; }
hr { border:none; border-top:1px solid #334; margin:18px 0; }
a { color:#4a8aca; text-decoration:none; }
a:hover { text-decoration:underline; }
.subtitle { color:#888; font-size:13px; margin-bottom:24px; }
.back { color:#4a8aca; font-size:13px; }
.admin-topbar { display:flex; align-items:center; gap:12px; padding:10px 0 16px; border-bottom:1px solid #334; margin-bottom:20px; }
.admin-topbar a { color:#aaa; text-decoration:none; font-size:12px; padding:6px 14px; border:1px solid #444; border-radius:6px; transition:all .15s; }
.admin-topbar a:hover { border-color:#f0c060; color:#f0c060; }
.admin-topbar a.active { border-color:#4a8aca; color:#4a8aca; background:#1a2a4a; }
.admin-topbar .back-btn { background:#222; border-color:#555; color:#ccc; font-weight:600; }
.admin-topbar .back-btn:hover { background:#333; border-color:#f0c060; color:#f0c060; }
.card-grid { display:grid; grid-template-columns:repeat(auto-fit,minmax(220px,1fr)); gap:14px; margin-top:18px; }
.card { background:#22224a; border:1px solid #334; border-radius:10px; padding:18px; text-decoration:none; color:#e0e0e0; transition:border-color .15s, transform .1s; }
.card:hover { border-color:#4a8aca; transform:translateY(-2px); text-decoration:none; }
.card .title { font-size:15px; font-weight:600; color:#f0c060; margin-bottom:4px; }
.card .desc { font-size:12px; color:#999; }
table { width:100%; border-collapse:collapse; margin:14px 0; background:#22224a; border-radius:8px; overflow:hidden; }
th { background:#2c2c5a; color:#9ec5fe; padding:10px 12px; text-align:left; font-size:13px; font-weight:600; }
td { padding:8px 12px; border-top:1px solid #2c2c4a; font-size:13px; color:#d0d0d0; vertical-align:top; }
.tag { display:inline-block; background:#1f4a1f; color:#a0d090; padding:1px 8px; border-radius:10px; font-size:11px; margin-right:4px; }
.tag-sub { background:#4a3a1f; color:#d0b090; }
.tag-orphan { background:#4a1f1f; color:#d09090; }
.status-ok { color:#7ec07e; }
.status-bad { color:#d07070; }
.metric { display:inline-block; margin-right:18px; color:#999; font-size:12px; }
.metric strong { color:#f0c060; font-size:14px; }
"""


def _admin_layout(title, body_html):
    topbar = (
        "<div class='admin-topbar'>"
        "<a class='back-btn' href='/'>&#8592; Chat</a>"
        "<a href='/admin'>Admin</a>"
        "<a href='/admin/access-control'>Access Control</a>"
        "<a href='/admin/permissions'>Berechtigungen</a>"
        "<a href='/admin/docs'>Docs</a>"
        "<a href='/admin/changelog'>Changelog</a>"
        "</div>"
    )
    return ("<!DOCTYPE html><html lang='de'><head><meta charset='UTF-8'>"
            f"<title>{_html_lib.escape(title)} — Admin</title>"
            f"<style>{_ADMIN_CSS}</style></head>"
            "<body><div class='container'>"
            f"{topbar}{body_html}"
            "</div></body></html>")


def _admin_status_check():
    """Quick status: web (8080) + clipper (8081) + watcher process.

    Matches both the Python-script form (`email_watcher.py`) and the compiled
    app-bundle binary name (`AssistantDev EmailWatcher`) introduced in 508796a.
    Without this, services running from the bundle show as offline.
    """
    import socket
    def port_alive(p):
        try:
            with socket.create_connection(("127.0.0.1", p), timeout=0.4):
                return True
        except Exception:
            return False
    web_ok = port_alive(8080)
    clip_ok = port_alive(8081)
    import subprocess as _sp
    def proc_alive(*patterns):
        for pat in patterns:
            try:
                out = _sp.run(["pgrep", "-f", pat], capture_output=True, text=True, timeout=2)
                if out.stdout.strip():
                    return True
            except Exception:
                pass
        return False
    watcher_ok = proc_alive("email_watcher.py", "AssistantDev EmailWatcher")
    kchat_ok = proc_alive("kchat_watcher.py")
    # WhatsApp Import: check last run from log file mtime
    wa_log = os.path.join(os.path.expanduser("~"), "AssistantDev", "logs", "whatsapp_import.log")
    wa_last_run = None
    wa_ok = False
    try:
        if os.path.exists(wa_log) and os.path.getsize(wa_log) > 0:
            mtime = os.path.getmtime(wa_log)
            import datetime as _dt
            wa_last_run = _dt.datetime.fromtimestamp(mtime).strftime('%d.%m.%Y %H:%M')
            # Consider "ok" if last run was within 25 minutes (interval is 20 min)
            import time as _tm
            wa_ok = (_tm.time() - mtime) < 1500
    except Exception:
        pass
    return {"web_8080": web_ok, "clipper_8081": clip_ok, "email_watcher": watcher_ok, "kchat_watcher": kchat_ok, "whatsapp_import": wa_ok, "whatsapp_last_run": wa_last_run}


@app.route('/api/services', methods=['GET'])
def api_services_status():
    """Returns status of all services as JSON."""
    st = _admin_status_check()
    services = [
        {"id": "web_server",      "name": "Web Server",      "port": 8080, "online": st["web_8080"]},
        {"id": "web_clipper",     "name": "Web Clipper",      "port": 8081, "online": st["clipper_8081"]},
        {"id": "email_watcher",   "name": "Email Watcher",    "port": None, "online": st["email_watcher"]},
        {"id": "kchat_watcher",   "name": "kChat Watcher",    "port": None, "online": st["kchat_watcher"]},
        {"id": "whatsapp_import", "name": "WhatsApp Import",  "port": None, "online": st["whatsapp_import"], "last_run": st.get("whatsapp_last_run"), "periodic": "20min"},
    ]
    return jsonify({"services": services})


@app.route('/api/services/restart', methods=['POST'])
def api_services_restart():
    """Restart a specific service by id."""
    import subprocess as _sp
    data = request.get_json(silent=True)
    service_id = (data or {}).get('service')
    # WhatsApp Import: trigger manually (not a daemon — runs once then exits)
    if service_id == 'whatsapp_import':
        script = os.path.join(os.path.expanduser("~"), "AssistantDev", "scripts", "whatsapp_db_import.py")
        log_path = os.path.join(os.path.expanduser("~"), "AssistantDev", "logs", "whatsapp_import.log")
        os.makedirs(os.path.dirname(log_path), exist_ok=True)
        lf = open(log_path, "a")
        _sp.Popen([sys.executable, script, "--agent", "privat", "--min-messages", "3"],
                   stdout=lf, stderr=_sp.STDOUT, cwd=os.path.join(os.path.expanduser("~"), "AssistantDev"))
        return jsonify({"ok": True, "restarted": "whatsapp_import", "note": "Import gestartet"})
    RESTART_MAP = {
        "web_server":    ("web_server.py", "AssistantDev WebServer"),
        "web_clipper":   ("web_clipper_server.py",),
        "email_watcher": ("email_watcher.py", "AssistantDev EmailWatcher"),
        "kchat_watcher": ("kchat_watcher.py",),
    }
    if service_id not in RESTART_MAP:
        return jsonify({"ok": False, "error": "Unbekannter Service: " + str(service_id)}), 400
    patterns = RESTART_MAP[service_id]
    # Kill
    for pat in patterns:
        try:
            _sp.run(["pkill", "-f", pat], capture_output=True, timeout=5)
        except Exception:
            pass
    import time as _time
    _time.sleep(2)
    # Start (except web_server — watchdog restarts it)
    if service_id != "web_server":
        script_name = patterns[0]
        script_path = os.path.join(BASE, "src", script_name)
        if os.path.exists(script_path):
            log_path = os.path.join(BASE, "logs", script_name.replace(".py", ".log"))
            os.makedirs(os.path.dirname(log_path), exist_ok=True)
            lf = open(log_path, "a")
            _sp.Popen(
                [sys.executable, script_path],
                stdout=lf, stderr=_sp.STDOUT,
                cwd=os.path.join(BASE, "src"),
            )
    return jsonify({"ok": True, "restarted": service_id})


@app.route('/api/open-window', methods=['POST'])
def api_open_window():
    """Open a new native pywebview window with a given path."""
    import subprocess as _sp
    data = request.get_json(silent=True)
    path = (data or {}).get('path', '/')
    # Sanitize path
    if not path.startswith('/'):
        path = '/' + path
    dashboard_script = os.path.join(BASE, "src", "dashboard_window.py")
    if not os.path.exists(dashboard_script):
        return jsonify({"ok": False, "error": "dashboard_window.py nicht gefunden"}), 500
    _sp.Popen(
        [sys.executable, dashboard_script, path],
        stdout=_sp.DEVNULL, stderr=_sp.DEVNULL,
    )
    return jsonify({"ok": True, "path": path})


@app.route('/admin', methods=['GET'])
def admin_root():
    st = _admin_status_check()
    def _badge(ok, label):
        cls = "status-ok" if ok else "status-bad"
        sign = "● online" if ok else "○ offline"
        return f"<span class='metric'><strong>{_html_lib.escape(label)}</strong> <span class='{cls}'>{sign}</span></span>"
    cards = [
        ("/admin/changelog", "📋 Changelog", "Alle Aenderungen am System"),
        ("/admin/docs", "📖 Technische Docs", "API Reference, Architektur, Workflows"),
        ("/admin/permissions", "🔐 Memory-Berechtigungen", "Wer hat Zugriff auf welches Memory"),
        ("/admin/access-control", "⚙ Access Control", "Cross-Agent Reads, Shared Memory"),
    ]
    cards_html = "".join(
        f"<a class='card' href='{href}'><div class='title'>{_html_lib.escape(t)}</div>"
        f"<div class='desc'>{_html_lib.escape(d)}</div></a>"
        for href, t, d in cards
    )
    body = (
        "<a href='/' class='back'>← Zurueck zum Chat</a>"
        "<h1>⚙ Admin-Bereich</h1>"
        "<div class='subtitle'>System-Verwaltung, Doku und Berechtigungen</div>"
        "<h2>System-Status</h2>"
        f"<div>{_badge(st['web_8080'], 'Web Server :8080')}{_badge(st['clipper_8081'], 'Web Clipper :8081')}{_badge(st['email_watcher'], 'Email Watcher')}</div>"
        "<h2>Bereiche</h2>"
        f"<div class='card-grid'>{cards_html}</div>"
    )
    return _admin_layout("Admin", body)


@app.route('/admin/changelog', methods=['GET'])
def admin_changelog():
    path = os.path.join(_ADMIN_BASE, "changelog.md")
    if not os.path.exists(path):
        return _admin_layout("Changelog", "<a href='/admin' class='back'>← Admin</a><h1>Changelog</h1><p>Datei nicht gefunden.</p>"), 404
    try:
        st = os.stat(path)
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            md = f.read()
    except Exception as e:
        return _admin_layout("Changelog", f"<p>Fehler: {_html_lib.escape(str(e))}</p>"), 500
    mtime = datetime.datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
    size_kb = st.st_size / 1024.0
    body = (
        "<a href='/admin' class='back'>← Admin</a>"
        "<h1>📋 Changelog</h1>"
        f"<div class='subtitle'>{path} — {size_kb:.1f} KB — zuletzt geaendert {mtime}</div>"
        f"{_admin_render_md(md)}"
    )
    return _admin_layout("Changelog", body)


@app.route('/admin/docs', methods=['GET'])
def admin_docs_list():
    if not os.path.isdir(_ADMIN_DOCS_DIR):
        return _admin_layout("Docs", "<a href='/admin' class='back'>← Admin</a><h1>Docs</h1><p>docs/ nicht gefunden.</p>"), 404
    items = []
    for f in sorted(os.listdir(_ADMIN_DOCS_DIR)):
        fp = os.path.join(_ADMIN_DOCS_DIR, f)
        if not os.path.isfile(fp):
            continue
        try:
            st = os.stat(fp)
            mtime = datetime.datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d")
            size_kb = st.st_size / 1024.0
        except Exception:
            mtime, size_kb = "?", 0.0
        items.append(
            f"<tr><td><a href='/admin/docs/{_html_lib.escape(f)}'>{_html_lib.escape(f)}</a></td>"
            f"<td>{size_kb:.1f} KB</td><td>{mtime}</td></tr>"
        )
    table = "<table><tr><th>Datei</th><th>Groesse</th><th>Geaendert</th></tr>" + "".join(items) + "</table>"
    body = (
        "<a href='/admin' class='back'>← Admin</a>"
        "<h1>📖 Technische Docs</h1>"
        f"<div class='subtitle'>{_ADMIN_DOCS_DIR}</div>"
        f"{table}"
    )
    return _admin_layout("Docs", body)


@app.route('/admin/docs/<path:filename>', methods=['GET'])
def admin_doc_view(filename):
    # Path traversal guard
    target = os.path.realpath(os.path.join(_ADMIN_DOCS_DIR, filename))
    docs_real = os.path.realpath(_ADMIN_DOCS_DIR)
    if not target.startswith(docs_real + os.sep) or not os.path.isfile(target):
        return _admin_layout("Docs", "<a href='/admin/docs' class='back'>← Docs</a><h1>Nicht gefunden</h1>"), 404
    try:
        st = os.stat(target)
        with open(target, 'r', encoding='utf-8', errors='replace') as f:
            md = f.read()
    except Exception as e:
        return _admin_layout("Docs", f"<p>Fehler: {_html_lib.escape(str(e))}</p>"), 500
    mtime = datetime.datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
    size_kb = st.st_size / 1024.0
    # Render markdown for .md, raw <pre> for .yaml/.yml/.txt
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".md":
        rendered = _admin_render_md(md)
    else:
        rendered = "<pre><code>" + _html_lib.escape(md) + "</code></pre>"
    body = (
        "<a href='/admin/docs' class='back'>← Docs</a>"
        f"<h1>{_html_lib.escape(filename)}</h1>"
        f"<div class='subtitle'>{size_kb:.1f} KB — zuletzt geaendert {mtime}</div>"
        f"{rendered}"
    )
    return _admin_layout(filename, body)


def _admin_collect_permissions():
    """Build agent-permission overview.
    Returns: list of dicts with keys: name, kind, memory_dir, file_count,
    parent, sub_agents, has_own_memory.
    Plus: list of orphan memory dirs (no agent definition)."""
    agents_dir = os.path.join(_ADMIN_DATALAKE, "config", "agents")
    if not os.path.isdir(agents_dir):
        return [], []

    # All agent definition files
    agent_files = sorted([f for f in os.listdir(agents_dir)
                          if f.endswith(".txt") and not f.startswith(".")])
    agent_names = [f[:-4] for f in agent_files]

    # Parent vs sub: parent has no underscore OR underscore-name where prefix is also an agent
    parents = set()
    subs_by_parent = {}
    for n in agent_names:
        if "_" in n:
            prefix = n.split("_", 1)[0]
            if prefix in agent_names:
                subs_by_parent.setdefault(prefix, []).append(n)
                continue
        parents.add(n)
    # Anything not in a parent's sub-list and not in `parents` is also a parent
    for n in agent_names:
        if n not in parents and n not in {s for subs in subs_by_parent.values() for s in subs}:
            parents.add(n)

    # Existing memory dirs in datalake root
    existing_memory_dirs = set()
    for entry in os.listdir(_ADMIN_DATALAKE):
        full = os.path.join(_ADMIN_DATALAKE, entry, "memory")
        if os.path.isdir(full):
            existing_memory_dirs.add(entry)

    rows = []
    for n in sorted(agent_names):
        is_parent = n in parents
        mem_dir = os.path.join(_ADMIN_DATALAKE, n, "memory") if is_parent else None
        file_count = 0
        has_own = False
        if mem_dir and os.path.isdir(mem_dir):
            try:
                file_count = sum(1 for f in os.listdir(mem_dir)
                                 if os.path.isfile(os.path.join(mem_dir, f)))
                has_own = True
            except Exception:
                pass
        # Sub agent: parent prefix and inherited memory
        parent = None
        if not is_parent and "_" in n:
            prefix = n.split("_", 1)[0]
            if prefix in agent_names:
                parent = prefix
        rows.append({
            "name": n,
            "kind": "Parent" if is_parent else "Sub-Agent",
            "has_own_memory": has_own,
            "memory_dir": mem_dir if has_own else None,
            "file_count": file_count,
            "parent": parent,
            "sub_agents": subs_by_parent.get(n, []),
        })

    # Orphan memory dirs: have memory/ but no agent .txt
    orphans = sorted(d for d in existing_memory_dirs if d not in agent_names)
    # Filter out non-agent system dirs
    system_dirs = {"config", "email_inbox", "claude_outputs", "global", "system_ward"}
    orphans = [d for d in orphans if d not in system_dirs]
    return rows, orphans


@app.route('/admin/permissions', methods=['GET'])
def admin_permissions():
    rows, orphans = _admin_collect_permissions()
    if not rows:
        return _admin_layout("Permissions",
            "<a href='/admin' class='back'>← Admin</a><h1>Memory-Berechtigungen</h1>"
            "<p>Keine Agent-Definitionen gefunden.</p>")
    tr = []
    for r in rows:
        name = r['name']
        esc_name = _html_lib.escape(name)
        # Working Memory path
        parent = r.get('parent', '')
        if r["kind"] == "Parent":
            wm_path = f"{esc_name}/working_memory/"
            wm_dir = os.path.join(BASE, name, 'working_memory')
        else:
            sub = name.split('_', 1)[1] if '_' in name else name
            wm_path = f"{_html_lib.escape(parent)}/working_memory/_{_html_lib.escape(sub)}/"
            wm_dir = os.path.join(BASE, parent, 'working_memory', '_' + sub)
        wm_exists = os.path.isdir(wm_dir)
        wm_count = len([f for f in os.listdir(wm_dir) if not f.startswith('_')]) if wm_exists else 0
        wm_html = f"<code>{wm_path}</code> ({wm_count} Dateien)" if wm_exists else f"<code>{wm_path}</code> <span style='color:#666;'>(leer)</span>"

        if r["kind"] == "Parent":
            kind_html = "<span class='tag'>Parent</span>"
            if r["has_own_memory"]:
                mem_html = f"<code>{esc_name}/memory/</code> ({r['file_count']} Dateien)"
            else:
                mem_html = "<span class='status-bad'>(kein memory/)</span>"
            access_html = "—"
            subs_html = ", ".join(_html_lib.escape(s) for s in r["sub_agents"]) if r["sub_agents"] else "—"
        else:
            kind_html = "<span class='tag tag-sub'>Sub-Agent</span>"
            mem_html = "—"
            access_html = (f"<code>{_html_lib.escape(parent)}/memory/</code>"
                           if parent else "(kein Parent)")
            subs_html = "—"
        tr.append(
            f"<tr><td><strong>{esc_name}</strong></td>"
            f"<td>{kind_html}</td>"
            f"<td>{mem_html}</td>"
            f"<td>{wm_html}</td>"
            f"<td>{access_html}</td>"
            f"<td>{subs_html}</td></tr>"
        )
    table = ("<table><tr><th>Agent</th><th>Typ</th><th>Eigenes Memory</th>"
             "<th>Working Memory</th>"
             "<th>Zugriff auf Parent-Memory</th><th>Sub-Agents</th></tr>"
             + "".join(tr) + "</table>")

    if orphans:
        orph_rows = "".join(f"<tr><td><code>{_html_lib.escape(d)}/memory/</code></td></tr>" for d in orphans)
        orph_html = ("<h2>Verwaiste Memory-Ordner</h2>"
                     "<p>Memory-Ordner ohne zugehoerige Agent-Definition (Datei in <code>config/agents/</code>):</p>"
                     "<table><tr><th>Pfad</th></tr>" + orph_rows + "</table>")
    else:
        orph_html = "<h2>Verwaiste Memory-Ordner</h2><p>Keine — alle Memory-Ordner haben einen Agent.</p>"

    n_parents = sum(1 for r in rows if r["kind"] == "Parent")
    n_subs = sum(1 for r in rows if r["kind"] == "Sub-Agent")
    # Shared Data Sources
    # Ordner-basierte Quellen (eigener Pfad)
    folder_sources = [
        ("webclips", "Webclips", os.path.join(_ADMIN_DATALAKE, "webclips")),
        ("calendar", "Kalender", os.path.join(_ADMIN_DATALAKE, "calendar")),
        ("whatsapp (privat)", "WhatsApp Chats", os.path.join(_ADMIN_DATALAKE, "privat", "memory", "whatsapp")),
    ]
    # Prefix-basierte Quellen (liegen verstreut in <agent>/memory/<prefix>_*.txt)
    def _count_prefix(prefix):
        total = 0
        agents_with = []
        try:
            for entry in os.listdir(_ADMIN_DATALAKE):
                mem = os.path.join(_ADMIN_DATALAKE, entry, "memory")
                if not os.path.isdir(mem):
                    continue
                try:
                    matches = [f for f in os.listdir(mem) if f.startswith(prefix) and f.endswith('.txt')]
                except OSError:
                    continue
                if matches:
                    total += len(matches)
                    agents_with.append(f"{entry} ({len(matches)})")
        except OSError:
            pass
        return total, agents_with

    kchat_total, kchat_agents = _count_prefix("kchat_")
    slack_total, slack_agents = _count_prefix("slack_")

    # E-Mails: IN_/OUT_ ueber alle Agent-Memorys aggregieren (nicht den
    # Staging-Ordner email_inbox/, der fast immer leer ist, weil der
    # Email-Watcher verarbeitet und wegraeumt).
    def _count_emails():
        total = 0
        per_agent = []
        try:
            for entry in os.listdir(_ADMIN_DATALAKE):
                mem = os.path.join(_ADMIN_DATALAKE, entry, "memory")
                if not os.path.isdir(mem):
                    continue
                try:
                    matches = [f for f in os.listdir(mem)
                               if f.endswith('.txt') and ('_IN_' in f or '_OUT_' in f)]
                except OSError:
                    continue
                if matches:
                    total += len(matches)
                    per_agent.append(f"{entry} ({len(matches)})")
        except OSError:
            pass
        return total, per_agent

    email_total, email_agents = _count_emails()
    staging_path = os.path.join(_ADMIN_DATALAKE, "email_inbox")
    staging_count = 0
    try:
        if os.path.isdir(staging_path):
            staging_count = sum(1 for x in os.listdir(staging_path) if x != 'processed')
    except OSError:
        pass

    shared_html = "<h2>Shared Data Sources</h2><table><tr><th>Quelle</th><th>Pfad</th><th>Status</th></tr>"
    # E-Mail-Zeile: Aggregat ueber alle Agent-Memorys
    email_details = ", ".join(email_agents)
    email_status = (
        f"<span class='status-ok'>● {email_total} Dateien</span>"
        f" <span style='color:#888;font-size:11px'>({_html_lib.escape(email_details)})</span>"
        f" <span style='color:#888;font-size:11px'>· Staging: {staging_count}</span>"
    ) if email_total else (
        "<span class='status-bad'>○ keine Dateien</span>"
    )
    shared_html += (f"<tr><td><strong>E-Mail Archive (IN_/OUT_)</strong></td>"
                    f"<td><code>&lt;agent&gt;/memory/*_IN_*.txt, *_OUT_*.txt</code></td>"
                    f"<td>{email_status}</td></tr>")

    for sid, slabel, spath in folder_sources:
        exists = os.path.isdir(spath)
        fcount = len(os.listdir(spath)) if exists else 0
        status = f"<span class='status-ok'>● {fcount} Dateien</span>" if exists else "<span class='status-bad'>○ nicht gefunden</span>"
        shared_html += f"<tr><td><strong>{_html_lib.escape(slabel)}</strong></td><td><code>{_html_lib.escape(sid)}/</code></td><td>{status}</td></tr>"

    # kChat-Log checken: Token-Fehler im UI sichtbar machen
    kchat_hint = ""
    try:
        with open("/tmp/kchat_watcher.log", "r", encoding="utf-8", errors="ignore") as _kl:
            _kl.seek(0, 2)
            _sz = _kl.tell()
            _kl.seek(max(0, _sz - 2000))
            _tail = _kl.read()
        if "401" in _tail or "ungueltig" in _tail.lower() or "unauthorized" in _tail.lower():
            kchat_hint = (" <span style='color:#d09090;font-size:11px'>"
                          "⚠ Watcher-Log zeigt Token-401: neuen Token in "
                          "<code>config/models.json</code> → <code>kchat.auth_token</code> eintragen</span>")
    except OSError:
        pass

    for sid, slabel, prefix, total, agents_list, hint in [
        ("kchat_*.txt pro Agent-Memory", "kChat Messages", "kchat_", kchat_total, kchat_agents, kchat_hint),
        ("slack_*.txt pro Agent-Memory", "Slack Messages", "slack_", slack_total, slack_agents, ""),
    ]:
        if total > 0:
            details = ", ".join(agents_list)
            status = (f"<span class='status-ok'>● {total} Dateien</span>"
                      f" <span style='color:#888;font-size:11px'>({_html_lib.escape(details)})</span>")
        else:
            status = "<span class='status-bad'>○ keine Dateien</span>"
        shared_html += (f"<tr><td><strong>{_html_lib.escape(slabel)}</strong></td>"
                        f"<td><code>{_html_lib.escape(sid)}</code></td>"
                        f"<td>{status}{hint}</td></tr>")
    shared_html += "</table>"
    body = (
        "<h1>🔐 Memory-Berechtigungen</h1>"
        "<div class='subtitle'>Wer hat Zugriff auf welches Memory? Datenbasis: "
        f"<code>{_ADMIN_DATALAKE}/config/agents/</code></div>"
        f"<div><span class='metric'><strong>{len(rows)}</strong> Agents</span>"
        f"<span class='metric'><strong>{n_parents}</strong> Parents</span>"
        f"<span class='metric'><strong>{n_subs}</strong> Sub-Agents</span>"
        f"<span class='metric'><strong>{len(orphans)}</strong> verwaist</span></div>"
        "<h2>Agents</h2>"
        f"{table}"
        f"{shared_html}"
        f"{orph_html}"
    )
    return _admin_layout("Permissions", body)


@app.route('/api/docs', methods=['GET'])
def api_docs():
    """Auto-generated API documentation page."""
    import re as _re_docs
    import inspect

    # Parse latest date from changelog
    changelog_path = os.path.join(os.path.expanduser("~"), "AssistantDev", "changelog.md")
    last_update = "unbekannt"
    try:
        with open(changelog_path, "r") as cf:
            for line in cf:
                m = _re_docs.search(r'\d{4}-\d{2}-\d{2}', line)
                if m:
                    last_update = m.group(0)
                    break
    except Exception:
        pass

    # Collect all routes
    rows = []
    for rule in sorted(app.url_map.iter_rules(), key=lambda r: r.rule):
        if rule.endpoint == "static":
            continue
        methods = ", ".join(sorted(rule.methods - {"OPTIONS", "HEAD"}))
        desc = ""
        view_fn = app.view_functions.get(rule.endpoint)
        if view_fn and view_fn.__doc__:
            desc = _html_lib.escape(view_fn.__doc__.strip().split("\n")[0])
        rows.append(
            f"<tr><td><code>{methods}</code></td>"
            f"<td><code>{_html_lib.escape(rule.rule)}</code></td>"
            f"<td>{desc}</td></tr>"
        )

    table = (
        "<table><tr><th>Method</th><th>Pfad</th><th>Beschreibung</th></tr>"
        + "\n".join(rows) + "</table>"
    )
    body = (
        '<a href="/" class="back">\u2190 Dashboard</a>'
        "<h1>AssistantDev API Documentation</h1>"
        f'<div class="subtitle">Letztes Update: {last_update} \u2014 '
        f"{len(rows)} Routen registriert</div>"
        f"{table}"
        '<hr><p style="color:#666;font-size:11px">AssistantDev \u2014 auto-generated</p>'
    )
    return _admin_layout("API Docs", body)

# ── ENDE ADMIN-BEREICH ────────────────────────────────────────────────────────


# ── MEMORY MANAGEMENT UI ─────────────────────────────────────────────────────

@app.route('/api/memory/list/<agent>', methods=['GET'])
def api_memory_list(agent):
    """List all memory files for an agent with metadata."""
    speicher = get_agent_speicher(agent)
    if not speicher:
        return jsonify([])
    memory_dir = os.path.join(speicher, 'memory')
    if not os.path.isdir(memory_dir):
        return jsonify([])
    result = []
    for fname in sorted(os.listdir(memory_dir)):
        fpath = os.path.join(memory_dir, fname)
        if not os.path.isfile(fpath):
            continue
        try:
            stat = os.stat(fpath)
            preview = ''
            try:
                with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                    preview = f.read(200)
            except Exception:
                pass
            result.append({
                'file': fname,
                'size': stat.st_size,
                'mtime': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M'),
                'preview': preview,
            })
        except Exception:
            continue
    return jsonify(result)


@app.route('/memory')
def memory_page():
    """Dedicated Memory Management UI."""
    return _admin_layout("Memory Management", _MEMORY_PAGE_HTML)


_MEMORY_PAGE_HTML = """
<a href="/" class="back">&larr; Dashboard</a>
<h1>Memory Management</h1>
<p class="subtitle">Working Memory und Memory-Files durchsuchen und verwalten</p>

<div style="margin-bottom:20px">
  <label style="color:#9ec5fe;font-size:13px;font-weight:600">Agent waehlen:</label>
  <select id="mm-agent-select" onchange="mmLoadAgent()" style="background:#22224a;color:#e0e0e0;border:1px solid #334;border-radius:6px;padding:6px 12px;font-size:14px;margin-left:8px;min-width:200px">
    <option value="">-- Agent waehlen --</option>
  </select>
</div>

<div id="mm-hint" style="color:#888;font-size:14px;margin:40px 0;text-align:center">Bitte zuerst einen Agenten waehlen.</div>

<div id="mm-content" style="display:none">
  <!-- Search bar -->
  <div style="margin-bottom:18px;display:flex;gap:10px">
    <input id="mm-search" type="text" placeholder="Dateisuche (min 2 Zeichen)..." oninput="mmFileSearch()" style="flex:1;background:#111;color:#e0e0e0;border:1px solid #334;border-radius:6px;padding:8px 12px;font-size:13px" />
    <button onclick="mmDeepSearch()" style="background:#4a8aca;color:#fff;border:none;border-radius:6px;padding:8px 16px;font-size:13px;cursor:pointer">Volltextsuche</button>
  </div>
  <div id="mm-search-results" style="margin-bottom:14px"></div>

  <!-- Working Memory -->
  <h2>Working Memory</h2>
  <div id="mm-working" style="margin-bottom:24px"><em style="color:#666">Lade...</em></div>

  <!-- All Memory Files -->
  <h2>Alle Memory-Files</h2>
  <div id="mm-files"><em style="color:#666">Lade...</em></div>
</div>

<!-- Deep search modal -->
<div id="mm-deep-modal" style="display:none;position:fixed;top:0;left:0;width:100%;height:100%;background:rgba(0,0,0,.7);z-index:999;justify-content:center;align-items:center">
  <div style="background:#1a1a2e;border:1px solid #334;border-radius:12px;padding:24px;width:90%;max-width:700px;max-height:80vh;overflow-y:auto">
    <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:14px">
      <h2 style="margin:0;color:#f0c060">Volltextsuche</h2>
      <button onclick="mmCloseDeep()" style="background:none;border:none;color:#888;font-size:20px;cursor:pointer">&times;</button>
    </div>
    <input id="mm-deep-query" type="text" placeholder="Suchbegriff..." style="width:100%;background:#111;color:#e0e0e0;border:1px solid #334;border-radius:6px;padding:8px 12px;font-size:14px;margin-bottom:10px" />
    <button onclick="mmRunDeep()" style="background:#4a8aca;color:#fff;border:none;border-radius:6px;padding:8px 16px;font-size:13px;cursor:pointer;margin-bottom:14px">Suchen</button>
    <div id="mm-deep-results"></div>
  </div>
</div>

<!-- File preview modal -->
<div id="mm-preview-modal" style="display:none;position:fixed;top:0;left:0;width:100%;height:100%;background:rgba(0,0,0,.7);z-index:999;justify-content:center;align-items:center">
  <div style="background:#1a1a2e;border:1px solid #334;border-radius:12px;padding:24px;width:90%;max-width:700px;max-height:80vh;overflow-y:auto">
    <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:14px">
      <h2 id="mm-preview-title" style="margin:0;color:#f0c060"></h2>
      <button onclick="mmClosePreview()" style="background:none;border:none;color:#888;font-size:20px;cursor:pointer">&times;</button>
    </div>
    <pre id="mm-preview-content" style="white-space:pre-wrap;word-break:break-word;max-height:60vh;overflow-y:auto"></pre>
  </div>
</div>

<script>
var mmAgent = '';

function mmLoadAgents() {
  fetch('/agents').then(function(r){return r.json()}).then(function(agents) {
    var sel = document.getElementById('mm-agent-select');
    agents.forEach(function(a) {
      var opt = document.createElement('option');
      opt.value = a.name;
      opt.textContent = a.label;
      sel.appendChild(opt);
      if (a.subagents) {
        a.subagents.forEach(function(s) {
          var sopt = document.createElement('option');
          sopt.value = s.name;
          sopt.textContent = '  \\u2514 ' + s.label;
          sel.appendChild(sopt);
        });
      }
    });
    var params = new URLSearchParams(window.location.search);
    if (params.get('agent')) {
      sel.value = params.get('agent');
      mmLoadAgent();
    }
  });
}

function mmLoadAgent() {
  var sel = document.getElementById('mm-agent-select');
  mmAgent = sel.value;
  if (!mmAgent) {
    document.getElementById('mm-hint').style.display = 'block';
    document.getElementById('mm-content').style.display = 'none';
    return;
  }
  document.getElementById('mm-hint').style.display = 'none';
  document.getElementById('mm-content').style.display = 'block';
  mmLoadWorking();
  mmLoadFiles();
}

function mmLoadWorking() {
  document.getElementById('mm-working').innerHTML = '<em style="color:#666">Lade...</em>';
  fetch('/api/working-memory/' + encodeURIComponent(mmAgent), {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({action: 'list'})
  }).then(function(r){return r.json()}).then(function(data) {
    var files = (data.manifest && data.manifest.files) || [];
    if (files.length === 0) {
      document.getElementById('mm-working').innerHTML = '<em style="color:#666">Keine Working-Memory-Dateien geladen.</em>';
      return;
    }
    var html = '<table><tr><th>Datei</th><th>Prioritaet</th><th>Beschreibung</th><th>Hinzugefuegt</th><th>Aktion</th></tr>';
    files.forEach(function(f) {
      html += '<tr><td><code>' + mmEsc(f.filename) + '</code></td>';
      html += '<td>' + (f.priority || '-') + '</td>';
      html += '<td>' + mmEsc(f.description || '-') + '</td>';
      html += '<td>' + mmEsc(f.added || '-') + '</td>';
      html += '<td><button onclick="mmRemoveWM(\\'' + mmEsc(f.filename).replace(/'/g, "\\\\'") + '\\')" style="background:#d07070;color:#fff;border:none;border-radius:4px;padding:3px 10px;cursor:pointer;font-size:12px">Entfernen</button></td></tr>';
    });
    html += '</table>';
    document.getElementById('mm-working').innerHTML = html;
  }).catch(function() {
    document.getElementById('mm-working').innerHTML = '<em style="color:#d07070">Fehler beim Laden.</em>';
  });
}

function mmRemoveWM(filename) {
  fetch('/api/working-memory/' + encodeURIComponent(mmAgent), {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({action: 'remove', filename: filename})
  }).then(function(){mmLoadWorking()});
}

function mmLoadFiles() {
  document.getElementById('mm-files').innerHTML = '<em style="color:#666">Lade...</em>';
  fetch('/api/memory/list/' + encodeURIComponent(mmAgent)).then(function(r){return r.json()}).then(function(files) {
    if (files.length === 0) {
      document.getElementById('mm-files').innerHTML = '<em style="color:#666">Keine Memory-Files gefunden.</em>';
      return;
    }
    var html = '<table><tr><th>Datei</th><th>Groesse</th><th>Geaendert</th><th>Vorschau</th></tr>';
    files.forEach(function(f) {
      html += '<tr><td><code>' + mmEsc(f.file) + '</code></td>';
      html += '<td>' + mmFormatSize(f.size) + '</td>';
      html += '<td>' + mmEsc(f.mtime) + '</td>';
      html += '<td><button onclick="mmShowPreview(\\'' + mmEsc(f.file).replace(/'/g, "\\\\'") + '\\')" style="background:#4a8aca;color:#fff;border:none;border-radius:4px;padding:3px 10px;cursor:pointer;font-size:12px">Anzeigen</button></td></tr>';
    });
    html += '</table>';
    html += '<p class="subtitle">' + files.length + ' Dateien</p>';
    document.getElementById('mm-files').innerHTML = html;
  }).catch(function() {
    document.getElementById('mm-files').innerHTML = '<em style="color:#d07070">Fehler beim Laden.</em>';
  });
}

function mmShowPreview(filename) {
  document.getElementById('mm-preview-title').textContent = filename;
  document.getElementById('mm-preview-content').textContent = 'Lade...';
  document.getElementById('mm-preview-modal').style.display = 'flex';
  fetch('/api/memory/list/' + encodeURIComponent(mmAgent)).then(function(r){return r.json()}).then(function(files) {
    var found = files.find(function(f){return f.file === filename});
    document.getElementById('mm-preview-content').textContent = found ? found.preview : 'Datei nicht gefunden.';
  });
}
function mmClosePreview() { document.getElementById('mm-preview-modal').style.display = 'none'; }

function mmFileSearch() {
  var q = document.getElementById('mm-search').value.trim();
  var el = document.getElementById('mm-search-results');
  if (q.length < 2) { el.innerHTML = ''; return; }
  fetch('/api/memory-files-search?agent=' + encodeURIComponent(mmAgent) + '&q=' + encodeURIComponent(q))
    .then(function(r){return r.json()}).then(function(results) {
      if (results.length === 0) { el.innerHTML = '<p style="color:#888;font-size:12px">Keine Treffer.</p>'; return; }
      var html = '<div style="background:#22224a;border:1px solid #334;border-radius:6px;padding:8px">';
      results.forEach(function(r) {
        html += '<div style="padding:4px 0;border-bottom:1px solid #2c2c4a"><code style="color:#9ec5fe">' + mmEsc(r.filename) + '</code> <span style="color:#888;font-size:11px">' + mmEsc(r.snippet) + '</span></div>';
      });
      html += '</div>';
      el.innerHTML = html;
    });
}

function mmDeepSearch() {
  document.getElementById('mm-deep-query').value = '';
  document.getElementById('mm-deep-results').innerHTML = '';
  document.getElementById('mm-deep-modal').style.display = 'flex';
  document.getElementById('mm-deep-query').focus();
}
function mmCloseDeep() { document.getElementById('mm-deep-modal').style.display = 'none'; }

function mmRunDeep() {
  var q = document.getElementById('mm-deep-query').value.trim();
  if (!q) return;
  var el = document.getElementById('mm-deep-results');
  el.innerHTML = '<em style="color:#666">Suche...</em>';
  fetch('/api/memory/search', {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({agent: mmAgent, query: q, max_results: 20})
  }).then(function(r){return r.json()}).then(function(data) {
    var results = data.results || [];
    if (results.length === 0) { el.innerHTML = '<p style="color:#888">Keine Ergebnisse.</p>'; return; }
    var html = '';
    results.forEach(function(r) {
      html += '<div style="background:#22224a;border:1px solid #334;border-radius:6px;padding:10px;margin-bottom:8px">';
      html += '<div style="font-weight:600;color:#9ec5fe;margin-bottom:4px">' + mmEsc(r.name || r.filename || 'Unbekannt') + '</div>';
      html += '<div style="color:#d0d0d0;font-size:12px;white-space:pre-wrap">' + mmEsc(r.preview || r.snippet || '') + '</div>';
      html += '</div>';
    });
    el.innerHTML = html;
  }).catch(function() {
    el.innerHTML = '<em style="color:#d07070">Fehler bei der Suche.</em>';
  });
}

function mmEsc(s) {
  if (!s) return '';
  var d = document.createElement('div');
  d.appendChild(document.createTextNode(s));
  return d.innerHTML;
}

function mmFormatSize(bytes) {
  if (bytes < 1024) return bytes + ' B';
  if (bytes < 1048576) return (bytes / 1024).toFixed(1) + ' KB';
  return (bytes / 1048576).toFixed(1) + ' MB';
}

mmLoadAgents();
</script>
"""

# ── ENDE MEMORY MANAGEMENT UI ────────────────────────────────────────────────


# ── MESSAGE DASHBOARD (Kanban-Posteingang) ──────────────────────────────────
# Aggregiert Nachrichten aus dem Datalake in Kanban-Spalten (eine pro Quelle).
# Wiederverwendet die Parser-Logik aus src/message_dashboard.py (PyQt6-App),
# exponiert sie aber als Web-API + HTML-Dashboard unter /messages und /api/messages/*.

import hashlib as _msgd_hashlib
import re as _msgd_re
import threading as _msgd_threading
import time as _msgd_time

_MSG_STATE_FILE = os.path.expanduser("~/.message_dashboard_state.json")
_MSG_PREVIEW_LEN = 200
_MSG_INBOX_WINDOW_DAYS = 90
_MSG_MAX_FILES_PER_SOURCE = 500
_MSG_PARSE_READ_BYTES = 8192
_MSG_CACHE_TTL_SECONDS = 60
_MSG_OWN_EMAILS = {
    "londoncityfox@gmail.com", "moritz.cremer@me.com", "moritz.cremer@icloud.com",
    "moritz@demoscapital.co", "moritz@vegatechnology.com.br",
    "moritz.cremer@signicat.com", "moritz.cremer@trustedcarrier.net",
    "moritz@casiopayaconsulting.io", "moritz@tangerina.com",
}

# Quellen — eine pro Spalte im Dashboard
_MSG_SOURCES = [
    {"key": "email_signicat",       "label": "Signicat",       "agent": "signicat",       "icon": "\U0001F4E7", "type": "email"},
    {"key": "email_privat",         "label": "Privat",         "agent": "privat",         "icon": "\U0001F4E7", "type": "email"},
    {"key": "email_trustedcarrier", "label": "TrustedCarrier", "agent": "trustedcarrier", "icon": "\U0001F4E7", "type": "email"},
    {"key": "email_standard",       "label": "Standard",       "agent": "standard",       "icon": "\U0001F4E7", "type": "email"},
    {"key": "email_systemward",     "label": "System Ward",    "agent": "system ward",    "icon": "\U0001F4E7", "type": "email"},
    {"key": "whatsapp",             "label": "WhatsApp",       "agent": None,             "icon": "\U0001F4F1", "type": "whatsapp"},
    {"key": "chat",                 "label": "Chat-Verlauf",   "agent": None,             "icon": "\U0001F4AC", "type": "chat"},
]

# Mapping Source -> empfohlener Agent fuer "Antworten"-Workflow
_MSG_SOURCE_TO_AGENT = {
    "email_signicat": "signicat",
    "email_privat": "privat",
    "email_trustedcarrier": "trustedcarrier",
    "email_standard": "standard",
    "email_systemward": "system ward",
    "whatsapp": "privat",
    "chat": "standard",
}

_MSG_EMAIL_FIELD_RE = _msgd_re.compile(
    r'^(Von|An|From|To|Betreff|Subject|Datum|Date|Kontakt|Richtung|Agent|Importiert|Quelle|Kanal|Cc|Kopie|Message-ID|Reply-To):\s*(.*)$',
    _msgd_re.IGNORECASE,
)
_MSG_EMAIL_ADDR_RE = _msgd_re.compile(r'[\w.+%-]+@[\w.-]+\.[a-zA-Z]{2,}')
_MSG_EMAIL_FILE_RE = _msgd_re.compile(r'^\d{4}-\d{2}-\d{2}_\d{2}-\d{2}-\d{2}_(IN|OUT)_.*\.txt$')
_MSG_LEGACY_EMAIL_RE = _msgd_re.compile(r'^email_.*\.txt$', _msgd_re.IGNORECASE)
_MSG_WHATSAPP_FILE_RE = _msgd_re.compile(r'^whatsapp_.*\.txt$', _msgd_re.IGNORECASE)
_MSG_KONV_FILE_RE = _msgd_re.compile(r'^konversation_.*\.txt$', _msgd_re.IGNORECASE)

_MSG_CACHE = {"messages": None, "ts": 0.0, "lock": _msgd_threading.Lock()}
_MSG_STATE_LOCK = _msgd_threading.Lock()


def _msg_hash_path(path):
    return _msgd_hashlib.sha1(path.encode("utf-8")).hexdigest()[:16]


def _msg_load_state():
    try:
        with open(_MSG_STATE_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
            if not isinstance(data, dict):
                return {"read_messages": []}
            data.setdefault("read_messages", [])
            return data
    except Exception:
        return {"read_messages": []}


def _msg_save_state(state):
    try:
        with open(_MSG_STATE_FILE, "w", encoding="utf-8") as f:
            json.dump(state, f, indent=2)
    except Exception as e:
        print(f"[MSGD] state save error: {e}", file=sys.stderr)


def _msg_strip_separator(text):
    sep = "\u2500" * 10
    idx = text.find(sep)
    if idx != -1:
        eol = text.find("\n", idx)
        if eol != -1:
            return text[eol+1:].lstrip()
    for marker in ["\n---\n", "\n----\n", "\n-----\n"]:
        idx = text.find(marker)
        if idx != -1:
            return text[idx+len(marker):].lstrip()
    return text


def _msg_parse_date(s):
    if not s:
        return None
    s = s.strip()
    try:
        from email.utils import parsedate_to_datetime
        dt = parsedate_to_datetime(s)
        if dt:
            if dt.tzinfo:
                dt = dt.astimezone().replace(tzinfo=None)
            return dt
    except Exception:
        pass
    try:
        return datetime.datetime.fromisoformat(
            s.replace("Z", "+00:00").split("+")[0].split(".")[0]
        )
    except Exception:
        pass
    try:
        return datetime.datetime.strptime(s[:19], "%Y-%m-%d_%H-%M-%S")
    except Exception:
        pass
    return None


def _msg_normalize_email_content(fpath, fname, source_key, agent_name, only_header=True):
    """Liest eine .txt Email-Datei und gibt normalisiertes Dict zurueck.
    only_header=True liest nur die ersten 8KB (Header + Preview), sonst bis zu 1MB."""
    try:
        mode_size = _MSG_PARSE_READ_BYTES if only_header else 1024 * 1024
        with open(fpath, "rb") as f:
            raw_bytes = f.read(mode_size)
        if not raw_bytes.strip():
            return None
        raw = raw_bytes.decode("utf-8", errors="replace")
    except Exception:
        return None

    fields = {}
    lines = raw.splitlines()
    header_end_idx = 0
    for i, line in enumerate(lines):
        if "\u2500\u2500\u2500" in line:
            header_end_idx = i + 1
            break
        if line.strip() == "" and fields:
            header_end_idx = i + 1
            break
        m = _MSG_EMAIL_FIELD_RE.match(line)
        if m:
            key = m.group(1).lower()
            value = m.group(2).strip()
            if key not in fields:
                fields[key] = value
        elif i > 30 and not fields:
            break

    body = (
        "\n".join(lines[header_end_idx:]).strip() if header_end_idx
        else _msg_strip_separator(raw).strip()
    )

    from_raw = fields.get("von", fields.get("from", ""))
    to_raw = fields.get("an", fields.get("to", fields.get("kanal", "")))
    subject = fields.get("betreff", fields.get("subject", "")).strip()
    if not subject:
        first_body_line = next((l for l in body.splitlines() if l.strip()), "").strip()
        subject = (first_body_line[:100] if first_body_line else fname)[:200]

    sender_email = ""
    m_addr = _MSG_EMAIL_ADDR_RE.search(from_raw)
    if m_addr:
        sender_email = m_addr.group(0).lower()

    sender_name = from_raw
    if "<" in from_raw and ">" in from_raw:
        sender_name = from_raw[:from_raw.index("<")].strip().strip('"')
    elif sender_email and sender_email in from_raw:
        sender_name = from_raw.replace(sender_email, "").strip().strip("<>").strip()
    if not sender_name:
        sender_name = sender_email or "(unbekannt)"

    direction = (fields.get("richtung", "") or "").strip().upper()
    if direction == "OUT":
        return None
    if sender_email and sender_email in _MSG_OWN_EMAILS:
        return None

    date_str = fields.get("datum", fields.get("date", ""))
    dt = _msg_parse_date(date_str)
    if not dt:
        m2 = _msgd_re.match(r"(\d{4}-\d{2}-\d{2})_(\d{2})-(\d{2})-(\d{2})_", fname)
        if m2:
            try:
                dt = datetime.datetime.strptime(
                    f"{m2.group(1)}_{m2.group(2)}-{m2.group(3)}-{m2.group(4)}",
                    "%Y-%m-%d_%H-%M-%S",
                )
            except Exception:
                pass
    if not dt:
        try:
            dt = datetime.datetime.fromtimestamp(os.path.getmtime(fpath))
        except Exception:
            dt = datetime.datetime.now()

    body_clean = body.strip()
    preview = " ".join(body_clean[:_MSG_PREVIEW_LEN].split())
    if len(body_clean) > _MSG_PREVIEW_LEN:
        preview = preview + "..."

    has_attachments = bool(_msgd_re.search(r"(anhang|attachment|\[cid:|\[Bild\]|\.pdf|\.docx|\.xlsx|\.png|\.jpg)", body_clean[:2000], _msgd_re.IGNORECASE))

    return {
        "id": _msg_hash_path(fpath),
        "source": source_key,
        "source_agent": agent_name,
        "sender_name": sender_name[:150],
        "sender_address": sender_email,
        "to": to_raw,
        "subject": subject[:200],
        "preview": preview,
        "full_content": body_clean if not only_header else "",
        "timestamp": dt.isoformat(),
        "timestamp_epoch": dt.timestamp(),
        "read": False,
        "has_attachments": has_attachments,
        "attachments": [],
        "raw_file_path": fpath,
        "message_id": fields.get("message-id", ""),
        "type": "email",
    }


def _msg_normalize_whatsapp_file(fpath, fname, source_key, agent_name, only_header=True):
    try:
        mode_size = _MSG_PARSE_READ_BYTES if only_header else 1024 * 512
        with open(fpath, "rb") as f:
            raw_bytes = f.read(mode_size)
        if not raw_bytes.strip():
            return None
        raw = raw_bytes.decode("utf-8", errors="replace")
    except Exception:
        return None

    # Header-Block parsen
    contact = ""
    period = ""
    msg_count = ""
    lines = raw.splitlines()
    body_start = 0
    for i, line in enumerate(lines[:10]):
        if line.lower().startswith("kontakt:"):
            contact = line.split(":", 1)[1].strip()
        elif line.lower().startswith("zeitraum:"):
            period = line.split(":", 1)[1].strip()
        elif line.lower().startswith("nachrichten:"):
            msg_count = line.split(":", 1)[1].strip()
        elif line.strip() == "" or line.startswith("["):
            body_start = i
            break

    body = "\n".join(lines[body_start:]).strip()

    # Letzte eingehende Nachricht finden (nicht "Ich:")
    last_in_line = ""
    last_ts = None
    for line in reversed(body.splitlines()):
        if line.startswith("[") and "]" in line:
            try:
                ts_str = line[1:line.index("]")]
                rest = line[line.index("]")+1:].strip()
                if rest.startswith("Ich:"):
                    continue
                last_in_line = rest
                try:
                    last_ts = datetime.datetime.strptime(ts_str, "%Y-%m-%d %H:%M")
                except Exception:
                    try:
                        last_ts = datetime.datetime.strptime(ts_str[:10], "%Y-%m-%d")
                    except Exception:
                        pass
                break
            except Exception:
                continue

    if not last_ts:
        try:
            last_ts = datetime.datetime.fromtimestamp(os.path.getmtime(fpath))
        except Exception:
            last_ts = datetime.datetime.now()

    preview = " ".join(last_in_line[:_MSG_PREVIEW_LEN].split()) if last_in_line else (body[:_MSG_PREVIEW_LEN].strip())
    if not preview:
        preview = f"{msg_count} Nachrichten" if msg_count else "(WhatsApp Chat)"

    subject = f"{contact or fname}"
    if period:
        subject += f" \u00b7 {period}"

    return {
        "id": _msg_hash_path(fpath),
        "source": source_key,
        "source_agent": agent_name,
        "sender_name": contact or fname,
        "sender_address": "",
        "to": "",
        "subject": subject[:200],
        "preview": preview,
        "full_content": body if not only_header else "",
        "timestamp": last_ts.isoformat(),
        "timestamp_epoch": last_ts.timestamp(),
        "read": False,
        "has_attachments": "[Bild]" in body or "[Video]" in body or "[Sprachnachricht]" in body,
        "attachments": [],
        "raw_file_path": fpath,
        "message_id": "",
        "type": "whatsapp",
    }


def _msg_normalize_chat_file(fpath, fname, source_key, agent_name, only_header=True):
    try:
        mode_size = _MSG_PARSE_READ_BYTES if only_header else 1024 * 512
        with open(fpath, "rb") as f:
            raw_bytes = f.read(mode_size)
        if not raw_bytes.strip():
            return None
        raw = raw_bytes.decode("utf-8", errors="replace")
    except Exception:
        return None

    preview = " ".join(raw[:_MSG_PREVIEW_LEN].split())
    try:
        dt = datetime.datetime.fromtimestamp(os.path.getmtime(fpath))
    except Exception:
        dt = datetime.datetime.now()

    # Versuche Agent/Datum aus Dateiname konversation_YYYY-MM-DD_HH-MM-SS_agent.txt
    agent_from_name = agent_name
    m = _msgd_re.match(r"konversation_(\d{4}-\d{2}-\d{2})_(\d{2})-(\d{2})-(\d{2})_(.+)\.txt$", fname, _msgd_re.IGNORECASE)
    if m:
        try:
            dt = datetime.datetime.strptime(
                f"{m.group(1)}_{m.group(2)}-{m.group(3)}-{m.group(4)}",
                "%Y-%m-%d_%H-%M-%S",
            )
        except Exception:
            pass
        agent_from_name = m.group(5)

    return {
        "id": _msg_hash_path(fpath),
        "source": source_key,
        "source_agent": agent_from_name or agent_name,
        "sender_name": agent_from_name or "Agent",
        "sender_address": "",
        "to": "",
        "subject": fname.replace(".txt", ""),
        "preview": preview,
        "full_content": raw if not only_header else "",
        "timestamp": dt.isoformat(),
        "timestamp_epoch": dt.timestamp(),
        "read": False,
        "has_attachments": False,
        "attachments": [],
        "raw_file_path": fpath,
        "message_id": "",
        "type": "chat",
    }


def _msg_scan_email_source(source):
    agent = source["agent"]
    speicher = get_agent_speicher(agent)
    memdir = os.path.join(speicher, "memory")
    out = []
    if not os.path.isdir(memdir):
        return out
    cutoff = _msgd_time.time() - (_MSG_INBOX_WINDOW_DAYS * 86400)
    candidates = []
    try:
        with os.scandir(memdir) as it:
            for entry in it:
                if not entry.is_file(follow_symlinks=False):
                    continue
                fname = entry.name
                if not fname.endswith(".txt"):
                    continue
                if fname.startswith("konversation_"):
                    continue
                if fname.startswith("whatsapp_"):
                    continue
                if fname.startswith("kchat_"):
                    continue
                if not (
                    _MSG_EMAIL_FILE_RE.match(fname)
                    or _MSG_LEGACY_EMAIL_RE.match(fname)
                ):
                    continue
                try:
                    st = entry.stat(follow_symlinks=False)
                    mtime = st.st_mtime
                except Exception:
                    mtime = 0
                # iCloud Drive setzt mtime beim Sync oft neu — deshalb ist das
                # Datum im Dateinamen (YYYY-MM-DD_HH-MM-SS_...) die verlaessliche
                # Quelle fuer Inbox-Alter.
                fname_ts = None
                m = _msgd_re.match(r"(\d{4}-\d{2}-\d{2})", fname)
                if m:
                    try:
                        fname_ts = datetime.datetime.strptime(m.group(1), "%Y-%m-%d").timestamp()
                    except Exception:
                        fname_ts = None
                if fname_ts is not None:
                    if fname_ts < cutoff:
                        continue
                    sort_ts = fname_ts
                else:
                    if mtime < cutoff:
                        continue
                    sort_ts = mtime
                candidates.append((sort_ts, entry.path, fname))
    except Exception:
        return out
    candidates.sort(key=lambda x: x[0], reverse=True)
    candidates = candidates[:_MSG_MAX_FILES_PER_SOURCE]
    for _, fpath, fname in candidates:
        try:
            msg = _msg_normalize_email_content(fpath, fname, source["key"], agent, only_header=True)
        except Exception:
            continue
        if msg:
            out.append(msg)
    return out


def _msg_scan_whatsapp_source(source):
    out = []
    cutoff = _msgd_time.time() - (_MSG_INBOX_WINDOW_DAYS * 86400)
    candidates = []
    for agent in ["signicat", "privat", "trustedcarrier", "standard", "system ward"]:
        speicher = get_agent_speicher(agent)
        for subdir in [os.path.join(speicher, "memory", "whatsapp"), os.path.join(speicher, "memory")]:
            if not os.path.isdir(subdir):
                continue
            try:
                with os.scandir(subdir) as it:
                    for entry in it:
                        if not entry.is_file(follow_symlinks=False):
                            continue
                        fname = entry.name
                        if not _MSG_WHATSAPP_FILE_RE.match(fname):
                            continue
                        try:
                            mtime = entry.stat(follow_symlinks=False).st_mtime
                        except Exception:
                            mtime = 0
                        # Datum aus Dateinamen (letztes YYYY-MM-DD vor .txt) hat
                        # Vorrang vor mtime (iCloud-Sync-Artefakte).
                        fname_ts = None
                        dm = _msgd_re.search(r"(\d{4}-\d{2}-\d{2})\.txt$", fname)
                        if dm:
                            try:
                                fname_ts = datetime.datetime.strptime(dm.group(1), "%Y-%m-%d").timestamp()
                            except Exception:
                                fname_ts = None
                        if fname_ts is not None:
                            if fname_ts < cutoff:
                                continue
                            sort_ts = fname_ts
                        else:
                            if mtime < cutoff:
                                continue
                            sort_ts = mtime
                        candidates.append((sort_ts, entry.path, fname, agent))
            except Exception:
                continue
    candidates.sort(key=lambda x: x[0], reverse=True)
    candidates = candidates[:_MSG_MAX_FILES_PER_SOURCE]
    for _, fpath, fname, agent in candidates:
        try:
            msg = _msg_normalize_whatsapp_file(fpath, fname, source["key"], agent, only_header=True)
        except Exception:
            continue
        if msg:
            out.append(msg)
    return out


def _msg_scan_chat_source(source):
    out = []
    cutoff = _msgd_time.time() - (_MSG_INBOX_WINDOW_DAYS * 86400)
    candidates = []
    for agent in ["signicat", "privat", "trustedcarrier", "standard", "system ward"]:
        speicher = get_agent_speicher(agent)
        memdir = os.path.join(speicher, "memory")
        if not os.path.isdir(memdir):
            continue
        try:
            with os.scandir(memdir) as it:
                for entry in it:
                    if not entry.is_file(follow_symlinks=False):
                        continue
                    fname = entry.name
                    if not _MSG_KONV_FILE_RE.match(fname):
                        continue
                    try:
                        st = entry.stat(follow_symlinks=False)
                        mtime = st.st_mtime
                        size = st.st_size
                    except Exception:
                        mtime = 0
                        size = 0
                    if size < 80:
                        continue
                    fname_ts = None
                    dm = _msgd_re.match(r"konversation_(\d{4}-\d{2}-\d{2})", fname)
                    if dm:
                        try:
                            fname_ts = datetime.datetime.strptime(dm.group(1), "%Y-%m-%d").timestamp()
                        except Exception:
                            fname_ts = None
                    if fname_ts is not None:
                        if fname_ts < cutoff:
                            continue
                        sort_ts = fname_ts
                    else:
                        if mtime < cutoff:
                            continue
                        sort_ts = mtime
                    candidates.append((sort_ts, entry.path, fname, agent))
        except Exception:
            continue
    candidates.sort(key=lambda x: x[0], reverse=True)
    candidates = candidates[:_MSG_MAX_FILES_PER_SOURCE]
    for _, fpath, fname, agent in candidates:
        try:
            msg = _msg_normalize_chat_file(fpath, fname, source["key"], agent, only_header=True)
        except Exception:
            continue
        if msg:
            out.append(msg)
    return out


def _msg_scan_source(source):
    stype = source.get("type")
    if stype == "email":
        return _msg_scan_email_source(source)
    if stype == "whatsapp":
        return _msg_scan_whatsapp_source(source)
    if stype == "chat":
        return _msg_scan_chat_source(source)
    return []


def _msg_get_all(force=False):
    """Liefert alle Nachrichten mit read-state Annotation. Nutzt einen
    einfachen In-Memory-Cache mit TTL damit die Kanban-UI fluessig scrollt."""
    with _MSG_CACHE["lock"]:
        now = _msgd_time.time()
        if (not force and _MSG_CACHE["messages"] is not None
                and (now - _MSG_CACHE["ts"]) < _MSG_CACHE_TTL_SECONDS):
            msgs = _MSG_CACHE["messages"]
        else:
            msgs = []
            for src in _MSG_SOURCES:
                try:
                    msgs.extend(_msg_scan_source(src))
                except Exception as e:
                    print(f"[MSGD] scan error source={src['key']}: {e}", file=sys.stderr)
            _MSG_CACHE["messages"] = msgs
            _MSG_CACHE["ts"] = now
        state = _msg_load_state()
        read_set = set(state.get("read_messages", []))
        # read-state frisch pro Request applizieren, ohne Cache zu invalidieren
        annotated = []
        for m in msgs:
            m2 = dict(m)
            m2["read"] = m2["id"] in read_set
            annotated.append(m2)
        return annotated


def _msg_find_by_id(msg_id):
    """Findet Nachricht inklusive full_content (laedt Volltext nach)."""
    msgs = _msg_get_all()
    hit = next((m for m in msgs if m["id"] == msg_id), None)
    if not hit:
        return None
    # Volltext nachladen
    fpath = hit.get("raw_file_path", "")
    if not fpath or not os.path.exists(fpath):
        return hit
    fname = os.path.basename(fpath)
    mtype = hit.get("type", "email")
    src_key = hit.get("source", "")
    agent = hit.get("source_agent", "")
    try:
        if mtype == "email":
            full = _msg_normalize_email_content(fpath, fname, src_key, agent, only_header=False)
        elif mtype == "whatsapp":
            full = _msg_normalize_whatsapp_file(fpath, fname, src_key, agent, only_header=False)
        else:
            full = _msg_normalize_chat_file(fpath, fname, src_key, agent, only_header=False)
        if full:
            full["read"] = hit["read"]
            return full
    except Exception as e:
        print(f"[MSGD] full-load error: {e}", file=sys.stderr)
    return hit


@app.route("/api/messages/sources")
def api_messages_sources():
    """Liefert Liste aller Quellen inkl. Counts + Unread.
    Felder pro Quelle: key, label, icon, type, available, count, unread,
    recommended_agent."""
    try:
        messages = _msg_get_all()
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500
    by_src = {}
    for m in messages:
        by_src.setdefault(m["source"], []).append(m)
    out = []
    for src in _MSG_SOURCES:
        items = by_src.get(src["key"], [])
        count = len(items)
        unread = sum(1 for m in items if not m["read"])
        out.append({
            "key": src["key"],
            "label": src["label"],
            "icon": src["icon"],
            "type": src["type"],
            "available": count > 0,
            "count": count,
            "unread": unread,
            "recommended_agent": _MSG_SOURCE_TO_AGENT.get(src["key"]),
        })
    return jsonify({"ok": True, "sources": out})


@app.route("/api/messages")
def api_messages():
    """Liefert alle Nachrichten (ohne full_content).
    Optional: ?source=<key>&limit=<n>."""
    try:
        messages = _msg_get_all(force=(request.args.get("refresh") == "1"))
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500
    source_filter = request.args.get("source")
    if source_filter:
        messages = [m for m in messages if m["source"] == source_filter]
    try:
        limit = int(request.args.get("limit", "2000"))
    except Exception:
        limit = 2000
    # Sortierung: ungelesen zuerst (aelteste oben), dann gelesen (neueste oben)
    messages.sort(key=lambda m: (
        0 if not m["read"] else 1,
        m["timestamp_epoch"] if not m["read"] else -m["timestamp_epoch"],
    ))
    messages = messages[:limit]
    slim = [{k: v for k, v in m.items() if k != "full_content"} for m in messages]
    return jsonify({"ok": True, "messages": slim, "count": len(slim)})


@app.route("/api/messages/<msg_id>")
def api_messages_detail(msg_id):
    msg = _msg_find_by_id(msg_id)
    if not msg:
        return jsonify({"ok": False, "error": "not found"}), 404
    return jsonify({"ok": True, "message": msg})


@app.route("/api/messages/mark-read", methods=["POST"])
def api_messages_mark_read():
    data = request.get_json(silent=True) or {}
    mid = (data.get("message_id") or "").strip()
    read = bool(data.get("read", True))
    if not mid:
        return jsonify({"ok": False, "error": "message_id required"}), 400
    with _MSG_STATE_LOCK:
        state = _msg_load_state()
        read_list = state.get("read_messages", [])
        read_set = set(read_list)
        if read:
            read_set.add(mid)
        else:
            read_set.discard(mid)
        state["read_messages"] = sorted(read_set)
        _msg_save_state(state)
    return jsonify({"ok": True, "read": read, "message_id": mid})


@app.route("/messages")
def messages_page():
    resp = make_response(_MSG_DASHBOARD_HTML)
    resp.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    resp.headers["Content-Type"] = "text/html; charset=utf-8"
    return resp


_MSG_DASHBOARD_HTML = r"""<!DOCTYPE html>
<html lang="de">
<head>
<meta charset="UTF-8">
<title>📬 AssistantDev — Messages</title>
<link rel="icon" href="data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 32 32'%3E%3Crect width='32' height='32' rx='6' fill='%23111'/%3E%3Ctext x='16' y='24' text-anchor='middle' font-family='system-ui' font-weight='700' font-size='22' fill='%23f0c060'%3EM%3C/text%3E%3C/svg%3E">
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap" rel="stylesheet">
<style>
  * { margin:0; padding:0; box-sizing:border-box; }
  html, body { height:100%; }
  body { font-family:'Inter','Helvetica Neue',sans-serif; background:#111; color:#e0e0e0; display:flex; flex-direction:column; overflow:hidden; }
  #md-header { background:#0c0c0c; border-bottom:1px solid #2a2a2a; padding:10px 18px; display:flex; align-items:center; gap:14px; flex-shrink:0; }
  #md-header h1 { font-size:11px; font-weight:700; color:#aaa; letter-spacing:2.5px; text-transform:uppercase; }
  #md-search { flex:1; max-width:520px; background:#1a1a1a; border:1px solid #333; color:#e0e0e0; border-radius:8px; padding:8px 14px; font-size:13px; font-family:inherit; outline:none; }
  #md-search:focus { border-color:#f0c060; }
  .md-hdr-btn { background:none; border:1px solid #333; color:#aaa; padding:7px 14px; font-size:12px; border-radius:7px; cursor:pointer; font-family:inherit; }
  .md-hdr-btn:hover { border-color:#f0c060; color:#f0c060; }
  .md-hdr-stat { font-size:11px; color:#777; }
  .md-hdr-stat strong { color:#f0c060; font-weight:700; }
  #md-board { flex:1; display:flex; gap:10px; padding:14px; overflow-x:auto; overflow-y:hidden; background:#111; }
  .md-col { flex:0 0 320px; min-width:320px; max-width:360px; background:#161616; border:1px solid #262626; border-radius:10px; display:flex; flex-direction:column; overflow:hidden; }
  .md-col-hdr { padding:10px 12px 8px; border-bottom:1px solid #232323; background:#1a1a1a; }
  .md-col-title { display:flex; align-items:center; gap:8px; font-size:13px; font-weight:600; color:#f0f0f0; }
  .md-col-title .md-col-icon { font-size:15px; }
  .md-col-title .md-col-count { font-size:11px; color:#777; font-weight:500; margin-left:auto; }
  .md-col-title .md-col-unread { background:#f0c060; color:#111; font-size:10px; font-weight:700; padding:1px 7px; border-radius:10px; }
  .md-col-search { margin-top:7px; width:100%; background:#111; border:1px solid #2a2a2a; color:#d0d0d0; font-size:11px; padding:5px 9px; border-radius:6px; outline:none; font-family:inherit; }
  .md-col-search:focus { border-color:#3a3a3a; }
  .md-col-search::placeholder { color:#555; }
  .md-col-body { flex:1; overflow-y:auto; padding:6px; }
  .md-col-empty { padding:24px 12px; text-align:center; color:#555; font-size:11px; }
  .md-card { background:#1d1d1d; border:1px solid #242424; border-radius:8px; padding:9px 11px; margin-bottom:6px; cursor:pointer; transition:border-color .12s, background .12s; position:relative; }
  .md-card:hover { border-color:#3a3a3a; background:#222; }
  .md-card.unread { background:#1c1a14; border-color:#3a3020; }
  .md-card.unread:hover { border-color:#f0c060; }
  .md-card.highlight { animation:md-flash 1.6s ease-out 1; }
  @keyframes md-flash { 0%,30% { border-color:#f0c060; box-shadow:0 0 0 3px rgba(240,192,96,0.25); } 100% { box-shadow:none; } }
  .md-card-top { display:flex; align-items:center; gap:6px; margin-bottom:3px; }
  .md-dot { width:7px; height:7px; border-radius:50%; flex-shrink:0; background:#333; }
  .md-dot.unread { background:#f0c060; }
  .md-card-sender { flex:1; font-size:12px; color:#fafafa; overflow:hidden; white-space:nowrap; text-overflow:ellipsis; }
  .md-card.unread .md-card-sender { font-weight:700; }
  .md-card-time { font-size:10px; color:#666; flex-shrink:0; }
  .md-card-subject { font-size:11.5px; color:#c9c9c9; margin-bottom:3px; overflow:hidden; text-overflow:ellipsis; white-space:nowrap; }
  .md-card.unread .md-card-subject { color:#e8e8e8; font-weight:600; }
  .md-card-preview { font-size:10.5px; color:#777; line-height:1.45; display:-webkit-box; -webkit-line-clamp:2; -webkit-box-orient:vertical; overflow:hidden; }
  .md-card-meta { display:flex; align-items:center; gap:6px; margin-top:5px; font-size:9.5px; color:#555; }
  .md-card-meta .md-attach { color:#9aa; }
  .md-card-expand { display:none; margin-top:10px; padding-top:9px; border-top:1px solid #2a2a2a; font-size:11px; line-height:1.55; color:#cfcfcf; }
  .md-card.expanded { cursor:default; background:#1a1a1a; border-color:#3a3a3a; }
  .md-card.expanded .md-card-expand { display:block; }
  .md-card-expand .md-full-body { background:#111; border:1px solid #262626; border-radius:6px; padding:8px; max-height:260px; overflow-y:auto; white-space:pre-wrap; word-break:break-word; font-family:'SF Mono',Menlo,monospace; font-size:10.5px; color:#c8c8c8; }
  .md-card-expand .md-meta-row { margin-bottom:4px; color:#888; font-size:10.5px; }
  .md-card-expand .md-meta-row b { color:#aaa; }
  .md-card-expand-actions { margin-top:8px; display:flex; gap:6px; flex-wrap:wrap; }
  .md-btn { background:#2a2a2a; border:1px solid #3a3a3a; color:#ddd; font-size:11px; padding:5px 11px; border-radius:6px; cursor:pointer; font-family:inherit; }
  .md-btn:hover { border-color:#f0c060; color:#f0c060; }
  .md-btn.primary { background:#f0c060; color:#111; border-color:#f0c060; font-weight:700; }
  .md-btn.primary:hover { background:#f5cc70; color:#000; }
  /* Agent Modal */
  #md-agent-modal { display:none; position:fixed; inset:0; background:rgba(0,0,0,0.72); z-index:90; justify-content:center; align-items:center; }
  #md-agent-modal.show { display:flex; }
  .md-modal-card { background:#181818; border:1px solid #2d2d2d; border-radius:12px; padding:22px 24px; width:min(560px,92%); max-height:80vh; overflow-y:auto; }
  .md-modal-card h2 { font-size:14px; font-weight:700; color:#f0c060; margin-bottom:4px; letter-spacing:1px; }
  .md-modal-card .md-modal-sub { font-size:11.5px; color:#888; margin-bottom:16px; }
  .md-agent-list { display:grid; grid-template-columns:1fr 1fr; gap:10px; margin-bottom:14px; }
  .md-agent-choice { background:#1e1e1e; border:1px solid #2e2e2e; border-radius:8px; padding:11px 13px; cursor:pointer; text-align:left; color:#e0e0e0; font-size:12px; font-family:inherit; transition:.12s; }
  .md-agent-choice:hover { border-color:#f0c060; color:#f0c060; background:#22201a; }
  .md-agent-choice.recommended { border-color:#f0c060; background:#22201a; color:#f0c060; font-weight:700; }
  .md-agent-choice .md-agent-desc { font-size:10px; font-weight:400; color:#888; margin-top:3px; display:block; }
  .md-agent-choice.recommended .md-agent-desc { color:#cba976; }
  .md-modal-actions { display:flex; justify-content:flex-end; gap:8px; }
  .md-close { background:none; border:none; color:#666; font-size:20px; cursor:pointer; float:right; margin-top:-6px; }
  .md-close:hover { color:#f0c060; }
  ::-webkit-scrollbar { width:9px; height:9px; }
  ::-webkit-scrollbar-track { background:#111; }
  ::-webkit-scrollbar-thumb { background:#2a2a2a; border-radius:5px; }
  ::-webkit-scrollbar-thumb:hover { background:#3d3d3d; }
  #md-toast { position:fixed; bottom:20px; right:20px; background:#2a2a2a; border:1px solid #3a3a3a; color:#f0c060; padding:10px 16px; border-radius:8px; font-size:12px; z-index:200; opacity:0; transition:opacity .2s; pointer-events:none; }
  #md-toast.show { opacity:1; }
</style>
</head>
<body>
<div id="md-header">
  <h1>📬 Messages</h1>
  <input id="md-search" type="text" placeholder="Alle Spalten durchsuchen..." autocomplete="off">
  <span class="md-hdr-stat" id="md-stat-text">Lade...</span>
  <button class="md-hdr-btn" id="md-btn-refresh">Aktualisieren</button>
  <button class="md-hdr-btn" onclick="location.href='/'">&larr; Chat</button>
</div>
<div id="md-board"></div>

<div id="md-agent-modal">
  <div class="md-modal-card">
    <button class="md-close" onclick="mdCloseAgentModal()">&times;</button>
    <h2>ANTWORTEN MIT AGENT</h2>
    <div class="md-modal-sub" id="md-modal-sub">Welcher Agent soll die Antwort formulieren?</div>
    <div class="md-agent-list" id="md-agent-list"></div>
    <div class="md-modal-actions">
      <button class="md-btn" onclick="mdCloseAgentModal()">Abbrechen</button>
    </div>
  </div>
</div>

<div id="md-toast"></div>

<script>
(function(){
  var BOARD = document.getElementById('md-board');
  var STATE = { sources:[], messages:[], byId:{}, globalQuery:'', colQuery:{}, selectedForReply:null };
  var REFRESH_MS = 60000;
  var agentsCache = null;

  function esc(s){
    if (s === null || s === undefined) return '';
    return String(s).replace(/[&<>"']/g, function(c){
      return ({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'})[c];
    });
  }

  function fmtTime(iso){
    try {
      var d = new Date(iso);
      var now = new Date();
      var diffMs = now - d;
      var diffH = diffMs / 3600000;
      if (diffH < 1) { var m = Math.max(1, Math.round(diffMs / 60000)); return m + 'min'; }
      if (diffH < 24) return Math.round(diffH) + 'h';
      if (diffH < 48) return 'gestern';
      if (diffH < 168) return Math.round(diffH/24) + 'd';
      return d.toLocaleDateString('de-DE', {day:'2-digit', month:'short'});
    } catch(e){ return ''; }
  }

  function fmtFullDate(iso){
    try { return new Date(iso).toLocaleString('de-DE'); } catch(e){ return iso; }
  }

  function showToast(msg){
    var t = document.getElementById('md-toast');
    t.textContent = msg;
    t.classList.add('show');
    clearTimeout(t._to);
    t._to = setTimeout(function(){ t.classList.remove('show'); }, 2200);
  }

  function filterForColumn(sourceKey){
    var q = (STATE.globalQuery || '').toLowerCase().trim();
    var qCol = (STATE.colQuery[sourceKey] || '').toLowerCase().trim();
    return STATE.messages.filter(function(m){
      if (m.source !== sourceKey) return false;
      var hay = (m.sender_name + ' ' + m.subject + ' ' + m.preview + ' ' + (m.sender_address||'')).toLowerCase();
      if (q && hay.indexOf(q) === -1) return false;
      if (qCol && hay.indexOf(qCol) === -1) return false;
      return true;
    }).sort(function(a,b){
      if (a.read !== b.read) return a.read ? 1 : -1;
      if (!a.read) return a.timestamp_epoch - b.timestamp_epoch;
      return b.timestamp_epoch - a.timestamp_epoch;
    });
  }

  function renderColumnCards(col, sourceKey){
    var list = filterForColumn(sourceKey);
    var body = col.querySelector('.md-col-body');
    body.innerHTML = '';
    if (!list.length){
      var e = document.createElement('div');
      e.className = 'md-col-empty';
      e.textContent = 'Keine Nachrichten.';
      body.appendChild(e);
      return;
    }
    list.forEach(function(m){
      var card = document.createElement('div');
      card.className = 'md-card' + (m.read ? '' : ' unread');
      card.setAttribute('data-id', m.id);
      card.innerHTML =
        '<div class="md-card-top">' +
          '<span class="md-dot ' + (m.read ? '' : 'unread') + '"></span>' +
          '<span class="md-card-sender">' + esc(m.sender_name) + '</span>' +
          '<span class="md-card-time">' + esc(fmtTime(m.timestamp)) + '</span>' +
        '</div>' +
        '<div class="md-card-subject">' + esc(m.subject || '(kein Betreff)') + '</div>' +
        '<div class="md-card-preview">' + esc(m.preview || '') + '</div>' +
        (m.has_attachments ? '<div class="md-card-meta"><span class="md-attach">📎 Anhang</span></div>' : '') +
        '<div class="md-card-expand"></div>';
      card.addEventListener('click', function(ev){
        if (ev.detail === 2) return;
        setTimeout(function(){ if (!card._dblClicked) onCardSingleClick(card, m); }, 160);
      });
      card.addEventListener('dblclick', function(){
        card._dblClicked = true;
        setTimeout(function(){ card._dblClicked = false; }, 400);
        openAgentModal(m);
      });
      body.appendChild(card);
    });
  }

  function updateColumnHeaderCounts(){
    var bySrc = {};
    STATE.messages.forEach(function(m){ (bySrc[m.source] = bySrc[m.source] || []).push(m); });
    STATE.sources.forEach(function(s){
      var col = BOARD.querySelector('[data-source="' + s.key + '"]');
      if (!col) return;
      var arr = bySrc[s.key] || [];
      var unread = arr.filter(function(x){ return !x.read; }).length;
      col.querySelector('.md-col-count').textContent = arr.length + ' gesamt';
      var ub = col.querySelector('.md-col-unread');
      if (unread > 0){ ub.style.display = 'inline-block'; ub.textContent = unread; }
      else { ub.style.display = 'none'; }
    });
    var total = STATE.messages.length;
    var totalUnread = STATE.messages.filter(function(x){ return !x.read; }).length;
    document.getElementById('md-stat-text').innerHTML =
      total + ' Nachrichten \u00b7 <strong>' + totalUnread + '</strong> ungelesen';
  }

  function renderBoard(){
    BOARD.innerHTML = '';
    STATE.sources.forEach(function(s){
      var col = document.createElement('div');
      col.className = 'md-col';
      col.setAttribute('data-source', s.key);
      col.innerHTML =
        '<div class="md-col-hdr">' +
          '<div class="md-col-title">' +
            '<span class="md-col-icon">' + s.icon + '</span>' +
            '<span>' + esc(s.label) + '</span>' +
            '<span class="md-col-unread" style="display:none"></span>' +
            '<span class="md-col-count">0 gesamt</span>' +
          '</div>' +
          '<input class="md-col-search" type="text" placeholder="In dieser Spalte..." autocomplete="off">' +
        '</div>' +
        '<div class="md-col-body"></div>';
      var inp = col.querySelector('.md-col-search');
      inp.addEventListener('input', function(){
        STATE.colQuery[s.key] = inp.value;
        renderColumnCards(col, s.key);
      });
      BOARD.appendChild(col);
      renderColumnCards(col, s.key);
    });
    updateColumnHeaderCounts();
  }

  function renderAllColumns(){
    STATE.sources.forEach(function(s){
      var col = BOARD.querySelector('[data-source="' + s.key + '"]');
      if (col) renderColumnCards(col, s.key);
    });
    updateColumnHeaderCounts();
  }

  async function onCardSingleClick(card, msg){
    if (card.classList.contains('expanded')){
      card.classList.remove('expanded');
      return;
    }
    // Alle anderen zuklappen in derselben Spalte
    var col = card.closest('.md-col');
    col.querySelectorAll('.md-card.expanded').forEach(function(c){ c.classList.remove('expanded'); });
    card.classList.add('expanded');
    var expandEl = card.querySelector('.md-card-expand');
    expandEl.innerHTML = '<em style="color:#666">Lade Vollansicht...</em>';
    try {
      var r = await fetch('/api/messages/' + encodeURIComponent(msg.id));
      var d = await r.json();
      if (!d.ok){ expandEl.innerHTML = '<span style="color:#c66">Fehler: ' + esc(d.error||'unbekannt') + '</span>'; return; }
      var m = d.message;
      expandEl.innerHTML =
        '<div class="md-meta-row"><b>Von:</b> ' + esc(m.sender_name) + (m.sender_address ? ' &lt;' + esc(m.sender_address) + '&gt;' : '') + '</div>' +
        (m.to ? '<div class="md-meta-row"><b>An:</b> ' + esc(m.to) + '</div>' : '') +
        '<div class="md-meta-row"><b>Datum:</b> ' + esc(fmtFullDate(m.timestamp)) + ' \u00b7 <b>Quelle:</b> ' + esc(m.source) + '</div>' +
        (m.message_id ? '<div class="md-meta-row"><b>Message-ID:</b> <code style="font-size:9px;color:#666">' + esc(m.message_id) + '</code></div>' : '') +
        '<div class="md-full-body">' + esc(m.full_content || m.preview || '(leer)') + '</div>' +
        '<div class="md-card-expand-actions">' +
          '<button class="md-btn primary" data-action="reply">\u21A9\uFE0E Mit Agent antworten</button>' +
          '<button class="md-btn" data-action="read">' + (msg.read ? 'Als ungelesen' : 'Als gelesen') + '</button>' +
          '<button class="md-btn" data-action="finder">📂 Im Finder</button>' +
        '</div>';
      expandEl.querySelector('[data-action="reply"]').addEventListener('click', function(e){ e.stopPropagation(); openAgentModal(m); });
      expandEl.querySelector('[data-action="read"]').addEventListener('click', function(e){ e.stopPropagation(); toggleRead(msg); });
      expandEl.querySelector('[data-action="finder"]').addEventListener('click', function(e){ e.stopPropagation(); openFinder(m.raw_file_path); });
      // Auto-mark as read
      if (!msg.read){ await markRead(msg.id, true); msg.read = true; card.classList.remove('unread'); card.querySelector('.md-dot').classList.remove('unread'); updateColumnHeaderCounts(); }
    } catch(e){
      expandEl.innerHTML = '<span style="color:#c66">Fehler: ' + esc(e.message) + '</span>';
    }
  }

  async function markRead(id, read){
    try {
      await fetch('/api/messages/mark-read', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({message_id:id, read:!!read})});
    } catch(e){}
  }

  async function toggleRead(msg){
    msg.read = !msg.read;
    await markRead(msg.id, msg.read);
    renderAllColumns();
    showToast(msg.read ? 'Als gelesen markiert' : 'Als ungelesen markiert');
  }

  async function openFinder(path){
    try {
      await fetch('/open_in_finder', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({path:path})});
    } catch(e){ showToast('Finder-Fehler'); }
  }

  async function ensureAgents(){
    if (agentsCache) return agentsCache;
    var r = await fetch('/agents');
    agentsCache = await r.json();
    return agentsCache;
  }

  async function openAgentModal(msg){
    STATE.selectedForReply = msg;
    document.getElementById('md-modal-sub').innerHTML =
      'Antwort auf: <strong>' + esc(msg.subject || '(kein Betreff)') + '</strong> von ' + esc(msg.sender_name);
    var list = document.getElementById('md-agent-list');
    list.innerHTML = '<em style="color:#666">Lade Agenten...</em>';
    try {
      var agents = await ensureAgents();
      var recommended = (STATE.sources.find(function(s){ return s.key === msg.source; }) || {}).recommended_agent;
      list.innerHTML = '';
      agents.forEach(function(a){
        var btn = document.createElement('button');
        btn.className = 'md-agent-choice' + (a.name === recommended ? ' recommended' : '');
        btn.innerHTML = '<span>' + esc(a.label || a.name) + (a.name === recommended ? ' \u2605' : '') + '</span>' +
                        (a.description ? '<span class="md-agent-desc">' + esc(a.description) + '</span>' : '');
        btn.addEventListener('click', function(){ openChatWithMessage(a.name, msg.id); });
        list.appendChild(btn);
      });
    } catch(e){
      list.innerHTML = '<span style="color:#c66">Fehler beim Laden der Agenten.</span>';
    }
    document.getElementById('md-agent-modal').classList.add('show');
  }

  function mdCloseAgentModal(){ document.getElementById('md-agent-modal').classList.remove('show'); }
  window.mdCloseAgentModal = mdCloseAgentModal;

  function openChatWithMessage(agentName, msgId){
    var url = '/?agent=' + encodeURIComponent(agentName) + '&preload_message=' + encodeURIComponent(msgId);
    try { window.open(url, '_blank'); } catch(e) { location.href = url; }
    mdCloseAgentModal();
  }

  async function loadSources(){
    var r = await fetch('/api/messages/sources');
    var d = await r.json();
    if (!d.ok) throw new Error(d.error || 'sources load failed');
    STATE.sources = d.sources;
  }

  async function loadMessages(force){
    var r = await fetch('/api/messages' + (force ? '?refresh=1' : ''));
    var d = await r.json();
    if (!d.ok) throw new Error(d.error || 'messages load failed');
    STATE.messages = d.messages;
    STATE.byId = {};
    STATE.messages.forEach(function(m){ STATE.byId[m.id] = m; });
  }

  async function fullReload(force){
    try {
      await loadSources();
      await loadMessages(force);
      renderBoard();
    } catch(e){
      BOARD.innerHTML = '<div style="padding:40px;color:#c66">Fehler: ' + esc(e.message) + '</div>';
    }
  }

  async function softRefresh(){
    try {
      var prevIds = new Set(STATE.messages.map(function(x){ return x.id; }));
      await loadSources();
      await loadMessages(false);
      renderAllColumns();
      // Neue Messages highlighten
      var newIds = STATE.messages.map(function(x){ return x.id; }).filter(function(x){ return !prevIds.has(x); });
      if (newIds.length){
        showToast(newIds.length + ' neue Nachricht(en)');
        newIds.forEach(function(id){
          var el = document.querySelector('[data-id="' + id + '"]');
          if (el){ el.classList.add('highlight'); setTimeout(function(){ el.classList.remove('highlight'); }, 2000); }
        });
      }
    } catch(e){ /* silent */ }
  }

  document.getElementById('md-search').addEventListener('input', function(e){
    STATE.globalQuery = e.target.value;
    renderAllColumns();
  });
  document.getElementById('md-btn-refresh').addEventListener('click', function(){
    var b = this; b.disabled = true; b.textContent = 'Lade...';
    fullReload(true).finally(function(){ b.disabled = false; b.textContent = 'Aktualisieren'; });
  });
  document.getElementById('md-agent-modal').addEventListener('click', function(e){
    if (e.target.id === 'md-agent-modal') mdCloseAgentModal();
  });

  fullReload(false);
  setInterval(softRefresh, REFRESH_MS);
})();
</script>
</body>
</html>
"""

# ── ENDE MESSAGE DASHBOARD ───────────────────────────────────────────────────


if __name__ == '__main__':
    # Browser wird NICHT mehr automatisch geoeffnet — native AssistantDev.app nutzt pywebview
    # Cleanup old sessions every hour
    def session_cleanup_loop():
        import time
        while True:
            time.sleep(3600)
            cleanup_old_sessions()
    threading.Thread(target=session_cleanup_loop, daemon=True).start()

    # Graceful Shutdown: SIGTERM abfangen, laufende Requests abwarten, dann sauber beenden
    import atexit
    def _save_all_sessions_on_exit():
        print('[AUTO-SAVE] Shutdown erkannt — sichere alle Sessions...')
        for sid in list(sessions.keys()):
            try:
                auto_save_session(sid)
                print(f'[AUTO-SAVE] Session {sid[:12]} gesichert')
            except Exception as e:
                print(f'[AUTO-SAVE] Fehler bei {sid[:12]}: {e}')

    def _graceful_shutdown(sig, frame):
        print(f'[SHUTDOWN] Signal {sig} empfangen — starte Graceful Shutdown...')
        _shutdown_event.set()
        import time as _shutdown_time
        for i in range(30):
            with _active_requests_lock:
                active = _active_requests
            if active == 0:
                print(f'[SHUTDOWN] Alle Requests abgeschlossen nach {i}s')
                break
            print(f'[SHUTDOWN] Warte auf {active} aktive(n) Request(s)... ({i+1}/30s)')
            _shutdown_time.sleep(1)
        else:
            with _active_requests_lock:
                active = _active_requests
            if active > 0:
                print(f'[SHUTDOWN] Timeout — {active} Request(s) noch aktiv, beende trotzdem')
        _save_all_sessions_on_exit()
        print('[SHUTDOWN] Sauber beendet.')
        sys.exit(0)

    atexit.register(_save_all_sessions_on_exit)
    signal.signal(signal.SIGTERM, _graceful_shutdown)
    # Build global search index in background at startup
    if build_global_index_async:
        build_global_index_async()
    # Background-backfill semantic embeddings for all agents (no-op without OpenAI key)
    if reindex_all_embeddings_async:
        try:
            reindex_all_embeddings_async()
        except Exception as e:
            print(f'[EMBEDDINGS] Initial backfill error: {e}')
    # Periodic incremental index update (every 5 minutes)
    def index_update_loop():
        import time
        time.sleep(60)  # Wait 1 min after startup before first run
        tick = 0
        while True:
            try:
                if update_all_indexes:
                    update_all_indexes()
                # Every 6th tick (~30 min) also catch new files for embeddings
                if reindex_all_embeddings_async and tick % 6 == 0:
                    reindex_all_embeddings_async()
            except Exception as e:
                print(f'[INDEX] Periodischer Update Fehler: {e}')
            tick += 1
            time.sleep(300)  # Every 5 minutes
    threading.Thread(target=index_update_loop, daemon=True).start()
    print('\nAssistant Web Interface laeuft auf http://localhost:8080')
    print('Zum Beenden: Control+C\n')
    app.run(host='127.0.0.1', port=8080, debug=False)
